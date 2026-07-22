"""
Evaluation script for canvas-enrollment-overview-gsheet-ppt task.

Checks:
1. Google Sheet "Fall 2013 Enrollment Overview" exists with correct enrollment data
2. PowerPoint file Enrollment_Overview_F2013.pptx exists with correct structure
3. Word document Enrollment_Report_F2013.docx exists with correct content
4. Email sent to academic.office@university.edu
"""

import argparse
import json
import os
import sys

import psycopg2

DB_CONFIG = {
    "host": os.environ.get("PGHOST", "localhost"),
    "port": 5432,
    "dbname": "toolathlon_gym",
    "user": "eigent",
    "password": "camel",
}

PASS_COUNT = 0
FAIL_COUNT = 0


def record(name, passed, detail=""):
    global PASS_COUNT, FAIL_COUNT
    if passed:
        PASS_COUNT += 1
        print(f"  [PASS] {name}")
    else:
        FAIL_COUNT += 1
        msg = f": {detail[:300]}" if detail else ""
        print(f"  [FAIL] {name}{msg}")


def num_close(a, b, tol=1.0):
    if a is None and b is None:
        return True
    if a is None or b is None:
        return False
    try:
        return abs(float(a) - float(b)) <= tol
    except (TypeError, ValueError):
        return str(a).strip().lower() == str(b).strip().lower()


# ============================================================================
# Check 1: Google Sheet
# ============================================================================

def _expected_fall_2013_codes():
    """Query canvas DB for Fall 2013 course codes (ending in '2013J')."""
    try:
        conn = psycopg2.connect(**DB_CONFIG)
        cur = conn.cursor()
        cur.execute(
            "SELECT course_code, total_students FROM canvas.courses "
            "WHERE course_code LIKE %s",
            ("%2013J",),
        )
        rows = cur.fetchall()
        cur.close()
        conn.close()
        return [(c.lower(), int(ts) if ts is not None else 0) for c, ts in rows]
    except Exception:
        # Fallback: known Fall 2013 codes
        return [("aaa-2013j", 383), ("bbb-2013j", 2237), ("ddd-2013j", 1938),
                ("eee-2013j", 1052), ("fff-2013j", 2283), ("ggg-2013j", 952)]


def check_gsheet():
    print("\n=== Checking Google Sheet ===")

    conn = psycopg2.connect(**DB_CONFIG)
    cur = conn.cursor()

    # Require exact title 'Fall 2013 Enrollment Overview'
    cur.execute(
        "SELECT id, title FROM gsheet.spreadsheets "
        "WHERE title ILIKE %s",
        ("Fall 2013 Enrollment Overview",),
    )
    sheets = cur.fetchall()
    record("Spreadsheet titled 'Fall 2013 Enrollment Overview' exists",
           len(sheets) >= 1, f"Found {len(sheets)} matching spreadsheets")

    if not sheets:
        cur.close()
        conn.close()
        return False

    ss_id, title = sheets[0]

    # Get cell data
    cur.execute("""
        SELECT c.row_index, c.col_index, c.value
        FROM gsheet.cells c
        JOIN gsheet.sheets s ON c.sheet_id = s.id
        WHERE c.spreadsheet_id = %s
        ORDER BY c.row_index, c.col_index
    """, (ss_id,))
    cells = cur.fetchall()

    grid = {}
    for row_i, col_i, val in cells:
        if row_i not in grid:
            grid[row_i] = {}
        grid[row_i][col_i] = val

    expected = _expected_fall_2013_codes()
    n_courses = len(expected)
    max_row = max(grid.keys()) if grid else 0
    record(f"Spreadsheet has header + {n_courses} courses",
           max_row >= n_courses,
           f"Max row index: {max_row}; expected at least {n_courses}")

    # Check for ALL expected Fall 2013 course codes in data
    all_vals = [str(v).lower() for row in grid.values() for v in row.values() if v]
    for code, _ts in expected:
        found_code = any(code in v for v in all_vals)
        record(f"Spreadsheet contains {code.upper()}", found_code)

    # Check for numeric student counts; require sum of integers >= sum(expected)
    numeric_vals = []
    for row in grid.values():
        for v in row.values():
            try:
                numeric_vals.append(float(v))
            except (TypeError, ValueError):
                pass
    expected_total = sum(ts for _, ts in expected)
    # Require all expected per-course student counts to appear OR the total
    # (Strict: cannot pass with placeholder "100").
    expected_counts = [ts for _, ts in expected]
    matched_counts = sum(
        1 for ec in expected_counts if any(num_close(v, ec, 5) for v in numeric_vals)
    )
    has_total = any(num_close(v, expected_total, 5) for v in numeric_vals)
    record("Spreadsheet has numeric enrollment data matching expected per-course counts",
           matched_counts >= max(1, len(expected_counts) - 1) or has_total,
           f"Matched {matched_counts}/{len(expected_counts)} counts; total found: {has_total}")

    cur.close()
    conn.close()
    return True


# ============================================================================
# Check 2: PowerPoint
# ============================================================================

def check_pptx(agent_workspace):
    print("\n=== Checking Enrollment_Overview_F2013.pptx ===")

    pptx_path = os.path.join(agent_workspace, "Enrollment_Overview_F2013.pptx")
    if not os.path.isfile(pptx_path):
        record("PPT file exists", False, f"Not found: {pptx_path}")
        return False
    record("PPT file exists", True)

    try:
        from pptx import Presentation
        prs = Presentation(pptx_path)
        slide_count = len(prs.slides)
        record("PPT has at least 8 slides (title + 6 courses + summary)",
               slide_count >= 8,
               f"Found {slide_count} slides")

        # Collect all text from slides
        all_text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    all_text += shape.text.lower() + " "

        record("PPT mentions Fall 2013",
               "fall 2013" in all_text or "2013" in all_text,
               "Missing year 2013")
        record("PPT mentions enrollment",
               "enrollment" in all_text or "enrolled" in all_text or "students" in all_text,
               "Missing enrollment content")

        # Check for ALL Fall 2013 course codes
        expected = _expected_fall_2013_codes()
        for code, _ts in expected:
            record(f"PPT mentions {code.upper()}",
                   code in all_text,
                   f"Missing {code} in slides")

        n_courses = len(expected)
        record(f"PPT has at least {n_courses + 2} slides (title + per-course + summary)",
               slide_count >= n_courses + 2,
               f"Found {slide_count} slides; expected at least {n_courses + 2}")

        return slide_count >= n_courses + 2

    except ImportError:
        size = os.path.getsize(pptx_path)
        record("PPT file has content (>5KB)", size > 5000, f"Size: {size} bytes")
        return size > 5000
    except Exception as e:
        record("PPT file readable", False, str(e))
        return False


# ============================================================================
# Check 3: Word document
# ============================================================================

def check_word(agent_workspace):
    print("\n=== Checking Enrollment_Report_F2013.docx ===")

    docx_path = os.path.join(agent_workspace, "Enrollment_Report_F2013.docx")
    if not os.path.isfile(docx_path):
        record("Word file exists", False, f"Not found: {docx_path}")
        return False
    record("Word file exists", True)

    try:
        from docx import Document
        doc = Document(docx_path)
        all_text = " ".join(p.text for p in doc.paragraphs).lower()
        headings = " ".join(p.text for p in doc.paragraphs
                           if p.style.name.startswith("Heading")).lower()

        record("Word doc has substantial content",
               len(all_text.strip()) >= 100,
               f"Content length: {len(all_text.strip())}")
        record("Word doc mentions enrollment",
               "enrollment" in all_text or "enroll" in all_text,
               "Missing 'enrollment' in document")
        record("Word doc mentions Fall 2013",
               "fall 2013" in all_text or "2013" in all_text,
               "Missing 2013")

        # Check for table
        tables = doc.tables
        record("Word doc has at least 1 table",
               len(tables) >= 1,
               f"Found {len(tables)} tables")

        # Check for total enrollment number (8845 = sum of Fall 2013 students)
        expected = _expected_fall_2013_codes()
        expected_total = sum(ts for _, ts in expected)
        total_strs = [str(expected_total), f"{expected_total:,}"]
        record(f"Word doc mentions total enrollment ({expected_total})",
               any(s in all_text for s in total_strs),
               f"Missing total {expected_total}; expected one of {total_strs}")

        return True

    except ImportError:
        size = os.path.getsize(docx_path)
        record("Word file has content (>3KB)", size > 3000, f"Size: {size} bytes")
        return size > 3000
    except Exception as e:
        record("Word file readable", False, str(e))
        return False


# ============================================================================
# Check 4: Email
# ============================================================================

def check_emails():
    print("\n=== Checking Emails ===")

    conn = psycopg2.connect(**DB_CONFIG)
    cur = conn.cursor()

    # Require recipient = academic.office@university.edu (AND not OR)
    cur.execute("""
        SELECT subject, from_addr, to_addr, body_text
        FROM email.messages
        WHERE to_addr::text ILIKE %s
    """, ("%academic.office@university.edu%",))
    addressed = cur.fetchall()
    record("Email to academic.office@university.edu found",
           len(addressed) >= 1, f"Found {len(addressed)} matching")

    if not addressed:
        cur.close()
        conn.close()
        return False

    # Among recipient-correct emails, require subject contains 'Fall 2013'
    # AND ('enrollment' or 'semester')
    valid = None
    for subject, from_addr, to_addr, body_text in addressed:
        sub = (subject or "").lower()
        if "fall 2013" in sub and ("enrollment" in sub or "semester" in sub):
            valid = (subject, from_addr, to_addr, body_text)
            break
    record("Email subject contains 'Fall 2013' AND ('enrollment' or 'semester')",
           valid is not None,
           f"Subjects observed: {[e[0] for e in addressed]}")

    if valid is None:
        cur.close()
        conn.close()
        return False

    body = (valid[3] or "").lower()
    expected = _expected_fall_2013_codes()
    expected_total = sum(ts for _, ts in expected)
    # Body must contain the actual numeric total enrollment.
    record(f"Email body mentions total enrollment number ({expected_total})",
           str(expected_total) in body or f"{expected_total:,}" in body,
           f"Body missing total {expected_total}")
    record("Email body mentions number of courses",
           "course" in body,
           "Body missing course count info")

    cur.close()
    conn.close()
    return True


# ============================================================================
# Main
# ============================================================================

def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--agent_workspace", required=False, default=".")
    parser.add_argument("--groundtruth_workspace", required=False, default=".")
    parser.add_argument("--launch_time", required=False)
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()

    gsheet_ok = check_gsheet()
    pptx_ok = check_pptx(args.agent_workspace)
    word_ok = check_word(args.agent_workspace)
    email_ok = check_emails()

    # Aggregation gate must include FAIL_COUNT==0 so per-row record() failures
    # propagate. Each check_* may return True at the end of its happy path while
    # individual record() calls (per-course code, numeric counts, body content)
    # have already incremented FAIL_COUNT.
    all_passed = (
        gsheet_ok and pptx_ok and word_ok and email_ok and FAIL_COUNT == 0
    )

    print(f"\n=== SUMMARY ===")
    print(f"  Passed: {PASS_COUNT}")
    print(f"  Failed: {FAIL_COUNT}")
    print(f"  Overall: {'PASS' if all_passed else 'FAIL'}")

    if args.res_log_file:
        result = {
            "passed": PASS_COUNT,
            "failed": FAIL_COUNT,
            "success": all_passed,
        }
        with open(args.res_log_file, "w") as f:
            json.dump(result, f, indent=2)

    sys.exit(0 if all_passed else 1)


if __name__ == "__main__":
    main()
