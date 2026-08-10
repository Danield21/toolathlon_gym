"""Evaluation for terminal-wc-yf-ppt-notion-email.
Checks:
1. Market_Strategy_Presentation.pptx (7 slides)
2. Notion "Market Strategy Tracker" database with 5 entries
3. Email to ceo@company.com
4. Email to marketing_team@company.com
5. market_correlation.py and category_analysis.py scripts exist
6. market_correlation.json output
7. category_market_analysis.json output
"""
import argparse
import json
import os
import sys

import psycopg2
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

DB_CONFIG = {
    "host": os.environ.get("PGHOST", "localhost"),
    "port": int(os.environ.get("PGPORT", "5432")),
    "dbname": os.environ.get("PGDATABASE", "toolathlon_gym"),
    "user": os.environ.get("PGUSER", "eigent"),
    "password": os.environ.get("PGPASSWORD", "camel"),
}

PASS_COUNT = 0
FAIL_COUNT = 0


def get_expected_from_db():
    """Query WC and YF schemas dynamically for expected values."""
    defaults = {
        "top_category": "electronics",
        "amzn_pct_change": 26.0,
        "jpm_pct_change": 62.0,
        "dji_pct_change": 24.0,
        # Alternate interpretation: pct change over the shared correlation
        # window (2025-03 .. 2026-02) instead of the full price history.
        # task.md explicitly anchors the email/ppt on the full-history values
        # (26/62/24), but a reasonable agent that reuses the correlation window
        # for the pct computation must not be failed.
        "amzn_pct_change_window": 2.4,
        "jpm_pct_change_window": 17.6,
        "dji_pct_change_window": 13.4,
    }
    try:
        conn = psycopg2.connect(**DB_CONFIG)
        cur = conn.cursor()

        # Top category by revenue
        cur.execute("""
            SELECT p.categories->0->>'name' as cat, SUM((li->>'total')::numeric) as rev
            FROM wc.orders o, jsonb_array_elements(o.line_items) li
            JOIN wc.products p ON (li->>'product_id')::int = p.id
            WHERE o.status NOT IN ('cancelled','refunded','failed')
            GROUP BY cat ORDER BY rev DESC LIMIT 1
        """)
        row = cur.fetchone()
        if row and row[0]:
            defaults["top_category"] = row[0].lower()

        for symbol, key, wkey in [("AMZN", "amzn_pct_change", "amzn_pct_change_window"),
                                  ("JPM", "jpm_pct_change", "jpm_pct_change_window"),
                                  ("^DJI", "dji_pct_change", "dji_pct_change_window")]:
            cur.execute("""
                SELECT
                    (SELECT close FROM yf.stock_prices WHERE symbol=%s ORDER BY date ASC LIMIT 1),
                    (SELECT close FROM yf.stock_prices WHERE symbol=%s ORDER BY date DESC LIMIT 1)
            """, (symbol, symbol))
            row = cur.fetchone()
            if row and row[0] and row[1]:
                defaults[key] = float((row[1] - row[0]) / row[0] * 100)
            # First/last trading day inside the correlation window 2025-03..2026-02
            cur.execute("""
                SELECT
                    (SELECT close FROM yf.stock_prices WHERE symbol=%s
                        AND date >= '2025-03-01' AND date < '2026-03-01' ORDER BY date ASC LIMIT 1),
                    (SELECT close FROM yf.stock_prices WHERE symbol=%s
                        AND date >= '2025-03-01' AND date < '2026-03-01' ORDER BY date DESC LIMIT 1)
            """, (symbol, symbol))
            row = cur.fetchone()
            if row and row[0] and row[1]:
                defaults[wkey] = float((row[1] - row[0]) / row[0] * 100)

        cur.close()
        conn.close()
    except Exception as e:
        print(f"  [WARN] DB query for expected values failed, using defaults: {e}")
    return defaults


EXPECTED = get_expected_from_db()


def check(name, condition, detail=""):
    global PASS_COUNT, FAIL_COUNT
    if condition:
        PASS_COUNT += 1
        print(f"  [PASS] {name}")
    else:
        FAIL_COUNT += 1
        print(f"  [FAIL] {name}: {str(detail)[:200]}")


def _db_conn():
    """Open a DB connection, or warn-and-return-None when the DB is unreachable
    so file-based checks can still complete instead of crashing the evaluator."""
    try:
        return psycopg2.connect(**DB_CONFIG)
    except Exception as e:
        print(f"  [WARN] DB connection failed, skipping DB checks: {e}")
        return None


def _to_float(v):
    """Robustly coerce a value to float. Returns None if not parseable."""
    if v is None:
        return None
    if isinstance(v, bool):
        return None
    if isinstance(v, (int, float)):
        return float(v)
    s = str(v).strip()
    if not s:
        return None
    # strip thousands separators, currency symbols, percent sign, spaces
    s = s.replace(",", "").replace("$", "").replace("¥", "").replace("€", "")
    s = s.replace("%", "").replace(" ", "").strip()
    try:
        return float(s)
    except ValueError:
        return None


def num_close(a, b, tol=2.0):
    fa, fb = _to_float(a), _to_float(b)
    if fa is not None and fb is not None:
        return abs(fa - fb) <= tol
    return str(a).strip().lower() == str(b).strip().lower()


def check_pptx(workspace):
    print("\n=== Check 1: Market_Strategy_Presentation.pptx ===")
    path = os.path.join(workspace, "Market_Strategy_Presentation.pptx")
    if not os.path.exists(path):
        check("PPTX file exists", False, f"Not found at {path}")
        return
    check("PPTX file exists", True)

    prs = Presentation(path)
    slides = list(prs.slides)
    check("Has 7 slides", len(slides) == 7, f"Found {len(slides)}")

    # Collect all text from all slides (including table cell text).
    # Table shapes have has_text_frame == False, so iterate cells explicitly.
    def _shape_text(shape):
        texts = []
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                if para.text:
                    texts.append(para.text)
        if getattr(shape, "has_table", False):
            try:
                for row in shape.table.rows:
                    for cell in row.cells:
                        if cell.text:
                            texts.append(cell.text)
            except Exception:
                pass
        # Recurse into grouped shapes: a group's own has_text_frame is False,
        # but its children may carry the real text.
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            try:
                for sub in shape.shapes:
                    texts.extend(_shape_text(sub))
            except Exception:
                pass
        # Charts (add_chart) have has_text_frame/has_table == False. Extract any
        # visible text runs (titles, axis titles, rich-text data labels) plus the
        # plotted data values (what data labels display). Each data value is
        # tagged with 'pct' so the percentage checks can see numbers that were
        # only shown as chart data labels (e.g. 26.18 plotted as a label).
        if getattr(shape, "has_chart", False):
            try:
                chart = shape.chart
                space = chart._chartSpace
                ns_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
                ns_c = "http://schemas.openxmlformats.org/drawingml/2006/chart"
                for t in space.iter(f"{{{ns_a}}}t"):
                    if t.text:
                        texts.append(t.text)
                for v in space.iter(f"{{{ns_c}}}v"):
                    if v.text and v.text.strip():
                        texts.append(f"{v.text.strip()} pct")
            except Exception:
                pass
        return texts

    all_texts = []
    for slide in slides:
        slide_text = []
        for shape in slide.shapes:
            slide_text.extend(_shape_text(shape))
        all_texts.append(" ".join(slide_text).lower())

    full_text = " ".join(all_texts)

    # Slide 1: Title
    check("Slide 1 has title 'Market Positioning'",
          "market" in all_texts[0] and "positioning" in all_texts[0] if len(all_texts) > 0 else False,
          all_texts[0][:100] if all_texts else "no slides")

    # Slide 2: Market Overview with AMZN, JPM, DJI
    if len(all_texts) > 1:
        check("Slide 2 mentions AMZN", "amzn" in all_texts[1] or "amazon" in all_texts[1],
              all_texts[1][:100])
        check("Slide 2 mentions JPM", "jpm" in all_texts[1] or "jpmorgan" in all_texts[1],
              all_texts[1][:100])
        check("Slide 2 mentions DJI", "dji" in all_texts[1] or "dow" in all_texts[1],
              all_texts[1][:100])

    # Slide 3: Revenue trends
    if len(all_texts) > 2:
        check("Slide 3 has revenue data",
              "revenue" in all_texts[2] or "2025" in all_texts[2],
              all_texts[2][:100])

    # Slide 4: Category performance
    if len(all_texts) > 3:
        check("Slide 4 mentions Electronics",
              "electronics" in all_texts[3], all_texts[3][:100])
        check("Slide 4 mentions category",
              "category" in all_texts[3] or "tv" in all_texts[3] or "audio" in all_texts[3],
              all_texts[3][:100])

    # Slide 5: Correlation
    if len(all_texts) > 4:
        check("Slide 5 mentions correlation",
              "correlation" in all_texts[4], all_texts[4][:100])

    # Slide 6: Recommendations
    if len(all_texts) > 5:
        check("Slide 6 has recommendations",
              "recommend" in all_texts[5] or "pricing" in all_texts[5] or "strategy" in all_texts[5],
              all_texts[5][:100])

    # Slide 7: Next steps
    if len(all_texts) > 6:
        check("Slide 7 has next steps",
              "next" in all_texts[6] or "action" in all_texts[6] or "step" in all_texts[6],
              all_texts[6][:100])

    # Check key numbers appear somewhere (dynamically computed)
    # Accept rounded-to-int OR actual 1-decimal value OR either +/-1 int
    import re
    def has_pct_near(text, target_pct, tol=2.0):
        """Return True if text contains a pct number within tol of target.
        Accepts '%', 'percent', 'per cent', 'per cent', 'pct' as the unit
        suffix so agents writing 'approximately 26 percent' also match."""
        for m in re.finditer(r"(-?\d+\.?\d*)\s*(?:%|percent|per\s*cent|pct)", text):
            try:
                v = float(m.group(1))
                if abs(v - target_pct) <= tol:
                    return True
            except Exception:
                pass
        return False

    # The pct-change anchor is the full price history (earliest-to-latest data
    # point, matching the email's 26/62/24 figures), but a correct agent may
    # consistently reuse the correlation window 2025-03..2026-02 (2.4/17.6/13.4).
    # Accept either interpretation so a genuinely correct computation is never
    # failed purely on which window was chosen.
    def pct_check(name, full_target, window_target):
        near_full = has_pct_near(full_text, full_target)
        near_win = has_pct_near(full_text, window_target)
        check(f"{name}",
              near_full or near_win,
              f"No percent near {full_target:.1f}% (full history) or {window_target:.1f}% (corr window)")

    pct_check("Mentions AMZN ~26% change", EXPECTED["amzn_pct_change"], EXPECTED["amzn_pct_change_window"])
    pct_check("Mentions JPM ~62% change", EXPECTED["jpm_pct_change"], EXPECTED["jpm_pct_change_window"])
    pct_check("Mentions DJI ~24% change", EXPECTED["dji_pct_change"], EXPECTED["dji_pct_change_window"])
    check(f"Mentions {EXPECTED['top_category'].title()} as top category",
          EXPECTED["top_category"] in full_text,
          f"No {EXPECTED['top_category']} mention")


def check_notion():
    print("\n=== Check 2: Notion Market Strategy Tracker ===")
    conn = _db_conn()
    if conn is None:
        return
    cur = conn.cursor()
    try:
        cur.execute("SELECT id, title FROM notion.databases")
        dbs = cur.fetchall()
        tracker_db = None
        for db_id, title in dbs:
            title_str = ""
            if isinstance(title, list):
                title_str = " ".join(
                    item.get("text", {}).get("content", "")
                    for item in title if isinstance(item, dict))
            elif isinstance(title, str):
                try:
                    parsed = json.loads(title)
                    if isinstance(parsed, list):
                        title_str = " ".join(
                            item.get("text", {}).get("content", "")
                            for item in parsed if isinstance(item, dict))
                    else:
                        title_str = str(title)
                except Exception:
                    title_str = str(title)
            else:
                title_str = str(title) if title else ""
            if "market" in title_str.lower() and "strategy" in title_str.lower() and "tracker" in title_str.lower():
                tracker_db = (db_id, title_str)
                break

        check("Market Strategy Tracker DB exists", tracker_db is not None,
              f"Databases found: {[d[1] for d in dbs]}")

        if tracker_db:
            # Check properties
            cur.execute("SELECT properties FROM notion.databases WHERE id = %s", (tracker_db[0],))
            props = cur.fetchone()[0]
            if isinstance(props, str):
                props = json.loads(props)
            prop_names = [k.lower().replace("_", " ") for k in props.keys()] if props else []
            check("Has Initiative property",
                  any("initiative" in p or "title" in str(props.get(k, {}).get("type", ""))
                      for k, p in [(k, k.lower()) for k in (props or {}).keys()]),
                  f"Props: {list((props or {}).keys())}")
            check("Has Category property",
                  any("category" in p for p in prop_names),
                  f"Props: {prop_names}")
            check("Has Market_Condition property",
                  any("market" in p and "condition" in p for p in prop_names),
                  f"Props: {prop_names}")

            # Check pages
            cur.execute("""
                SELECT id, properties FROM notion.pages
                WHERE parent->>'database_id' = %s AND NOT archived
            """, (tracker_db[0],))
            pages = cur.fetchall()
            # >= 5 (not == 5): in multi-agent parallel runs a duplicate creation
            # would produce more rows, but all 5 required entries must still exist.
            check("Has at least 5 initiative entries", len(pages) >= 5,
                  f"Found {len(pages)}")

            # Aggregate all page property texts for content validation
            page_titles = []
            all_page_props_text = []
            for pid, props in pages:
                if isinstance(props, str):
                    props = json.loads(props)
                all_page_props_text.append(json.dumps(props).lower())
                for k, v in (props or {}).items():
                    if isinstance(v, dict) and "title" in v:
                        title_items = v.get("title", [])
                        if isinstance(title_items, list):
                            for item in title_items:
                                if isinstance(item, dict):
                                    page_titles.append(item.get("text", {}).get("content", ""))
            all_titles = " ".join(page_titles).lower()
            combined_props = " ".join(all_page_props_text)

            check("Has TV & Home Theater initiative",
                  "tv" in all_titles or "home theater" in all_titles, f"Titles: {page_titles}")
            check("Has Electronics initiative",
                  "electronics" in all_titles, f"Titles: {page_titles}")
            check("Has Audio initiative",
                  "audio" in all_titles, f"Titles: {page_titles}")

            # Verify Expected_Impact values appear (task.md specifies 5000/3000/2500/1500/2000)
            expected_impacts = ["5000", "3000", "2500", "1500", "2000"]
            matched = sum(1 for v in expected_impacts if v in combined_props)
            check("Notion entries include ALL 5 expected impact values",
                  matched == 5, f"Only {matched}/5 impact values found")

            # Verify expected categories (Pricing/Marketing/Expansion/Inventory/Cost Management)
            expected_cats = ["pricing", "marketing", "expansion", "inventory", "cost"]
            cat_matched = sum(1 for c in expected_cats if c in combined_props)
            check("Notion entries include ALL 5 expected categories",
                  cat_matched == 5, f"Only {cat_matched}/5 categories found")

            # Verify Timeline dates (per task.md: 2026-04-15, 2026-04-01, 2026-05-01, 2026-04-10, 2026-05-15)
            expected_timelines = ["2026-04-15", "2026-04-01", "2026-05-01", "2026-04-10", "2026-05-15"]
            tl_matched = sum(1 for tl in expected_timelines if tl in combined_props)
            check("Notion entries include ALL 5 expected timeline dates",
                  tl_matched == 5, f"Only {tl_matched}/5 timeline dates found")

            # Verify exact per-entry mapping (Initiative -> Category -> Impact -> Timeline)
            # Per task.md: initiative title -> (category, impact, timeline)
            expected_entries = [
                ("tv", "pricing", "5000", "2026-04-15"),
                ("electronics promotion", "marketing", "3000", "2026-04-01"),
                ("audio", "expansion", "2500", "2026-05-01"),
                ("camera", "inventory", "1500", "2026-04-10"),
                ("home appliances", "cost", "2000", "2026-05-15"),
            ]
            for title_kw, cat_kw, impact_val, tl_val in expected_entries:
                matched_entry = False
                for entry_text in all_page_props_text:
                    if title_kw in entry_text and cat_kw in entry_text and impact_val in entry_text and tl_val in entry_text:
                        matched_entry = True
                        break
                check(f"Notion entry '{title_kw}' has matching category/impact/timeline",
                      matched_entry,
                      f"Expected {title_kw}+{cat_kw}+{impact_val}+{tl_val}")
    except Exception as e:
        check("Notion check", False, str(e))
    finally:
        cur.close()
        conn.close()


def check_emails():
    print("\n=== Check 3: Emails ===")
    conn = _db_conn()
    if conn is None:
        return
    cur = conn.cursor()
    try:
        # Check CEO email
        cur.execute("""
            SELECT subject, body_text, to_addr FROM email.messages
            WHERE to_addr::text LIKE '%%ceo@company.com%%'
            ORDER BY date DESC LIMIT 1
        """)
        ceo_email = cur.fetchone()
        check("CEO email exists", ceo_email is not None,
              "No email to ceo@company.com found")
        if ceo_email:
            subj = (ceo_email[0] or "").lower()
            body = (ceo_email[1] or "").lower()
            check("CEO email subject mentions strategy or summary",
                  "strategy" in subj or "summary" in subj or "executive" in subj,
                  f"Subject: {ceo_email[0]}")
            check("CEO email body mentions AMZN",
                  "amzn" in body or "amazon" in body, "No AMZN mention")
            check("CEO email body mentions Electronics",
                  "electronics" in body, "No Electronics mention")
            check("CEO email body mentions recommendation",
                  "recommend" in body or "pricing" in body or "independent" in body or "strategy" in body,
                  "No recommendation")

        # Check marketing team email
        cur.execute("""
            SELECT subject, body_text, to_addr FROM email.messages
            WHERE to_addr::text LIKE '%%marketing_team@company.com%%'
            ORDER BY date DESC LIMIT 1
        """)
        mkt_email = cur.fetchone()
        check("Marketing team email exists", mkt_email is not None,
              "No email to marketing_team@company.com found")
        if mkt_email:
            subj = (mkt_email[0] or "").lower()
            body = (mkt_email[1] or "").lower()
            check("Marketing email subject mentions category",
                  "category" in subj or "performance" in subj or "q2" in subj,
                  f"Subject: {mkt_email[0]}")
            check("Marketing email body mentions categories",
                  "electronics" in body and ("camera" in body or "watch" in body),
                  "Missing category details")
    except Exception as e:
        check("Email check", False, str(e))
    finally:
        cur.close()
        conn.close()


def check_reverse_validation():
    print("\n=== Reverse Validation ===")
    conn = _db_conn()
    if conn is None:
        return
    cur = conn.cursor()
    try:
        # Check Notion tracker has exactly 5 entries, no noise
        cur.execute("SELECT id, title FROM notion.databases")
        dbs = cur.fetchall()
        tracker_db = None
        for db_id, title in dbs:
            title_str = ""
            if isinstance(title, list):
                title_str = " ".join(
                    item.get("text", {}).get("content", "")
                    for item in title if isinstance(item, dict))
            elif isinstance(title, str):
                try:
                    parsed = json.loads(title)
                    if isinstance(parsed, list):
                        title_str = " ".join(
                            item.get("text", {}).get("content", "")
                            for item in parsed if isinstance(item, dict))
                    else:
                        title_str = str(title)
                except Exception:
                    title_str = str(title)
            else:
                title_str = str(title) if title else ""
            if "market" in title_str.lower() and "strategy" in title_str.lower() and "tracker" in title_str.lower():
                tracker_db = (db_id, title_str)
                break

        if tracker_db:
            cur.execute("""
                SELECT id, properties FROM notion.pages
                WHERE parent->>'database_id' = %s AND NOT archived
            """, (tracker_db[0],))
            pages = cur.fetchall()
            # At least the 5 required entries must exist (>= rather than == so
            # parallel multi-agent duplicate creation is not treated as noise).
            check("Notion tracker has at least the 5 required entries",
                  len(pages) >= 5,
                  f"Found {len(pages)} entries, expected at least 5")

            # Check no unrelated categories in notion entries
            all_props_text = " ".join(json.dumps(p[1]).lower() for p in pages if p[1])
            noise_categories = ["healthcare", "food", "real estate", "energy"]
            for cat in noise_categories:
                check(f"Notion entries do not contain noise category '{cat}'",
                      cat not in all_props_text,
                      f"Found '{cat}' in Notion data")

        # Check no emails sent to noise recipients
        noise_recipients = [
            "all-staff@company.com",
            "hr@company.com",
            "newsletter@company.com",
            "sales_team@company.com",
        ]
        for addr in noise_recipients:
            cur.execute(
                "SELECT COUNT(*) FROM email.messages WHERE to_addr::text ILIKE %s",
                (f"%{addr}%",),
            )
            cnt = cur.fetchone()[0]
            check(f"No email sent to noise recipient {addr}", cnt == 0,
                  f"Found {cnt} emails to {addr}")
    except Exception as e:
        check("Reverse validation", False, str(e))
    finally:
        cur.close()
        conn.close()


def check_scripts(workspace):
    print("\n=== Check 4: Scripts and JSON outputs ===")
    check("market_correlation.py exists",
          os.path.exists(os.path.join(workspace, "market_correlation.py")))
    check("category_analysis.py exists",
          os.path.exists(os.path.join(workspace, "category_analysis.py")))

    # Check market_correlation.json
    corr_path = os.path.join(workspace, "market_correlation.json")
    if os.path.exists(corr_path):
        check("market_correlation.json exists", True)
        try:
            with open(corr_path) as f:
                data = json.load(f)
            check("market_correlation.json has recommendation",
                  "recommendation" in data or "strategy" in str(data).lower(),
                  f"Keys: {list(data.keys()) if isinstance(data, dict) else 'not dict'}")
            data_str = json.dumps(data).lower()
            # Accept either the ticker (amzn/jpm/dji, as task.md specifies) or the
            # full company/index name (Amazon / JPMorgan / Dow Jones), so a correct
            # agent that keys its JSON by full name is not failed on naming.
            check("market_correlation.json mentions AMZN correlation",
                  "amzn" in data_str or "amazon" in data_str,
                  f"Content: {data_str[:200]}")
            check("market_correlation.json mentions JPM correlation",
                  "jpm" in data_str or "jpmorgan" in data_str,
                  f"Content: {data_str[:200]}")
            check("market_correlation.json mentions DJI correlation",
                  "dji" in data_str or "dow jones" in data_str,
                  f"Content: {data_str[:200]}")
            # Validate recommendation matches correlation rules
            import re
            reco_txt = str(data.get('recommendation', '')).lower() if isinstance(data, dict) else data_str
            has_valid_rule_usage = any(term in reco_txt for term in [
                'market-aligned', 'aligned', 'counter-cyclical', 'counter', 'independent'])
            check("market_correlation.json recommendation uses rule terminology",
                  has_valid_rule_usage, f"reco: {reco_txt[:150]}")
            # Recommendation consistency check: parse AMZN correlation and verify rule.
            # Robust parser: only trust a value whose key/path clearly identifies it
            # as a correlation/coefficient (corr|coefficient|coeff|r), OR the sole
            # AMZN leaf *whose magnitude is a plausible correlation* (|r|<=1). It
            # never blindly returns the first number inside an AMZN block (e.g.
            # avg_close), so a correctly reported correlation is not misread as a
            # different metric. Numeric strings ("-0.02") are accepted as well as
            # real numbers, because agents frequently format round(r, 2) into JSON.
            amzn_corr = None
            if isinstance(data, dict):
                corr_kw = ("corr", "coefficient", "coeff")
                corr_leaf = ("r", "r2", "r_squared", "corr", "correlation",
                             "coefficient", "coeff")

                def _is_corr_path(path):
                    pl = path.lower()
                    leaf = pl.split(".")[-1].split("[")[0].strip()
                    return any(w in pl for w in corr_kw) or leaf in corr_leaf

                def _leaf_num(v):
                    # correlation values can legitimately be stored as JSON
                    # strings (e.g. "-0.02" from an f-string), so coerce those too.
                    if isinstance(v, bool):
                        return None
                    if isinstance(v, (int, float)):
                        return float(v)
                    if isinstance(v, str):
                        try:
                            return float(v.strip())
                        except ValueError:
                            return None
                    return None

                def _amzn_leaves(obj, path=""):
                    found = []
                    if isinstance(obj, dict):
                        for k, v in obj.items():
                            found.extend(_amzn_leaves(v, f"{path}.{k}" if path else str(k)))
                    elif isinstance(obj, list):
                        for i, item in enumerate(obj):
                            found.extend(_amzn_leaves(item, f"{path}[{i}]"))
                    else:
                        fv = _leaf_num(obj)
                        if fv is not None:
                            pl = path.lower()
                            if "amzn" in pl or "amazon" in pl:
                                found.append((path, fv))
                    return found

                leaves = _amzn_leaves(data)
                if leaves:
                    corr = [v for p, v in leaves if _is_corr_path(p)]
                    if corr:
                        amzn_corr = corr[0]
                    elif len(leaves) == 1:
                        # Sole AMZN leaf is only a correlation when its magnitude is
                        # plausible (|r| <= 1). A lone avg_close=206.45 must NOT be
                        # treated as a correlation, otherwise a correct
                        # 'independent pricing strategy' recommendation gets failed.
                        single = leaves[0][1]
                        if abs(single) <= 1.01:
                            amzn_corr = single
            if amzn_corr is not None:
                abs_corr = abs(amzn_corr)
                # Per task.md rules:
                #  abs(AMZN) > 0.5 -> "market-aligned pricing"
                #  AMZN < -0.3 -> "counter-cyclical promotions"
                #  else -> "independent pricing strategy"
                if abs_corr > 0.5:
                    expected_rule = "aligned"
                elif amzn_corr < -0.3:
                    expected_rule = "counter"
                else:
                    expected_rule = "independent"
                check(f"Recommendation consistent with AMZN corr={amzn_corr:.3f} (expect '{expected_rule}')",
                      expected_rule in reco_txt,
                      f"reco: {reco_txt[:150]}")
        except Exception as e:
            check("market_correlation.json valid JSON", False, str(e))
    else:
        check("market_correlation.json exists", False, f"Not found at {corr_path}")

    # Check category_market_analysis.json
    cat_path = os.path.join(workspace, "category_market_analysis.json")
    if os.path.exists(cat_path):
        check("category_market_analysis.json exists", True)
        try:
            with open(cat_path) as f:
                data = json.load(f)
            data_str = json.dumps(data).lower()
            check("category_market_analysis.json has category data",
                  "electronics" in data_str, f"Content: {data_str[:200]}")
        except Exception as e:
            check("category_market_analysis.json valid JSON", False, str(e))
    else:
        check("category_market_analysis.json exists", False, f"Not found at {cat_path}")


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--agent_workspace", required=False, default=".")
    parser.add_argument("--groundtruth_workspace", required=False, default=".")
    parser.add_argument("--launch_time", required=False)
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()

    check_pptx(args.agent_workspace)
    check_notion()
    check_emails()
    check_scripts(args.agent_workspace)
    check_reverse_validation()

    total = PASS_COUNT + FAIL_COUNT
    if total == 0:
        print("\nFAIL: No checks performed.")
        sys.exit(1)

    accuracy = PASS_COUNT / total * 100
    print(f"\nOverall: {PASS_COUNT}/{total} checks passed ({accuracy:.1f}%)")

    result = {"total_passed": PASS_COUNT, "total_checks": total, "accuracy": accuracy}
    if args.res_log_file:
        with open(args.res_log_file, "w") as f:
            json.dump(result, f, indent=2)

    # Tightened: require all checks to pass. Strict per-entry notion + pct
    # checks are defanged by 70% threshold — bad cases produce 5-9 failures
    # and still pass at ~85-92% accuracy.
    if FAIL_COUNT == 0:
        print("PASS")
        sys.exit(0)
    else:
        print("FAIL")
        sys.exit(1)


if __name__ == "__main__":
    main()
