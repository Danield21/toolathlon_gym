"""Evaluation for yf-sector-analysis-ppt-email.

Blocking checks: Sector_Analysis.xlsx (Excel data) and Sector_Analysis.pptx (PPT structure).
Non-blocking: Email DB check.
"""
import argparse
import os
import sys
import openpyxl
from pptx import Presentation


TICKERS = ["AMZN", "GOOGL", "JNJ", "JPM", "XOM"]
SECTORS = {"AMZN": "Technology", "GOOGL": "Technology", "XOM": "Energy",
           "JNJ": "Healthcare", "JPM": "Financials"}


def num_close(a, b, tol=1.0):
    if a is None and b is None:
        return True
    if a is None or b is None:
        return False
    try:
        return abs(float(a) - float(b)) <= tol
    except (TypeError, ValueError):
        return str(a).strip().lower() == str(b).strip().lower()


def str_match(a, b):
    if a is None or b is None:
        return a is None and b is None
    return str(a).strip().lower() == str(b).strip().lower()


def load_sheet_rows(wb, sheet_name):
    for name in wb.sheetnames:
        if name.strip().lower() == sheet_name.strip().lower():
            return [[cell.value for cell in row] for row in wb[name].iter_rows()]
    return None


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--agent_workspace", required=False, default=".")
    parser.add_argument("--groundtruth_workspace", required=False)
    parser.add_argument("--launch_time", required=False)
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()

    task_root = os.path.abspath(os.path.join(os.path.dirname(__file__), ".."))
    gt_dir = args.groundtruth_workspace or os.path.join(task_root, "groundtruth_workspace")
    all_errors = []

    # ---- Check Excel ----
    agent_file = os.path.join(args.agent_workspace, "Sector_Analysis.xlsx")
    gt_file = os.path.join(gt_dir, "Sector_Analysis.xlsx")

    if not os.path.exists(agent_file):
        all_errors.append("Agent output Sector_Analysis.xlsx not found")
    elif not os.path.exists(gt_file):
        all_errors.append("Groundtruth Sector_Analysis.xlsx not found")
    else:
        agent_wb = openpyxl.load_workbook(agent_file, data_only=True)
        gt_wb = openpyxl.load_workbook(gt_file, data_only=True)

        # Check Stock Performance sheet
        print("  Checking Stock Performance...")
        a_rows = load_sheet_rows(agent_wb, "Stock Performance")
        g_rows = load_sheet_rows(gt_wb, "Stock Performance")
        if a_rows is None:
            all_errors.append("Sheet 'Stock Performance' not found in agent output")
        elif g_rows is None:
            all_errors.append("Sheet 'Stock Performance' not found in groundtruth")
        else:
            a_data = a_rows[1:] if len(a_rows) > 1 else []
            g_data = g_rows[1:] if len(g_rows) > 1 else []

            a_lookup = {}
            for row in a_data:
                if row and row[0] is not None:
                    a_lookup[str(row[0]).strip().upper()] = row
            for g_row in g_data:
                if not g_row or g_row[0] is None:
                    continue
                key = str(g_row[0]).strip().upper()
                a_row = a_lookup.get(key)
                if a_row is None:
                    all_errors.append(f"Missing stock: {key}")
                    continue
                # Col 1: Sector
                if len(a_row) > 1 and len(g_row) > 1:
                    if not str_match(a_row[1], g_row[1]):
                        all_errors.append(f"{key}.Sector: {a_row[1]} vs {g_row[1]}")
                # Col 2: Latest_Close - tighten from 5.0 to 1.0 (specific date close)
                if len(a_row) > 2 and len(g_row) > 2:
                    if not num_close(a_row[2], g_row[2], 1.0):
                        all_errors.append(f"{key}.Latest_Close: {a_row[2]} vs {g_row[2]} (tol=1.0)")
                # Col 3: YTD_Return_Pct - tighten from 3.0 to 0.5 percentage points
                if len(a_row) > 3 and len(g_row) > 3:
                    if not num_close(a_row[3], g_row[3], 0.5):
                        all_errors.append(f"{key}.YTD_Return_Pct: {a_row[3]} vs {g_row[3]} (tol=0.5)")
                # Col 4: One_Year_Return_Pct - tighten from 3.0 to 0.5
                if len(a_row) > 4 and len(g_row) > 4:
                    if not num_close(a_row[4], g_row[4], 0.5):
                        all_errors.append(f"{key}.One_Year_Return_Pct: {a_row[4]} vs {g_row[4]} (tol=0.5)")

            # Verify alphabetical sort by Symbol
            symbols_in_order = [str(r[0]).strip().upper() for r in a_data if r and r[0]]
            if symbols_in_order != sorted(symbols_in_order):
                all_errors.append(f"Stock Performance not sorted alphabetically by Symbol. Got {symbols_in_order}")
            if not all_errors:
                print("    PASS")

        # Check Sector Summary sheet
        print("  Checking Sector Summary...")
        a_rows = load_sheet_rows(agent_wb, "Sector Summary")
        g_rows = load_sheet_rows(gt_wb, "Sector Summary")
        prev_errors = len(all_errors)
        if a_rows is None:
            all_errors.append("Sheet 'Sector Summary' not found in agent output")
        elif g_rows is None:
            all_errors.append("Sheet 'Sector Summary' not found in groundtruth")
        else:
            a_data = a_rows[1:] if len(a_rows) > 1 else []
            g_data = g_rows[1:] if len(g_rows) > 1 else []

            a_lookup = {}
            for row in a_data:
                if row and row[0] is not None:
                    a_lookup[str(row[0]).strip().lower()] = row
            for g_row in g_data:
                if not g_row or g_row[0] is None:
                    continue
                key = str(g_row[0]).strip().lower()
                a_row = a_lookup.get(key)
                if a_row is None:
                    all_errors.append(f"Missing sector: {g_row[0]}")
                    continue
                # Col 1: Num_Stocks
                if len(a_row) > 1 and len(g_row) > 1:
                    if not num_close(a_row[1], g_row[1], 0):
                        all_errors.append(f"{key}.Num_Stocks: {a_row[1]} vs {g_row[1]}")
                # Col 2: Avg_YTD_Return_Pct - tighten 3.0 to 0.5
                if len(a_row) > 2 and len(g_row) > 2:
                    if not num_close(a_row[2], g_row[2], 0.5):
                        all_errors.append(f"{key}.Avg_YTD_Return: {a_row[2]} vs {g_row[2]} (tol=0.5)")
                # Col 3: Avg_One_Year_Return_Pct (was missing!)
                if len(a_row) > 3 and len(g_row) > 3:
                    if not num_close(a_row[3], g_row[3], 0.5):
                        all_errors.append(f"{key}.Avg_One_Year_Return: {a_row[3]} vs {g_row[3]} (tol=0.5)")

            # Verify alphabetical sort by Sector
            sectors_in_order = [str(r[0]).strip() for r in a_data if r and r[0]]
            if sectors_in_order != sorted(sectors_in_order, key=lambda s: s.lower()):
                all_errors.append(f"Sector Summary not sorted alphabetically. Got {sectors_in_order}")
            new_errors = len(all_errors) - prev_errors
            if new_errors == 0:
                print("    PASS")

    # ---- Check PowerPoint ----
    agent_ppt = os.path.join(args.agent_workspace, "Sector_Analysis.pptx")
    if not os.path.exists(agent_ppt):
        all_errors.append("Agent output Sector_Analysis.pptx not found")
    else:
        print("  Checking Sector_Analysis.pptx...")
        prs = Presentation(agent_ppt)
        slides = list(prs.slides)
        if len(slides) < 4:
            all_errors.append(f"PPT has {len(slides)} slides, expected at least 4")
        else:
            # Check title slide title + subtitle
            slide1_text = ""
            for shape in slides[0].shapes:
                if shape.has_text_frame:
                    slide1_text += shape.text_frame.text + "\n"
            sl1_lower = slide1_text.lower()
            if "sector performance analysis" not in sl1_lower:
                all_errors.append(f"Slide 1 title missing 'Sector Performance Analysis'. Found: {slide1_text[:200]}")
            if "market report" not in sl1_lower or "2026-03-06" not in sl1_lower:
                all_errors.append(f"Slide 1 subtitle missing 'Market Report - 2026-03-06'. Found: {slide1_text[:200]}")

            # Slide 2: Individual Stock Performance
            slide2_text = ""
            for shape in slides[1].shapes:
                if shape.has_text_frame:
                    slide2_text += shape.text_frame.text + "\n"
            sl2_lower = slide2_text.lower()
            if "individual stock performance" not in sl2_lower:
                all_errors.append(f"Slide 2 title not 'Individual Stock Performance'. Found: {slide2_text[:200]}")

            # Slide 3: Sector Comparison
            slide3_text = ""
            for shape in slides[2].shapes:
                if shape.has_text_frame:
                    slide3_text += shape.text_frame.text + "\n"
            sl3_lower = slide3_text.lower()
            if "sector comparison" not in sl3_lower:
                all_errors.append(f"Slide 3 title not 'Sector Comparison'. Found: {slide3_text[:200]}")

            # Slide 4: Key Findings
            slide4_text = ""
            for shape in slides[3].shapes:
                if shape.has_text_frame:
                    slide4_text += shape.text_frame.text + "\n"
            sl4_lower = slide4_text.lower()
            if "key findings" not in sl4_lower:
                all_errors.append(f"Slide 4 title not 'Key Findings'. Found: {slide4_text[:200]}")
            # Key Findings should mention best/worst sector and stock
            if not (("best" in sl4_lower and "worst" in sl4_lower)
                    or ("highest" in sl4_lower and "lowest" in sl4_lower)):
                all_errors.append(f"Slide 4 'Key Findings' must identify best AND worst. Found: {slide4_text[:300]}")

            # All PPT text contains all tickers
            all_ppt_text = ""
            for slide in slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        all_ppt_text += shape.text_frame.text.lower() + " "

            for ticker in TICKERS:
                if ticker.lower() not in all_ppt_text:
                    all_errors.append(f"PPT missing ticker: {ticker}")

        if all_errors:
            print(f"    FAIL: {[e for e in all_errors if 'PPT' in e or 'lide' in e or 'ticker' in e][:3]}")

    # ---- Email check (BLOCKING) ----
    print("  Checking Email...")
    try:
        import psycopg2
        conn = psycopg2.connect(host=os.environ.get("PGHOST", "localhost"), port=5432, dbname="toolathlon_gym",
                                user="eigent", password="camel")
        cur = conn.cursor()
        cur.execute("""
            SELECT subject, from_addr, to_addr, body_text FROM email.messages
            WHERE to_addr::text ILIKE '%%investments@company.com%%'
        """)
        emails = cur.fetchall()
        cur.close()
        conn.close()

        if not emails:
            all_errors.append("No email sent to investments@company.com")
        else:
            ok_subj = False
            ok_body = False
            for subject, from_addr, to_addr, body in emails:
                sl = (subject or "").strip().lower()
                bl = (body or "").lower()
                # Subject: exact 'Sector Analysis Report - 2026-03-06'
                if sl == "sector analysis report - 2026-03-06":
                    ok_subj = True
                # Body must mention best/worst sector + best/worst stock
                if (("best" in bl or "highest" in bl)
                        and ("worst" in bl or "lowest" in bl)
                        and any(t.lower() in bl for t in TICKERS)):
                    ok_body = True
            if not ok_subj:
                all_errors.append(
                    f"No email subject 'Sector Analysis Report - 2026-03-06'. "
                    f"Subjects: {[e[0] for e in emails]}"
                )
            if not ok_body:
                all_errors.append(
                    f"No email body mentions best/worst sector and best/worst stock. "
                    f"Bodies (first 200): {[(e[3] or '')[:200] for e in emails][:1]}"
                )
    except Exception as e:
        all_errors.append(f"Email check error: {e}")

    if all_errors:
        print(f"\n=== RESULT: FAIL ({len(all_errors)} errors) ===")
        for e in all_errors[:15]:
            print(f"  {e}")
        sys.exit(1)
    else:
        print("\n=== RESULT: PASS ===")
        sys.exit(0)


if __name__ == "__main__":
    main()
