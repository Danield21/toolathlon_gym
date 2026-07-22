"""Evaluation for yf-portfolio-stress-test."""
import argparse, os, sys


def num_close(a, b, abs_tol=1.0, rel_tol=0.05):
    try:
        a_f, b_f = float(a), float(b)
        return abs(a_f - b_f) <= max(abs_tol, abs(b_f) * rel_tol)
    except (TypeError, ValueError):
        return str(a).strip().lower() == str(b).strip().lower()


def load_sheet_rows(wb, sheet_name):
    for name in wb.sheetnames:
        if name.strip().lower() == sheet_name.strip().lower():
            return [[cell.value for cell in row] for row in wb[name].iter_rows()]
    return None


EXPECTED_SCENARIOS = {
    "market crash": {"equity": -30, "gold": 15, "loss": -27.75, "exceeds": "yes", "prob": 5},
    "rate hike": {"equity": -10, "gold": 3, "loss": -9.35, "exceeds": "no", "prob": 15},
    "recession": {"equity": -20, "gold": 20, "loss": -18, "exceeds": "no", "prob": 10},
    "stagflation": {"equity": -15, "gold": 25, "loss": -13, "exceeds": "no", "prob": 8},
}

EXPECTED_POSITIONS = {
    "GOOGL": 25,
    "AMZN": 25,
    "JPM": 20,
    "JNJ": 15,
    "XOM": 10,
    "GC=F": 5,
}


def check_excel(agent_workspace):
    errors = []
    import openpyxl
    path = os.path.join(agent_workspace, "Stress_Test.xlsx")
    if not os.path.exists(path):
        return ["Stress_Test.xlsx not found"]
    try:
        wb = openpyxl.load_workbook(path, data_only=True)

        # Scenario Analysis
        rows = load_sheet_rows(wb, "Scenario Analysis")
        if rows is None:
            errors.append("Sheet 'Scenario Analysis' not found")
        else:
            data_rows = [r for r in rows[1:] if r and r[0] is not None]
            if len(data_rows) != 4:
                errors.append(f"Scenario Analysis has {len(data_rows)} rows, expected 4")
            scenarios = {str(r[0]).strip().lower(): r for r in data_rows if r[0]}
            for sname, expected in EXPECTED_SCENARIOS.items():
                if sname not in scenarios:
                    errors.append(f"Scenario '{sname}' missing")
                    continue
                r = scenarios[sname]
                # Equity_Shock_Pct (col 1)
                if len(r) > 1 and not num_close(r[1], expected["equity"], abs_tol=0.5):
                    errors.append(f"{sname} Equity_Shock_Pct={r[1]}, expected {expected['equity']}")
                # Gold_Shock_Pct (col 2)
                if len(r) > 2 and not num_close(r[2], expected["gold"], abs_tol=0.5):
                    errors.append(f"{sname} Gold_Shock_Pct={r[2]}, expected {expected['gold']}")
                # Portfolio_Loss_Pct (col 3)
                if len(r) > 3 and not num_close(r[3], expected["loss"], abs_tol=0.5):
                    errors.append(f"{sname} Portfolio_Loss_Pct={r[3]}, expected {expected['loss']}")
                # Exceeds_Limit (col 4)
                if len(r) > 4 and str(r[4] or "").strip().lower() != expected["exceeds"]:
                    errors.append(f"{sname} Exceeds_Limit={r[4]}, expected {expected['exceeds'].title()}")
                # Probability_Pct (col 5)
                if len(r) > 5 and not num_close(r[5], expected["prob"], abs_tol=0.5):
                    errors.append(f"{sname} Probability_Pct={r[5]}, expected {expected['prob']}")

        # Position Impact - validate every (scenario, symbol) row
        rows2 = load_sheet_rows(wb, "Position Impact")
        if rows2 is None:
            errors.append("Sheet 'Position Impact' not found")
        else:
            data_rows2 = [r for r in rows2[1:] if r and r[0] is not None]
            if len(data_rows2) != 24:
                errors.append(f"Position Impact has {len(data_rows2)} rows, expected 24")
            # Build lookup by (scenario_lc, symbol_uc)
            pi_lookup = {}
            for r in data_rows2:
                if r and r[0] and r[1]:
                    k = (str(r[0]).strip().lower(), str(r[1]).strip().upper())
                    pi_lookup[k] = r
            for sname, sexp in EXPECTED_SCENARIOS.items():
                for sym, weight in EXPECTED_POSITIONS.items():
                    k = (sname, sym.upper())
                    r = pi_lookup.get(k)
                    if r is None:
                        errors.append(f"Position Impact missing ({sname}, {sym})")
                        continue
                    # Weight_Pct (col 2)
                    if len(r) > 2 and not num_close(r[2], weight, abs_tol=0.5):
                        errors.append(f"({sname},{sym}) Weight_Pct={r[2]}, expected {weight}")
                    # Shock_Pct (col 3): for equities use equity, for GC=F use gold
                    expected_shock = sexp["gold"] if sym == "GC=F" else sexp["equity"]
                    if len(r) > 3 and not num_close(r[3], expected_shock, abs_tol=0.5):
                        errors.append(f"({sname},{sym}) Shock_Pct={r[3]}, expected {expected_shock}")
                    # Position_Loss_Pct (col 4) = weight * shock / 100
                    expected_loss = round(weight * expected_shock / 100, 2)
                    if len(r) > 4 and not num_close(r[4], expected_loss, abs_tol=0.1):
                        errors.append(f"({sname},{sym}) Position_Loss_Pct={r[4]}, expected {expected_loss}")

        # Risk Summary
        rows3 = load_sheet_rows(wb, "Risk Summary")
        if rows3 is None:
            errors.append("Sheet 'Risk Summary' not found")
        else:
            data_rows3 = [r for r in rows3[1:] if r and r[0] is not None]
            lookup = {str(r[0]).strip().lower(): r[1] for r in data_rows3 if r[0]}
            if "worst_case_loss" in lookup:
                # Allow signed or unsigned (both 27.75 and -27.75)
                v = lookup["worst_case_loss"]
                if not (num_close(v, 27.75, abs_tol=0.5) or num_close(v, -27.75, abs_tol=0.5)):
                    errors.append(f"Worst_Case_Loss={v}, expected ~27.75 or ~-27.75")
            else:
                errors.append("Worst_Case_Loss not found")
            if "scenarios_exceeding_limit" in lookup:
                if not num_close(lookup["scenarios_exceeding_limit"], 1, abs_tol=0):
                    errors.append(f"Scenarios_Exceeding_Limit={lookup['scenarios_exceeding_limit']}, expected 1")
            if "expected_loss" in lookup:
                v = lookup["expected_loss"]
                if not (num_close(v, -5.63, abs_tol=0.3) or num_close(v, 5.63, abs_tol=0.3)):
                    errors.append(f"Expected_Loss={v}, expected ~-5.63")
            if "risk_rating" in lookup:
                if str(lookup["risk_rating"]).strip().lower() != "elevated":
                    errors.append(f"Risk_Rating={lookup['risk_rating']}, expected Elevated")

    except Exception as e:
        errors.append(f"Error reading Excel: {e}")
    return errors


def check_pptx(agent_workspace):
    errors = []
    path = os.path.join(agent_workspace, "Risk_Presentation.pptx")
    if not os.path.exists(path):
        return ["Risk_Presentation.pptx not found"]
    try:
        from pptx import Presentation
        prs = Presentation(path)
        if len(prs.slides) < 5:
            errors.append(f"Presentation has {len(prs.slides)} slides, expected at least 5")
            return errors

        # Per-slide text
        slide_texts = []
        for slide in prs.slides:
            stext = ""
            for shape in slide.shapes:
                if shape.has_text_frame:
                    stext += " " + shape.text_frame.text
            slide_texts.append(stext.lower())

        # Slide 1: Portfolio Stress Test Results title
        if not ("portfolio stress test" in slide_texts[0] or
                ("stress test" in slide_texts[0] and "portfolio" in slide_texts[0])):
            errors.append("Slide 1 missing 'Portfolio Stress Test Results' title")

        # Slide 2: portfolio composition / weights
        if not ("composition" in slide_texts[1] or "weight" in slide_texts[1] or
                any(t.lower() in slide_texts[1] for t in ["googl", "amzn", "jpm"])):
            errors.append("Slide 2 missing portfolio composition/weights content")

        # Slide 3: scenarios and impact
        scenarios_found = sum(1 for s in EXPECTED_SCENARIOS.keys() if s in slide_texts[2])
        if scenarios_found < 3:
            errors.append(f"Slide 3 missing scenarios (found {scenarios_found} of 4)")

        # Slide 4: worst case + risk policy / limit
        if not ("worst" in slide_texts[3] or "27.75" in slide_texts[3] or "market crash" in slide_texts[3]):
            errors.append("Slide 4 missing worst-case scenario content")
        if not ("limit" in slide_texts[3] or "policy" in slide_texts[3] or
                "exceed" in slide_texts[3] or "breach" in slide_texts[3]):
            errors.append("Slide 4 missing risk policy limit reference")

        # Slide 5: risk summary metrics + recommendation
        if not ("recommend" in slide_texts[4] or "elevated" in slide_texts[4]):
            errors.append("Slide 5 missing risk summary/recommendation")
    except Exception as e:
        errors.append(f"Error reading PPTX: {e}")
    return errors


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--agent_workspace", required=False)
    parser.add_argument("--groundtruth_workspace", required=False)
    parser.add_argument("--launch_time", required=False)
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()
    agent_ws = args.agent_workspace or os.path.join(os.path.dirname(__file__), "..", "groundtruth_workspace")

    all_errors = []

    print("  Checking Excel file...")
    errs = check_excel(agent_ws)
    if errs:
        all_errors.extend(errs)
        for e in errs[:5]:
            print(f"    ERROR: {e}")
    else:
        print("    PASS")

    print("  Checking PowerPoint...")
    errs = check_pptx(agent_ws)
    if errs:
        all_errors.extend(errs)
        for e in errs[:3]:
            print(f"    ERROR: {e}")
    else:
        print("    PASS")

    if all_errors:
        print(f"\n=== RESULT: FAIL ({len(all_errors)} errors) ===")
        for e in all_errors[:10]:
            print(f"  {e}")
        sys.exit(1)
    else:
        print("\n=== RESULT: PASS ===")
        sys.exit(0)


if __name__ == "__main__":
    main()
