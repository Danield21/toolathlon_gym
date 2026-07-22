"""Evaluation script for wc-yf-pricing-strategy-excel-ppt-email."""
import os
import argparse, json, os, sys
import openpyxl

def num_close(a, b, rel_tol=0.15, abs_tol=0.5):
    return abs(float(a) - float(b)) <= max(abs_tol, abs(float(b)) * rel_tol)


DB_CONFIG = {
    "host": os.environ.get("PGHOST", "localhost"), "port": 5432,
    "dbname": os.environ.get("PGDATABASE", "toolathlon_gym"),
    "user": "eigent", "password": "camel"
}

PASS_COUNT = 0
FAIL_COUNT = 0

def check(name, condition, detail=""):
    global PASS_COUNT, FAIL_COUNT
    if condition:
        PASS_COUNT += 1
        print(f"  [PASS] {name}")
    else:
        FAIL_COUNT += 1
        detail_str = str(detail)[:200] if detail else ""
        print(f"  [FAIL] {name}: {detail_str}")

def safe_float(val, default=None):
    try:
        if val is None: return default
        return float(str(val).replace(",", "").replace("%", "").replace("$", "").strip())
    except (ValueError, TypeError):
        return default

def get_conn():
    import psycopg2
    return psycopg2.connect(**DB_CONFIG)

def run_evaluation(agent_workspace, groundtruth_workspace, launch_time, res_log_file):
    global PASS_COUNT, FAIL_COUNT
    PASS_COUNT = 0
    FAIL_COUNT = 0

    # Check Dynamic_Pricing_Strategy.xlsx
    excel_path = os.path.join(agent_workspace, "Dynamic_Pricing_Strategy.xlsx")
    check("Dynamic_Pricing_Strategy.xlsx exists", os.path.exists(excel_path))
    if os.path.exists(excel_path):
        wb = openpyxl.load_workbook(excel_path)
        gt_path = os.path.join(groundtruth_workspace, "Dynamic_Pricing_Strategy.xlsx")
        gt_wb = openpyxl.load_workbook(gt_path) if os.path.exists(gt_path) else None

        if gt_wb:
            for sheet_name in gt_wb.sheetnames:
                check(f"{sheet_name} sheet exists", sheet_name in wb.sheetnames)
                if sheet_name in wb.sheetnames:
                    ws = wb[sheet_name]
                    gt_ws = gt_wb[sheet_name]
                    # Check headers
                    gt_headers = [str(c.value).strip().lower() if c.value else "" for c in gt_ws[1]]
                    headers = [str(c.value).strip().lower() if c.value else "" for c in ws[1]]
                    for h in gt_headers:
                        if h:
                            check(f"{sheet_name} has {h} column", h in headers, f"headers: {headers[:10]}")
                    # Check row count - exact match (filter blank trailing rows)
                    gt_rows = [r for r in gt_ws.iter_rows(min_row=2, values_only=True)
                               if any(c is not None for c in r)]
                    data_rows = [r for r in ws.iter_rows(min_row=2, values_only=True)
                                 if any(c is not None for c in r)]
                    # For Price_Recommendations and Product_Analysis, agent may
                    # legitimately recommend different products. Require >= GT
                    # row count instead of exact.
                    if sheet_name.lower() in ("price_recommendations", "product_analysis"):
                        check(f"{sheet_name} row count >= {len(gt_rows)}",
                              len(data_rows) >= len(gt_rows), f"got {len(data_rows)}")
                    else:
                        check(f"{sheet_name} row count == {len(gt_rows)}",
                              len(data_rows) == len(gt_rows), f"got {len(data_rows)}")

                    # Cell value comparison against groundtruth - ALL rows, not just first 3
                    header_map = {h: i for i, h in enumerate(headers)}
                    gt_header_map = {h: i for i, h in enumerate(gt_headers)}
                    # Build agent lookup by first column (typically Product_Name or Indicator/Metric)
                    if data_rows:
                        a_lookup = {}
                        for r in data_rows:
                            if r and r[0] is not None:
                                a_lookup[str(r[0]).strip().lower()] = r
                        for ri, gt_row in enumerate(gt_rows):
                            if not gt_row or gt_row[0] is None:
                                continue
                            gkey = str(gt_row[0]).strip().lower()
                            agent_row = a_lookup.get(gkey)
                            if agent_row is None:
                                # For Price_Recommendations / Product_Analysis, agent
                                # may pick different products. Don't fail just because
                                # a specific GT product wasn't selected by the agent.
                                if sheet_name.lower() in ("price_recommendations", "product_analysis"):
                                    continue
                                check(f"{sheet_name} row '{gt_row[0]}' present", False, f"row missing in agent output")
                                continue
                            for ci, gt_h in enumerate(gt_headers):
                                if not gt_h or ci >= len(gt_row):
                                    continue
                                gv = gt_row[ci]
                                agent_ci = header_map.get(gt_h)
                                if agent_ci is None or agent_ci >= len(agent_row):
                                    continue
                                av = agent_row[agent_ci]
                                gf = safe_float(gv)
                                af = safe_float(av)
                                if gf is not None and af is not None:
                                    # Per-column tolerance: counts exact, prices/values 5%, percentages 0.5
                                    gh_l = gt_h.lower()
                                    if "count" in gh_l or gh_l in ("total_sales", "stock_level"):
                                        tol = 1
                                    elif gh_l == "recommended_price" or gh_l == "change_pct":
                                        # Subjective recommendation; relax to 15% / 5.0 abs
                                        tol = max(5.0, abs(gf) * 0.15)
                                    elif "pct" in gh_l or "percentage" in gh_l:
                                        # Other percentages (Estimated_Revenue_Impact, Avg_Price_Change)
                                        # also depend on agent's pricing logic — relax to 2.0
                                        tol = 2.0
                                    elif "products_" in gh_l:
                                        # Integer counts of agent classification — relax to 5
                                        tol = 5
                                    else:
                                        tol = max(0.5, abs(gf) * 0.05)
                                    check(f"{sheet_name} '{gt_row[0]}' {gt_h}",
                                          abs(gf - af) <= tol, f"expected ~{gf}, got {af}")
                                elif gv is not None and av is not None:
                                    gs = str(gv).strip().lower()
                                    avs = str(av).strip().lower()
                                    if gs:
                                        # Subjective categorical fields: agent's classification
                                        # may legitimately differ from GT since rules aren't
                                        # explicitly specified. Accept any value in the valid set.
                                        gh_l = gt_h.lower()
                                        if gh_l == "price_elasticity_estimate":
                                            check(f"{sheet_name} '{gt_row[0]}' {gt_h} valid",
                                                  avs in ("high", "medium", "low"),
                                                  f"expected high|medium|low, got '{avs[:50]}'")
                                        elif gh_l == "trend":
                                            check(f"{sheet_name} '{gt_row[0]}' {gt_h} valid",
                                                  avs in ("up", "down", "flat"),
                                                  f"expected up|down|flat, got '{avs[:50]}'")
                                        elif gh_l == "impact_on_pricing":
                                            check(f"{sheet_name} '{gt_row[0]}' {gt_h} valid",
                                                  avs in ("direct", "indirect", "none"),
                                                  f"expected direct|indirect|none, got '{avs[:50]}'")
                                        elif gh_l == "rationale":
                                            # Non-empty subjective text
                                            check(f"{sheet_name} '{gt_row[0]}' {gt_h} non-empty",
                                                  len(avs) > 0,
                                                  "got empty/missing")
                                        else:
                                            check(f"{sheet_name} '{gt_row[0]}' {gt_h} text",
                                                  gs == avs,  # exact match for non-subjective
                                                  f"expected '{gs[:50]}', got '{avs[:50]}'")

    # Check Pricing_Strategy_Deck.pptx
    pptx_path = os.path.join(agent_workspace, "Pricing_Strategy_Deck.pptx")
    check("Pricing_Strategy_Deck.pptx exists", os.path.exists(pptx_path))
    if os.path.exists(pptx_path):
        from pptx import Presentation
        prs = Presentation(pptx_path)
        gt_pptx_path = os.path.join(groundtruth_workspace, "Pricing_Strategy_Deck.pptx")
        if os.path.exists(gt_pptx_path):
            gt_prs = Presentation(gt_pptx_path)
            gt_slide_count = len(gt_prs.slides)
            check(f"Pricing_Strategy_Deck.pptx has {gt_slide_count} slides", len(prs.slides) == gt_slide_count, f"got {len(prs.slides)} slides")
            # Compare slide titles
            gt_titles = []
            for s in gt_prs.slides:
                if s.shapes.title:
                    gt_titles.append(s.shapes.title.text.strip().lower())
            agent_titles = []
            for s in prs.slides:
                if s.shapes.title:
                    agent_titles.append(s.shapes.title.text.strip().lower())
            def _norm_title(t):
                return (t or "").strip().rstrip(":.- ").lower()
            agent_titles_norm = [_norm_title(t) for t in agent_titles]
            for gt in gt_titles:
                gt_n = _norm_title(gt)
                # Require either exact normalized equality OR strong overlap
                # (GT title appears as a leading substring of agent title, or
                # agent title is a leading substring of GT, with shared >= 12
                # chars to avoid trivial matches like 'introduction').
                def _match(at_n):
                    if not at_n:
                        return False
                    if at_n == gt_n:
                        return True
                    if len(gt_n) >= 12 and at_n.startswith(gt_n):
                        return True
                    if len(at_n) >= 12 and gt_n.startswith(at_n):
                        return True
                    return False
                found = any(_match(at) for at in agent_titles_norm)
                check(f"Pricing_Strategy_Deck.pptx has slide \"{gt[:40]}\"",
                      found, f"agent titles: {agent_titles[:5]}")
        else:
            check("Pricing_Strategy_Deck.pptx has >= 3 slides", len(prs.slides) >= 3, f"got {len(prs.slides)} slides")

    # Check Python script exists (terminal usage) - require specific filename per task.md
    py_files = [f for f in os.listdir(agent_workspace) if f.endswith(".py")]
    check("Python script 'pricing_strategist.py' exists",
          "pricing_strategist.py" in py_files,
          f"found: {py_files}")

    # Database checks - email check is BLOCKING (required deliverable)
    try:
        conn = get_conn()
        cur = conn.cursor()
        # Try Sent folder first; fall back to any folder with subject pattern
        cur.execute("""
            SELECT subject, to_addr, body_text FROM email.messages
            WHERE subject ILIKE '%dynamic pricing strategy%'
               OR subject ILIKE '%pricing strategy proposal%'
            ORDER BY id DESC
        """)
        rows = cur.fetchall()
        email_found = len(rows) > 0
        check("Email 'Dynamic Pricing Strategy Proposal' sent", email_found, f"no matching email found")
        if email_found:
            r = rows[0]
            subj, to_addr, body = r[0], r[1], r[2] if len(r) > 2 else ""
            # Validate recipient contains pricing-committee@company.com
            recipient_ok = to_addr is not None and "pricing-committee@company.com" in str(to_addr).lower()
            check("Email recipient is pricing-committee@company.com", recipient_ok, f"to_addr: {to_addr}")
            # Validate body mentions revenue impact: require BOTH a revenue/impact term
            # AND a numeric/percentage indicator. Also require body length > 80 chars
            # to reject trivial 'thanks' bodies.
            body_lower = str(body or "").lower()
            has_topic = any(kw in body_lower for kw in ("revenue", "impact"))
            has_quant = ("%" in body_lower) or ("percent" in body_lower) or any(
                c.isdigit() for c in body_lower
            )
            has_pricing = any(
                kw in body_lower for kw in ("pricing", "price", "recommend")
            )
            body_ok = has_topic and has_quant and has_pricing and len(body_lower.strip()) >= 80
            check(
                "Email body discusses revenue impact with quantitative pricing detail",
                body_ok,
                f"body excerpt ({len(body_lower)} chars): {body_lower[:150]}"
            )
        conn.close()
    except Exception as e:
        check("DB email check", False, str(e))

    return FAIL_COUNT == 0, f"Passed {PASS_COUNT}/{PASS_COUNT + FAIL_COUNT} checks"

def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--agent_workspace", required=False, default=".")
    parser.add_argument("--groundtruth_workspace", required=False, default=".")
    parser.add_argument("--launch_time", required=False, default="2026-03-07 10:00:00")
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()

    success, message = run_evaluation(
        args.agent_workspace, args.groundtruth_workspace,
        args.launch_time, args.res_log_file
    )
    print(message)
    sys.exit(0 if success else 1)

if __name__ == "__main__":
    main()