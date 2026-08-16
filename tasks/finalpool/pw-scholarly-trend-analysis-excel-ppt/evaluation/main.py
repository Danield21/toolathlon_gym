"""Evaluation script for pw-scholarly-trend-analysis-excel-ppt."""
import os
import argparse, json, os, sys
import openpyxl

# --- verify_v2 smart primitives ---
_EVAL_DIR = os.path.dirname(os.path.abspath(__file__))
_PACK_ROOT = os.path.abspath(os.path.join(_EVAL_DIR, "..", "..", "..", ".."))
if _PACK_ROOT not in sys.path:
    sys.path.insert(0, _PACK_ROOT)
try:
    from utils.verify_v2 import smart_column_exists
    from utils.verify_v2.eval_helpers import get_sheet_rows_as_dicts, get_gt_column_values
    _HAS_VERIFY_V2 = True
except Exception:
    _HAS_VERIFY_V2 = False

TASK_NAME = "pw-scholarly-trend-analysis-excel-ppt"


DB_CONFIG = {
    "host": os.environ.get("PGHOST", "localhost"), "port": int(os.environ.get("PGPORT", "5432")),
    "dbname": os.environ.get("PGDATABASE", "toolathlon_gym"),
    "user": "eigent", "password": "camel"
}

PASS_COUNT = 0
FAIL_COUNT = 0

def check(name, condition, detail=""):
    global PASS_COUNT, FAIL_COUNT
    if condition:
        PASS_COUNT += 1
        print(f"  [PASS] {name}")
    else:
        FAIL_COUNT += 1
        detail_str = str(detail)[:200] if detail else ""
        print(f"  [FAIL] {name}: {detail_str}")

def safe_float(val, default=None):
    try:
        if val is None:
            return default
        return float(str(val).replace(',', '').replace('%', '').replace('$', '').strip())
    except (ValueError, TypeError):
        return default

def get_conn():
    import psycopg2
    return psycopg2.connect(**DB_CONFIG)

def run_evaluation(agent_workspace, groundtruth_workspace, launch_time, res_log_file):
    global PASS_COUNT, FAIL_COUNT
    PASS_COUNT = 0
    FAIL_COUNT = 0

    
    excel_path = os.path.join(agent_workspace, "Trend_Analysis_Report.xlsx")
    check("Trend_Analysis_Report.xlsx exists", os.path.exists(excel_path))
    if os.path.exists(excel_path):
        wb = openpyxl.load_workbook(excel_path)
        gt_path = os.path.join(groundtruth_workspace, "Trend_Analysis_Report.xlsx")
        gt_wb = openpyxl.load_workbook(gt_path) if os.path.exists(gt_path) else None

        def check_columns(sheet_name, expected_cols, min_rows):
            """Verify sheet exists, >= min_rows, contains required columns.
            Uses LLM semantic mapping via verify_v2, falls back to strict
            header match."""
            check(f"{sheet_name} sheet exists", sheet_name in wb.sheetnames)
            if sheet_name not in wb.sheetnames:
                return
            _ws = wb[sheet_name]
            _data_rows = list(_ws.iter_rows(min_row=2, values_only=True))
            check(f"{sheet_name} has >= {min_rows} rows",
                  len(_data_rows) >= min_rows, f"got {len(_data_rows)}")
            if _HAS_VERIFY_V2 and gt_wb is not None:
                _raw_headers, _agent_rows = get_sheet_rows_as_dicts(wb, sheet_name)
                for _exp in expected_cols:
                    _gt_vals = get_gt_column_values(gt_wb, sheet_name, _exp)
                    _ok, _matched, _reason = smart_column_exists(
                        expected_col=_exp, agent_headers=_raw_headers,
                        gt_samples=_gt_vals[:3], agent_rows=_agent_rows,
                        task_name=TASK_NAME,
                    )
                    _detail = _reason
                    if _ok and _matched and _matched.lower() != _exp.lower():
                        _detail = f"LLM-mapped to {_matched!r}"
                    check(f"{sheet_name} has {_exp} column", _ok, _detail)
            else:
                _headers = [str(c.value).strip().lower() if c.value else "" for c in _ws[1]]
                for _exp in expected_cols:
                    check(f"{sheet_name} has {_exp} column",
                          _exp.lower() in _headers, f"headers: {_headers[:8]}")

        check("Data_Analysis sheet exists", "Data_Analysis" in wb.sheetnames)
        if "Data_Analysis" in wb.sheetnames:
            ws = wb["Data_Analysis"]
            data_rows = list(ws.iter_rows(min_row=2, values_only=True))
            check("Data_Analysis has >= 4 rows", len(data_rows) >= 4, f"got {len(data_rows)}")

            check_columns('Data_Analysis', ['Paper_ID', 'Citations', 'Relevance_Score'], 4)

            # Paper_ID overlap check: agent rows must mention at least 2 of the 4
            # GT papers (preprocess seeds these IDs into scholarly.arxiv_papers).
            expected_ids = {'2301.01234', '2302.05678', '2303.09012', '2304.03456'}
            seen_ids = set()
            for row in data_rows:
                for cell in row:
                    if cell is None:
                        continue
                    s = str(cell).strip()
                    for pid in expected_ids:
                        if pid in s:
                            seen_ids.add(pid)
            check("Data_Analysis references at least 2 of the seeded arxiv paper IDs",
                  len(seen_ids) >= 2,
                  f"seen_ids={seen_ids}, expected any of {expected_ids}")
        check("Metrics sheet exists", "Metrics" in wb.sheetnames)
        if "Metrics" in wb.sheetnames:
            ws = wb["Metrics"]
            data_rows = list(ws.iter_rows(min_row=2, values_only=True))
            check("Metrics has >= 4 rows", len(data_rows) >= 4, f"got {len(data_rows)}")

            check_columns('Metrics', ['Metric', 'Value'], 4)
        check("Recommendations sheet exists", "Recommendations" in wb.sheetnames)
        if "Recommendations" in wb.sheetnames:
            ws = wb["Recommendations"]
            data_rows = list(ws.iter_rows(min_row=2, values_only=True))
            check("Recommendations has >= 2 rows", len(data_rows) >= 2, f"got {len(data_rows)}")

            check_columns('Recommendations', ['Priority', 'Action'], 2)
    # PPT and script checks should run regardless of Excel state.
    # Per task.md: 'Trend_Analysis_Presentation.pptx with at least 4 slides'
    pptx_path = os.path.join(agent_workspace,
                             "Trend_Analysis_Presentation.pptx")
    if not os.path.exists(pptx_path):
        import glob as globmod
        pptx_files = globmod.glob(os.path.join(agent_workspace, "*.pptx"))
        if pptx_files:
            pptx_path = pptx_files[0]
    check("Trend_Analysis_Presentation.pptx (or any *.pptx) exists",
          os.path.exists(pptx_path),
          f"checked {pptx_path}")
    if os.path.exists(pptx_path):
        try:
            from pptx import Presentation
            prs = Presentation(pptx_path)
            slide_count = len(prs.slides)
            check("PPT has at least 4 slides",
                  slide_count >= 4,
                  f"got {slide_count} slides")
            # Aggregate all slide text and check 'finding' keyword (task says
            # 'covering the key findings')
            all_text = ""
            for s in prs.slides:
                for sh in s.shapes:
                    if sh.has_text_frame:
                        for p in sh.text_frame.paragraphs:
                            all_text += " " + p.text
            all_text = all_text.lower()
            check("PPT mentions findings/trend keywords",
                  any(k in all_text for k in ["finding", "trend", "key", "result"]),
                  f"text len {len(all_text)}")
        except Exception as e:
            check("PPT readable", False, str(e))

    proc_path = os.path.join(agent_workspace, "scholarly_trend_processor.py")
    proc_exists = os.path.exists(proc_path)
    check("scholarly_trend_processor.py exists", proc_exists)
    if proc_exists:
        try:
            proc_size = os.path.getsize(proc_path)
            check("scholarly_trend_processor.py is non-trivial (>=200 bytes)",
                  proc_size >= 200, f"size={proc_size} bytes")
        except OSError as e:
            check("scholarly_trend_processor.py size check", False, str(e))

    return FAIL_COUNT == 0, f"Passed {PASS_COUNT}/{PASS_COUNT + FAIL_COUNT} checks"

def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--agent_workspace", required=False, default=".")
    parser.add_argument("--groundtruth_workspace", required=False, default=".")
    parser.add_argument("--launch_time", required=False, default="2026-03-07 10:00:00")
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()

    success, message = run_evaluation(
        args.agent_workspace, args.groundtruth_workspace,
        args.launch_time, args.res_log_file
    )
    print(message)
    sys.exit(0 if success else 1)

if __name__ == "__main__":
    main()
