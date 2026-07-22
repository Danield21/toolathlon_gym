#!/usr/bin/env python3
"""Evaluation script for sales-territory-performance-audit task validation"""

from argparse import ArgumentParser
import json
import os
import sys

PASS_COUNT = 0
FAIL_COUNT = 0


def record(name, passed, detail=""):
    global PASS_COUNT, FAIL_COUNT
    if passed:
        PASS_COUNT += 1
        print(f"  [PASS] {name}")
    else:
        FAIL_COUNT += 1
        msg = f": {detail[:300]}" if detail else ""
        print(f"  [FAIL] {name}{msg}")


def num_close(a, b, tol=1.0):
    try: return abs(float(a) - float(b)) <= tol
    except (TypeError, ValueError): return False


def str_match(a, b):
    if a is None or b is None: return a is None and b is None
    return str(a).strip().lower() == str(b).strip().lower()


def _find_xlsx_by_keywords(workspace, keywords, exclude=None):
    """Score xlsx files by keyword matches: filename match weighs 10x more than content.
    Returns sorted list of (path) by descending score, ties broken alphabetically."""
    import glob
    import openpyxl
    exclude = exclude or set()
    scored = []
    for path in sorted(glob.glob(os.path.join(workspace, "*.xlsx"))):
        if path in exclude:
            continue
        fname_low = os.path.basename(path).lower()
        if fname_low.startswith("~$") or fname_low.startswith("."):
            continue
        score = 0
        # Filename match (each keyword found in filename adds 10)
        for kw in keywords:
            if kw in fname_low:
                score += 10
        # Content match (each keyword found in first 5 rows + sheetnames adds 1)
        try:
            wb = openpyxl.load_workbook(path, data_only=True, read_only=True)
            content_text = " ".join(s.lower() for s in wb.sheetnames)
            for ws in wb.worksheets:
                row_count = 0
                for row in ws.iter_rows(values_only=True):
                    if row_count > 5:
                        break
                    row_count += 1
                    for cell in row:
                        if cell is not None:
                            content_text += " " + str(cell).lower()
            wb.close()
            for kw in keywords:
                if kw in content_text:
                    score += 1
        except Exception:
            pass
        if score > 0:
            scored.append((score, path))
    scored.sort(key=lambda x: (-x[0], x[1]))
    return [p for _, p in scored]


def _compare_xlsx_to_gt(agent_xlsx_path, gt_xlsx_path, label):
    import openpyxl
    fname = os.path.basename(agent_xlsx_path)
    try:
        wb = openpyxl.load_workbook(agent_xlsx_path, data_only=True)
    except Exception as e:
        record(f"xlsx {fname} readable", False, str(e))
        return

    gt_wb = openpyxl.load_workbook(gt_xlsx_path, data_only=True)
    for gt_sname in gt_wb.sheetnames:
        gt_ws = gt_wb[gt_sname]
        a_ws = None
        for asn in wb.sheetnames:
            if asn.strip().lower() == gt_sname.strip().lower():
                a_ws = wb[asn]; break
        if a_ws is None and len(wb.sheetnames) == 1 and len(gt_wb.sheetnames) == 1:
            a_ws = wb[wb.sheetnames[0]]
        if a_ws is None:
            record(f"GT {label} sheet '{gt_sname}' exists in {fname}", False,
                   f"Available: {wb.sheetnames}")
            continue
        gt_rows = [r for r in gt_ws.iter_rows(min_row=2, values_only=True) if any(c is not None for c in r)]
        a_rows = [r for r in a_ws.iter_rows(min_row=2, values_only=True) if any(c is not None for c in r)]
        # Tolerant row count: allow ±25% (min ±2) — territory groupings may vary.
        row_tol = max(2, int(len(gt_rows) * 0.25))
        row_diff = abs(len(a_rows) - len(gt_rows))
        record(f"GT {label} '{gt_sname}' row count (within tolerance)",
               row_diff <= row_tol and len(a_rows) >= max(2, len(gt_rows) - row_tol),
               f"Expected ~{len(gt_rows)} (±{row_tol}), got {len(a_rows)}")
        for ri in range(min(3, len(gt_rows))):
            if ri >= len(a_rows): break
            ok = True
            for ci in range(min(len(gt_rows[ri]), len(a_rows[ri]) if ri < len(a_rows) else 0)):
                gv, av = gt_rows[ri][ci], a_rows[ri][ci]
                if gv is None: continue
                if isinstance(gv, (int, float)):
                    if not num_close(av, gv, max(abs(gv)*0.1, 1.0)): ok = False; break
                else:
                    if not str_match(av, gv): ok = False; break
            record(f"GT {label} '{gt_sname}' row {ri+1} values", ok,
                   f"gt={gt_rows[ri][:4]}, agent={a_rows[ri][:4] if ri < len(a_rows) else 'missing'}")
    gt_wb.close()
    wb.close()


def check_xlsx_content(workspace, groundtruth_workspace="."):
    """Locate two xlsx deliverables (performance scorecard + coaching plans)
    by filename or content keyword, compare each to GT."""
    print("\n=== Check: XLSX files ===")
    import openpyxl

    targets = [
        ("performance_scorecard.xlsx", "performance scorecard xlsx",
         ["performance_scorecard", "performance-scorecard", "performance scorecard",
          "performance", "scorecard", "metric"]),
        ("coaching_plans.xlsx", "coaching plans xlsx",
         ["coaching_plans", "coaching-plans", "coaching plans",
          "coaching", "development", "training"]),
    ]

    used_paths = set()
    for gt_fname, label, keywords in targets:
        gt_path = os.path.join(groundtruth_workspace, gt_fname)
        if not os.path.isfile(gt_path):
            continue
        exact = os.path.join(workspace, gt_fname)
        chosen = None
        if os.path.isfile(exact) and exact not in used_paths:
            chosen = exact
        else:
            cands = _find_xlsx_by_keywords(workspace, keywords, exclude=used_paths)
            if cands: chosen = cands[0]
        if chosen is None:
            record(f"xlsx for {label} exists", False,
                   f"No xlsx with keywords {keywords[:3]} found")
            continue
        used_paths.add(chosen)
        record(f"xlsx for {label} exists ({os.path.basename(chosen)})", True)
        try:
            wb = openpyxl.load_workbook(chosen, data_only=True)
            for ws in wb.worksheets:
                rows = list(ws.iter_rows(values_only=True))
                record(f"xlsx {os.path.basename(chosen)} '{ws.title}' has data",
                       len(rows) >= 2, f"{len(rows)} rows")
            wb.close()
        except Exception as e:
            record(f"xlsx {os.path.basename(chosen)} readable", False, str(e))
            continue
        _compare_xlsx_to_gt(chosen, gt_path, label)


def _find_docx_by_keywords(workspace, keywords, exclude=None):
    import glob
    from docx import Document
    exclude = exclude or set()
    candidates = []
    for path in sorted(glob.glob(os.path.join(workspace, "*.docx"))):
        if path in exclude:
            continue
        fname_low = os.path.basename(path).lower()
        if fname_low.startswith("~$") or fname_low.startswith("."):
            continue
        if any(kw in fname_low for kw in keywords):
            candidates.append(path)
            continue
        try:
            doc = Document(path)
            text_low = "\n".join(p.text for p in doc.paragraphs).lower()
            if any(kw in text_low for kw in keywords):
                candidates.append(path)
        except Exception:
            continue
    return candidates


def check_docx_content(workspace):
    """Locate coaching analysis docx by filename or content keyword."""
    print("\n=== Check: DOCX files ===")
    from docx import Document
    keywords = ["coaching", "territory", "performance", "analysis", "audit"]
    exact = os.path.join(workspace, "coaching_analysis.docx")
    chosen = exact if os.path.isfile(exact) else None
    if chosen is None:
        cands = _find_docx_by_keywords(workspace, keywords)
        if cands: chosen = cands[0]
    if chosen is None:
        record("docx for coaching analysis exists", False,
               f"No docx with keywords {keywords} found")
        return False
    record(f"docx for coaching analysis exists ({os.path.basename(chosen)})", True)
    try:
        doc = Document(chosen)
        text_low = " ".join(p.text for p in doc.paragraphs).lower()
        record(f"docx {os.path.basename(chosen)} substantive content (>=400 chars)",
               len(text_low) >= 400, f"text len: {len(text_low)}")
        for term in ["territory", "coaching", "performance"]:
            record(f"docx mentions '{term}'", term in text_low, "missing")
    except Exception as e:
        record(f"docx {os.path.basename(chosen)} readable", False, str(e))
    return True


def _find_pptx_by_keywords(workspace, keywords):
    import glob
    from pptx import Presentation
    candidates = []
    for path in sorted(glob.glob(os.path.join(workspace, "*.pptx"))):
        fname_low = os.path.basename(path).lower()
        if fname_low.startswith("~$") or fname_low.startswith("."):
            continue
        if any(kw in fname_low for kw in keywords):
            candidates.append(path)
            continue
        try:
            prs = Presentation(path)
            text_low = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text_low += " " + shape.text.lower()
            if any(kw in text_low for kw in keywords):
                candidates.append(path)
        except Exception:
            continue
    return candidates


def check_pptx_content(workspace):
    """Locate leadership presentation pptx by filename or content keyword."""
    print("\n=== Check: PPTX leadership presentation ===")
    from pptx import Presentation
    keywords = ["leadership", "presentation", "executive", "territory", "performance"]
    exact = os.path.join(workspace, "leadership_presentation.pptx")
    chosen = exact if os.path.isfile(exact) else None
    if chosen is None:
        cands = _find_pptx_by_keywords(workspace, keywords)
        if cands: chosen = cands[0]
    if chosen is None:
        record("pptx for leadership presentation exists", False,
               f"No pptx with keywords {keywords} found")
        return False
    record(f"pptx for leadership presentation exists ({os.path.basename(chosen)})", True)
    try:
        prs = Presentation(chosen)
        record(f"pptx {os.path.basename(chosen)} has >= 4 slides",
               len(prs.slides) >= 4, f"{len(prs.slides)} slides")
        all_text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    all_text += " " + shape.text
        all_text_low = all_text.lower()
        for term in ["territory", "performance"]:
            record(f"pptx mentions '{term}'", term in all_text_low, "missing")
    except Exception as e:
        record(f"pptx {os.path.basename(chosen)} readable", False, str(e))
    return True


def check_email_phase6():
    print("\n=== Check: Email to senior management (phase 6) ===")
    try:
        import psycopg2
        DB = {
            "host": os.environ.get("PGHOST", "localhost"),
            "port": int(os.environ.get("PGPORT", "5432")),
            "dbname": os.environ.get("PGDATABASE", "toolathlon_gym"),
            "user": "eigent", "password": "camel",
        }
        conn = psycopg2.connect(**DB); cur = conn.cursor()
        cur.execute("SELECT subject, to_addr, COALESCE(body_text, body_html, '') FROM email.messages")
        rows = cur.fetchall(); cur.close(); conn.close()
    except Exception as e:
        record("email DB reachable", False, str(e)); return
    matched = None
    for subj, to_addr, body in rows:
        text = ((subj or "") + " " + (body or "")).lower()
        if any(k in text for k in ["territory audit", "sales territory", "performance audit",
                                   "territory realignment", "audit report"]):
            matched = (subj, to_addr); break
    record("Email to senior management about territory/audit exists",
           matched is not None, f"checked {len(rows)} emails")


def check_gsheet_dashboard():
    """Phase 4 requires designing performance dashboards. The task_config
    includes google_sheet, so verify a territory-related cloud spreadsheet
    has been created."""
    print("\n=== Check: Google Sheet territory dashboard (phase 4) ===")
    try:
        import psycopg2
        DB = {
            "host": os.environ.get("PGHOST", "localhost"),
            "port": int(os.environ.get("PGPORT", "5432")),
            "dbname": os.environ.get("PGDATABASE", "toolathlon_gym"),
            "user": "eigent", "password": "camel",
        }
        conn = psycopg2.connect(**DB); cur = conn.cursor()
        cur.execute("SELECT id, title FROM gsheet.spreadsheets")
        sheets = cur.fetchall(); cur.close(); conn.close()
    except Exception as e:
        record("gsheet DB reachable", False, str(e))
        return
    matched = []
    for sid, title in sheets:
        title_low = (title or "").lower()
        if any(k in title_low for k in
               ["territory", "sales performance", "sales dashboard",
                "performance dashboard", "sales territory"]):
            matched.append((sid, title))
    record("Cloud spreadsheet territory/performance dashboard exists",
           len(matched) >= 1,
           f"checked {len(sheets)} sheets, matched={[m[1] for m in matched][:3]}")


def check_calendar_phase6():
    print("\n=== Check: Strategic review meeting (phase 6) ===")
    try:
        import psycopg2
        DB = {
            "host": os.environ.get("PGHOST", "localhost"),
            "port": int(os.environ.get("PGPORT", "5432")),
            "dbname": os.environ.get("PGDATABASE", "toolathlon_gym"),
            "user": "eigent", "password": "camel",
        }
        conn = psycopg2.connect(**DB); cur = conn.cursor()
        cur.execute("SELECT id, summary, COALESCE(description, '') FROM gcal.events")
        rows = cur.fetchall(); cur.close(); conn.close()
    except Exception as e:
        record("calendar DB reachable", False, str(e)); return
    matched = None
    for _id, summary, desc in rows:
        text = ((summary or "") + " " + (desc or "")).lower()
        if any(k in text for k in ["strategic review", "territory review",
                                   "sales territory", "audit review"]):
            matched = (summary, desc); break
    record("Calendar event scheduled for strategic review meeting",
           matched is not None, f"checked {len(rows)} events")


def main():
    parser = ArgumentParser()
    parser.add_argument("--agent_workspace", required=False, default=".")
    parser.add_argument("--groundtruth_workspace", required=False, default=".")
    parser.add_argument("--res_log_file", required=False)
    parser.add_argument("--launch_time", required=False)
    args = parser.parse_args()

    ws = args.agent_workspace
    if not os.path.isdir(ws):
        print(f"Agent workspace not found: {ws}")
        sys.exit(1)

    # GT self-test toleration: when agent==groundtruth (V1 GT-only smoke
    # test), the email/gcal checks would fail because no email/event was
    # actually sent.
    is_gt_self_test = (
        args.agent_workspace and args.groundtruth_workspace
        and os.path.abspath(args.agent_workspace) == os.path.abspath(args.groundtruth_workspace)
    )

    check_xlsx_content(ws, args.groundtruth_workspace)
    check_docx_content(ws)
    check_pptx_content(ws)
    if not is_gt_self_test:
        check_email_phase6()
        check_calendar_phase6()
        check_gsheet_dashboard()
    else:
        print("  [SKIP] email/calendar/gsheet checks (GT self-test mode)")

    total = PASS_COUNT + FAIL_COUNT
    if total == 0:
        print("\nFAIL: No checks were performed.")
        sys.exit(1)

    print(f"\nOverall: {PASS_COUNT}/{total} checks passed; FAIL_COUNT={FAIL_COUNT}")

    result = {
        "total_passed": PASS_COUNT,
        "total_checks": total,
        "fail_count": FAIL_COUNT,
    }

    if args.res_log_file:
        with open(args.res_log_file, "w") as f:
            json.dump(result, f, indent=2)

    if FAIL_COUNT == 0:
        print("PASS")
        sys.exit(0)
    else:
        print("FAIL")
        sys.exit(1)


if __name__ == "__main__":
    main()
