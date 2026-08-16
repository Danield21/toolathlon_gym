"""
Evaluation script for canvas-assignment-effectiveness-ppt-notion-email task.

Checks:
1. Excel file Assessment_Effectiveness.xlsx with 3 sheets and correct data
2. PowerPoint Curriculum_Review.pptx with 6+ slides
3. Notion database "Assignment Improvement Tracker" with revision entries
4. Email to curriculum_committee@university.edu

Robustness notes (see NewBenchmark audit T3__canvas-assignment-effectiveness-ppt-notion-email):
- DB connection reads PGHOST/PGPORT/PGDATABASE/PGUSER/PGPASSWORD (per-case worker DBs supported).
- Excel is read with data_only=False so formula cells are not silently read as None; a robust
  _to_float handles numbers, strings (currency/thousands/percent/space stripping) and treats
  formula cells as unparseable.
- Numeric columns (Points_Possible, Submission_Count, Completion_Rate, Mean_Score,
  Score_Std_Dev, Discrimination_Index) are compared against the groundtruth per-row with
  per-column tolerances. A blank/formula/unparseable cell counts as a MISMATCH (never a
  silent skip), and each comparison requires a minimum number of matched rows before it runs
  (else it FAILs) -- so a model cannot bypass the core metric computation by leaving cells
  empty or writing Excel formulas. These values are all reproducible from the db seed, so a
  model that follows the documented methodology matches within tolerance.
- The DI comparison uses a tolerant hybrid tolerance (absolute floor + relative) that absorbs
  the small methodological deviations a correct-but-differently-rounded agent can produce
  (observed <= 0.007 for the documented method), while NOT absorbing the much larger errors
  from using the wrong denominator (e.g. the enrollments table count instead of distinct
  graded submitters), so the task still discriminates.
- Effectiveness labels are verified per-row against the classification rule applied to the
  agent's own Discrimination Index (Good if > 0.3, Acceptable if 0.15..0.3, Poor if < 0.15),
  so labeling every assignment "Good" is caught.
- Row-count checks are lenient (>= thresholds) and GT comparison is content-based (keyed by
  course+assignment, with assignment-name fallback), so multi-agent duplicate writes or extra
  assignments with zero submissions cannot cause false FAILs. Columns are anchored by header
  keywords (with a positional fallback to the documented column layout).
- Course Summary uses a UNIFIED eligible-assignment set with the Assignment Metrics layer
  (case-study 2026-08-12, case #15). Previously the summary compared the agent's written
  Avg_Completion_Rate (computed over whatever rows the agent chose to average, e.g. its full
  57-row set including 6 extra zero-submission assignments) against a GT summary computed over
  a 51-row set -- two different eligible sets that cannot be compared. The summary's numeric
  aggregates (Avg_Completion_Rate, Avg_DI, Good/Acceptable/Poor counts, Total) are now
  RE-COMPUTED from the agent's own per-assignment Assignment-Metrics rows restricted to the
  GT-keyed eligible set, then compared to the GT summary. This guarantees both layers reference
  the same eligible set, so an agent that writes correct per-assignment values also produces a
  correct summary, and an agent that invents per-assignment values still fails. The agent's
  literal Course Summary cells are still read for the presence/row-count checks but are no
  longer the basis of the numeric GT comparison.
"""
import argparse
import json
import os
import sys

import psycopg2

DB_CONFIG = {
    "host": os.environ.get("PGHOST", "localhost"),
    "port": int(os.environ.get("PGPORT", "5432")),
    "dbname": os.environ.get("PGDATABASE", "toolathlon_gym"),
    "user": os.environ.get("PGUSER", "eigent"),
    "password": os.environ.get("PGPASSWORD", "camel"),
}

PASS_COUNT = 0
FAIL_COUNT = 0

# Tolerances / thresholds for the GT comparison. The GT (51 assignments, 7 courses, 17
# revision-needed) is fully reproducible from the db seed with the documented methodology
# (raw percentage scores, enrolled = unique students with >=1 graded non-null submission,
# 27% group size = floor(0.27*N), tie-break by user id, population std dev). The DI tolerance
# absorbs the small deviations a correct-but-differently-rounded agent can produce (<= 0.007
# for the documented method), while leaving the larger errors of a wrong denominator intact.
DI_TOL_ABS = 0.05        # absolute floor for DI tolerance
DI_TOL_REL = 0.05        # relative tolerance (fraction of the GT value)
REQUIRED_GT_PRESENT = 48  # of the 51 GT assignments that must appear in the agent sheet
MIN_COMPARE_ROWS = 40    # min matched GT rows before a numeric comparison runs (else FAIL)
MATCH_RATIO = 0.85       # fraction of matched rows that must be within tolerance
MIN_REVISION_NAMES = 14  # of the 17 GT revision-needed assignments that must appear
MIN_REV_COMPARE = 10     # min matched revision assignments before revision numeric checks run
MIN_SUMMARY_COMPARE = 5  # min matched courses before summary numeric checks run

# Per-column absolute tolerances for the Assignment Metrics sheet (a correct model reproduces
# these exactly; tolerances absorb rounding-method differences, e.g. a model rounding the
# 1-decimal columns to whole numbers, while still catching fabricated values by a wide margin).
COL_TOL = {
    "Points_Possible": 0.05,
    "Submission_Count": 1.0,
    "Completion_Rate": 0.5,
    "Mean_Score": 0.5,
    "Score_Std_Dev": 0.5,
}
SUM_COUNT_TOL = 2.0      # Course Summary count columns (Total/Good/Acceptable/Poor)
SUM_COMP_TOL = 0.5       # Course Summary Avg_Completion_Rate
SUM_DI_TOL = 0.1         # Course Summary Avg_DI (wider than DI_TOL so an averaged value
                         #   produced by per-assignment DIs at the tolerance edge still passes)
REV_COMP_TOL = 0.5       # Revision Needed Completion_Rate

# Assignment Metrics numeric columns: (name, GT column index, header keyword groups).
# Keyword groups are tried in order; within a group ALL substrings must appear in a header.
METRIC_NUMERIC_COLS = [
    ("Points_Possible", 2, (("points",),)),
    ("Submission_Count", 3, (("submission",),)),
    ("Completion_Rate", 4, (("completion",),)),
    ("Mean_Score", 5, (("mean",),)),
    ("Score_Std_Dev", 6, (("std",), ("deviation",))),
    ("Discrimination_Index", 7, (("discrimination",), ("di",))),
]


def check(name, condition, detail=""):
    global PASS_COUNT, FAIL_COUNT
    if condition:
        PASS_COUNT += 1
        print(f"  [PASS] {name}")
    else:
        FAIL_COUNT += 1
        msg = f": {str(detail)[:300]}" if detail else ""
        print(f"  [FAIL] {name}{msg}")


def _to_float(v):
    """Robustly coerce a cell value to float.

    Supports int/float/str (strips thousands separators, currency symbols, '%' and spaces).
    Returns None for None, booleans, formula strings and any other unparseable value.
    """
    if v is None:
        return None
    if isinstance(v, bool):
        return None
    if isinstance(v, (int, float)):
        return float(v)
    if isinstance(v, str):
        s = v.strip()
        if s.startswith("="):
            return None  # Excel formula cell -> value not stored as literal
        for ch in (",", "$", "€", "¥", "%", " "):
            s = s.replace(ch, "")
        try:
            return float(s)
        except ValueError:
            return None
    return None


def num_close(a, b, tol=1.0):
    """Numeric close with fallback to case-insensitive string equality for unparseable values."""
    fa, fb = _to_float(a), _to_float(b)
    if fa is not None and fb is not None:
        return abs(fa - fb) <= tol
    if a is None or b is None:
        return False
    return str(a).strip().lower() == str(b).strip().lower()


def di_close(a, b):
    """Tolerant DI comparison: hybrid absolute + relative tolerance."""
    fa, fb = _to_float(a), _to_float(b)
    if fa is not None and fb is not None:
        return abs(fa - fb) <= max(DI_TOL_ABS, DI_TOL_REL * abs(fb))
    if a is None or b is None:
        return False
    return str(a).strip().lower() == str(b).strip().lower()


def _classify(di):
    """Effectiveness classification from a DI value per task.md / guide.md."""
    if di > 0.3:
        return "good"
    if di >= 0.15:
        return "acceptable"
    return "poor"


def _non_empty_rows(ws):
    """Return rows that have a non-empty first cell (skips blank trailing rows)."""
    return [r for r in ws.iter_rows(values_only=True)
            if r and r[0] is not None and str(r[0]).strip() != ""]


def _find_sheet(wb, *substrings):
    for s in wb.sheetnames:
        if all(sub in s.lower() for sub in substrings):
            return wb[s]
    return None


def _build_header_map(ws, min_cells=3):
    """Return {lowercased_header: column_index} from the first row that looks like a header.

    Skips single-cell (e.g. merged-title) rows so a stray title row does not confuse the map.
    """
    for row in ws.iter_rows(values_only=True):
        vals = [c for c in row if c is not None and str(c).strip() != ""]
        if len(vals) < min_cells:
            continue
        hmap = {}
        for j, h in enumerate(row):
            if h is not None:
                s = str(h).strip().lower()
                if s and s not in hmap:
                    hmap[s] = j
        return hmap
    return {}


def _agent_col(hmap, default, keyword_groups):
    """Find an agent column index by header keywords; fall back to `default` (GT layout)."""
    for group in keyword_groups:
        for key, idx in hmap.items():
            if all(k in key for k in group):
                return idx
    return default


def _compare_column(name, gt_idx, agent_idx, tol, gt_row, agent_row, stats, use_di=False):
    """Compare one cell. A blank/formula/unparseable agent cell counts as a MISMATCH."""
    gt_v = gt_row[gt_idx] if len(gt_row) > gt_idx else None
    ag_v = agent_row[agent_idx] if agent_idx is not None and len(agent_row) > agent_idx else None
    stats[name]["checked"] += 1
    if use_di:
        if di_close(ag_v, gt_v):
            stats[name]["matched"] += 1
    else:
        if num_close(ag_v, gt_v, tol):
            stats[name]["matched"] += 1


def _emit_numeric_checks(label, stats, min_checked, missing):
    """Emit PASS/FAIL per numeric column. Never silently passes on too few matched rows."""
    for name, s in stats.items():
        if s["checked"] < min_checked:
            check(f"{label} {name} matches groundtruth",
                  False,
                  f"only {s['checked']} matched rows with values; expected literal numbers")
        else:
            ratio = s["matched"] / s["checked"] if s["checked"] else 0.0
            check(f"{label} {name} matches groundtruth",
                  ratio >= MATCH_RATIO,
                  f"{s['matched']}/{s['checked']} within tolerance; missing rows: {missing[:3]}")


def _recompute_course_summary(rows_for_course, comp_idx, di_idx, eff_idx):
    """Re-aggregate one course's summary from its per-assignment rows.

    Returns (total, good, acceptable, poor, avg_completion, avg_di). rows_for_course is an
    iterable of Assignment-Metrics rows (tuples) for a single course. This is used by the
    Course Summary check so that both the Assignment Metrics layer and the Summary layer
    aggregate over the SAME eligible-assignment set (case-study 2026-08-12, case #15).
    """
    total = 0
    good = acceptable = poor = 0
    comps = []
    dis = []
    for r in rows_for_course:
        total += 1
        eff = str(r[eff_idx]).strip().lower() if eff_idx is not None and len(r) > eff_idx and r[eff_idx] else ""
        if "good" in eff:
            good += 1
        elif "accept" in eff:
            acceptable += 1
        elif "poor" in eff:
            poor += 1
        cv = _to_float(r[comp_idx]) if comp_idx is not None and len(r) > comp_idx else None
        if cv is not None:
            comps.append(cv)
        dv = _to_float(r[di_idx]) if di_idx is not None and len(r) > di_idx else None
        if dv is not None:
            dis.append(dv)
    avg_comp = round(sum(comps) / len(comps), 1) if comps else None
    avg_di = round(sum(dis) / len(dis), 3) if dis else None
    return total, good, acceptable, poor, avg_comp, avg_di


# ============================================================
# Check 1: Excel
# ============================================================
def check_excel(agent_workspace, gt_workspace):
    print("\n=== Checking Excel ===")
    xlsx_path = os.path.join(agent_workspace, "Assessment_Effectiveness.xlsx")
    if not os.path.isfile(xlsx_path):
        check("Assessment_Effectiveness.xlsx exists", False, f"Not found: {xlsx_path}")
        return

    check("Assessment_Effectiveness.xlsx exists", True)

    try:
        import openpyxl
        wb = openpyxl.load_workbook(xlsx_path, data_only=False)
        sheet_names = wb.sheetnames

        check("Sheet 'Assignment Metrics' exists",
              any("assignment" in s.lower() and "metric" in s.lower() for s in sheet_names),
              f"Sheets: {sheet_names}")
        check("Sheet 'Course Summary' exists",
              any("course" in s.lower() and "summary" in s.lower() for s in sheet_names),
              f"Sheets: {sheet_names}")
        check("Sheet 'Revision Needed' exists",
              any("revision" in s.lower() for s in sheet_names),
              f"Sheets: {sheet_names}")

        metrics_ws = _find_sheet(wb, "assignment", "metric") or wb[sheet_names[0]]
        rows = _non_empty_rows(metrics_ws)
        check("Assignment Metrics has header + at least 51 data rows",
              len(rows) >= 51,
              f"Found {len(rows)} non-empty rows")

        if len(rows) > 1:
            all_text = " ".join(str(c) for row in rows for c in row if c).lower()
            check("Contains Fall 2014 course names",
                  "applied analytics" in all_text or "biochemistry" in all_text
                  or "foundations of finance" in all_text or "creative computing" in all_text,
                  f"Sample: {all_text[:200]}")

        # Compare with groundtruth (content-based, tolerant, per-column numeric)
        gt_xlsx = os.path.join(gt_workspace, "Assessment_Effectiveness.xlsx")
        # Function-scope state shared with the Course Summary check below so BOTH layers
        # aggregate over the SAME eligible-assignment set (the GT-keyed rows). Without this,
        # the summary compares the agent's written Avg_Completion_Rate (averaged over whatever
        # rows the agent happened to include, e.g. extra zero-submission assignments) against a
        # GT summary computed over a different set -- two non-comparable eligible sets
        # (case-study 2026-08-12, case #15).
        gt_rows_by_course = {}
        agent_gtkeyed_by_course = {}
        agent_shared_hmap = {}
        agent_shared_idx = {}
        if os.path.isfile(gt_xlsx):
            gt_wb = openpyxl.load_workbook(gt_xlsx, data_only=False)
            gt_metrics = _find_sheet(gt_wb, "assignment", "metric") or gt_wb[gt_wb.sheetnames[0]]
            gt_rows = _non_empty_rows(gt_metrics)

            gt_by_key = {}
            for row in gt_rows[1:]:
                if len(row) >= 2 and row[0] and row[1]:
                    gt_by_key[(str(row[0]).strip(), str(row[1]).strip())] = row
                    gt_rows_by_course.setdefault(str(row[0]).strip(), []).append(row)

            agent_hmap = _build_header_map(metrics_ws)
            agent_course_idx = _agent_col(agent_hmap, 0, (("course",),))
            agent_assign_idx = _agent_col(agent_hmap, 1, (("assignment",),))

            agent_rows = rows[1:]
            agent_by_key = {}
            agent_by_assign = {}
            for row in agent_rows:
                ac = row[agent_course_idx] if agent_course_idx is not None and len(row) > agent_course_idx else None
                aa = row[agent_assign_idx] if agent_assign_idx is not None and len(row) > agent_assign_idx else None
                if ac and aa:
                    agent_by_key[(str(ac).strip(), str(aa).strip())] = row
                if aa is not None:
                    agent_by_assign.setdefault(str(aa).strip(), row)

            # Agent numeric column indices (header-anchored with GT-layout fallback)
            agent_idx_of = {}
            for name, gt_idx, kwgroups in METRIC_NUMERIC_COLS:
                agent_idx_of[name] = _agent_col(agent_hmap, gt_idx, kwgroups)
            eff_idx = _agent_col(agent_hmap, 8, (("effectiveness",), ("effect",)))
            di_agent_idx = agent_idx_of["Discrimination_Index"]

            # Share the agent header map + column indices with the Course Summary check so it
            # can re-aggregate the agent's per-assignment rows over the GT-keyed eligible set.
            agent_shared_hmap.update(agent_hmap)
            agent_shared_idx.update(agent_idx_of)
            agent_shared_idx["__eff__"] = eff_idx

            stats = {name: {"checked": 0, "matched": 0} for name, _, _ in METRIC_NUMERIC_COLS}
            eff_checked = 0
            eff_matched = 0
            missing = []
            checked = 0
            for key, gt_row in gt_by_key.items():
                agent_row = agent_by_key.get(key) or agent_by_assign.get(str(key[1]).strip())
                if agent_row is None:
                    missing.append(key)
                    continue
                checked += 1
                # Collect the agent rows that match a GT key, grouped by course, so the Course
                # Summary check can re-aggregate over the SAME eligible set (case-study #15).
                agent_gtkeyed_by_course.setdefault(str(key[0]).strip(), []).append(agent_row)
                for name, gt_idx, kwgroups in METRIC_NUMERIC_COLS:
                    use_di = (name == "Discrimination_Index")
                    tol = 0.0 if use_di else COL_TOL[name]
                    _compare_column(name, gt_idx, agent_idx_of[name], tol,
                                    gt_row, agent_row, stats, use_di=use_di)

                # Effectiveness label: consistent with the classification rule applied to the
                # agent's own DI (or, as a fallback, to the GT DI) -- catches all-'Good' etc.
                ag_di_v = agent_row[di_agent_idx] if di_agent_idx is not None and len(agent_row) > di_agent_idx else None
                fa = _to_float(ag_di_v)
                if fa is not None:
                    ag_label = agent_row[eff_idx] if eff_idx is not None and len(agent_row) > eff_idx else None
                    label = str(ag_label).strip().lower() if ag_label is not None else ""
                    ok = _classify(fa) in label
                    if not ok:
                        gt_di_v = gt_row[7] if len(gt_row) > 7 else None
                        fb = _to_float(gt_di_v)
                        if fb is not None:
                            ok = _classify(fb) in label
                    eff_checked += 1
                    if ok:
                        eff_matched += 1

            check("Most GT assignments present in agent sheet",
                  checked >= REQUIRED_GT_PRESENT,
                  f"{checked}/{len(gt_by_key)}; missing: {missing[:5]}")

            # Per-column numeric comparison (blank/formula cells are mismatches, never a pass)
            _emit_numeric_checks("Assignment Metrics", stats, MIN_COMPARE_ROWS, missing)

            if eff_checked < MIN_COMPARE_ROWS:
                check("Effectiveness labels match Discrimination Index classification",
                      False,
                      f"only {eff_checked} rows with parseable DI; cannot verify labels")
            else:
                ratio = eff_matched / eff_checked
                check("Effectiveness labels match Discrimination Index classification",
                      ratio >= MATCH_RATIO,
                      f"{eff_matched}/{eff_checked} labels consistent with DI values")

        # Check Course Summary sheet
        summary_ws = _find_sheet(wb, "course", "summary")
        if summary_ws:
            summary_rows = _non_empty_rows(summary_ws)
            check("Course Summary has header + at least 7 course rows",
                  len(summary_rows) >= 8,
                  f"Found {len(summary_rows)} rows")

            # Presence of the majority of the 7 Fall 2014 courses (partial names, tolerant
            # of course-name formatting differences such as a missing term suffix).
            gt_course_hints = [
                "applied analytics",
                "biochemistry",
                "creative computing",
                "data-driven design",
                "environmental economics",
                "foundations of finance",
                "global governance",
            ]
            if len(summary_rows) > 1:
                summary_text = " ".join(str(c) for row in summary_rows for c in row if c).lower()
                present = sum(1 for hint in gt_course_hints if hint in summary_text)
                check("Course Summary lists most Fall 2014 courses",
                      present >= 5,
                      f"{present}/7 course names found")

            # Numeric comparison of the summary values vs GT.
            #
            # CRITICAL (case-study 2026-08-12, case #15): both the Assignment Metrics layer and
            # this Summary layer MUST aggregate over the SAME eligible-assignment set. We do NOT
            # compare the agent's literal Course-Summary cells against the GT summary cells,
            # because the agent may have written its summary over a different row set (e.g. it
            # included extra zero-submission assignments that GT excludes). Instead we
            # re-aggregate the agent's OWN per-assignment Assignment-Metrics values, restricted
            # to the GT-keyed eligible set, and compare the re-aggregated summary to the GT
            # summary. This way an agent that writes correct per-assignment values also produces
            # a correct summary, and an agent that fabricates per-assignment values still fails.
            if os.path.isfile(gt_xlsx) and agent_gtkeyed_by_course:
                gt_wb3 = openpyxl.load_workbook(gt_xlsx, data_only=False)
                gt_sum_ws = _find_sheet(gt_wb3, "course", "summary")
                if gt_sum_ws and len(summary_rows) > 1:
                    gt_sum_rows = _non_empty_rows(gt_sum_ws)[1:]

                    # Column indices into the agent's Assignment-Metrics rows (shared from the
                    # metrics check above). These define how to read Completion_Rate, DI, and
                    # Effectiveness from each per-assignment row.
                    a_comp_idx = agent_shared_idx.get("Completion_Rate")
                    a_di_idx = agent_shared_idx.get("Discrimination_Index")
                    a_eff_idx = agent_shared_idx.get("__eff__")

                    sum_stats = {
                        "Total_Assignments": {"checked": 0, "matched": 0},
                        "Good_Count": {"checked": 0, "matched": 0},
                        "Acceptable_Count": {"checked": 0, "matched": 0},
                        "Poor_Count": {"checked": 0, "matched": 0},
                        "Avg_Completion_Rate": {"checked": 0, "matched": 0},
                        "Avg_DI": {"checked": 0, "matched": 0},
                    }
                    sum_missing = []
                    for gs in gt_sum_rows:
                        gname = str(gs[0]).strip().lower()
                        hint = next((h for h in gt_course_hints if h in gname), None)
                        if hint is None:
                            continue
                        # Match the GT course to an agent course name. The agent's per-assignment
                        # rows were grouped under the EXACT course string from gt_by_key keys, so
                        # find the agent course whose name contains this hint.
                        ag_course = next(
                            (ac for ac in agent_gtkeyed_by_course
                             if hint in ac.lower()),
                            None,
                        )
                        if ag_course is None:
                            sum_missing.append(gname)
                            continue
                        ag_rows = agent_gtkeyed_by_course[ag_course]
                        # GT summary counts/values (authoritative). gs layout:
                        # [Course, Total, Good, Acceptable, Poor, AvgCompletion, AvgDI]
                        gt_total = _to_float(gs[1]) if len(gs) > 1 else None
                        gt_good = _to_float(gs[2]) if len(gs) > 2 else None
                        gt_acc = _to_float(gs[3]) if len(gs) > 3 else None
                        gt_poor = _to_float(gs[4]) if len(gs) > 4 else None
                        gt_avgcomp = _to_float(gs[5]) if len(gs) > 5 else None
                        gt_avgdi = _to_float(gs[6]) if len(gs) > 6 else None
                        # Re-aggregate the agent's eligible rows over the unified set.
                        (a_total, a_good, a_acc, a_poor,
                         a_avgcomp, a_avgdi) = _recompute_course_summary(
                            ag_rows, a_comp_idx, a_di_idx, a_eff_idx,
                        )

                        def _cmp(name, ag_v, gt_v, tol):
                            sum_stats[name]["checked"] += 1
                            if num_close(ag_v, gt_v, tol):
                                sum_stats[name]["matched"] += 1

                        _cmp("Total_Assignments", a_total, gt_total, 1.0)
                        _cmp("Good_Count", a_good, gt_good, SUM_COUNT_TOL)
                        _cmp("Acceptable_Count", a_acc, gt_acc, SUM_COUNT_TOL)
                        _cmp("Poor_Count", a_poor, gt_poor, SUM_COUNT_TOL)
                        _cmp("Avg_Completion_Rate", a_avgcomp, gt_avgcomp, SUM_COMP_TOL)
                        _cmp("Avg_DI", a_avgdi, gt_avgdi, SUM_DI_TOL)

                    _emit_numeric_checks("Course Summary", sum_stats,
                                         MIN_SUMMARY_COMPARE, sum_missing)

        # Check Revision Needed sheet
        revision_ws = _find_sheet(wb, "revision")
        if revision_ws:
            revision_rows = _non_empty_rows(revision_ws)
            check("Revision Needed has header + at least 14 entries",
                  len(revision_rows) >= 15,
                  f"Found {len(revision_rows)} rows")

            rev_hmap = _build_header_map(revision_ws)
            rev_course_idx = _agent_col(rev_hmap, 0, (("course",),))
            rev_assign_idx = _agent_col(rev_hmap, 1, (("assignment",),))
            rev_di_idx = _agent_col(rev_hmap, 2, (("current", "di"), ("di",)))
            rev_comp_idx = _agent_col(rev_hmap, 3, (("completion",),))

            agent_rev_map = {}
            for r in revision_rows[1:]:
                aa = r[rev_assign_idx] if rev_assign_idx is not None and len(r) > rev_assign_idx else None
                if aa is not None:
                    agent_rev_map.setdefault(str(aa).strip(), r)

            # Presence of the majority of the GT revision-needed assignment names
            gt_rev_names = set()
            gt_rev_rows = []
            if os.path.isfile(gt_xlsx):
                gt_wb2 = openpyxl.load_workbook(gt_xlsx, data_only=False)
                gt_rev_ws = _find_sheet(gt_wb2, "revision")
                if gt_rev_ws:
                    gt_rev_rows = _non_empty_rows(gt_rev_ws)[1:]
                    for row in gt_rev_rows:
                        if len(row) >= 2 and row[1] is not None:
                            gt_rev_names.add(str(row[1]).strip())
            if gt_rev_names and len(revision_rows) > 1:
                present = len(gt_rev_names & set(agent_rev_map.keys()))
                check("Revision Needed contains most flagged assignments",
                      present >= MIN_REVISION_NAMES,
                      f"{present}/{len(gt_rev_names)} flagged assignments present")

            # Numeric comparison of Current_DI / Completion_Rate for the matched assignments
            if gt_rev_names and os.path.isfile(gt_xlsx):
                rev_stats = {
                    "Current_DI": {"checked": 0, "matched": 0},
                    "Completion_Rate": {"checked": 0, "matched": 0},
                }
                rev_matched = 0
                for gr in gt_rev_rows:
                    gname = str(gr[1]).strip()
                    ar = agent_rev_map.get(gname)
                    if ar is None:
                        continue
                    rev_matched += 1
                    _compare_column("Current_DI", 2, rev_di_idx, 0.0, gr, ar, rev_stats, use_di=True)
                    _compare_column("Completion_Rate", 3, rev_comp_idx, REV_COMP_TOL, gr, ar, rev_stats)

                if rev_matched >= MIN_REV_COMPARE:
                    _emit_numeric_checks("Revision Needed", rev_stats, MIN_REV_COMPARE, [])
                else:
                    check("Revision Needed numeric values match groundtruth",
                          False,
                          f"only {rev_matched} matched revision assignments")

    except ImportError:
        check("openpyxl available", False, "Cannot parse Excel without openpyxl")
    except Exception as e:
        check("Excel parsing", False, str(e))


# ============================================================
# Check 2: PowerPoint
# ============================================================
def check_pptx(agent_workspace):
    print("\n=== Checking PowerPoint ===")
    pptx_path = os.path.join(agent_workspace, "Curriculum_Review.pptx")
    if not os.path.isfile(pptx_path):
        check("Curriculum_Review.pptx exists", False, f"Not found: {pptx_path}")
        return

    check("Curriculum_Review.pptx exists", True)

    try:
        from pptx import Presentation
        prs = Presentation(pptx_path)
        slide_count = len(prs.slides)
        check("PPT has at least 6 slides", slide_count >= 6,
              f"Found {slide_count} slides")

        all_text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    all_text += shape.text.lower() + " "

        check("PPT mentions discrimination index or DI",
              "discrimination" in all_text or " di " in all_text or "di=" in all_text,
              "Missing DI content")
        check("PPT mentions methodology or formula",
              "methodol" in all_text or "formula" in all_text or "27%" in all_text,
              "Missing methodology")
        check("PPT mentions revision or improvement",
              "revision" in all_text or "improv" in all_text or "needs" in all_text,
              "Missing revision content")
        check("PPT mentions recommendations",
              "recommend" in all_text or "suggest" in all_text or "action" in all_text,
              "Missing recommendations")
        check("PPT mentions Fall 2014 or course names",
              "fall 2014" in all_text or "2014" in all_text,
              "Missing term reference")

    except ImportError:
        size = os.path.getsize(pptx_path)
        check("PPT file has content (>5KB)", size > 5000, f"Size: {size}")
    except Exception as e:
        check("PPT parsing", False, str(e))


# ============================================================
# Check 3: Notion
# ============================================================
def check_notion():
    print("\n=== Checking Notion ===")
    try:
        conn = psycopg2.connect(**DB_CONFIG)
        cur = conn.cursor()

        cur.execute("""
            SELECT id, title, properties FROM notion.databases
            WHERE title::text ILIKE '%%improvement%%tracker%%'
               OR title::text ILIKE '%%assignment%%tracker%%'
        """)
        dbs = cur.fetchall()
        check("Notion database 'Assignment Improvement Tracker' exists",
              len(dbs) >= 1,
              f"Found {len(dbs)} matching databases")

        if dbs:
            props_raw = dbs[0][2]
            if isinstance(props_raw, str):
                try:
                    props = json.loads(props_raw)
                except Exception:
                    props = {}
            elif isinstance(props_raw, dict):
                props = props_raw
            else:
                props = {}
            if not isinstance(props, dict):
                props = {}

            prop_names = [k.lower() for k in props.keys()]
            check("Database has 'Assignment' property",
                  any("assignment" in p for p in prop_names),
                  f"Properties: {prop_names}")
            check("Database has 'Course' property",
                  any("course" in p for p in prop_names),
                  f"Properties: {prop_names}")
            check("Database has 'Status' property",
                  any("status" in p for p in prop_names),
                  f"Properties: {prop_names}")

            # Collect pages across all matching databases (multi-agent may create several).
            db_ids = [d[0] for d in dbs]
            seen = set()
            pages = []
            for db_id in db_ids:
                cur.execute("""
                    SELECT id, properties FROM notion.pages
                    WHERE parent::text LIKE %s
                      AND (archived IS NULL OR archived = false)
                      AND (in_trash IS NULL OR in_trash = false)
                """, (f'%{db_id}%',))
                for pid, propval in cur.fetchall():
                    if pid not in seen:
                        seen.add(pid)
                        pages.append((pid, propval))
            check("Notion has entries for all revision-needed assignments (>= 17)",
                  len(pages) >= 17,
                  f"Found {len(pages)} pages")

            if pages:
                all_props_text = " ".join(
                    json.dumps(p[1]) if isinstance(p[1], dict) else str(p[1])
                    for p in pages
                ).lower()
                check("Pages contain assignment names",
                      "cma" in all_props_text or "tma" in all_props_text or "final" in all_props_text,
                      f"Sample: {all_props_text[:200]}")
                check("Pages contain course references",
                      "foundation" in all_props_text or "creative" in all_props_text or "finance" in all_props_text,
                      f"Sample: {all_props_text[:200]}")

        cur.close()
        conn.close()
    except Exception as e:
        check("Notion check", False, str(e))


# ============================================================
# Check 4: Email
# ============================================================
def check_email():
    print("\n=== Checking Email ===")
    try:
        conn = psycopg2.connect(**DB_CONFIG)
        cur = conn.cursor()

        cur.execute("""
            SELECT id, subject, from_addr, to_addr, body_text
            FROM email.messages
            WHERE to_addr::text ILIKE '%%curriculum_committee@university.edu%%'
        """)
        emails = cur.fetchall()
        check("Email to curriculum_committee@university.edu found",
              len(emails) >= 1,
              f"Found {len(emails)} matching emails")

        if emails:
            email = emails[0]
            subject = str(email[1] or "").lower()
            body = str(email[4] or "").lower()

            check("Email subject mentions assignment or effectiveness",
                  "assignment" in subject or "effectiveness" in subject or "fall 2014" in subject,
                  f"Subject: {email[1]}")
            check("Email body has substantive content",
                  len(body) > 50,
                  f"Body length: {len(body)}")
            check("Email body mentions key findings",
                  any(term in body for term in ["revision", "poor", "completion", "discrimination", "good", "acceptable"]),
                  f"Body sample: {body[:200]}")

        cur.close()
        conn.close()
    except Exception as e:
        check("Email check", False, str(e))


# ============================================================
# Main
# ============================================================
def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--agent_workspace", required=True)
    parser.add_argument("--groundtruth_workspace", required=False, default="")
    parser.add_argument("--launch_time", required=False)
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()

    gt_ws = args.groundtruth_workspace or args.agent_workspace

    check_excel(args.agent_workspace, gt_ws)
    check_pptx(args.agent_workspace)
    check_notion()
    check_email()

    total = PASS_COUNT + FAIL_COUNT
    print(f"\n=== Results: {PASS_COUNT}/{total} passed ===")

    if args.res_log_file:
        result = {"passed": PASS_COUNT, "failed": FAIL_COUNT, "total": total}
        with open(args.res_log_file, "w") as f:
            json.dump(result, f, indent=2)

    if FAIL_COUNT > 0:
        print(f"{FAIL_COUNT} checks failed")
        sys.exit(1)
    else:
        print("All checks passed!")
        sys.exit(0)


if __name__ == "__main__":
    main()
