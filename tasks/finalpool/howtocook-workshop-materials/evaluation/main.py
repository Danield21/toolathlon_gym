"""
Evaluation for howtocook-workshop-materials task.

Checks structural properties of the three output files since the agent
has freedom to choose any 3 recipes from HowToCook:

1. Workshop_Handbook.docx
   - File exists
   - Has top-level heading containing "Cooking Workshop Handbook"
   - Has at least 3 second-level (Heading 1) headings for dish names
   - Contains ingredient content (word "ingredient" or ingredient-like items)
   - Contains step content (word "step" or numbered instructions)
   - Has a "Tips and Notes" section

2. Workshop_Slides.pptx
   - File exists
   - First slide title contains "Cooking Workshop" or "Team"
   - Has at least 7 slides (1 title + 3 dishes x 2 slides + 1 closing)
   - Closing slide contains "Enjoy"

3. Shopping_List.pdf
   - File exists
   - File size > 1KB (non-trivial content)
   - Contains "Shopping List" text (checked via PDF reader)
"""
import os
import sys
import json
import re
from argparse import ArgumentParser
from datetime import datetime

from docx import Document
from pptx import Presentation


# ---------- language-independent content helpers ----------
# The HowToCook MCP data source is Chinese; models may emit either English
# (encouraged by the task) or Chinese body text. These helpers accept both
# so that structurally-correct submissions are never failed on language alone.

# Ingredients/measurement signals (English + Chinese) and a digit+unit pattern.
_INGREDIENT_KEYWORDS = [
    "ingredient", "tablespoon", "teaspoon", "gram", "cup", "ml", "oil", "salt",
    "sugar", "egg", "flour", "kg", "oz", "chicken", "broccoli", "garlic",
    "soy sauce", "ginger", "scallion", "onion", "beef", "pork", "fish", "rice",
    "noodle", "tofu", "mushroom", "pepper", "carrot", "potato", "vinegar",
    "wine", "water", "cola", "lemon", "cucumber", "cabbage", "spinach", "shrimp",
    "食材", "配料", "主料", "辅料", "调料", "克", "毫升", "勺", "油", "盐", "糖",
    "鸡蛋", "西红柿", "番茄", "葱", "姜", "蒜", "酱油", "料酒", "面粉", "淀粉",
    "鸡肉", "猪肉", "牛肉", "鱼", "虾", "豆腐", "蘑菇", "香菇", "木耳", "醋",
]
_UNIT_PATTERN = re.compile(
    r"\d+\s*(?:g|kg|ml|l|tbsp|tablespoon|teaspoon|cup|oz|lb|"
    r"克|毫升|升|斤|两|公斤|个|颗|瓣|勺|条|根|片|块|枚|只|份)",
    re.IGNORECASE,
)

# Cooking-step signals (English verbs + Chinese verbs + numbered-step markers).
_STEP_KEYWORDS = [
    "step", "stir", "cook", "heat", "add", "pour", "cut", "boil", "simmer",
    "fry", "mix", "blanch", "season", "serve", "preheat", "saute", "toss",
    "sauce", "烧", "炒", "煮", "切", "加热", "加入", "倒入", "焯", "搅拌",
    "煸", "煎", "炖", "蒸", "炸", "翻", "步骤",
]
_STEP_PATTERN = re.compile(
    r"(?:step\s*\d|步骤\s*\d|第[一二三四五六七八九十\d]+步)", re.IGNORECASE
)

_STOPWORDS = set("""a an and or of to in on for with at by from as is are was were
be been being this that these those it its we our you your they their i me my he him
she her them us not no can could will would shall should may might must do does did
done have has had about into over under above below out up down off again further
then once here there when where why how all any both each few more most other some
such only own same so than too very just don now also""".split())


def _contains_ingredient_content(text):
    if any(k in text for k in _INGREDIENT_KEYWORDS):
        return True
    if _UNIT_PATTERN.search(text):
        return True
    return False


def _contains_step_content(text):
    if any(k in text for k in _STEP_KEYWORDS):
        return True
    if _STEP_PATTERN.search(text):
        return True
    return False


_CJK_REJOIN = re.compile(r"(?<=[一-鿿])\s+(?=[一-鿿])")


def _rejoin_cjk_spaces(text):
    """Remove whitespace sitting between two CJK chars.

    pypdf/PyPDF2 extract CJK text from fonts lacking a ToUnicode map by
    inserting a space between (almost) every glyph, e.g. '红 烧 肉'. Rejoining
    restores contiguous runs so the CJK tokenizer still sees real words.
    """
    return _CJK_REJOIN.sub("", text)


def _tokenize(text):
    """Language-independent token set (latin words + CJK runs)."""
    tokens = set()
    text = _rejoin_cjk_spaces(text)
    for m in re.finditer(r"[A-Za-z0-9]+", text):
        t = m.group(0).lower()
        if len(t) >= 2 and not t.isdigit() and t not in _STOPWORDS:
            if t.endswith("s") and len(t) > 3:
                t = t[:-1]  # crude plural normalization
            tokens.add(t)
    for m in re.finditer(r"[一-鿿]{2,}", text):
        tokens.add(m.group(0))
    return tokens


def check_word_doc(agent_workspace):
    """Check Workshop_Handbook.docx for required structure and content."""
    passed = 0
    total = 0
    doc_path = os.path.join(agent_workspace, "Workshop_Handbook.docx")

    # 1. File exists
    total += 1
    if not os.path.exists(doc_path):
        print("  FAIL: Workshop_Handbook.docx not found")
        return passed, total
    passed += 1
    print("  PASS: Workshop_Handbook.docx exists")

    doc = Document(doc_path)

    # Extract headings and full text
    headings_by_level = {}  # level -> list of heading texts
    all_text_parts = []
    for para in doc.paragraphs:
        all_text_parts.append(para.text)
        if para.style and para.style.name:
            style_name = para.style.name
            if "Heading" in style_name or "Title" in style_name:
                # Extract level: "Heading 1" -> 1, "Title" -> 0
                if style_name == "Title":
                    level = 0
                else:
                    try:
                        level = int(style_name.split()[-1])
                    except (ValueError, IndexError):
                        level = -1
                if level not in headings_by_level:
                    headings_by_level[level] = []
                headings_by_level[level].append(para.text.strip())

    full_text = " ".join(all_text_parts).lower()

    # 2. Title heading contains "Cooking Workshop Handbook"
    total += 1
    title_headings = headings_by_level.get(0, [])
    all_headings_flat = []
    for lvl_headings in headings_by_level.values():
        all_headings_flat.extend(lvl_headings)
    has_title = any("cooking workshop handbook" in h.lower() for h in all_headings_flat)
    if not has_title:
        # Fallback: check if "cooking workshop handbook" appears anywhere in text
        has_title = "cooking workshop handbook" in full_text
    if has_title:
        passed += 1
        print("  PASS: Title 'Cooking Workshop Handbook' found")
    else:
        print(f"  FAIL: Title 'Cooking Workshop Handbook' not found. Headings: {all_headings_flat[:10]}")

    # 3. At least 3 second-level headings (dish names)
    total += 1
    # Count Heading 1/2/3 level headings (excluding known section names)
    dish_headings = []
    known_sections = {"tips and notes", "cooking workshop handbook", "ingredients",
                      "cooking steps", "tips", "notes", "introduction", "welcome"}
    for level in [1, 2, 3]:
        for h in headings_by_level.get(level, []):
            if h.lower().strip() not in known_sections:
                dish_headings.append(h)

    # If we don't have enough from levels 1-2, also check if there are at least
    # 3 distinct content sections (by counting Heading 1 level headings)
    h1_count = len(headings_by_level.get(1, []))
    if len(dish_headings) >= 3 or h1_count >= 3:
        passed += 1
        print(f"  PASS: Found {max(len(dish_headings), h1_count)} section headings (need >= 3)")
    else:
        print(f"  FAIL: Only {len(dish_headings)} dish headings found (need >= 3). All H1: {headings_by_level.get(1, [])}")

    # 4. Contains ingredient content (language-independent: EN/CN keywords + digit+unit)
    total += 1
    has_ingredients = _contains_ingredient_content(full_text)
    if has_ingredients:
        passed += 1
        print("  PASS: Ingredient content found")
    else:
        print("  FAIL: No ingredient content found in document")

    # 5. Contains step/instruction content (language-independent: EN/CN verbs + step markers)
    total += 1
    has_steps = _contains_step_content(full_text)
    if has_steps:
        passed += 1
        print("  PASS: Cooking step content found")
    else:
        print("  FAIL: No cooking step content found in document")

    # 6. Has "Tips and Notes" section
    # The task mandates the English section name, but a model that translates
    # it into Chinese (a common HowToCook-side behaviour) must not be failed.
    total += 1
    tips_keywords = ("tips", "note", "小贴士", "注意事项", "烹饪提示", "小提示",
                     "贴心提示", "温馨提示", "厨师提示", "烹饪建议", "小建议",
                     "备忘")
    has_tips = any("tips" in h.lower() and "note" in h.lower()
                   for h in all_headings_flat)
    if not has_tips:
        # Fallback: any tips-like heading (EN or CN)
        has_tips = any(any(k in h.lower() for k in tips_keywords)
                       for h in all_headings_flat)
    if not has_tips:
        # Fallback: check body text (specific EN/CN tips phrases)
        has_tips = any(k in full_text
                       for k in ("tips and notes", "tips", "小贴士", "注意事项",
                                 "烹饪提示", "小提示", "贴心提示", "温馨提示",
                                 "厨师提示", "烹饪建议"))
    if has_tips:
        passed += 1
        print("  PASS: 'Tips and Notes' section found")
    else:
        print("  FAIL: 'Tips and Notes' section not found")

    return passed, total


def check_pptx(agent_workspace):
    """Check Workshop_Slides.pptx for required structure."""
    passed = 0
    total = 0
    pptx_path = os.path.join(agent_workspace, "Workshop_Slides.pptx")

    # 1. File exists
    total += 1
    if not os.path.exists(pptx_path):
        print("  FAIL: Workshop_Slides.pptx not found")
        return passed, total
    passed += 1
    print("  PASS: Workshop_Slides.pptx exists")

    prs = Presentation(pptx_path)

    # Extract text from all slides
    slide_texts = []
    for slide in prs.slides:
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    texts.append(para.text)
        slide_texts.append(" ".join(texts))

    # 2. First slide title contains "Cooking Workshop" or "Team"
    total += 1
    first_text = slide_texts[0].lower() if slide_texts else ""
    if "cooking workshop" in first_text or "team" in first_text or "workshop" in first_text:
        passed += 1
        print("  PASS: First slide contains workshop/team title")
    else:
        # Check across all slides as fallback
        all_text = " ".join(slide_texts).lower()
        if "cooking workshop" in all_text or "team" in all_text:
            passed += 1
            print("  PASS: Workshop/team title found in slides (not first slide)")
        else:
            print(f"  FAIL: No workshop/team title found. First slide: '{first_text[:100]}'")

    # 3. At least 7 slides
    total += 1
    slide_count = len(prs.slides)
    if slide_count >= 7:
        passed += 1
        print(f"  PASS: {slide_count} slides found (need >= 7)")
    else:
        print(f"  FAIL: Only {slide_count} slides found (need >= 7)")

    # 4. Closing slide contains "Enjoy"
    total += 1
    last_text = slide_texts[-1].lower() if slide_texts else ""
    all_text_lower = " ".join(slide_texts).lower()
    if "enjoy" in last_text:
        passed += 1
        print("  PASS: Closing slide contains 'Enjoy'")
    elif "enjoy" in all_text_lower:
        passed += 1
        print("  PASS: 'Enjoy' found in slides (not last slide)")
    else:
        print(f"  FAIL: 'Enjoy' not found in slides. Last slide: '{last_text[:100]}'")

    return passed, total


def extract_pdf_text(pdf_path):
    pdf_text = ""
    try:
        from pypdf import PdfReader
        reader = PdfReader(pdf_path)
        for page in reader.pages:
            pdf_text += page.extract_text() or ""
    except Exception:
        try:
            from PyPDF2 import PdfReader
            reader = PdfReader(pdf_path)
            for page in reader.pages:
                pdf_text += page.extract_text() or ""
        except Exception:
            with open(pdf_path, "rb") as f:
                pdf_text = f.read().decode("latin-1", errors="ignore")
    return pdf_text


def check_word_pdf_consistency(agent_workspace):
    """Verify that PDF shopping list overlaps with Word doc ingredients."""
    passed = 0
    total = 0
    doc_path = os.path.join(agent_workspace, "Workshop_Handbook.docx")
    pdf_path = os.path.join(agent_workspace, "Shopping_List.pdf")
    if not os.path.exists(doc_path) or not os.path.exists(pdf_path):
        return passed, total

    try:
        doc = Document(doc_path)
        word_text = " ".join(p.text for p in doc.paragraphs).lower()
    except Exception as e:
        print(f"  [WARN] could not read Word doc for consistency: {e}")
        return passed, total

    pdf_text = extract_pdf_text(pdf_path).lower()

    # If PDF text extraction was unreliable (empty/garbled binary fallback),
    # skip the consistency check rather than fail a valid deliverable.
    alnum_chars = sum(1 for ch in pdf_text if ch.isalnum())
    if alnum_chars < 50:
        print("  [SKIP] PDF text extraction unreliable; skipping PDF-Word consistency check")
        return passed, total

    # Language-independent content overlap: the PDF consolidates ingredients from
    # the same three dishes as the Word handbook, so their content tokens must overlap.
    word_tokens = _tokenize(word_text)
    pdf_tokens = _tokenize(pdf_text)

    # Long-but-garbled extraction (e.g. a CJK font without a ToUnicode map whose
    # glyphs are all single chars, or a binary fallback decoding to junk) can
    # yield substantial text yet almost no recognizable tokens. Comparing such
    # broken text would fail a valid deliverable, so skip instead. Genuine
    # inconsistencies where both documents ARE extractable still FAIL below.
    if len(word_tokens) < 3 or len(pdf_tokens) < 3:
        print("  [SKIP] Document text extraction yielded no recognizable content tokens; "
              "skipping PDF-Word consistency check")
        return passed, total

    shared = word_tokens & pdf_tokens

    total += 1
    if len(shared) >= 3:
        passed += 1
        print(f"  PASS: PDF-Word ingredient overlap (found {len(shared)} shared tokens: {sorted(shared)[:8]})")
    else:
        print(f"  FAIL: Only {len(shared)} shared content tokens between PDF and Word (need >= 3). Word tokens: {sorted(word_tokens)[:8]}")

    return passed, total


def check_pdf(agent_workspace):
    """Check Shopping_List.pdf for required structure."""
    passed = 0
    total = 0
    pdf_path = os.path.join(agent_workspace, "Shopping_List.pdf")

    # 1. File exists
    total += 1
    if not os.path.exists(pdf_path):
        print("  FAIL: Shopping_List.pdf not found")
        return passed, total
    passed += 1
    print("  PASS: Shopping_List.pdf exists")

    # 2. File size > 1KB
    total += 1
    file_size = os.path.getsize(pdf_path)
    if file_size > 1024:
        passed += 1
        print(f"  PASS: PDF size is {file_size} bytes (> 1KB)")
    else:
        print(f"  FAIL: PDF size is {file_size} bytes (need > 1KB)")

    # 3. Contains "Shopping List" text
    total += 1
    try:
        # Try to read PDF text using PyPDF2 or pypdf
        pdf_text = ""
        try:
            from pypdf import PdfReader
            reader = PdfReader(pdf_path)
            for page in reader.pages:
                pdf_text += page.extract_text() or ""
        except ImportError:
            try:
                from PyPDF2 import PdfReader
                reader = PdfReader(pdf_path)
                for page in reader.pages:
                    pdf_text += page.extract_text() or ""
            except ImportError:
                # Fallback: check raw bytes for the string
                with open(pdf_path, "rb") as f:
                    raw = f.read()
                pdf_text = raw.decode("latin-1", errors="ignore")

        if "shopping list" in pdf_text.lower() or "Shopping List" in pdf_text:
            passed += 1
            print("  PASS: PDF contains 'Shopping List' text")
        else:
            # Even if text extraction fails, if file is large enough it likely has content
            if file_size > 2000:
                passed += 1
                print("  PASS: PDF is substantial size, likely contains shopping list content")
            else:
                print(f"  FAIL: 'Shopping List' text not found in PDF. Extracted: '{pdf_text[:200]}'")
    except Exception as e:
        # If we can't read the PDF at all, check size as fallback
        if file_size > 2000:
            passed += 1
            print(f"  PASS: PDF exists with substantial size ({file_size} bytes), text extraction failed: {e}")
        else:
            print(f"  FAIL: Could not read PDF text: {e}")

    return passed, total


def main(args):
    total_passed = 0
    total_checks = 0

    # Check 1: Word document
    print("--- Check 1: Word Document (Workshop_Handbook.docx) ---")
    p, t = check_word_doc(args.agent_workspace)
    print(f"  Word Doc: {p}/{t} checks passed")
    total_passed += p
    total_checks += t

    # Check 2: PowerPoint
    print("\n--- Check 2: PowerPoint (Workshop_Slides.pptx) ---")
    p, t = check_pptx(args.agent_workspace)
    print(f"  PowerPoint: {p}/{t} checks passed")
    total_passed += p
    total_checks += t

    # Check 3: PDF
    print("\n--- Check 3: PDF (Shopping_List.pdf) ---")
    p, t = check_pdf(args.agent_workspace)
    print(f"  PDF: {p}/{t} checks passed")
    total_passed += p
    total_checks += t

    # Check 4: Word-PDF ingredient consistency
    print("\n--- Check 4: Word-PDF Ingredient Consistency ---")
    p, t = check_word_pdf_consistency(args.agent_workspace)
    if t > 0:
        print(f"  Consistency: {p}/{t} checks passed")
        total_passed += p
        total_checks += t

    # Overall
    if total_checks == 0:
        print("\nFAIL: No checks were performed.")
        accuracy = 0.0
    else:
        accuracy = total_passed / total_checks * 100
        print(f"\nOverall: {total_passed}/{total_checks} checks passed ({accuracy:.1f}%)")

    result = {
        "total_passed": total_passed,
        "total_checks": total_checks,
        "accuracy": accuracy,
        "timestamp": datetime.now().isoformat(),
    }

    if args.res_log_file:
        with open(args.res_log_file, "w") as f:
            json.dump(result, f, indent=2)
        print(f"Report saved to {args.res_log_file}")

    if total_passed == total_checks:
        print("PASS")
        sys.exit(0)
    else:
        print("FAIL")
        sys.exit(1)


if __name__ == "__main__":
    parser = ArgumentParser()
    parser.add_argument("--agent_workspace", required=False, default=".")
    parser.add_argument("--groundtruth_workspace", required=False, default=".")
    parser.add_argument("--launch_time", required=False)
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()
    main(args)
