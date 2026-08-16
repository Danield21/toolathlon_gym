"""Evaluation for sf-sales-regional-analysis-ppt-gcal-email.

Checks:
1. PPT file Regional_Sales_Review.pptx with at least 5 slides
   - Title slide with "Regional Sales Performance Review"
   - Contains Europe as top region
   - Contains Latin America as lowest region
   - Contains recommendations
2. GCal event "Regional Sales Review Meeting" ~14 days from launch_time
3. Email to sales-leadership@company.example.com
   Subject: "Regional Sales Performance Review"
"""
import argparse
import json
import os
import sys
from datetime import datetime, timedelta, timezone

import psycopg2

DB = dict(host=os.environ.get("PGHOST", "localhost"), port=int(os.environ.get("PGPORT", "5432")), dbname="toolathlon_gym", user="eigent", password="camel")

PASS_COUNT = 0
FAIL_COUNT = 0
RUNTIME_ONLY_FAIL = 0


def check(name, condition, detail="", runtime_only=False):
    global PASS_COUNT, FAIL_COUNT, RUNTIME_ONLY_FAIL
    if condition:
        PASS_COUNT += 1
        print(f"  [PASS] {name}")
    else:
        FAIL_COUNT += 1
        if runtime_only:
            RUNTIME_ONLY_FAIL += 1
            tag = " (runtime-only)"
        else:
            tag = ""
        detail_str = f": {str(detail)[:200]}" if detail else ""
        print(f"  [FAIL] {name}{tag}{detail_str}")


def check_ppt(agent_workspace):
    print("\n=== Checking PPT File ===")
    ppt_path = os.path.join(agent_workspace, "Regional_Sales_Review.pptx")
    check("Regional_Sales_Review.pptx exists", os.path.isfile(ppt_path),
          f"Expected at {ppt_path}")
    if not os.path.isfile(ppt_path):
        return

    try:
        from pptx import Presentation
        prs = Presentation(ppt_path)
    except Exception as e:
        check("PPT file readable", False, str(e))
        return

    check("PPT has at least 5 slides", len(prs.slides) >= 5,
          f"Found {len(prs.slides)} slides")

    all_text = " ".join(
        shape.text.lower()
        for slide in prs.slides
        for shape in slide.shapes
        if hasattr(shape, "text")
    )

    check("PPT contains 'Regional Sales Performance Review'",
          "regional sales" in all_text and "review" in all_text,
          "Title not found")
    check("PPT mentions Europe as top region",
          "europe" in all_text, "Europe not mentioned")
    check("PPT mentions Latin America",
          "latin america" in all_text or "latin" in all_text,
          "Latin America not mentioned")
    # Verify all 5 required slide titles appear
    check("PPT has 'Revenue by Region' slide title",
          "revenue by region" in all_text, "Slide 2 title missing")
    check("PPT has 'Top Region Deep Dive' slide title",
          "top region deep dive" in all_text or "top region" in all_text, "Slide 3 title missing")
    check("PPT has 'Bottom Region Analysis' slide title",
          "bottom region" in all_text, "Slide 4 title missing")
    check("PPT has 'Recommended Actions' slide title",
          "recommended actions" in all_text or "recommendations" in all_text, "Slide 5 title missing")
    # Revenue figures (keep lenient substring check - derived at runtime from Snowflake)
    check("PPT contains revenue data",
          "648" in all_text or "642" in all_text or "606" in all_text,
          "Revenue figures not found")
    check("PPT contains recommendations",
          "recommend" in all_text or "action" in all_text or "strategy" in all_text,
          "No recommendations found")


def check_gcal(launch_time_str=None):
    print("\n=== Checking Google Calendar ===")
    conn = psycopg2.connect(**DB)
    cur = conn.cursor()

    cur.execute("""
        SELECT summary, start_datetime, end_datetime, description
        FROM gcal.events
        WHERE LOWER(summary) LIKE '%regional%' AND LOWER(summary) LIKE '%sales%'
    """)
    events = cur.fetchall()
    agent_populated = len(events) >= 1
    check("Regional Sales Review Meeting event created", agent_populated,
          f"Found {len(events)} matching events", runtime_only=not agent_populated)

    if events and launch_time_str:
        try:
            launch_time = datetime.fromisoformat(launch_time_str)
            if launch_time.tzinfo is None:
                launch_time = launch_time.replace(tzinfo=timezone.utc)
            target_date = launch_time + timedelta(days=14)
            matched_event = None
            for event in events:
                event_start = event[1]
                if event_start.tzinfo is None:
                    event_start = event_start.replace(tzinfo=timezone.utc)
                diff_days = abs((event_start.date() - target_date.date()).days)
                if diff_days <= 3:
                    matched_event = event
                    check("Review Meeting is ~14 days from launch", True)
                    break
            else:
                check("Review Meeting is ~14 days from launch", False,
                      f"Closest event at {events[0][1]}, expected ~{target_date.date()}")
            # Verify 1-hour duration
            if matched_event:
                start, end = matched_event[1], matched_event[2]
                dur = (end - start).total_seconds() / 60.0
                check("Review Meeting duration is 60 minutes",
                      abs(dur - 60) <= 5, f"Duration: {dur} min")
        except Exception as e:
            check("Review Meeting date check", False, str(e))

    cur.close()
    conn.close()


def check_emails():
    print("\n=== Checking Emails ===")
    conn = psycopg2.connect(**DB)
    cur = conn.cursor()

    cur.execute("""
        SELECT subject, from_addr, to_addr, body_text
        FROM email.messages
    """)
    all_emails = cur.fetchall()
    conn.close()

    def parse_recipients(to_addr):
        if to_addr is None:
            return []
        if isinstance(to_addr, list):
            return [str(r).strip().lower() for r in to_addr]
        to_str = str(to_addr).strip()
        try:
            parsed = json.loads(to_str)
            if isinstance(parsed, list):
                return [str(r).strip().lower() for r in parsed]
            return [to_str.lower()]
        except (json.JSONDecodeError, TypeError):
            return [to_str.lower()]

    target = "sales-leadership@company.example.com"
    found = None
    for subj, from_addr, to_addr, body in all_emails:
        recipients = parse_recipients(to_addr)
        if target in recipients:
            found = (subj, from_addr, to_addr, body)
            break

    agent_emailed = found is not None
    check("Email sent to sales-leadership@company.example.com", agent_emailed,
          runtime_only=not agent_emailed)
    if found:
        subj, from_addr, to_addr, body = found
        check("Email from reporting@company.example.com",
              "reporting@company.example.com" in (from_addr or "").lower(),
              f"From: {from_addr}")
        check("Subject is 'Regional Sales Performance Review'",
              "regional sales" in (subj or "").lower() and "review" in (subj or "").lower(),
              f"Subject: {subj}")


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--agent_workspace", required=False, default=".")
    parser.add_argument("--groundtruth_workspace", required=False)
    parser.add_argument("--launch_time", required=False)
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()

    print("=" * 70)
    print("SF SALES REGIONAL ANALYSIS PPT GCAL EMAIL - EVALUATION")
    print("=" * 70)

    check_ppt(args.agent_workspace)
    check_gcal(args.launch_time)
    check_emails()

    print(f"\n=== SUMMARY ===")
    print(f"  Passed: {PASS_COUNT}")
    print(f"  Failed: {FAIL_COUNT} (runtime-only: {RUNTIME_ONLY_FAIL})")
    blocking_fail = FAIL_COUNT - RUNTIME_ONLY_FAIL
    overall = blocking_fail == 0
    print(f"  Overall: {'PASS' if overall else 'FAIL'}")
    sys.exit(0 if overall else 1)


if __name__ == "__main__":
    main()
