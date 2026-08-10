"""
Evaluation script for ppt-snowflake-executive task.

Dynamically queries PostgreSQL to compute expected Q4 2024 values,
then checks agent output files for correctness.
"""

from argparse import ArgumentParser
import sys
import os
import re
from pathlib import Path


def _norm_num(s):
    """Reduce a number/string to a canonical form by stripping thousands
    separators, currency symbols, and spaces. Lets '2,579', '$411,857.18',
    and '411 857' all be compared against plain '2579' / '411857.18'."""
    return re.sub(r"[\s,$€£¥]", "", str(s))


def get_expected_data():
    """Query PostgreSQL to compute expected Q4 2024 values."""
    import psycopg2

    conn = psycopg2.connect(
        host=os.environ.get("PGHOST", "localhost"), port=5432, dbname="toolathlon_gym",
        user="eigent", password="camel"
    )
    cur = conn.cursor()

    # Regional summary
    cur.execute('''
        SELECT c."REGION", COUNT(*) as orders,
               ROUND(SUM(o."TOTAL_AMOUNT"::float)::numeric, 2) as revenue
        FROM sf_data."SALES_DW__PUBLIC__ORDERS" o
        JOIN sf_data."SALES_DW__PUBLIC__CUSTOMERS" c
            ON o."CUSTOMER_ID" = c."CUSTOMER_ID"
        WHERE o."ORDER_DATE" >= '2024-10-01' AND o."ORDER_DATE" < '2025-01-01'
        GROUP BY c."REGION" ORDER BY c."REGION"
    ''')
    regions = [(r[0], int(r[1]), float(r[2])) for r in cur.fetchall()]

    # Segment analysis
    cur.execute('''
        SELECT c."SEGMENT", COUNT(*) as orders,
               ROUND(SUM(o."TOTAL_AMOUNT"::float)::numeric, 2) as revenue
        FROM sf_data."SALES_DW__PUBLIC__ORDERS" o
        JOIN sf_data."SALES_DW__PUBLIC__CUSTOMERS" c
            ON o."CUSTOMER_ID" = c."CUSTOMER_ID"
        WHERE o."ORDER_DATE" >= '2024-10-01' AND o."ORDER_DATE" < '2025-01-01'
        GROUP BY c."SEGMENT" ORDER BY c."SEGMENT"
    ''')
    segments = [(r[0], int(r[1]), float(r[2])) for r in cur.fetchall()]

    # Totals
    cur.execute('''
        SELECT COUNT(*) as orders,
               ROUND(SUM(o."TOTAL_AMOUNT"::float)::numeric, 2) as revenue
        FROM sf_data."SALES_DW__PUBLIC__ORDERS" o
        JOIN sf_data."SALES_DW__PUBLIC__CUSTOMERS" c
            ON o."CUSTOMER_ID" = c."CUSTOMER_ID"
        WHERE o."ORDER_DATE" >= '2024-10-01' AND o."ORDER_DATE" < '2025-01-01'
    ''')
    total_orders, total_revenue = cur.fetchone()
    total_orders = int(total_orders)
    total_revenue = float(total_revenue)

    conn.close()
    return regions, segments, total_orders, total_revenue


def check_excel(workspace, regions, segments, total_revenue):
    """Check Q4_Executive_Analysis.xlsx for correctness."""
    import openpyxl

    xlsx_path = Path(workspace) / "Q4_Executive_Analysis.xlsx"
    if not xlsx_path.exists():
        return False, f"Q4_Executive_Analysis.xlsx not found in {workspace}"

    wb = openpyxl.load_workbook(xlsx_path, data_only=True)

    # Check sheet names
    if "Regional Summary" not in wb.sheetnames:
        return False, f"Missing 'Regional Summary' sheet. Found: {wb.sheetnames}"
    if "Segment Analysis" not in wb.sheetnames:
        return False, f"Missing 'Segment Analysis' sheet. Found: {wb.sheetnames}"

    # --- Check Regional Summary ---
    ws1 = wb["Regional Summary"]
    rows1 = list(ws1.iter_rows(values_only=True))
    if len(rows1) < 2:
        return False, "Regional Summary sheet has no data rows"

    # Check header
    header1 = [str(h).strip() if h else "" for h in rows1[0]]
    expected_cols = ["Region", "Orders", "Revenue", "Avg_Order_Value"]
    for col in expected_cols:
        if col not in header1:
            return False, f"Regional Summary missing column '{col}'. Found: {header1}"

    # Map column indices
    idx = {col: header1.index(col) for col in expected_cols}

    data_rows = rows1[1:]
    if len(data_rows) != len(regions):
        return False, f"Regional Summary: expected {len(regions)} data rows, got {len(data_rows)}"

    for i, (exp_region, exp_orders, exp_revenue) in enumerate(regions):
        row = data_rows[i]
        region_val = str(row[idx["Region"]]).strip() if row[idx["Region"]] else ""
        orders_val = row[idx["Orders"]]
        revenue_val = row[idx["Revenue"]]
        avg_val = row[idx["Avg_Order_Value"]]

        if region_val != exp_region:
            return False, f"Regional Summary row {i+1}: expected region '{exp_region}', got '{region_val}'"

        if orders_val is None or int(orders_val) != exp_orders:
            return False, f"Regional Summary '{exp_region}': expected orders {exp_orders}, got {orders_val}"

        if revenue_val is None or abs(float(revenue_val) - exp_revenue) > 1.0:
            return False, f"Regional Summary '{exp_region}': expected revenue {exp_revenue}, got {revenue_val}"

        exp_avg = round(exp_revenue / exp_orders, 2)
        if avg_val is None or abs(float(avg_val) - exp_avg) > 0.02:
            return False, f"Regional Summary '{exp_region}': expected avg {exp_avg}, got {avg_val}"

    print("  [PASS] Regional Summary data correct")

    # --- Check Segment Analysis ---
    ws2 = wb["Segment Analysis"]
    rows2 = list(ws2.iter_rows(values_only=True))
    if len(rows2) < 2:
        return False, "Segment Analysis sheet has no data rows"

    header2 = [str(h).strip() if h else "" for h in rows2[0]]
    expected_cols2 = ["Segment", "Orders", "Revenue", "Revenue_Share_Pct"]
    for col in expected_cols2:
        if col not in header2:
            return False, f"Segment Analysis missing column '{col}'. Found: {header2}"

    idx2 = {col: header2.index(col) for col in expected_cols2}

    data_rows2 = rows2[1:]
    if len(data_rows2) != len(segments):
        return False, f"Segment Analysis: expected {len(segments)} data rows, got {len(data_rows2)}"

    for i, (exp_segment, exp_orders, exp_revenue) in enumerate(segments):
        row = data_rows2[i]
        seg_val = str(row[idx2["Segment"]]).strip() if row[idx2["Segment"]] else ""
        orders_val = row[idx2["Orders"]]
        revenue_val = row[idx2["Revenue"]]
        pct_val = row[idx2["Revenue_Share_Pct"]]

        if seg_val != exp_segment:
            return False, f"Segment Analysis row {i+1}: expected '{exp_segment}', got '{seg_val}'"

        if orders_val is None or int(orders_val) != exp_orders:
            return False, f"Segment Analysis '{exp_segment}': expected orders {exp_orders}, got {orders_val}"

        if revenue_val is None or abs(float(revenue_val) - exp_revenue) > 1.0:
            return False, f"Segment Analysis '{exp_segment}': expected revenue {exp_revenue}, got {revenue_val}"

        exp_pct = round(exp_revenue / total_revenue * 100, 1)
        if pct_val is None or abs(float(pct_val) - exp_pct) > 0.2:
            return False, f"Segment Analysis '{exp_segment}': expected pct {exp_pct}, got {pct_val}"

    print("  [PASS] Segment Analysis data correct")
    wb.close()
    return True, "Excel file checks passed"


def check_pptx(workspace, regions, total_orders, total_revenue):
    """Check Q4_Executive_Briefing.pptx for correctness."""
    from pptx import Presentation

    pptx_path = Path(workspace) / "Q4_Executive_Briefing.pptx"
    if not pptx_path.exists():
        return False, f"Q4_Executive_Briefing.pptx not found in {workspace}"

    prs = Presentation(str(pptx_path))
    slides = list(prs.slides)

    # At least 7 slides: 1 title + 5 regions + 1 summary
    if len(slides) < 7:
        return False, f"Expected at least 7 slides, got {len(slides)}"
    print(f"  Slide count: {len(slides)}")

    # Collect all text from all slides
    all_text = []
    for slide in slides:
        slide_texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    slide_texts.append(paragraph.text)
        all_text.append("\n".join(slide_texts))

    full_text = "\n".join(all_text).lower()

    # Check title slide contains "Q4 2024"
    first_slide_text = all_text[0].lower()
    if "q4 2024" not in first_slide_text:
        return False, f"Title slide does not contain 'Q4 2024'. Text: {all_text[0][:200]}"
    print("  [PASS] Title slide contains 'Q4 2024'")

    # Check all 5 regions appear somewhere in the presentation
    region_names = [r[0] for r in regions]
    for region in region_names:
        if region.lower() not in full_text:
            return False, f"Region '{region}' not found in presentation text"
    print("  [PASS] All 5 regions present in presentation")

    # Verify each region has a dedicated slide 2..6 with title matching that region
    # slides 2..6 (indices 1..5) should each contain one region name in the title/first text
    region_slides_ok = 0
    for i in range(1, min(6, len(all_text))):
        slide_text = all_text[i].lower()
        matching_regions = [r for r in region_names if r.lower() in slide_text]
        if len(matching_regions) == 1:
            region_slides_ok += 1
    if region_slides_ok < 4:
        return False, f"Expected each of slides 2-6 to cover exactly one region; only {region_slides_ok}/5 matched"
    print(f"  [PASS] {region_slides_ok}/5 region slides have unique region coverage")

    # Check total orders and revenue appear in the summary
    # Look in last few slides for totals. Normalize away thousands separators,
    # currency symbols, and spaces so "2,579" / "$411,857.18" / "411 857" match.
    found_total_orders = False
    found_total_revenue = False
    orders_candidates = [str(total_orders), f"{total_orders:,}"]
    rev_candidates = [
        f"{total_revenue:,.2f}", f"{total_revenue:,.0f}",
        f"{total_revenue:.2f}", f"{total_revenue:.0f}",
        f"{total_revenue:,.1f}",
    ]
    for slide_text in all_text[-3:]:
        norm_slide = _norm_num(slide_text)
        if any(_norm_num(c) in norm_slide for c in orders_candidates):
            found_total_orders = True
        if any(_norm_num(c) in norm_slide for c in rev_candidates):
            found_total_revenue = True

    if not found_total_orders:
        return False, f"Total orders ({total_orders}) not found in summary slides"
    if not found_total_revenue:
        return False, f"Total revenue ({total_revenue}) not found in summary slides"
    print("  [PASS] Summary slide contains total orders and revenue")

    return True, "PPTX file checks passed"


def check_pdf(workspace, total_orders, total_revenue):
    """Check Q4_Executive_Briefing.pdf exists and carries Q4 content (not just a
    placeholder). A real PDF-to-PDF export of the briefing must mention Q4 and the
    total revenue; a blank/reportlab placeholder would fail here."""
    pdf_path = Path(workspace) / "Q4_Executive_Briefing.pdf"
    if not pdf_path.exists():
        return False, f"Q4_Executive_Briefing.pdf not found in {workspace}"

    size = pdf_path.stat().st_size
    if size < 1024:
        return False, f"PDF file too small ({size} bytes), likely invalid"

    # Extract text. Require a PDF reader so a non-text placeholder cannot pass.
    text = ""
    reader_ok = False
    try:
        import PyPDF2
        with open(pdf_path, "rb") as f:
            reader = PyPDF2.PdfReader(f)
            for page in reader.pages:
                text += page.extract_text() or ""
        reader_ok = True
    except ImportError:
        try:
            import pdfplumber
            with pdfplumber.open(pdf_path) as p:
                for page in p.pages:
                    text += page.extract_text() or ""
            reader_ok = True
        except ImportError:
            pass
    if not reader_ok:
        return False, "No PDF reader library (PyPDF2/pdfplumber) available in eval env"
    if not text.strip():
        return False, "PDF contains no extractable text (possible placeholder)"

    text_norm = _norm_num(text).lower()
    has_q4 = "q4" in text_norm or "quarter 4" in text_norm or "q42024" in text_norm or ("2024" in text_norm and "quarter" in text_norm)
    if not has_q4:
        return False, "PDF text does not reference Q4 / fourth quarter 2024"

    # Total revenue must appear (tolerant to formatting via _norm_num).
    rev_canon = _norm_num(f"{total_revenue:.2f}")
    nums = re.findall(r"\d+(?:\.\d+)?", text_norm)
    rev_hit = any(abs(float(n) - total_revenue) < max(10.0, total_revenue * 0.01) for n in nums)
    if not rev_hit:
        return False, f"PDF text does not mention total revenue (~${total_revenue:,.2f})"

    print(f"  [PASS] PDF exists (size={size}) and references Q4 totals")
    return True, "PDF file check passed"


if __name__ == "__main__":
    parser = ArgumentParser()
    parser.add_argument("--agent_workspace", required=False)
    parser.add_argument("--groundtruth_workspace", required=False)
    parser.add_argument("--res_log_file", required=False)
    parser.add_argument("--launch_time", required=False, help="Launch time")
    args = parser.parse_args()

    workspace = args.agent_workspace
    if not workspace:
        print("Error: --agent_workspace is required")
        sys.exit(1)

    print("Fetching expected data from database...")
    try:
        regions, segments, total_orders, total_revenue = get_expected_data()
        print(f"  Regions: {len(regions)}, Segments: {len(segments)}")
        print(f"  Total: {total_orders} orders, ${total_revenue:,.2f} revenue")
    except Exception as e:
        print(f"Error querying database: {e}")
        sys.exit(1)

    all_passed = True

    # Check Excel
    print("\n--- Check 1: Excel File ---")
    try:
        ok, msg = check_excel(workspace, regions, segments, total_revenue)
        if not ok:
            print(f"  [FAIL] {msg}")
            all_passed = False
        else:
            print(f"  {msg}")
    except Exception as e:
        print(f"  [FAIL] Excel check error: {e}")
        all_passed = False

    # Check PPTX
    print("\n--- Check 2: PowerPoint File ---")
    try:
        ok, msg = check_pptx(workspace, regions, total_orders, total_revenue)
        if not ok:
            print(f"  [FAIL] {msg}")
            all_passed = False
        else:
            print(f"  {msg}")
    except Exception as e:
        print(f"  [FAIL] PPTX check error: {e}")
        all_passed = False

    # Check PDF
    print("\n--- Check 3: PDF File ---")
    try:
        ok, msg = check_pdf(workspace, total_orders, total_revenue)
        if not ok:
            print(f"  [FAIL] {msg}")
            all_passed = False
        else:
            print(f"  {msg}")
    except Exception as e:
        print(f"  [FAIL] PDF check error: {e}")
        all_passed = False

    if all_passed:
        print("\nPass all tests!")
    else:
        print("\nSome checks failed.")
        sys.exit(1)
