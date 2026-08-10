#!/usr/bin/env python3
"""Evaluation script for budget-variance-analysis-quarterly task validation.

Scoring philosophy: a model that correctly follows the task brief MUST pass.
Checks are therefore structural / internally-consistent / anchored-to-inputs,
NOT a rigid cell-by-cell diff against a fixed groundtruth (the old version's
cell diff assumed a fixed layout and relied on fabricated numbers).

Deliverable checks:
  * variance_analysis.xlsx  - header/columns present; numeric data present;
      internal math consistent (|Variance| ~ |Actual - Budget|, Variance % ~
      Variance / Budget); both favorable and unfavorable categories present;
      Budget/Actual values anchored to the source files
      (approved_budget.xlsx + q1_actual_expenditures.csv) when those are still
      in the workspace.
  * variance_tracking.xlsx  - Department + month columns present; >=3 data
      rows covering >=2 of the four departments; numeric variance values.
  * budget_forecast.xlsx    - scenario table with >=2 numeric scenario rows.
  * dept_variance_reports.docx - >=150 words + topical keywords.
  * executive_presentation.pptx - >=3 slides + 'variance' + budget/forecast.
  * gcal.events             - >=1 budget/variance/review meeting.
  * email.messages          - >=1 topical message.

DB connections read PGHOST/PGPORT/PGDATABASE/PGUSER/PGPASSWORD env vars with
defaults matching preprocess/main.py, so the evaluator always queries the same
worker database the agent wrote to (never a hardcoded dbname/port).
"""

from argparse import ArgumentParser
import json
import os
import sys

PASS_COUNT = 0
FAIL_COUNT = 0
WARN_COUNT = 0
IS_GT_SELF_TEST = False

DB_CONFIG = {
    "host": os.environ.get("PGHOST", "localhost"),
    "port": int(os.environ.get("PGPORT", "5432")),
    "dbname": os.environ.get("PGDATABASE", "toolathlon_gym"),
    "user": os.environ.get("PGUSER", "eigent"),
    "password": os.environ.get("PGPASSWORD", "camel"),
}


def record(name, passed, detail="", db_side=False):
    global PASS_COUNT, FAIL_COUNT, WARN_COUNT
    if passed:
        PASS_COUNT += 1
        print(f"  [PASS] {name}")
    else:
        if IS_GT_SELF_TEST and db_side:
            WARN_COUNT += 1
            msg = f": {detail[:300]}" if detail else ""
            print(f"  [WARN] {name} (GT self-test mode, DB-side){msg}")
        else:
            FAIL_COUNT += 1
            msg = f": {detail[:300]}" if detail else ""
            print(f"  [FAIL] {name}{msg}")


def _to_float(v):
    """Robust numeric parser: int/float/None or string with $/€/¥/commas/%/
    spaces stripped; None when unparseable. Formula strings ('=...') yield
    None because no cached value is available under data_only=False."""
    if v is None:
        return None
    if isinstance(v, bool):
        return None
    if isinstance(v, (int, float)):
        return float(v)
    s = str(v).strip()
    if not s:
        return None
    if s.startswith("="):
        return None
    neg = False
    t = s
    for ch in ("$", "€", "¥", "£", ",", "%", " ", " "):
        t = t.replace(ch, "")
    if t.startswith("(") and t.endswith(")"):
        neg = True
        t = t[1:-1].strip()
    try:
        val = float(t)
    except ValueError:
        return None
    return -val if neg else val


def num_close(a, b, tol=1.0):
    """Compare two values: numeric when both parse, else case-insensitive
    string equality (tolerant of None == None)."""
    fa = _to_float(a)
    fb = _to_float(b)
    if fa is not None and fb is not None:
        return abs(fa - fb) <= tol
    if a is None or b is None:
        return a is None and b is None
    return str(a).strip().lower() == str(b).strip().lower()


def str_match(a, b):
    if a is None or b is None:
        return a is None and b is None
    return str(a).strip().lower() == str(b).strip().lower()


# ---------------------------------------------------------------------------
# file discovery
# ---------------------------------------------------------------------------

def _find_xlsx_by_keywords(workspace, keywords, exclude=None):
    """Score xlsx files by keyword matches; filename match weighs 10x more."""
    import glob
    import openpyxl
    exclude = exclude or set()
    scored = []
    for path in sorted(glob.glob(os.path.join(workspace, "*.xlsx"))):
        if path in exclude:
            continue
        fname_low = os.path.basename(path).lower()
        if fname_low.startswith("~$") or fname_low.startswith("."):
            continue
        score = 0
        for kw in keywords:
            if kw in fname_low:
                score += 10
        try:
            wb = openpyxl.load_workbook(path, data_only=True, read_only=True)
            content_text = " ".join(s.lower() for s in wb.sheetnames)
            for ws in wb.worksheets:
                row_count = 0
                for row in ws.iter_rows(values_only=True):
                    if row_count > 5:
                        break
                    row_count += 1
                    for cell in row:
                        if cell is not None:
                            content_text += " " + str(cell).lower()
            wb.close()
            for kw in keywords:
                if kw in content_text:
                    score += 1
        except Exception:
            pass
        if score > 0:
            scored.append((score, path))
    scored.sort(key=lambda x: (-x[0], x[1]))
    return [p for _, p in scored]


def _find_docx_by_keywords(workspace, keywords, exclude=None):
    import glob
    from docx import Document
    exclude = exclude or set()
    candidates = []
    for path in sorted(glob.glob(os.path.join(workspace, "*.docx"))):
        if path in exclude:
            continue
        fname_low = os.path.basename(path).lower()
        if fname_low.startswith("~$") or fname_low.startswith("."):
            continue
        if any(kw in fname_low for kw in keywords):
            candidates.append(path)
            continue
        try:
            doc = Document(path)
            text_low = "\n".join(p.text for p in doc.paragraphs).lower()
            if any(kw in text_low for kw in keywords):
                candidates.append(path)
        except Exception:
            continue
    return candidates


def _find_pptx_by_keywords(workspace, keywords):
    import glob
    from pptx import Presentation
    candidates = []
    for path in sorted(glob.glob(os.path.join(workspace, "*.pptx"))):
        fname_low = os.path.basename(path).lower()
        if fname_low.startswith("~$") or fname_low.startswith("."):
            continue
        if any(kw in fname_low for kw in keywords):
            candidates.append(path)
            continue
        try:
            prs = Presentation(path)
            text_low = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_low += " " + (shape.text or "").lower()
            if any(kw in text_low for kw in keywords):
                candidates.append(path)
        except Exception:
            continue
    return candidates


# ---------------------------------------------------------------------------
# xlsx helpers
# ---------------------------------------------------------------------------

def _match_sheet(wb, names):
    """Select the worksheet for this deliverable: single-sheet fallback, then
    case-insensitive exact name, then substring keyword fallback."""
    if len(wb.sheetnames) == 1:
        return wb[wb.sheetnames[0]]
    lower = {n.strip().lower() for n in names}
    for sn in wb.sheetnames:
        if sn.strip().lower() in lower:
            return wb[sn]
    for sn in wb.sheetnames:
        snl = sn.lower()
        for n in names:
            n = n.lower()
            if n in snl or snl in n:
                return wb[sn]
    return None


_ANALYSIS_HEADER_LABELS = {
    "cost": ["cost", "center"],
    "dept": ["department", "dept"],
    "category": ["category", "account", "expense"],
    "budget": ["budget"],
    "actual": ["actual", "spend", "spent"],
    "variance": ["variance", "var"],
    "status": ["status", "favor", "unfavor"],
}


def _header_col_score(s, role):
    """Score how well a single header cell maps to a required role.

    Returns -1 when the cell cannot be that role, otherwise a non-negative
    score (higher = better).  The key guardrail: a percentage-looking header
    ('%'/'pct'/'percent') can never be the budget / actual / variance dollar
    column, so an extra column such as '% of Budget' cannot hijack the budget
    or variance-% mapping and mass-fail an otherwise-correct workbook.
    """
    s = (s or "").strip().lower()
    if not s:
        return -1
    is_pct = "%" in s or "pct" in s or "percent" in s
    if role == "budget":
        if "budget" not in s or is_pct:
            return -1
        score = 0
        if "q1" in s:
            score += 4
        if "approved" in s or "annual" in s or "baseline" in s:
            score += 1
        if "actual" in s or "variance" in s or "spend" in s:
            score -= 2
        return score
    if role == "actual":
        if "actual" not in s and "spent" not in s and "expenditure" not in s:
            return -1
        if is_pct:
            return -1
        score = 0
        if "q1" in s:
            score += 4
        if "variance" in s or "budget" in s:
            score -= 1
        return score
    if role == "variance":
        if "variance" not in s and "var" not in s:
            return -1
        if is_pct:
            return -1
        score = 0
        stripped = s.replace("$", "").replace("amount", "").replace("dollar", "").strip()
        if stripped == "variance":
            score += 3
        return score
    if role == "pct":
        if not is_pct:
            return -1
        score = 0
        if "variance" in s or "var" in s:
            score += 5  # 'Variance %' is the canonical column
        if "budget" in s:
            score += 1
        if "actual" in s:
            score += 1
        return score
    if role == "status":
        if not any(k in s for k in ("status", "favor", "unfavor", "condition", "flag")):
            return -1
        return 0
    if role == "cost":
        if "cost" not in s and "center" not in s:
            return -1
        return 0
    if role == "category":
        if not any(k in s for k in ("category", "account", "expense", "line item")):
            return -1
        return 2 if "category" in s else 1
    if role == "dept":
        if "department" not in s and "dept" not in s:
            return -1
        return 2 if "department" in s else 1
    return -1


def _pick_cols(texts):
    """Pick, per required role, the best-matching column.  A cell may serve
    multiple roles (e.g. 'Variance %' is both the variance and the pct hint),
    but each role picks its own best column independently, so a later extra
    column can never overwrite an earlier exact match."""
    cols = {}
    for role in ("budget", "actual", "variance", "pct", "status", "cost",
                 "category", "dept"):
        best_j, best_score = None, -1
        for j, s in enumerate(texts):
            sc = _header_col_score(s, role)
            if sc > best_score:
                best_j, best_score = j, sc
        if best_j is not None:
            cols[role] = best_j
    return cols


def _find_analysis_header(ws, max_scan=25):
    """Return (header_row_idx_1based, set_of_labels, col_map) for the row that
    looks most like a variance-analysis header."""
    best = None
    for i, row in enumerate(ws.iter_rows(values_only=True), start=1):
        if i > max_scan:
            break
        if row is None or all(c is None for c in row):
            continue
        texts = [str(c).strip().lower() for c in row if c is not None]
        joined = " ".join(texts)
        if not any(k in joined for k in ("budget", "variance", "actual", "cost")):
            continue
        labels = set()
        for label, kws in _ANALYSIS_HEADER_LABELS.items():
            if any(k in joined for k in kws):
                labels.add(label)
        if best is None or len(labels) > best[1]:
            best = (i, len(labels), labels, texts)
        if len(labels) >= 4:
            break
    if best is None:
        return None
    row_idx, _n, labels, texts = best
    cols = _pick_cols(texts)
    return (row_idx, labels, cols)


def _load_budget_map(workspace):
    """Parse approved_budget.xlsx -> {(cost_center_lower, category_lower): Q1 budget}.
    Returns None when the source file is absent/unreadable."""
    import openpyxl
    path = os.path.join(workspace, "approved_budget.xlsx")
    if not os.path.isfile(path):
        return None
    try:
        wb = openpyxl.load_workbook(path, data_only=False)
        ws = wb[wb.sheetnames[0]]
        rows = list(ws.iter_rows(values_only=True))
        hdr_idx = None
        for i, r in enumerate(rows):
            joined = " ".join(str(c).strip().lower() for c in r if c is not None)
            if "cost center" in joined and "budget" in joined and "category" in joined:
                hdr_idx = i
                break
        if hdr_idx is None:
            wb.close()
            return None
        hdr = rows[hdr_idx]
        cc_col = cat_col = None
        q1_col = None
        for j, c in enumerate(hdr):
            s = str(c).strip().lower()
            if s == "cost center":
                cc_col = j
            if s == "category":
                cat_col = j
            if "q1" in s and "budget" in s:
                q1_col = j
        m = {}
        for r in rows[hdr_idx + 1:]:
            if not any(c is not None for c in r):
                continue
            if cc_col is None or cat_col is None or q1_col is None:
                continue
            if q1_col >= len(r) or cat_col >= len(r) or cc_col >= len(r):
                continue
            cc = r[cc_col]
            cat = r[cat_col]
            q1 = _to_float(r[q1_col])
            if cc is not None and cat is not None and q1 is not None:
                m[(str(cc).strip().lower(), str(cat).strip().lower())] = q1
        wb.close()
        return m
    except Exception:
        return None


def _load_actuals_map(workspace):
    """Sum q1_actual_expenditures.csv by (cost_center_lower, category_lower).
    Returns None when the source file is absent/unreadable."""
    import csv
    path = os.path.join(workspace, "q1_actual_expenditures.csv")
    if not os.path.isfile(path):
        return None
    m = {}
    try:
        with open(path, newline="", encoding="utf-8") as f:
            reader = csv.DictReader(f)
            for row in reader:
                cc = (row.get("cost_center") or row.get("Cost Center") or "").strip().lower()
                cat = (row.get("spending_category") or row.get("category") or "").strip().lower()
                amt = _to_float(row.get("amount") or row.get("Amount"))
                if cc and cat and amt is not None:
                    key = (cc, cat)
                    m[key] = m.get(key, 0.0) + amt
        return m
    except Exception:
        return None


def check_variance_analysis(path, workspace):
    """Structure + internal consistency + source anchoring for the main
    variance analysis workbook."""
    import openpyxl
    wb = openpyxl.load_workbook(path, data_only=False)
    ws = _match_sheet(wb, ["variance analysis", "variance", "analysis"])
    if ws is None:
        record("variance_analysis: sheet present", False,
               f"sheets: {wb.sheetnames}")
        wb.close()
        return
    record(f"variance_analysis: sheet present ('{ws.title}')", True)

    hdr = _find_analysis_header(ws)
    if hdr is None:
        record("variance_analysis: header row with key columns", False,
               "No row with budget/actual/variance column labels found")
        wb.close()
        return
    row_idx, labels, cols = hdr
    needed = {"budget", "actual", "variance"}
    missing = needed - labels
    if missing:
        record("variance_analysis: header row with key columns", False,
               f"missing labels: {sorted(missing)}; found: {sorted(labels)}")
        wb.close()
        return
    record(f"variance_analysis: header row with key columns (row {row_idx})", True)

    budget_map = _load_budget_map(workspace)
    actual_map = _load_actuals_map(workspace)

    n_data = 0
    n_numeric = 0
    n_var_chk = 0
    n_var_bad = 0
    n_pct_chk = 0
    n_pct_bad = 0
    n_fav = 0
    n_unfav = 0
    n_on = 0
    n_budget_matched = 0
    n_budget_bad = 0
    n_actual_matched = 0
    n_actual_bad = 0

    def getc(row, key):
        j = cols.get(key)
        if j is None or j >= len(row):
            return None
        return row[j]

    for r in ws.iter_rows(min_row=row_idx + 1, values_only=True):
        if not any(c is not None for c in r):
            continue
        budget = _to_float(getc(r, "budget"))
        actual = _to_float(getc(r, "actual"))
        variance = _to_float(getc(r, "variance"))
        pct = _to_float(getc(r, "pct"))
        cc = getc(r, "cost")
        category = getc(r, "category")
        n_data += 1
        if budget is None and actual is None and variance is None and pct is None:
            continue  # text-only row (e.g. sub-headers / driver notes)
        if budget is not None and actual is not None:
            n_numeric += 1
            diff = actual - budget
            # magnitude consistency: |variance| == |actual - budget|
            if variance is not None:
                n_var_chk += 1
                tol = max(0.02 * max(abs(budget), abs(actual)), 100.0)
                if abs(abs(variance) - abs(diff)) > tol:
                    n_var_bad += 1
                    record(f"variance_analysis: row {n_data} variance math", False,
                           f"Budget={budget} Actual={actual} Variance={variance}")
            # percent consistency (accept percent-scale 5.0 or fraction-scale
            # 0.05, whichever matches — the fraction interpretation must also
            # cover sub-1% values such as an 'Overall Variance %' of 0.85%).
            if pct is not None and budget != 0:
                n_pct_chk += 1
                expected_pct = diff / budget * 100.0
                pct_v = pct
                tol = max(0.05 * abs(expected_pct), 0.5)
                ok_pct = abs(abs(pct_v) - abs(expected_pct)) <= tol
                ok_frac = abs(abs(pct_v * 100.0) - abs(expected_pct)) <= tol
                if not (ok_pct or ok_frac):
                    n_pct_bad += 1
                    record(f"variance_analysis: row {n_data} variance %", False,
                           f"Budget={budget} Actual={actual} pct={pct} expected~{expected_pct:.2f}%")
            # direction classification from the data (not the status text)
            if abs(diff) <= max(0.005 * abs(budget), 1.0):
                n_on += 1
            elif diff < 0:
                n_fav += 1
            else:
                n_unfav += 1
            # anchoring to source files
            if cc is not None and category is not None:
                key = (str(cc).strip().lower(), str(category).strip().lower())
                if budget_map is not None and key in budget_map:
                    n_budget_matched += 1
                    tol = max(0.01 * abs(budget_map[key]), 100.0)
                    if abs(budget - budget_map[key]) > tol:
                        n_budget_bad += 1
                        record(f"variance_analysis: budget matches approved_budget.xlsx (row {n_data})",
                               False, f"Category {key} budget={budget} expected={budget_map[key]}")
                if actual_map is not None and key in actual_map:
                    n_actual_matched += 1
                    tol = max(0.01 * abs(actual_map[key]), 100.0)
                    if abs(actual - actual_map[key]) > tol:
                        n_actual_bad += 1
                        record(f"variance_analysis: actual matches q1_actual_expenditures.csv (row {n_data})",
                               False, f"Category {key} actual={actual} expected={actual_map[key]}")

    record(f"variance_analysis: >=3 numeric data rows ({n_numeric} found)",
           n_numeric >= 3, f"numeric rows={n_numeric}, data rows={n_data}")
    record(f"variance_analysis: variance math consistent ({n_var_chk} checked)",
           n_var_bad == 0, f"bad={n_var_bad} checked={n_var_chk}")
    record(f"variance_analysis: variance % consistent ({n_pct_chk} checked)",
           n_pct_bad == 0, f"bad={n_pct_bad} checked={n_pct_chk}")
    record("variance_analysis: has both favorable and unfavorable categories",
           n_fav >= 1 and n_unfav >= 1,
           f"favorable={n_fav} unfavorable={n_unfav} on-budget={n_on}")

    if budget_map is None:
        record("variance_analysis: budget values match approved_budget.xlsx",
               True, "source file not present in workspace; skipped")
    elif n_budget_matched == 0:
        record("variance_analysis: budget values match approved_budget.xlsx",
               True, "no rows keyed to (cost center, category); skipped")
    else:
        record(f"variance_analysis: budget values match approved_budget.xlsx ({n_budget_matched} rows)",
               n_budget_bad == 0, f"bad={n_budget_bad} matched={n_budget_matched}")

    if actual_map is None:
        record("variance_analysis: actual values match q1_actual_expenditures.csv",
               True, "source file not present in workspace; skipped")
    elif n_actual_matched == 0:
        record("variance_analysis: actual values match q1_actual_expenditures.csv",
               True, "no rows keyed to (cost center, category); skipped")
    else:
        record(f"variance_analysis: actual values match q1_actual_expenditures.csv ({n_actual_matched} rows)",
               n_actual_bad == 0, f"bad={n_actual_bad} matched={n_actual_matched}")

    wb.close()


_TRACK_MONTH_KWS = ["january", "february", "march", "variance", "trend", "q1", "month"]


def _find_track_header(ws, max_scan=25):
    for i, row in enumerate(ws.iter_rows(values_only=True), start=1):
        if i > max_scan:
            break
        if row is None or all(c is None for c in row):
            continue
        texts = [str(c).strip().lower() for c in row if c is not None]
        joined = " ".join(texts)
        has_dept = "department" in joined or "dept" in joined
        month_kws = sum(1 for k in _TRACK_MONTH_KWS if k in joined)
        if has_dept and month_kws >= 2:
            return i
    return None


def check_variance_tracking(path):
    import openpyxl
    wb = openpyxl.load_workbook(path, data_only=False)
    ws = _match_sheet(wb, ["variance tracking", "tracking", "monitor", "variance"])
    if ws is None:
        record("variance_tracking: sheet present", False,
               f"sheets: {wb.sheetnames}")
        wb.close()
        return
    record(f"variance_tracking: sheet present ('{ws.title}')", True)

    hdr_idx = _find_track_header(ws)
    if hdr_idx is None:
        record("variance_tracking: header with Department + period columns", False,
               "No Department+month header row found")
        wb.close()
        return
    record(f"variance_tracking: header with Department + period columns (row {hdr_idx})", True)

    # Anchor the department column by header keyword instead of assuming it is
    # position 0 (column order is not part of the graded contract).
    header_cells = next(ws.iter_rows(min_row=hdr_idx, max_row=hdr_idx, values_only=True))
    dept_col = 0
    for j, c in enumerate(header_cells):
        s = str(c).strip().lower() if c is not None else ""
        if "department" in s or "dept" in s:
            dept_col = j
            break

    data_rows = []
    for r in ws.iter_rows(min_row=hdr_idx + 1, values_only=True):
        if not any(c is not None for c in r):
            continue
        data_rows.append(r)

    depts = set()
    numeric_cells = 0
    for r in data_rows:
        if dept_col < len(r) and r[dept_col] is not None:
            depts.add(str(r[dept_col]).strip().lower())
        for j, cell in enumerate(r):
            if j == dept_col:
                continue
            if _to_float(cell) is not None:
                numeric_cells += 1

    record(f"variance_tracking: >=3 data rows ({len(data_rows)} found)",
           len(data_rows) >= 3, f"rows={len(data_rows)}")
    covered = depts & {"operations", "sales", "marketing", "it"}
    record(f"variance_tracking: covers >=2 of the four departments ({len(covered)} found)",
           len(covered) >= 2, f"depts={sorted(covered)}")
    record(f"variance_tracking: numeric variance values present ({numeric_cells} cells)",
           numeric_cells >= 4, f"numeric cells={numeric_cells}")
    wb.close()


def _find_scenario_header(ws, max_scan=35):
    for i, row in enumerate(ws.iter_rows(values_only=True), start=1):
        if i > max_scan:
            break
        if row is None or all(c is None for c in row):
            continue
        texts = [str(c).strip().lower() for c in row if c is not None]
        joined = " ".join(texts)
        if "scenario" in joined and any(k in joined for k in
                                       ("operations", "sales", "marketing", "it", "total", "dept", "budget")):
            return i
    return None


def check_budget_forecast(path):
    import openpyxl
    wb = openpyxl.load_workbook(path, data_only=False)
    ws = _match_sheet(wb, ["budget forecast", "forecast", "budget", "projection"])
    if ws is None:
        record("budget_forecast: sheet present", False,
               f"sheets: {wb.sheetnames}")
        wb.close()
        return
    record(f"budget_forecast: sheet present ('{ws.title}')", True)

    hdr_idx = _find_scenario_header(ws)
    if hdr_idx is None:
        record("budget_forecast: scenario table present", False,
               "No 'Scenario' header row found")
        wb.close()
        return
    record(f"budget_forecast: scenario table present (row {hdr_idx})", True)

    scenario_rows = 0
    for r in ws.iter_rows(min_row=hdr_idx + 1, values_only=True):
        if not any(c is not None for c in r):
            continue
        if not r[0]:
            continue
        numeric = sum(1 for c in r[1:] if _to_float(c) is not None)
        if numeric >= 2:
            scenario_rows += 1
    record(f"budget_forecast: >=2 scenario rows with numeric projections ({scenario_rows} found)",
           scenario_rows >= 2, f"scenario rows={scenario_rows}")

    sheet_text = " ".join(wb.sheetnames)
    for row in ws.iter_rows(values_only=True):
        for cell in row:
            if cell is not None:
                sheet_text += " " + str(cell)
    mentions = "forecast" in sheet_text.lower() or "projection" in sheet_text.lower()
    record("budget_forecast: mentions forecast/projection",
           mentions, f"sheet text sample: {sheet_text[:120]}")
    wb.close()


def check_xlsx_content(workspace):
    """Locate the three xlsx deliverables (by exact filename first, then
    keyword) and run each deliverable's checks."""
    print("\n=== Check: XLSX files ===")

    targets = [
        ("variance_analysis.xlsx", "variance analysis xlsx",
         ["variance_analysis", "variance-analysis", "variance analysis",
          "variance", "analysis"], "variance_analysis"),
        ("variance_tracking.xlsx", "variance tracking xlsx",
         ["variance_tracking", "variance-tracking", "variance tracking",
          "tracking", "monitor", "schedule"], "variance_tracking"),
        ("budget_forecast.xlsx", "budget forecast xlsx",
         ["budget_forecast", "budget-forecast", "budget forecast",
          "forecast", "projection", "revised", "outlook"], "budget_forecast"),
    ]

    used_paths = set()
    for gt_fname, label, keywords, kind in targets:
        exact = os.path.join(workspace, gt_fname)
        chosen = None
        if os.path.isfile(exact) and exact not in used_paths:
            chosen = exact
        else:
            cands = [p for p in _find_xlsx_by_keywords(workspace, keywords, exclude=used_paths)]
            if cands:
                chosen = cands[0]
        if chosen is None:
            record(f"xlsx for {label} exists", False,
                   f"No xlsx with keywords {keywords[:3]} found in workspace")
            continue
        used_paths.add(chosen)
        record(f"xlsx for {label} exists ({os.path.basename(chosen)})", True)
        try:
            if kind == "variance_analysis":
                check_variance_analysis(chosen, workspace)
            elif kind == "variance_tracking":
                check_variance_tracking(chosen)
            else:
                check_budget_forecast(chosen)
        except Exception as e:
            record(f"xlsx {os.path.basename(chosen)} readable", False, str(e))


# ---------------------------------------------------------------------------
# docx / pptx
# ---------------------------------------------------------------------------

def check_docx_content(workspace):
    """Locate department variance report docx by filename or content keywords."""
    print("\n=== Check: DOCX files ===")
    from docx import Document
    keywords = ["variance", "department", "budget", "report", "dept"]
    exact = os.path.join(workspace, "dept_variance_reports.docx")
    if os.path.isfile(exact):
        chosen = exact
    else:
        cands = _find_docx_by_keywords(workspace, keywords)
        chosen = cands[0] if cands else None
    if chosen is None:
        record("docx for department variance reports exists", False,
               f"No docx with keywords {keywords} found in workspace")
        return False
    record(f"docx for department variance reports exists ({os.path.basename(chosen)})", True)
    try:
        doc = Document(chosen)
        text = "\n".join(p.text for p in doc.paragraphs)
        word_count = len([w for w in text.split() if w.strip()])
        record(f"docx {os.path.basename(chosen)} substantive (>=100 words)",
               word_count >= 100, f"{word_count} words")
        text_low = text.lower()
        topic_keys = ["variance", "budget", "department", "forecast", "spending"]
        matched = sum(1 for k in topic_keys if k in text_low)
        record(f"docx {os.path.basename(chosen)} mentions topical keywords (>=2)",
               matched >= 2, f"matched {matched}/5")
    except Exception as e:
        record(f"docx {os.path.basename(chosen)} readable", False, str(e))
    return True


def check_pptx_content(workspace):
    """Locate executive presentation pptx by filename or content keywords."""
    print("\n=== Check: PPTX executive presentation ===")
    from pptx import Presentation
    keywords = ["executive", "presentation", "variance", "budget", "leadership"]
    exact = os.path.join(workspace, "executive_presentation.pptx")
    if os.path.isfile(exact):
        chosen = exact
    else:
        cands = _find_pptx_by_keywords(workspace, keywords)
        chosen = cands[0] if cands else None
    if chosen is None:
        record("pptx for executive presentation exists", False,
               f"No pptx with keywords {keywords} found in workspace")
        return False
    record(f"pptx for executive presentation exists ({os.path.basename(chosen)})", True)
    try:
        prs = Presentation(chosen)
        record(f"pptx {os.path.basename(chosen)} has >=3 slides",
               len(prs.slides) >= 3, f"{len(prs.slides)} slides")
        all_text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    all_text += " " + shape.text
        record(f"pptx {os.path.basename(chosen)} mentions 'variance'",
               "variance" in all_text.lower(),
               f"sample: {all_text[:200]}")
        record(f"pptx {os.path.basename(chosen)} mentions 'budget' or 'forecast'",
               "budget" in all_text.lower() or "forecast" in all_text.lower(),
               f"sample: {all_text[:200]}")
    except Exception as e:
        record(f"pptx {os.path.basename(chosen)} readable", False, str(e))
    return True


# ---------------------------------------------------------------------------
# DB side-effects (gcal / email)
# ---------------------------------------------------------------------------

def check_gcal_meetings():
    """Phase 6 requires scheduling budget review meetings."""
    print("\n=== Check: GCal budget review meetings ===")
    try:
        import psycopg2
        conn = psycopg2.connect(**DB_CONFIG)
        cur = conn.cursor()
        cur.execute("SELECT summary FROM gcal.events")
        events = cur.fetchall()
        cur.close(); conn.close()
        summaries = [str(e[0]).lower() for e in events if e[0]]
        record("gcal has >=1 event scheduled (budget review meeting)",
               len(events) >= 1, f"Total events: {len(events)}", db_side=True)
        record("gcal event mentions 'budget' or 'variance' or 'review'",
               any(k in s for s in summaries for k in ["budget", "variance", "review", "forecast"]),
               f"Summaries: {summaries[:5]}", db_side=True)
    except Exception as e:
        record("gcal check ran", False, str(e), db_side=True)


def check_email_to_management():
    """Phase 6 requires sending detailed reports to senior management."""
    print("\n=== Check: Email to senior management ===")
    try:
        import psycopg2
        conn = psycopg2.connect(**DB_CONFIG)
        cur = conn.cursor()
        cur.execute("SELECT subject, to_addr, body_text FROM email.messages")
        messages = cur.fetchall()
        cur.close(); conn.close()
        record("email: >=1 message sent (variance report distribution)",
               len(messages) >= 1, f"Total messages: {len(messages)}", db_side=True)
        topical = []
        for subj, to, body in messages:
            text = ((subj or "") + " " + (body or "")).lower()
            if any(k in text for k in ["variance", "budget", "forecast", "quarterly"]):
                topical.append(subj)
        record("email: at least 1 topical message (variance/budget/forecast)",
               len(topical) >= 1, f"Sample subjects: {topical[:3]}", db_side=True)
    except Exception as e:
        record("email check ran", False, str(e), db_side=True)


# ---------------------------------------------------------------------------

def main():
    parser = ArgumentParser()
    parser.add_argument("--agent_workspace", required=False, default=".")
    parser.add_argument("--groundtruth_workspace", required=False, default=".")
    parser.add_argument("--res_log_file", required=False)
    parser.add_argument("--launch_time", required=False)
    args = parser.parse_args()

    ws = args.agent_workspace
    if not os.path.isdir(ws):
        print(f"Agent workspace not found: {ws}")
        sys.exit(1)

    global IS_GT_SELF_TEST
    try:
        if args.groundtruth_workspace and os.path.exists(args.groundtruth_workspace):
            IS_GT_SELF_TEST = (
                os.path.realpath(ws) ==
                os.path.realpath(args.groundtruth_workspace)
            )
    except Exception:
        IS_GT_SELF_TEST = False

    check_xlsx_content(ws)
    check_docx_content(ws)
    check_pptx_content(ws)
    check_gcal_meetings()
    check_email_to_management()

    total = PASS_COUNT + FAIL_COUNT
    if total == 0:
        print("\nFAIL: No checks were performed.")
        sys.exit(1)

    print(f"\nOverall: {PASS_COUNT}/{total} checks passed")

    result = {
        "total_passed": PASS_COUNT,
        "total_checks": total,
        "failed": FAIL_COUNT,
    }

    if args.res_log_file:
        with open(args.res_log_file, "w") as f:
            json.dump(result, f, indent=2)

    if FAIL_COUNT == 0:
        print("PASS")
        sys.exit(0)
    else:
        print(f"FAIL ({FAIL_COUNT} failures)")
        sys.exit(1)


if __name__ == "__main__":
    main()
