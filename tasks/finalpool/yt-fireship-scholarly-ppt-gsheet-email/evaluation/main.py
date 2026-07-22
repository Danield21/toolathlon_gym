"""
Evaluation for yt-fireship-scholarly-ppt-gsheet-email task.

Checks:
1. AI_Tech_Overview_2025.pptx exists with >= 4 slides and AI-related content
2. GSheet has spreadsheet with Videos sheet with >= 4 rows
3. GSheet has Papers sheet with >= 3 rows
4. GSheet has Technology_Map sheet with >= 3 rows
5. Email sent to research@ai-group.org
6. Email sent to team@lab.edu
"""
import json
import os
import sys
from argparse import ArgumentParser

import psycopg2

DB_CONFIG = {
    "host": os.environ.get("PGHOST", "localhost"),
    "port": 5432,
    "dbname": os.environ.get("PGDATABASE", "toolathlon_gym"),
    "user": "eigent",
    "password": "camel",
}

PASS_COUNT = 0
FAIL_COUNT = 0


def record(name, passed, detail=""):
    global PASS_COUNT, FAIL_COUNT
    if passed:
        PASS_COUNT += 1
        print(f"  [PASS] {name}")
    else:
        FAIL_COUNT += 1
        msg = f": {detail[:300]}" if detail else ""
        print(f"  [FAIL] {name}{msg}")


def num_close(a, b, tol=1.0):
    try: return abs(float(a) - float(b)) <= tol
    except (TypeError, ValueError): return False


def str_match(a, b):
    if a is None or b is None: return a is None and b is None
    return str(a).strip().lower() == str(b).strip().lower()


def check_pptx(agent_workspace):
    print("\n=== Check 1: PPT AI_Tech_Overview_2025.pptx ===")

    pptx_path = os.path.join(agent_workspace, "AI_Tech_Overview_2025.pptx")
    if not os.path.exists(pptx_path):
        record("AI_Tech_Overview_2025.pptx exists", False, f"Not found at {pptx_path}")
        return
    record("AI_Tech_Overview_2025.pptx exists", True)

    try:
        from pptx import Presentation
        prs = Presentation(pptx_path)
        slide_count = len(prs.slides)
        record("PPT has >= 4 slides", slide_count >= 4, f"Found {slide_count} slides")

        all_text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    all_text += " " + shape.text
        all_text_lower = all_text.lower()

        # Require at least 5 distinct AI/ML technical terms across slides (not just generic 'ai')
        ai_terms = [
            "deepseek", "openai", "claude", "grok", "llm",
            "transformer", "language model", "machine learning", "neural",
            "gpt", "bert", "vibe coding", "generative", "openrouter",
            "llama", "anthropic", "cuda", "rag",
        ]
        found_terms = [t for t in ai_terms if t in all_text_lower]
        record("PPT contains >= 5 distinct AI/ML technical terms", len(found_terms) >= 5,
               f"Found terms ({len(found_terms)}): {found_terms}")

        # Slide titles/headings should reference each pillar
        # Slide 2: videos with view counts; slide 3: papers/citations; slide 4: technology landscape; slide 5: insights
        record("PPT mentions 'video' (slide 2 context)", "video" in all_text_lower,
               f"Sample: {all_text_lower[:200]}")
        record("PPT mentions 'paper' (slide 3 context)", "paper" in all_text_lower,
               "Missing 'paper'")
        record("PPT mentions 'technology' (slide 4 context)",
               "technology" in all_text_lower or "technologies" in all_text_lower,
               "Missing 'technology'")
        record("PPT mentions 'insight' or 'recommendation' (slide 5)",
               "insight" in all_text_lower or "recommendation" in all_text_lower,
               "Missing slide 5 content marker")

    except ImportError:
        record("python-pptx available", False, "python-pptx not installed")
    except Exception as e:
        record("PPT readable", False, str(e))


def check_gsheet():
    print("\n=== Check 2: GSheet AI_Research_Overview ===")

    conn = psycopg2.connect(**DB_CONFIG)
    cur = conn.cursor()

    # Find the spreadsheet — require strict workbook title 'AI_Research_Overview'
    cur.execute("SELECT id, title FROM gsheet.spreadsheets")
    spreadsheets = cur.fetchall()

    ai_spreadsheet = None
    for ss_id, title in spreadsheets:
        if title and title.lower().strip() == "ai_research_overview":
            ai_spreadsheet = (ss_id, title)
            break
    # Allow case-insensitive match on the canonical name (no fallback to ANY spreadsheet)
    if not ai_spreadsheet:
        for ss_id, title in spreadsheets:
            if title and "ai_research_overview" in title.lower().strip():
                ai_spreadsheet = (ss_id, title)
                break

    record("GSheet 'AI_Research_Overview' exists", ai_spreadsheet is not None,
           f"Spreadsheets found: {[s[1] for s in spreadsheets]}")

    if not ai_spreadsheet:
        cur.close()
        conn.close()
        return

    ss_id = ai_spreadsheet[0]

    # Check sheets
    cur.execute("SELECT id, title FROM gsheet.sheets WHERE spreadsheet_id = %s", (ss_id,))
    sheets = cur.fetchall()
    sheet_titles_lower = [(s[0], s[1].lower() if s[1] else "") for s in sheets]

    has_videos = any("video" in t for _, t in sheet_titles_lower)
    has_papers = any("paper" in t for _, t in sheet_titles_lower)
    has_tech_map = any("tech" in t or "map" in t or "technology" in t for _, t in sheet_titles_lower)

    record("GSheet has Videos sheet", has_videos, f"Sheets: {[s[1] for s in sheets]}")
    record("GSheet has Papers sheet", has_papers, f"Sheets: {[s[1] for s in sheets]}")
    record("GSheet has Technology_Map sheet", has_tech_map, f"Sheets: {[s[1] for s in sheets]}")

    def count_data_rows(sheet_id):
        cur.execute("""
            SELECT COUNT(DISTINCT row_index) FROM gsheet.cells
            WHERE sheet_id = %s AND row_index > 0
        """, (sheet_id,))
        return cur.fetchone()[0]

    if has_videos:
        video_sheet_id = next(sid for sid, t in sheet_titles_lower if "video" in t)
        video_rows = count_data_rows(video_sheet_id)
        record("Videos sheet has >= 4 data rows", video_rows >= 4,
               f"Found {video_rows} rows")

    if has_papers:
        paper_sheet_id = next(sid for sid, t in sheet_titles_lower if "paper" in t)
        paper_rows = count_data_rows(paper_sheet_id)
        record("Papers sheet has >= 5 data rows", paper_rows >= 5,
               f"Found {paper_rows} rows")
        # Verify at least one well-known paper title appears
        cur.execute("""
            SELECT value FROM gsheet.cells WHERE sheet_id = %s
        """, (paper_sheet_id,))
        cell_values = " ".join((str(r[0]) if r[0] else "") for r in cur.fetchall()).lower()
        # The injected scholar_papers contain titles such as "Attention Is All You Need", "BERT", "GPT", "deepseek"
        known_papers = ["attention is all you need", "bert", "gpt", "deepseek", "llama"]
        found_paper = [p for p in known_papers if p in cell_values]
        record("Papers sheet references >= 1 well-known AI paper",
               len(found_paper) >= 1, f"Found: {found_paper}; sample cells: {cell_values[:200]}")

    if has_tech_map:
        tech_sheet_id = next(sid for sid, t in sheet_titles_lower if "tech" in t or "map" in t or "technology" in t)
        tech_rows = count_data_rows(tech_sheet_id)
        record("Technology_Map sheet has >= 3 data rows", tech_rows >= 3,
               f"Found {tech_rows} rows")

    cur.close()
    conn.close()


def check_xlsx_content(workspace, groundtruth_workspace="."):
    print("\n=== Check: XLSX AI_Research_Overview_local.xlsx ===")
    import openpyxl
    xlsx_path = os.path.join(workspace, "AI_Research_Overview_local.xlsx")
    if not os.path.isfile(xlsx_path):
        record("xlsx AI_Research_Overview_local.xlsx exists", False, f"Not found at {xlsx_path}")
        return False
    record("xlsx AI_Research_Overview_local.xlsx exists", True)
    try:
        wb = openpyxl.load_workbook(xlsx_path, data_only=True)
        for ws in wb.worksheets:
            rows = list(ws.iter_rows(values_only=True))
            record(f"xlsx '{ws.title}' has data", len(rows) >= 2, f"{len(rows)} rows")

        # --- Groundtruth XLSX value comparison ---
        gt_path = os.path.join(groundtruth_workspace, "AI_Research_Overview_local.xlsx")
        if os.path.isfile(gt_path):
            gt_wb = openpyxl.load_workbook(gt_path, data_only=True)
            for gt_sname in gt_wb.sheetnames:
                gt_ws = gt_wb[gt_sname]
                a_ws = None
                for asn in wb.sheetnames:
                    if asn.strip().lower() == gt_sname.strip().lower():
                        a_ws = wb[asn]; break
                if a_ws is None:
                    record(f"GT sheet '{gt_sname}' exists in agent xlsx", False, f"Available: {wb.sheetnames}")
                    continue
                gt_rows = [r for r in gt_ws.iter_rows(min_row=2, values_only=True) if any(c is not None for c in r)]
                a_rows = [r for r in a_ws.iter_rows(min_row=2, values_only=True) if any(c is not None for c in r)]
                record(f"GT '{gt_sname}' row count", len(a_rows) == len(gt_rows),
                       f"Expected {len(gt_rows)}, got {len(a_rows)}")
                for ri in range(min(3, len(gt_rows))):
                    if ri >= len(a_rows): break
                    ok = True
                    for ci in range(min(len(gt_rows[ri]), len(a_rows[ri]))):
                        gv, av = gt_rows[ri][ci], a_rows[ri][ci]
                        if gv is None: continue
                        if isinstance(gv, (int, float)):
                            if not num_close(av, gv, max(abs(gv)*0.1, 1.0)): ok = False; break
                        else:
                            if not str_match(av, gv): ok = False; break
                    record(f"GT '{gt_sname}' row {ri+1} values", ok,
                           f"gt={gt_rows[ri][:4]}, agent={a_rows[ri][:4] if ri < len(a_rows) else 'missing'}")
            gt_wb.close()

        wb.close()
    except Exception as e:
        record("xlsx readable", False, str(e))
    return True


def check_emails():
    print("\n=== Check 3: Emails sent ===")

    conn = psycopg2.connect(**DB_CONFIG)
    cur = conn.cursor()
    cur.execute("SELECT subject, from_addr, to_addr, body_text FROM email.messages")
    messages = cur.fetchall()
    cur.close()
    conn.close()

    def to_addresses(to_addr):
        if isinstance(to_addr, list):
            return " ".join(str(r).lower() for r in to_addr)
        elif to_addr:
            try:
                parsed = json.loads(str(to_addr))
                return " ".join(str(r).lower() for r in parsed) if isinstance(parsed, list) else str(to_addr).lower()
            except Exception:
                return str(to_addr).lower()
        return ""

    to_research = [m for m in messages if "research@ai-group.org" in to_addresses(m[2])]
    to_team = [m for m in messages if "team@lab.edu" in to_addresses(m[2])]

    record("Email sent to research@ai-group.org", len(to_research) >= 1,
           f"Total messages: {len(messages)}")
    record("Email sent to team@lab.edu", len(to_team) >= 1,
           f"Total messages: {len(messages)}")

    if to_research:
        subj, _, _, body = to_research[0]
        content = ((subj or "") + " " + (body or "")).lower()
        # Subject must reference overview/research; body must mention presentation
        record("Research email subject mentions overview/research/AI",
               any(k in (subj or "").lower() for k in ["overview", "research", "ai", "technology"]),
               f"Subject: {subj}")
        record("Research email body mentions 'presentation' (deliverable handoff)",
               "presentation" in content,
               f"Body sample: {(body or '')[:200]}")
        # Specific tech terms required
        tech_count = sum(1 for kw in ["deepseek", "openai", "claude", "llm", "transformer", "machine learning", "language model"] if kw in content)
        record("Research email body mentions >= 2 specific AI tech terms",
               tech_count >= 2,
               f"Found {tech_count} terms in body")

    if to_team:
        subj_t, _, _, body_t = to_team[0]
        content_t = ((subj_t or "") + " " + (body_t or "")).lower()
        record("Team email body mentions 'spreadsheet'",
               "spreadsheet" in content_t or "sheet" in content_t,
               f"Subject: {subj_t}; body: {(body_t or '')[:200]}")


def main():
    parser = ArgumentParser()
    parser.add_argument("--agent_workspace", required=False, default=".")
    parser.add_argument("--groundtruth_workspace", required=False, default=".")
    parser.add_argument("--launch_time", required=False)
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()

    check_pptx(args.agent_workspace)
    # Note: AI_Research_Overview_local.xlsx is not required by task.md — skip local xlsx check.
    check_gsheet()
    check_emails()

    total = PASS_COUNT + FAIL_COUNT
    if total == 0:
        print("\nFAIL: No checks were performed.")
        sys.exit(1)

    print(f"\nOverall: {PASS_COUNT}/{total} checks passed")

    result = {
        "total_passed": PASS_COUNT,
        "total_checks": total,
        "failed": FAIL_COUNT,
    }

    if args.res_log_file:
        with open(args.res_log_file, "w") as f:
            json.dump(result, f, indent=2)

    if FAIL_COUNT == 0:
        print("PASS")
        sys.exit(0)
    else:
        print(f"FAIL ({FAIL_COUNT} failures)")
        sys.exit(1)


if __name__ == "__main__":
    main()
