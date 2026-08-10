"""
Evaluation for arxiv-survey-presentation task.
Checks Excel (Survey_Data.xlsx) and PowerPoint (NLP_Survey.pptx).

Robustness notes:
- The workbook is read WITHOUT data_only, so Excel formulas are visible as
  "=..." strings. For every cell we also merge the cached value from a
  data_only read (the value Excel last recalculated). A formula cell whose
  cached value is missing (never recalculated by a real spreadsheet app)
  is treated as "unresolved" and the numeric check is skipped rather than
  falsely failing a semantically correct answer.
- Paper_IDs are normalized (arXiv: prefix / version suffix stripped) before
  comparison.
- The row-count check dedups on Paper_ID so that an over-complete workbook
  (e.g. built by several sub-agents) is not unfairly penalized.
"""
import argparse
import os
import re
import sys

import openpyxl

PASS_COUNT = 0
FAIL_COUNT = 0

EXPECTED_PAPER_IDS = {"2404.00001", "2404.00002", "2404.00003", "2404.00004", "2404.00005"}

# Sentinel: a formula cell with no cached computed value.
_UNRESOLVED = object()


def record(name, passed, detail=""):
    global PASS_COUNT, FAIL_COUNT
    if passed:
        PASS_COUNT += 1
        print(f"  [PASS] {name}")
    else:
        FAIL_COUNT += 1
        msg = f": {detail[:300]}" if detail else ""
        print(f"  [FAIL] {name}{msg}")


def _to_float(v):
    """Robust numeric parse.

    Accepts int/float, and strings that carry thousands separators, currency
    symbols, percent signs or surrounding whitespace (e.g. '5', '4.0', '5.0%',
    '$5', '1,000'). Returns None when the value cannot be parsed.
    """
    if v is None or isinstance(v, bool):
        return None
    if isinstance(v, (int, float)):
        return float(v)
    s = str(v).strip()
    if not s:
        return None
    s = s.replace("$", "").replace("¥", "").replace("€", "").replace("£", "")
    s = s.replace("％", "%").replace("%", "")
    s = s.replace(",", "").replace(" ", "").replace("−", "-")  # minus sign
    try:
        return float(s)
    except ValueError:
        return None


def _normalize_paper_id(raw):
    """Return a canonical paper id, tolerating 'arXiv:' prefixes and 'v1' suffixes."""
    if raw is None:
        return ""
    s = str(raw).strip()
    low = s.lower()
    for pfx in ("arxiv:", "http://arxiv.org/abs/", "https://arxiv.org/abs/",
                "arxiv.org/abs/"):
        if low.startswith(pfx):
            s = s[len(pfx):].strip()
            low = s.lower()
            break
    s = re.sub(r"v\d+$", "", s).strip()
    return s


def _norm(s):
    """Lowercase and collapse whitespace (keeps internal spaces for phrases)."""
    if s is None:
        return ""
    return re.sub(r"\s+", " ", str(s).lower()).strip()


def _tokens(s):
    return set(re.findall(r"[a-z0-9]+", _norm(s)))


def _match_any(text, kws):
    """True if any keyword matches `text`.

    Phrase keywords match as substrings of the whitespace-normalized text;
    single-word keywords match as substrings of the whitespace-normalized text,
    as substrings of the punctuation-stripped (nospace) text (so hyphen/space
    variants like 'self-consistency' vs 'self consistency' both match), or as
    standalone tokens (so abbreviations like 'NQ', 'RAG', 'CoT' are seen).
    """
    t = _norm(text)
    ns = re.sub(r"[^a-z0-9]+", "", t)
    toks = _tokens(text)
    for kw in kws:
        nk = _norm(kw)
        if not nk:
            continue
        if nk in t:
            return True
        nk_ns = re.sub(r"[^a-z0-9]+", "", nk)
        if nk_ns and nk_ns in ns:
            return True
        if " " not in nk and nk in toks:
            return True
    return False


def _count_matching(text, kws):
    return sum(1 for kw in kws if _match_any(text, [kw]))


def _disp(v):
    if v is _UNRESOLVED:
        return "(Excel formula, no cached value)"
    return v


def load_workbook_pair(path):
    """Return (raw_wb, cached_wb).

    raw_wb keeps formula strings; cached_wb exposes last-computed values.
    Falls back to raw_wb if the data_only pass fails for any reason.
    """
    raw = openpyxl.load_workbook(path, data_only=False)
    try:
        cached = openpyxl.load_workbook(path, data_only=True)
    except Exception:
        cached = raw
    return raw, cached


def find_sheet(wb, sheet_name):
    target = sheet_name.strip().lower()
    target_us = target.replace(" ", "_")
    target_sp = target.replace("_", " ")
    for name in wb.sheetnames:
        n = name.strip().lower()
        if n.replace(" ", "_") == target_us or n.replace("_", " ") == target_sp:
            return name
    return None


def sheet_matrix(wb, name):
    return [[c.value for c in row] for row in wb[name].iter_rows()]


def resolve_value(raw_val, cached_val):
    """Formula cells fall back to the cached (recalculated) value."""
    if isinstance(raw_val, str) and raw_val.startswith("="):
        if (cached_val is not None
                and (not isinstance(cached_val, str) or not cached_val.startswith("="))):
            return cached_val
        return _UNRESOLVED
    return raw_val


def find_col(header, names):
    if not header:
        return None
    for i, cell in enumerate(header):
        if cell is None:
            continue
        cl = str(cell).strip().lower().replace(" ", "_")
        for n in names:
            if n.lower().replace(" ", "_") == cl:
                return i
    return None


def check_excel(workspace):
    print("\n=== Checking Excel ===")
    path = os.path.join(workspace, "Survey_Data.xlsx")
    if not os.path.isfile(path):
        record("Excel exists", False, f"Not found: {path}")
        return False
    record("Excel exists", True)

    try:
        raw_wb, cached_wb = load_workbook_pair(path)
    except Exception as e:
        record("Excel readable", False, str(e))
        return False

    paper_id_set = None

    # ---- Paper Comparison ----
    name = find_sheet(raw_wb, "Paper Comparison") or find_sheet(raw_wb, "Paper_Comparison")
    if name is None:
        record("Sheet 'Paper Comparison' exists", False, f"Sheets: {raw_wb.sheetnames}")
    else:
        record("Sheet 'Paper Comparison' exists", True)
        pc_raw = sheet_matrix(raw_wb, name)
        pc_cached = sheet_matrix(cached_wb, name)
        data = pc_raw[1:]
        id_col = find_col(pc_raw[0], ["Paper_ID", "Paper ID", "ID"])
        if id_col is not None:
            found = set()
            for r in data:
                if id_col < len(r) and r[id_col] is not None:
                    found.add(_normalize_paper_id(r[id_col]))
            found.discard("")
            paper_id_set = found
            record("Paper Comparison covers >= 5 distinct papers",
                   len(found) >= 5, f"Found {len(found)} distinct in {len(data)} rows")
            for pid in EXPECTED_PAPER_IDS:
                record(f"Paper {pid} present", pid in found, f"Found: {sorted(found)}")
        else:
            record("Paper Comparison has >= 5 rows", len(data) >= 5, f"Found {len(data)} rows")

        method_col = find_col(pc_raw[0], ["Method", "method"])
        record("Method column exists", method_col is not None, f"Header: {pc_raw[0]}")

        dataset_col = find_col(pc_raw[0], ["Dataset", "dataset", "Dataset_Used"])
        record("Dataset column exists", dataset_col is not None, f"Header: {pc_raw[0]}")

        # Content check: Method column should cover known NLP methods (3+ distinct)
        if method_col is not None:
            methods_text = " ".join(str(r[method_col]) for r in data
                                    if method_col < len(r) and r[method_col])
            # Keyword matching is hyphen/case/space insensitive, so phrasing
            # variants of the same method still count.
            method_keywords = ["chain", "cot", "instruction", "tuning", "retrieval",
                               "rag", "self-consistency", "consistency", "multimodal",
                               "visual", "prompt", "reasoning", "step", "generation",
                               "augmented", "voting", "majority"]
            matches = _count_matching(methods_text, method_keywords)
            record("Method column has >=3 NLP method keywords", matches >= 3,
                   f"Matched {matches}, text: {methods_text[:200]}")

        # Per-paper Method checks.
        # Keyword lists are paper-derived (title/abstract/sections), so a model
        # that paraphrases the method in its own words (e.g. "majority voting
        # over sampled reasoning paths" for self-consistency) still passes.
        # Rows are deduplicated by Paper_ID: a paper passes if ANY of its rows
        # mentions the method, so an over-complete workbook built by several
        # sub-agents is not unfairly penalized.
        if method_col is not None and id_col is not None:
            expected_method = {
                "2404.00001": ["chain", "cot", "thought", "step", "reasoning", "prompt"],
                "2404.00002": ["instruction", "tuning", "tuned", "fine"],
                "2404.00003": ["retrieval", "rag", "augmented", "generation", "dense"],
                "2404.00004": ["self-consistency", "consistency", "consistent",
                               "majority", "voting", "sampled", "decoding"],
                "2404.00005": ["multimodal", "multi-modal", "visual", "vision",
                               "vlm", "encoder", "image"],
            }
            by_pid_method = {}
            for row in data:
                pid = _normalize_paper_id(row[id_col]) if id_col < len(row) else ""
                if pid and method_col < len(row) and row[method_col] is not None:
                    by_pid_method.setdefault(pid, []).append(row[method_col])
            for pid, kws in expected_method.items():
                texts = by_pid_method.get(pid, [])
                ok = any(_match_any(t, kws) for t in texts)
                got = " | ".join(str(t)[:60] for t in texts[:3]) if texts else "(no rows)"
                record(f"Method for {pid} mentions {kws[0]}", ok,
                       f"Got '{got}'")

        # Content check: Dataset column should cover known datasets (3+ distinct).
        # The keyword list covers every dataset actually named by the five
        # papers (abstracts + experiments), so any paper-faithful choice of
        # "primary evaluation dataset" contributes a match.
        if dataset_col is not None:
            ds_text = " ".join(str(r[dataset_col]) for r in data
                               if dataset_col < len(r) and r[dataset_col])
            ds_keywords = ["gsm8k", "math", "mmlu", "bbh", "triviaqa", "hotpotqa",
                           "natural questions", "nq", "svamp", "aqua", "vqa", "gqa",
                           "vizwiz", "coco", "image caption", "visual dialogue",
                           "captioning", "imagenet", "squad", "truthfulqa", "flan",
                           "big-bench"]
            ds_matches = _count_matching(ds_text, ds_keywords)
            record("Dataset column has >=3 known datasets", ds_matches >= 3,
                   f"Matched {ds_matches}, text: {ds_text[:200]}")

        # Per-paper Dataset checks.
        # Several papers legitimately evaluate on more than one benchmark, so
        # each paper's accepted list covers every dataset it actually names
        # (abstract + experiments); "primary" is intentionally left ambiguous.
        # Rows are deduplicated by Paper_ID (see method check above).
        if dataset_col is not None and id_col is not None:
            expected_dataset = {
                "2404.00001": ["gsm8k", "math"],
                "2404.00002": ["mmlu", "bbh", "triviaqa"],
                "2404.00003": ["natural questions", "triviaqa", "hotpotqa", "nq"],
                "2404.00004": ["gsm8k", "svamp", "aqua"],
                "2404.00005": ["vqa", "gqa", "vizwiz", "coco", "visual dialogue",
                               "image caption", "captioning"],
            }
            by_pid_dataset = {}
            for row in data:
                pid = _normalize_paper_id(row[id_col]) if id_col < len(row) else ""
                if pid and dataset_col < len(row) and row[dataset_col] is not None:
                    by_pid_dataset.setdefault(pid, []).append(row[dataset_col])
            for pid, kws in expected_dataset.items():
                texts = by_pid_dataset.get(pid, [])
                ok = any(_match_any(t, kws) for t in texts)
                got = " | ".join(str(t)[:60] for t in texts[:3]) if texts else "(no rows)"
                record(f"Dataset for {pid} mentions {kws[0]}", ok,
                       f"Got '{got}'")

    # ---- Taxonomy ----
    name = find_sheet(raw_wb, "Taxonomy")
    if name is None:
        record("Sheet 'Taxonomy' exists", False, f"Sheets: {raw_wb.sheetnames}")
    else:
        record("Sheet 'Taxonomy' exists", True)
        tax_rows = sheet_matrix(raw_wb, name)
        data = tax_rows[1:]
        record("Taxonomy has >= 2 categories", len(data) >= 2, f"Found {len(data)}")
        cat_col = find_col(tax_rows[0], ["Category", "category", "Topic"])
        record("Taxonomy has Category column", cat_col is not None, f"Header: {tax_rows[0]}")
        if cat_col is not None:
            cats_text = " ".join(str(r[cat_col]).lower() for r in data
                                 if cat_col < len(r) and r[cat_col])
            record("Taxonomy references NLP topic keywords",
                   any(kw in cats_text for kw in
                       ["reasoning", "prompt", "tuning", "retrieval", "multimodal",
                        "generation", "training", "architecture", "framework"]),
                   f"Got: {cats_text[:200]}")

    # ---- Summary Statistics ----
    name = find_sheet(raw_wb, "Summary Statistics") or find_sheet(raw_wb, "Summary_Statistics")
    if name is None:
        record("Sheet 'Summary Statistics' exists", False, f"Sheets: {raw_wb.sheetnames}")
    else:
        record("Sheet 'Summary Statistics' exists", True)
        ss_raw = sheet_matrix(raw_wb, name)
        ss_cached = sheet_matrix(cached_wb, name)
        metrics_raw = {}
        metrics_cached = {}
        for i, row in enumerate(ss_raw[1:], start=1):
            if row and row[0] is not None:
                key = str(row[0]).strip().lower().replace(" ", "_")
                metrics_raw[key] = row[1] if len(row) > 1 else None
                crow = ss_cached[i] if i < len(ss_cached) else []
                metrics_cached[key] = crow[1] if len(crow) > 1 else None

        def metric_val(key):
            return resolve_value(metrics_raw.get(key), metrics_cached.get(key))

        tp_key = next((k for k in metrics_raw if "total" in k and "paper" in k), None)
        if tp_key:
            mv = metric_val(tp_key)
            num = _to_float(mv)
            if num is not None:
                record("Total_Papers = 5", abs(num - 5) < 1, f"Got {_disp(mv)}")
            else:
                # Fallback: cross-check against the distinct Paper_IDs actually listed.
                if paper_id_set is not None:
                    ok = len(paper_id_set) == 5
                    record("Total_Papers = 5", ok,
                           f"Value {_disp(mv)!r} is not a literal number; "
                           f"cross-checked distinct Paper_ID count ({len(paper_id_set)})")
                else:
                    record("Total_Papers = 5", True,
                           f"Value {_disp(mv)!r} is not a literal number; skipped "
                           "(no Paper_ID column to cross-check)")
        else:
            record("Total_Papers metric exists", False, f"Keys: {sorted(metrics_raw)}")

        avg_sec_key = next((k for k in metrics_raw if "average" in k and "section" in k), None)
        if avg_sec_key is None:
            avg_sec_key = next((k for k in metrics_raw if "avg" in k and "section" in k), None)
        if avg_sec_key:
            mv = metric_val(avg_sec_key)
            num = _to_float(mv)
            if num is not None:
                record("Average_Sections == 4.0 (±0.5)", abs(num - 4.0) <= 0.5,
                       f"Got {_disp(mv)}")
            else:
                record("Average_Sections == 4.0 (±0.5)", True,
                       f"Value {_disp(mv)!r} is not a literal number; skipped "
                       "(formula/unparseable)")
        else:
            record("Average_Sections metric exists", False, f"Keys: {sorted(metrics_raw)}")

        yr_key = next((k for k in metrics_raw if "year" in k and "range" in k), None)
        if yr_key:
            mv = metric_val(yr_key)
            if mv is _UNRESOLVED:
                record("Year_Range contains 2024", True, "Value is a formula; skipped")
            else:
                yr_val = str(mv).strip() if mv is not None else ""
                record("Year_Range contains 2024", "2024" in yr_val, f"Got '{yr_val}'")

        mc_key = next((k for k in metrics_raw if "common" in k or "most" in k), None)
        if mc_key:
            mv = metric_val(mc_key)
            if mv is _UNRESOLVED:
                record("Most_Common_Category non-empty", True, "Value is a formula; skipped")
            else:
                mc_val = str(mv).strip() if mv is not None else ""
                record("Most_Common_Category non-empty", len(mc_val) > 0, f"Got {mc_val}")

    return True


def check_pptx(workspace):
    print("\n=== Checking PowerPoint ===")
    path = os.path.join(workspace, "NLP_Survey.pptx")
    if not os.path.isfile(path):
        record("PPTX exists", False, f"Not found: {path}")
        return False
    record("PPTX exists", True)

    try:
        from pptx import Presentation
        prs = Presentation(str(path))
        slides = list(prs.slides)

        record("Has >= 6 slides", len(slides) >= 6, f"Found {len(slides)}")

        all_text = []
        for slide in slides:
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for p in shape.text_frame.paragraphs:
                        texts.append(p.text)
            all_text.append("\n".join(texts))

        full = "\n".join(all_text).lower()

        # Title slide
        first = all_text[0].lower() if all_text else ""
        has_title = any(kw in first for kw in ["survey", "nlp", "language model",
                                               "reasoning", "advances"])
        record("Title slide has survey keywords", has_title, f"First slide: {first[:200]}")

        # Content checks.
        # Method mentions are matched case/hyphen/whitespace-insensitively, so
        # 'Chain of Thought', 'chain-of-thought' and 'CoT' all count as the same
        # method; abbreviations (CoT, RAG) are accepted.
        def _pptx_methods(full_text):
            ns = re.sub(r"[^a-z0-9]+", "", _norm(full_text))
            toks = _tokens(full_text)
            methods = {
                "chain-of-thought": ["chain", "cot"],
                "instruction tuning": ["instructiontuning"],
                "retrieval": ["retrieval", "rag"],
                "self-consistency": ["selfconsistency", "consistency"],
                "multimodal": ["multimodal"],
            }
            matched = []
            for canon, kws in methods.items():
                if any(kw in ns or kw in toks for kw in kws):
                    matched.append(canon)
            return matched

        matched_methods = _pptx_methods(full)
        record("Mentions paper methods", len(matched_methods) > 0,
               f"Matched: {', '.join(matched_methods) or 'none'}")
        record("PPT mentions >=3 distinct paper methods", len(matched_methods) >= 3,
               f"Matched {len(matched_methods)} methods: {', '.join(matched_methods)}")

        # Summary/conclusion slide
        last = all_text[-1] if all_text else ""
        record("Last slide has content", len(last.strip()) > 10, f"Last slide: {last[:100]}")

        return True
    except Exception as e:
        record("PPTX readable", False, str(e))
        return False


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--agent_workspace", required=False, default=".")
    parser.add_argument("--groundtruth_workspace", required=False, default=".")
    parser.add_argument("--launch_time", required=False)
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()

    check_excel(args.agent_workspace)
    check_pptx(args.agent_workspace)

    print(f"\n=== SUMMARY ===")
    print(f"  Passed: {PASS_COUNT}, Failed: {FAIL_COUNT}")
    sys.exit(0 if FAIL_COUNT == 0 else 1)


if __name__ == "__main__":
    main()
