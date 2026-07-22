"""Evaluation for wc-shipping-zone-ppt."""
import argparse
import os
import sys

import openpyxl
import psycopg2

DB = {"host": os.environ.get("PGHOST", "localhost"), "port": 5432, "dbname": "toolathlon_gym", "user": "eigent", "password": "camel"}

PASS_COUNT = 0
FAIL_COUNT = 0


def record(name, passed, detail=""):
    global PASS_COUNT, FAIL_COUNT
    if passed:
        PASS_COUNT += 1
        print(f"  [PASS] {name}")
    else:
        FAIL_COUNT += 1
        msg = f": {str(detail)[:300]}" if detail else ""
        print(f"  [FAIL] {name}{msg}")


def num_close(a, b, tol=0.5):
    try:
        return abs(float(a) - float(b)) <= tol
    except (TypeError, ValueError):
        return False


def get_expected_zone_data():
    """Compute expected zone performance from read-only DB."""
    conn = psycopg2.connect(**DB)
    cur = conn.cursor()
    cur.execute("""
        SELECT
          CASE
            WHEN shipping->>'state' = 'CA' AND shipping->>'country' = 'US' THEN 'California'
            WHEN shipping->>'country' = 'US' THEN 'Domestic US'
            ELSE 'International'
          END as zone_name,
          COUNT(*) as order_count,
          ROUND(SUM(shipping_total)::numeric, 2)::float as total_shipping,
          ROUND(AVG(shipping_total)::numeric, 2)::float as avg_shipping
        FROM wc.orders
        GROUP BY 1
        ORDER BY COUNT(*) DESC
    """)
    rows = cur.fetchall()
    cur.close()
    conn.close()
    return rows


def check_excel(agent_workspace):
    print("\n=== Checking Shipping_Zone_Report.xlsx ===")
    agent_file = os.path.join(agent_workspace, "Shipping_Zone_Report.xlsx")
    if not os.path.isfile(agent_file):
        record("Excel file exists", False, f"Not found: {agent_file}")
        return
    record("Excel file exists", True)

    try:
        wb = openpyxl.load_workbook(agent_file, data_only=True)
    except Exception as e:
        record("Excel readable", False, str(e))
        return

    expected_zones = get_expected_zone_data()

    # Find Zone Performance sheet (exact preferred)
    zp_sheet = None
    for name in wb.sheetnames:
        if name.strip().lower() == "zone performance":
            zp_sheet = wb[name]; break
    if zp_sheet is None:
        for name in wb.sheetnames:
            if "zone" in name.lower() and "perform" in name.lower():
                zp_sheet = wb[name]; break

    if zp_sheet is None:
        record("Sheet 'Zone Performance' exists", False, f"Sheets: {wb.sheetnames}")
    else:
        record("Sheet 'Zone Performance' exists", True)

        # Header check
        header_row = list(zp_sheet[1])
        a_header = [str(h.value or "").strip().lower().replace(" ", "_") for h in header_row]
        expected_headers = ["zone_name", "order_count", "total_shipping_cost", "avg_shipping_cost"]
        for i, eh in enumerate(expected_headers):
            ah = a_header[i] if i < len(a_header) else "MISSING"
            record(f"Zone Performance header[{i}] = '{eh}'", ah == eh, f"Got '{ah}'")

        rows = [r for r in zp_sheet.iter_rows(min_row=2, values_only=True) if r and r[0] is not None]

        # Row count must equal expected_zones (zones with orders)
        record(f"Zone Performance has {len(expected_zones)} rows",
               len(rows) == len(expected_zones),
               f"Got {len(rows)} rows, expected {len(expected_zones)}")

        a_lookup = {str(r[0]).strip().lower(): r for r in rows}
        for ez in expected_zones:
            zone_name, exp_count, exp_total, exp_avg = ez
            r = a_lookup.get(zone_name.lower())
            if not r:
                record(f"Zone '{zone_name}' present", False, f"Got: {list(a_lookup.keys())}")
                continue
            record(f"Zone '{zone_name}' present", True)
            # Order_Count tol=0
            record(f"'{zone_name}' Order_Count = {exp_count}",
                   num_close(r[1], exp_count, 0),
                   f"Got {r[1]}")
            # Total_Shipping_Cost tol 0.05 (2-decimal)
            record(f"'{zone_name}' Total_Shipping_Cost ~ {float(exp_total)}",
                   num_close(r[2], float(exp_total), 0.05),
                   f"Got {r[2]}")
            # Avg_Shipping_Cost tol 0.02 (2-decimal)
            record(f"'{zone_name}' Avg_Shipping_Cost ~ {float(exp_avg)}",
                   num_close(r[3], float(exp_avg), 0.02),
                   f"Got {r[3]}")

        # Sort: Order_Count descending
        try:
            counts = [int(r[1]) for r in rows]
            record("Zone Performance sorted Order_Count descending",
                   counts == sorted(counts, reverse=True),
                   f"Counts: {counts}")
        except (TypeError, ValueError):
            pass

    # Summary sheet
    sum_sheet = None
    for name in wb.sheetnames:
        if name.strip().lower() == "summary":
            sum_sheet = wb[name]; break
    if sum_sheet is None:
        record("Sheet 'Summary' exists", False, f"Sheets: {wb.sheetnames}")
    else:
        record("Sheet 'Summary' exists", True)
        summary = {}
        for row in sum_sheet.iter_rows(min_row=2, values_only=True):
            if row and row[0]:
                summary[str(row[0]).strip().lower()] = row[1]

        total_orders = sum(int(ez[1]) for ez in expected_zones)
        total_revenue = round(sum(float(ez[2]) for ez in expected_zones), 2)
        zones_count = len(expected_zones)

        record(f"Summary Total_Orders = {total_orders}",
               num_close(summary.get("total_orders"), total_orders, 0),
               f"Got {summary.get('total_orders')}")
        record(f"Summary Total_Shipping_Revenue ~ {total_revenue}",
               num_close(summary.get("total_shipping_revenue"), total_revenue, 0.05),
               f"Got {summary.get('total_shipping_revenue')}")
        record(f"Summary Zones_Count = {zones_count}",
               num_close(summary.get("zones_count"), zones_count, 0),
               f"Got {summary.get('zones_count')}")


def check_pptx(agent_workspace, expected_zones):
    print("\n=== Checking Shipping_Review.pptx ===")
    from pptx import Presentation

    pptx_file = os.path.join(agent_workspace, "Shipping_Review.pptx")
    if not os.path.isfile(pptx_file):
        record("PPTX file exists", False, f"Not found: {pptx_file}")
        return
    record("PPTX file exists", True)

    try:
        prs = Presentation(pptx_file)
    except Exception as e:
        record("PPTX readable", False, str(e))
        return

    slide_count = len(prs.slides)
    record("PPTX has >= 4 slides", slide_count >= 4, f"Got {slide_count}")

    slide_texts = []
    for slide in prs.slides:
        text = ""
        for shape in slide.shapes:
            if shape.has_text_frame:
                text += shape.text_frame.text + "\n"
        slide_texts.append(text)
    full = " ".join(slide_texts).lower()

    # Slide 1: title contains 'Shipping Zone Performance Review' and date 2026-03-06
    if slide_texts:
        s1 = slide_texts[0].lower()
        record("Slide 1 contains 'Shipping Zone Performance Review'",
               "shipping zone performance review" in s1,
               f"Slide 1: {slide_texts[0][:200]}")
        record("Slide 1 contains date 2026-03-06",
               "2026-03-06" in s1 or "2026/03/06" in s1 or "march 6, 2026" in s1.lower(),
               f"Slide 1: {slide_texts[0][:200]}")

    # Per-zone slides
    for zone_name, exp_count, exp_total, exp_avg in expected_zones:
        zone_slide_text = None
        for st in slide_texts:
            if zone_name.lower() in st.lower() and "summary" not in st.lower()[:30]:
                zone_slide_text = st.lower()
                break
        record(f"Slide for zone '{zone_name}' present",
               zone_slide_text is not None,
               "Not found")
        if zone_slide_text:
            # order count present
            record(f"'{zone_name}' slide mentions order count {exp_count}",
                   str(exp_count) in zone_slide_text,
                   f"Looking for '{exp_count}'")
            # total shipping cost
            total_str = f"{float(exp_total):,.2f}"
            record(f"'{zone_name}' slide mentions total cost {total_str}",
                   total_str in zone_slide_text or f"{float(exp_total):.2f}" in zone_slide_text,
                   f"Looking for '{total_str}'")

    # Summary slide with total orders
    total_orders = sum(int(ez[1]) for ez in expected_zones)
    has_summary_slide = any("summary" in st.lower() for st in slide_texts)
    record("Has a 'Summary' slide", has_summary_slide, "No 'summary' keyword found in any slide")
    if has_summary_slide:
        record(f"Summary slide mentions total orders {total_orders}",
               str(total_orders) in full,
               f"Looking for '{total_orders}' in PPTX")


def check_gsheet(expected_zones):
    print("\n=== Checking Google Sheet ===")
    try:
        conn = psycopg2.connect(**DB)
        cur = conn.cursor()
        cur.execute("SELECT id, title FROM gsheet.spreadsheets")
        rows = cur.fetchall()

        target_title = "shipping performance dashboard"
        ss_id = None
        for sid, t in rows:
            if t and target_title in t.strip().lower():
                ss_id = sid
                break

        record(f"GSheet titled '{target_title}'",
               ss_id is not None,
               f"Titles: {[t for _, t in rows]}")
        if ss_id is None:
            cur.close(); conn.close()
            return

        # Sheet 'Zone Data'
        cur.execute("SELECT id, title FROM gsheet.sheets WHERE spreadsheet_id = %s", (ss_id,))
        subs = cur.fetchall()
        zd_id = None
        for sid, t in subs:
            if t and t.strip().lower() == "zone data":
                zd_id = sid
                break
        record("Worksheet 'Zone Data' exists",
               zd_id is not None,
               f"Worksheets: {[t for _, t in subs]}")
        if zd_id is None:
            cur.close(); conn.close()
            return

        cur.execute("""SELECT row_index, col_index, COALESCE(formatted_value, value)
                       FROM gsheet.cells WHERE spreadsheet_id = %s AND sheet_id = %s""",
                    (ss_id, zd_id))
        cells = cur.fetchall()
        cur.close(); conn.close()

        grid = {}
        for r, c, v in cells:
            grid.setdefault(r, {})[c] = v

        # Header row
        header = grid.get(0, {})
        h_cols = [str(header.get(i, "") or "").strip().lower().replace(" ", "_") for i in sorted(header.keys())]
        expected_h = ["zone_name", "order_count", "total_shipping_cost", "avg_shipping_cost"]
        for i, eh in enumerate(expected_h):
            ah = h_cols[i] if i < len(h_cols) else "MISSING"
            record(f"GSheet header[{i}] = '{eh}'",
                   ah == eh,
                   f"Got '{ah}'")

        # Data rows match expected_zones
        data_rows = sorted(r for r in grid if r > 0)
        record(f"GSheet has {len(expected_zones)} data rows",
               len(data_rows) == len(expected_zones),
               f"Got {len(data_rows)}")

        # Per-zone validation
        for ez in expected_zones:
            zone_name, exp_count, exp_total, exp_avg = ez
            found = False
            for r in data_rows:
                rd = grid[r]
                v0 = str(rd.get(0, "") or "").strip().lower()
                if zone_name.lower() == v0:
                    found = True
                    try:
                        n = int(float(str(rd.get(1, ""))))
                        record(f"GSheet '{zone_name}' Order_Count = {exp_count}",
                               n == int(exp_count),
                               f"Got {n}")
                    except (TypeError, ValueError):
                        record(f"GSheet '{zone_name}' Order_Count numeric", False, f"Got {rd.get(1)}")
                    try:
                        t = float(str(rd.get(2, "")))
                        record(f"GSheet '{zone_name}' Total ~ {float(exp_total)}",
                               num_close(t, float(exp_total), 0.05),
                               f"Got {t}")
                    except (TypeError, ValueError):
                        record(f"GSheet '{zone_name}' Total numeric", False, f"Got {rd.get(2)}")
                    break
            if not found:
                record(f"GSheet zone '{zone_name}' present", False, "Missing")
    except Exception as e:
        record("GSheet check error", False, str(e))


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--agent_workspace", required=False, default=".")
    parser.add_argument("--groundtruth_workspace", required=False, default=".")
    parser.add_argument("--launch_time", required=False)
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()

    expected_zones = get_expected_zone_data()

    check_excel(args.agent_workspace)
    check_pptx(args.agent_workspace, expected_zones)
    check_gsheet(expected_zones)

    print(f"\n=== SUMMARY ===")
    print(f"  Passed: {PASS_COUNT}")
    print(f"  Failed: {FAIL_COUNT}")
    overall = FAIL_COUNT == 0
    print(f"  Overall: {'PASS' if overall else 'FAIL'}")
    sys.exit(0 if overall else 1)


if __name__ == "__main__":
    main()
