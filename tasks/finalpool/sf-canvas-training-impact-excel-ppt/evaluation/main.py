"""Evaluation script for sf-canvas-training-impact-excel-ppt."""
import os
import argparse, json, os, sys
import openpyxl

def num_close(a, b, rel_tol=0.15, abs_tol=0.5):
    return abs(float(a) - float(b)) <= max(abs_tol, abs(float(b)) * rel_tol)


DB_CONFIG = {
    "host": os.environ.get("PGHOST", "localhost"), "port": int(os.environ.get("PGPORT", "5432")),
    "dbname": os.environ.get("PGDATABASE", "toolathlon_gym"),
    "user": "eigent", "password": "camel"
}

PASS_COUNT = 0
FAIL_COUNT = 0

def check(name, condition, detail=""):
    global PASS_COUNT, FAIL_COUNT
    if condition:
        PASS_COUNT += 1
        print(f"  [PASS] {name}")
    else:
        FAIL_COUNT += 1
        detail_str = str(detail)[:200] if detail else ""
        print(f"  [FAIL] {name}: {detail_str}")

def safe_float(val, default=None):
    try:
        if val is None: return default
        return float(str(val).replace(",", "").replace("%", "").replace("$", "").strip())
    except (ValueError, TypeError):
        return default

def get_conn():
    import psycopg2
    return psycopg2.connect(**DB_CONFIG)

def run_evaluation(agent_workspace, groundtruth_workspace, launch_time, res_log_file):
    global PASS_COUNT, FAIL_COUNT
    PASS_COUNT = 0
    FAIL_COUNT = 0

    # Pre-load groundtruth Excel once (shared by Excel + PPT sections)
    gt_path = os.path.join(groundtruth_workspace, "Training_Impact_Analysis.xlsx")
    gt_wb = None
    try:
        if os.path.exists(gt_path):
            gt_wb = openpyxl.load_workbook(gt_path)
    except Exception:
        gt_wb = None

    # Check Training_Impact_Analysis.xlsx
    excel_path = os.path.join(agent_workspace, "Training_Impact_Analysis.xlsx")
    check("Training_Impact_Analysis.xlsx exists", os.path.exists(excel_path))
    if os.path.exists(excel_path):
        wb = openpyxl.load_workbook(excel_path)

        if gt_wb:
            for sheet_name in gt_wb.sheetnames:
                check(f"{sheet_name} sheet exists", sheet_name in wb.sheetnames)
                if sheet_name in wb.sheetnames:
                    ws = wb[sheet_name]
                    gt_ws = gt_wb[sheet_name]
                    # Check headers
                    gt_headers = [str(c.value).strip().lower() if c.value else "" for c in gt_ws[1]]
                    headers = [str(c.value).strip().lower() if c.value else "" for c in ws[1]]
                    for h in gt_headers:
                        if h:
                            check(f"{sheet_name} has {h} column", h in headers, f"headers: {headers[:10]}")
                    # Check row count - must equal GT (deterministic data)
                    gt_rows = [r for r in gt_ws.iter_rows(min_row=2, values_only=True)
                               if any(c is not None for c in r)]
                    data_rows = [r for r in ws.iter_rows(min_row=2, values_only=True)
                                 if any(c is not None for c in r)]
                    # Course_Overview is a representative course subset;
                    # accept >= GT row count (agent may legitimately list all
                    # courses from the LMS, e.g. 22 instead of GT's 4).
                    if sheet_name.strip().lower() == "course_overview":
                        check(f"{sheet_name} has >= {len(gt_rows)} data rows",
                              len(data_rows) >= len(gt_rows), f"got {len(data_rows)}")
                    else:
                        check(f"{sheet_name} has exactly {len(gt_rows)} data rows",
                              len(data_rows) == len(gt_rows), f"got {len(data_rows)}")

                    # Cell value comparison against groundtruth - iterate ALL gt rows
                    header_map = {h: i for i, h in enumerate(headers)}
                    # Match by first column key to handle row-order differences
                    gt_keys = {}
                    for r in gt_rows:
                        if r and r[0] is not None:
                            gt_keys[str(r[0]).strip().lower()] = r
                    a_keys = {}
                    for r in data_rows:
                        if r and r[0] is not None:
                            a_keys[str(r[0]).strip().lower()] = r

                    # Determine: which numeric columns in this sheet are
                    # plausibly tied across all GT rows (so accept any single
                    # value from the tied set for "Highest_Performing_Dept" /
                    # "Lowest_Performing_Dept")
                    is_impact_summary = sheet_name.strip().lower() == "impact_summary"
                    # Build allowed-values sets for impact-summary text fields.
                    # Highest_Performing_Dept: any dept tied at max rating.
                    # Lowest_Performing_Dept: any dept tied at min rating.
                    allowed_high_dept = set()
                    allowed_low_dept = set()
                    allowed_popular_course = set()
                    if is_impact_summary and gt_wb and "Department_Performance" in gt_wb.sheetnames:
                        dp_ws = gt_wb["Department_Performance"]
                        dp_rows = list(dp_ws.iter_rows(min_row=2, values_only=True))
                        ratings = []
                        for r in dp_rows:
                            if r and len(r) > 2:
                                rf = safe_float(r[2])
                                if rf is not None:
                                    ratings.append((r[0], rf))
                        if ratings:
                            # Round to 2 decimals to handle agent rounding diff
                            r2 = [(d, round(v, 2)) for d, v in ratings]
                            max_r = max(v for _, v in r2)
                            min_r = min(v for _, v in r2)
                            allowed_high_dept = {str(d).strip().lower()
                                                 for d, v in r2 if v == max_r}
                            allowed_low_dept = {str(d).strip().lower()
                                                for d, v in r2 if v == min_r}
                    # Most popular course tie set: collect any course with max enrollment
                    if is_impact_summary and gt_wb and "Course_Overview" in gt_wb.sheetnames:
                        co_ws = gt_wb["Course_Overview"]
                        co_rows = list(co_ws.iter_rows(min_row=2, values_only=True))
                        enrollments = []
                        for r in co_rows:
                            if r and len(r) > 2:
                                ec = safe_float(r[2])
                                if ec is not None:
                                    enrollments.append((r[0], ec))
                        if enrollments:
                            max_e = max(v for _, v in enrollments)
                            allowed_popular_course = {
                                str(name).strip().lower()
                                for name, v in enrollments if v == max_e
                            }

                    for k, gt_row in gt_keys.items():
                        if k not in a_keys:
                            check(f"{sheet_name} row '{k[:40]}' present", False,
                                  f"missing key {k}")
                            continue
                        agent_row = a_keys[k]
                        for ci, gt_h in enumerate(gt_headers):
                            if not gt_h or ci >= len(gt_row):
                                continue
                            gv = gt_row[ci]
                            agent_ci = header_map.get(gt_h)
                            if agent_ci is None or agent_ci >= len(agent_row):
                                continue
                            av = agent_row[agent_ci]
                            gf = safe_float(gv)
                            af = safe_float(av)
                            if gf is not None and af is not None:
                                # Tighter tolerance: 1% rel, 0.1 abs (was 5%)
                                tol = max(0.1, abs(gf) * 0.01)
                                check(f"{sheet_name} '{k[:30]}' {gt_h} ~{gf}",
                                      abs(gf - af) <= tol, f"got {af}")
                            elif gv is not None and av is not None:
                                gs = str(gv).strip().lower()
                                avs = str(av).strip().lower()
                                if not gs:
                                    continue
                                # For impact_summary VALUE column on
                                # highest_performing_dept / lowest_performing_dept
                                # rows, accept any tied dept (by rating).
                                # For most_popular_course, accept any tied course.
                                is_value_col = (gt_h == "value")
                                handled = False
                                if is_impact_summary and is_value_col:
                                    if "highest_performing_dept" in k and allowed_high_dept:
                                        check(f"{sheet_name} '{k[:30]}' {gt_h} text",
                                              avs in allowed_high_dept,
                                              f"got '{avs[:50]}', allowed {sorted(allowed_high_dept)}")
                                        handled = True
                                    elif "lowest_performing_dept" in k and allowed_low_dept:
                                        check(f"{sheet_name} '{k[:30]}' {gt_h} text",
                                              avs in allowed_low_dept,
                                              f"got '{avs[:50]}', allowed {sorted(allowed_low_dept)}")
                                        handled = True
                                    elif "most_popular_course" in k and allowed_popular_course:
                                        check(f"{sheet_name} '{k[:30]}' {gt_h} text",
                                              avs in allowed_popular_course,
                                              f"got '{avs[:50]}', allowed {sorted(allowed_popular_course)}")
                                        handled = True
                                if not handled:
                                    check(f"{sheet_name} '{k[:30]}' {gt_h} text",
                                          gs == avs,
                                          f"expected '{gs[:50]}', got '{avs[:50]}'")

    # Check Training_Impact_Presentation.pptx
    pptx_path = os.path.join(agent_workspace, "Training_Impact_Presentation.pptx")
    check("Training_Impact_Presentation.pptx exists", os.path.exists(pptx_path))
    if os.path.exists(pptx_path):
        from pptx import Presentation
        prs = Presentation(pptx_path)
        gt_pptx_path = os.path.join(groundtruth_workspace, "Training_Impact_Presentation.pptx")
        if os.path.exists(gt_pptx_path):
            gt_prs = Presentation(gt_pptx_path)
            gt_slide_count = len(gt_prs.slides)
            check(f"Training_Impact_Presentation.pptx has >= {gt_slide_count} slides", len(prs.slides) >= gt_slide_count, f"got {len(prs.slides)} slides")
            # First slide title must contain 'Training Impact Analysis Report' substring
            first_title = ""
            if prs.slides[0].shapes.title:
                first_title = prs.slides[0].shapes.title.text.strip().lower()
            else:
                # fallback: scan first slide for any title text
                for shape in prs.slides[0].shapes:
                    if hasattr(shape, "text") and shape.text:
                        first_title = shape.text.lower()
                        break
            check("First slide titled 'Training Impact Analysis Report'",
                  "training impact analysis report" in first_title or
                  ("training impact" in first_title and "report" in first_title),
                  f"first slide title: {first_title[:80]}")

            # Compare slide titles - use tighter equality
            gt_titles = []
            for s in gt_prs.slides:
                if s.shapes.title:
                    gt_titles.append(s.shapes.title.text.strip().lower())
            agent_titles = []
            for s in prs.slides:
                if s.shapes.title:
                    agent_titles.append(s.shapes.title.text.strip().lower())
            for gt in gt_titles:
                # Require either exact match or that GT title is contained in an agent title
                found = any(gt == at or gt in at for at in agent_titles)
                check(f"Training_Impact_Presentation.pptx has slide \"{gt[:40]}\"", found, f"agent titles: {agent_titles[:5]}")

            # PPT mentions a Highest_Performing_Dept and a Lowest_Performing_Dept.
            # Because all dept ratings tie at 3 in GT, accept any dept name from
            # the GT department list — but require the PPT text references at
            # least 2 distinct department names (not just any one).
            all_ppt_text = []
            for s in prs.slides:
                for shape in s.shapes:
                    if hasattr(shape, "text"):
                        all_ppt_text.append(shape.text.lower())
            ppt_full = " ".join(all_ppt_text)
            gt_dept_set = set()
            if gt_wb and "Department_Performance" in gt_wb.sheetnames:
                for r in gt_wb["Department_Performance"].iter_rows(min_row=2, values_only=True):
                    if r and r[0]:
                        gt_dept_set.add(str(r[0]).strip().lower())
            mentioned_depts = [d for d in gt_dept_set if d in ppt_full]
            check("PPT mentions at least 2 department names",
                  len(mentioned_depts) >= 2,
                  f"depts mentioned: {mentioned_depts[:5]}")
        else:
            check("Training_Impact_Presentation.pptx has >= 3 slides", len(prs.slides) >= 3, f"got {len(prs.slides)} slides")

    # Check Python script exists (terminal usage) - require name 'training_impact.py'
    py_files = [f for f in os.listdir(agent_workspace) if f.endswith(".py")]
    check("Python script 'training_impact.py' exists",
          "training_impact.py" in py_files,
          f"py_files: {py_files}")

    # Check impact_analysis.json exists with valid structure and topical keys
    impact_path = os.path.join(agent_workspace, "impact_analysis.json")
    if os.path.exists(impact_path):
        check("impact_analysis.json exists", True)
        try:
            with open(impact_path) as f:
                impact_data = json.load(f)
            check("impact_analysis.json is valid JSON", True)
            # Inspect a flat dump of all string keys/values to verify the file
            # actually mentions both department + course/training topics.
            blob_text = json.dumps(impact_data).lower()
            has_dept_topic = any(t in blob_text for t in [
                "department", "engineering", "operations", "sales", "hr",
                "finance", "support", "r&d"])
            has_course_topic = any(t in blob_text for t in [
                "course", "training", "enrollment", "curriculum"])
            check("impact_analysis.json mentions a department-related topic",
                  has_dept_topic, blob_text[:300])
            check("impact_analysis.json mentions a course/training topic",
                  has_course_topic, blob_text[:300])
            check("impact_analysis.json has substantive content",
                  bool(impact_data) and (
                      isinstance(impact_data, dict) and len(impact_data) >= 2
                      or isinstance(impact_data, list) and len(impact_data) >= 1
                  ),
                  f"content keys/len: {list(impact_data.keys())[:5] if isinstance(impact_data, dict) else len(impact_data)}")
        except Exception as e:
            check("impact_analysis.json is valid JSON", False, str(e))
    else:
        check("impact_analysis.json exists", False, "file not found")

    return FAIL_COUNT == 0, f"Passed {PASS_COUNT}/{PASS_COUNT + FAIL_COUNT} checks"

def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--agent_workspace", required=False, default=".")
    parser.add_argument("--groundtruth_workspace", required=False, default=".")
    parser.add_argument("--launch_time", required=False, default="2026-03-07 10:00:00")
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()

    success, message = run_evaluation(
        args.agent_workspace, args.groundtruth_workspace,
        args.launch_time, args.res_log_file
    )
    print(message)
    sys.exit(0 if success else 1)

if __name__ == "__main__":
    main()