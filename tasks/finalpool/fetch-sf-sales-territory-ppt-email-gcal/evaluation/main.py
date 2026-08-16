"""Evaluation for fetch-sf-sales-territory-ppt-email-gcal."""
import argparse
import os
import sys

import psycopg2

# ──────────────────────────────────────────────────────────────────────────
# EVALUATION GROUND TRUTH SPEC (gcal tz root-fix v3, case-study 2026-08-13)
# gcal.events.start_datetime is TIMESTAMPTZ; bare r[2].hour silently compares
# wrong in non-UTC PG sessions. Use gcal_helpers.
# ──────────────────────────────────────────────────────────────────────────
# task.md line 19: "March 28, 2026 from 09:00 to 10:30" → SF company PT
EXPECTED_TIMEZONE = "America/Los_Angeles"

from utils.evaluation.gcal_helpers import get_zone_components  # noqa: E402

DB = dict(host=os.environ.get("PGHOST", "localhost"), port=int(os.environ.get("PGPORT", "5432")),
          dbname=os.environ.get("PGDATABASE", "toolathlon_gym"),
          user="eigent", password="camel")


def num_close(a, b, tol=0.5):
    try:
        return abs(float(a) - float(b)) <= tol
    except (TypeError, ValueError):
        return str(a).strip().lower() == str(b).strip().lower()


def check_pptx(agent_workspace):
    errors = []
    path = os.path.join(agent_workspace, "Territory_Scorecard.pptx")
    if not os.path.exists(path):
        return ["Territory_Scorecard.pptx not found"]
    try:
        from pptx import Presentation
        prs = Presentation(path)
        slides = list(prs.slides)
        if len(slides) < 5:
            errors.append(f"Expected 5 slides, found {len(slides)}")

        # Gather all text from the presentation
        all_text = ""
        for slide in slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    all_text += shape.text_frame.text + "\n"
        all_lower = all_text.lower()

        # Check title slide
        if "territory performance" not in all_lower and "q1 2026" not in all_lower:
            errors.append("Title slide missing expected title text")

        # Check revenue data for ALL 5 regions (not just 3)
        for region in ["asia pacific", "europe", "latin america", "middle east", "north america"]:
            if region not in all_lower:
                errors.append(f"Region '{region}' not mentioned in presentation")

        # Check quota attainment values
        if "103.7" not in all_text and "103.6" not in all_text and "103.8" not in all_text:
            errors.append("Asia Pacific attainment ~103.7% not found")
        if "94.7" not in all_text and "94.6" not in all_text and "94.8" not in all_text:
            errors.append("Latin America attainment ~94.7% not found")

        # Check ALL 4 segments
        for seg in ["consumer", "enterprise", "smb", "government"]:
            if seg not in all_lower:
                errors.append(f"Segment '{seg}' not mentioned in presentation")

        # Check pipeline coverage
        if "pipeline" not in all_lower and "coverage" not in all_lower:
            errors.append("Pipeline coverage not discussed")

        # Check recommendations slide
        if "recommendation" not in all_lower:
            errors.append("Recommendations slide not found")

    except Exception as e:
        errors.append(f"Error reading PPTX: {e}")
    return errors


def check_gcal():
    errors = []
    try:
        conn = psycopg2.connect(**DB)
        cur = conn.cursor()
        cur.execute("""
            SELECT summary, description, start_datetime, end_datetime FROM gcal.events
            WHERE start_datetime::date = '2026-03-28'
        """)
        rows = cur.fetchall()
        cur.close()
        conn.close()
        if not rows:
            errors.append("No GCal event found on 2026-03-28")
        else:
            matched = False
            for r in rows:
                sum_lower = (r[0] or "").lower()
                if ("territory" in sum_lower or "review" in sum_lower or "executive" in sum_lower):
                    # Verify 09:00-10:30 window in PT (session-tz-independent).
                    # gcal.events.start_datetime is TIMESTAMPTZ; r[2] is psycopg2
                    # datetime in session tz — use helper to extract PT components.
                    _sd, sh = get_zone_components(r[2], EXPECTED_TIMEZONE)[:2]
                    _sd2, eh, em = get_zone_components(r[3], EXPECTED_TIMEZONE)
                    if sh == 9 and eh == 10 and em == 30:
                        matched = True
                        break
            if not matched:
                errors.append(f"No territory review event at 09:00-10:30 on 2026-03-28 (found: {[r[0] for r in rows]})")
    except Exception as e:
        errors.append(f"Error checking GCal: {e}")
    return errors


def check_email():
    errors = []
    try:
        conn = psycopg2.connect(**DB)
        cur = conn.cursor()
        cur.execute("""
            SELECT subject, body_text FROM email.messages
            WHERE to_addr::text ILIKE '%%executive_team@company.com%%'
            ORDER BY id DESC LIMIT 5
        """)
        rows = cur.fetchall()
        cur.close()
        conn.close()
        if not rows:
            errors.append("No email found to executive_team@company.com")
        else:
            subjects = [r[0].lower() if r[0] else "" for r in rows]
            if not any("territory" in s or "q1" in s or "performance" in s for s in subjects):
                errors.append(f"Email subject doesn't match (found: {[r[0] for r in rows]})")
    except Exception as e:
        errors.append(f"Error checking email: {e}")
    return errors


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--agent_workspace", required=False)
    parser.add_argument("--groundtruth_workspace", required=False)
    parser.add_argument("--launch_time", required=False)
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()
    agent_ws = args.agent_workspace or os.path.join(os.path.dirname(__file__), "..", "groundtruth_workspace")

    all_errors = []

    print("  Checking PowerPoint...")
    errs = check_pptx(agent_ws)
    if errs:
        all_errors.extend(errs)
        for e in errs[:5]:
            print(f"    ERROR: {e}")
    else:
        print("    PASS")

    print("  Checking GCal event...")
    errs = check_gcal()
    if errs:
        all_errors.extend(errs)
        for e in errs[:3]:
            print(f"    ERROR: {e}")
    else:
        print("    PASS")

    print("  Checking email...")
    errs = check_email()
    if errs:
        all_errors.extend(errs)
        for e in errs[:3]:
            print(f"    ERROR: {e}")
    else:
        print("    PASS")

    if all_errors:
        print(f"\n=== RESULT: FAIL ({len(all_errors)} errors) ===")
        for e in all_errors[:10]:
            print(f"  {e}")
        sys.exit(1)
    else:
        print("\n=== RESULT: PASS ===")
        sys.exit(0)


if __name__ == "__main__":
    main()
