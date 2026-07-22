"""Evaluation script for yf-sf-market-sales-excel-ppt-gcal."""
import os
import argparse, json, os, sys
import openpyxl

def num_close(a, b, rel_tol=0.15, abs_tol=0.5):
    return abs(float(a) - float(b)) <= max(abs_tol, abs(float(b)) * rel_tol)


DB_CONFIG = {
    "host": os.environ.get("PGHOST", "localhost"), "port": 5432,
    "dbname": os.environ.get("PGDATABASE", "toolathlon_gym"),
    "user": "eigent", "password": "camel"
}

PASS_COUNT = 0
FAIL_COUNT = 0

def check(name, condition, detail=""):
    global PASS_COUNT, FAIL_COUNT
    if condition:
        PASS_COUNT += 1
        print(f"  [PASS] {name}")
    else:
        FAIL_COUNT += 1
        detail_str = str(detail)[:200] if detail else ""
        print(f"  [FAIL] {name}: {detail_str}")

def safe_float(val, default=None):
    try:
        if val is None: return default
        return float(str(val).replace(",", "").replace("%", "").replace("$", "").strip())
    except (ValueError, TypeError):
        return default

def get_conn():
    import psycopg2
    return psycopg2.connect(**DB_CONFIG)

def run_evaluation(agent_workspace, groundtruth_workspace, launch_time, res_log_file):
    global PASS_COUNT, FAIL_COUNT
    PASS_COUNT = 0
    FAIL_COUNT = 0

    # Check Market_Sales_Correlation.xlsx
    excel_path = os.path.join(agent_workspace, "Market_Sales_Correlation.xlsx")
    check("Market_Sales_Correlation.xlsx exists", os.path.exists(excel_path))
    if os.path.exists(excel_path):
        wb = openpyxl.load_workbook(excel_path)
        gt_path = os.path.join(groundtruth_workspace, "Market_Sales_Correlation.xlsx")
        gt_wb = openpyxl.load_workbook(gt_path) if os.path.exists(gt_path) else None

        if gt_wb:
            for sheet_name in gt_wb.sheetnames:
                check(f"{sheet_name} sheet exists", sheet_name in wb.sheetnames)
                if sheet_name in wb.sheetnames:
                    ws = wb[sheet_name]
                    gt_ws = gt_wb[sheet_name]
                    # Check headers
                    gt_headers = [str(c.value).strip().lower() if c.value else "" for c in gt_ws[1]]
                    headers = [str(c.value).strip().lower() if c.value else "" for c in ws[1]]
                    for h in gt_headers:
                        if h:
                            check(f"{sheet_name} has {h} column", h in headers, f"headers: {headers[:10]}")
                    # Check row count
                    gt_rows = list(gt_ws.iter_rows(min_row=2, values_only=True))
                    data_rows = list(ws.iter_rows(min_row=2, values_only=True))
                    min_rows = max(1, len(gt_rows) - 2)
                    check(f"{sheet_name} has >= {min_rows} data rows", len(data_rows) >= min_rows, f"got {len(data_rows)}")

                    # For Stock_Overview, validate alphabetical sort by Symbol
                    if sheet_name.lower() == "stock_overview":
                        non_empty = [r for r in data_rows if r and r[0]]
                        symbols = [str(r[0]).strip().upper() for r in non_empty]
                        check(
                            f"{sheet_name} sorted alphabetically by Symbol",
                            symbols == sorted(symbols),
                            f"order: {symbols}",
                        )

                    # Cell value comparison against groundtruth — validate ALL rows
                    header_map = {h: i for i, h in enumerate(headers)}
                    gt_header_map = {h: i for i, h in enumerate(gt_headers)}

                    # Special handling: Key_Insights uses free-form 'Finding' text
                    # which the agent cannot reproduce verbatim. Validate structure
                    # and substantive content rather than per-row PK match.
                    if sheet_name.lower() == "key_insights":
                        non_empty_rows = [
                            r for r in data_rows
                            if r and any(c is not None and str(c).strip() for c in r)
                        ]
                        check(
                            f"{sheet_name} has at least 4 insights",
                            len(non_empty_rows) >= 4,
                            f"got {len(non_empty_rows)} non-empty rows",
                        )
                        # All required columns must have non-empty content for each row
                        for ri, row in enumerate(non_empty_rows[:8]):
                            for ci, gt_h in enumerate(gt_headers):
                                if not gt_h:
                                    continue
                                agent_ci = header_map.get(gt_h)
                                if agent_ci is None or agent_ci >= len(row):
                                    continue
                                cell = row[agent_ci]
                                check(
                                    f"{sheet_name} row {ri+1} {gt_h} non-empty",
                                    cell is not None and len(str(cell).strip()) > 0,
                                    f"row {ri+1} {gt_h} is empty",
                                )
                        # Check that insights mention key concepts
                        all_text = " ".join(
                            str(c).lower() for r in non_empty_rows for c in r
                            if c is not None
                        )
                        check(
                            f"{sheet_name} insights mention 'correlat' concept",
                            "correlat" in all_text,
                            f"text excerpt: {all_text[:150]}",
                        )
                        # Skip the per-row PK lookup for Key_Insights
                        continue

                    # Special handling: Sales_Trends pulls from a live DB whose
                    # exact month range is not pinned by the task. Validate
                    # structure (sorted Month column, positive Total_Revenue,
                    # plausible Order_Count and Avg_Order_Value) and a minimum
                    # row count, but do not enforce exact PK match for every GT
                    # month (agent may pick a different but valid window).
                    if sheet_name.lower() == "sales_trends":
                        non_empty_rows = [
                            r for r in data_rows
                            if r and any(c is not None and str(c).strip() for c in r)
                        ]
                        check(
                            f"{sheet_name} has at least 6 monthly rows",
                            len(non_empty_rows) >= 6,
                            f"got {len(non_empty_rows)} rows",
                        )
                        month_idx = header_map.get("month")
                        rev_idx = header_map.get("total_revenue")
                        ord_idx = header_map.get("order_count")
                        aov_idx = header_map.get("avg_order_value")
                        if month_idx is not None:
                            months = []
                            for r in non_empty_rows:
                                if month_idx < len(r) and r[month_idx] is not None:
                                    months.append(str(r[month_idx]).strip())
                            check(
                                f"{sheet_name} months sorted ascending",
                                months == sorted(months),
                                f"order: {months[:6]}",
                            )
                        if rev_idx is not None:
                            revs = [safe_float(r[rev_idx]) for r in non_empty_rows
                                    if rev_idx < len(r)]
                            revs = [v for v in revs if v is not None]
                            ok_rev = revs and all(v > 0 for v in revs)
                            check(
                                f"{sheet_name} Total_Revenue positive",
                                bool(ok_rev),
                                f"sample: {revs[:3]}",
                            )
                        if ord_idx is not None:
                            ords = [safe_float(r[ord_idx]) for r in non_empty_rows
                                    if ord_idx < len(r)]
                            ords = [v for v in ords if v is not None]
                            ok_ord = ords and all(v > 0 for v in ords)
                            check(
                                f"{sheet_name} Order_Count positive",
                                bool(ok_ord),
                                f"sample: {ords[:3]}",
                            )
                        if aov_idx is not None:
                            aovs = [safe_float(r[aov_idx]) for r in non_empty_rows
                                    if aov_idx < len(r)]
                            aovs = [v for v in aovs if v is not None]
                            ok_aov = aovs and all(v > 0 for v in aovs)
                            check(
                                f"{sheet_name} Avg_Order_Value positive",
                                bool(ok_aov),
                                f"sample: {aovs[:3]}",
                            )
                        continue

                    # Special handling: Correlation_Matrix is a computed analysis
                    # whose exact values depend on the agent's chosen methodology
                    # (e.g. price returns vs. revenue trends, window length).
                    # Validate structure (PK presence + diagonal=1.0 + value range)
                    # rather than exact numerical equality.
                    if sheet_name.lower() == "correlation_matrix":
                        # All GT row PKs must be present in agent rows
                        pk_col_cm = gt_headers[0] if gt_headers else None
                        if pk_col_cm and pk_col_cm in header_map:
                            gt_pk_idx_cm = gt_header_map[pk_col_cm]
                            agent_pk_idx_cm = header_map[pk_col_cm]
                            agent_pks = set()
                            for r in data_rows:
                                if r and len(r) > agent_pk_idx_cm and r[agent_pk_idx_cm] is not None:
                                    agent_pks.add(str(r[agent_pk_idx_cm]).strip().lower())
                            for gt_row in gt_rows:
                                if not gt_row or len(gt_row) <= gt_pk_idx_cm or gt_row[gt_pk_idx_cm] is None:
                                    continue
                                pk_val = str(gt_row[gt_pk_idx_cm]).strip().lower()
                                check(
                                    f"{sheet_name} row pk={pk_val} present",
                                    pk_val in agent_pks,
                                    f"missing symbol {pk_val}",
                                )
                            # Validate values lie in [-1.0, 1.0]; diagonal cells (Symbol == column header) ~ 1.0
                            for r in data_rows:
                                if not r or r[agent_pk_idx_cm] is None:
                                    continue
                                row_pk = str(r[agent_pk_idx_cm]).strip().upper()
                                for ci, h in enumerate(headers):
                                    if ci == agent_pk_idx_cm or not h:
                                        continue
                                    cv = safe_float(r[ci]) if ci < len(r) else None
                                    if cv is None:
                                        continue
                                    in_range = -1.05 <= cv <= 1.05
                                    check(
                                        f"{sheet_name} value at ({row_pk}, {h.upper()}) in [-1, 1]",
                                        in_range,
                                        f"value: {cv}",
                                    )
                                    if h.upper() == row_pk:
                                        # Diagonal: a symbol's correlation with itself should be ~1.0
                                        check(
                                            f"{sheet_name} diagonal ({row_pk}) ~ 1.0",
                                            abs(cv - 1.0) <= 0.1,
                                            f"diagonal={cv}",
                                        )
                        continue

                    # Match by primary key (first non-empty column) when possible.
                    pk_col = gt_headers[0] if gt_headers else None
                    if pk_col and pk_col in header_map:
                        gt_pk_idx = gt_header_map[pk_col]
                        agent_pk_idx = header_map[pk_col]
                        agent_lookup = {}
                        for r in data_rows:
                            if r and len(r) > agent_pk_idx and r[agent_pk_idx] is not None:
                                agent_lookup[str(r[agent_pk_idx]).strip().lower()] = r
                        for gt_row in gt_rows:
                            if not gt_row or len(gt_row) <= gt_pk_idx or gt_row[gt_pk_idx] is None:
                                continue
                            pk_val = str(gt_row[gt_pk_idx]).strip().lower()
                            agent_row = agent_lookup.get(pk_val)
                            check(f"{sheet_name} row pk={pk_val} present", agent_row is not None)
                            if agent_row is None:
                                continue
                            for ci, gt_h in enumerate(gt_headers):
                                if not gt_h or ci >= len(gt_row):
                                    continue
                                gv = gt_row[ci]
                                agent_ci = header_map.get(gt_h)
                                if agent_ci is None or agent_ci >= len(agent_row):
                                    continue
                                av = agent_row[agent_ci]
                                gf = safe_float(gv)
                                af = safe_float(av)
                                if gf is not None and af is not None:
                                    # Tighter tolerance for read-only data
                                    tol = max(0.05, abs(gf) * 0.05)
                                    check(f"{sheet_name} pk={pk_val} {gt_h} ~{gf}",
                                          abs(gf - af) <= tol, f"got {af}")
                                elif gv is not None and av is not None:
                                    gs = str(gv).strip().lower()
                                    avs = str(av).strip().lower()
                                    if gs:
                                        check(f"{sheet_name} pk={pk_val} {gt_h} text",
                                              gs == avs,
                                              f"expected {gs[:50]}, got {avs[:50]}")

    # Check Market_Sales_Strategy.pptx
    pptx_path = os.path.join(agent_workspace, "Market_Sales_Strategy.pptx")
    check("Market_Sales_Strategy.pptx exists", os.path.exists(pptx_path))
    if os.path.exists(pptx_path):
        from pptx import Presentation
        prs = Presentation(pptx_path)
        gt_pptx_path = os.path.join(groundtruth_workspace, "Market_Sales_Strategy.pptx")
        if os.path.exists(gt_pptx_path):
            gt_prs = Presentation(gt_pptx_path)
            gt_slide_count = len(gt_prs.slides)
            check(f"Market_Sales_Strategy.pptx has >= {gt_slide_count} slides", len(prs.slides) >= gt_slide_count, f"got {len(prs.slides)} slides")
            # Compare slide titles
            gt_titles = []
            for s in gt_prs.slides:
                if s.shapes.title:
                    gt_titles.append(s.shapes.title.text.strip().lower())
            agent_titles = []
            for s in prs.slides:
                if s.shapes.title:
                    agent_titles.append(s.shapes.title.text.strip().lower())
            for gt in gt_titles:
                # Substring would let "High" match "Highest"; use exact match (still case-insensitive).
                found = any(gt == at for at in agent_titles)
                check(f"Market_Sales_Strategy.pptx has slide \"{gt[:40]}\"", found, f"agent titles: {agent_titles[:5]}")
        else:
            check("Market_Sales_Strategy.pptx has >= 3 slides", len(prs.slides) >= 3, f"got {len(prs.slides)} slides")

    # Check the specific Python script exists (terminal usage)
    py_files = [f for f in os.listdir(agent_workspace) if f.endswith(".py")]
    check(
        "market_sales_correlator.py exists",
        "market_sales_correlator.py" in py_files,
        f"found: {py_files}",
    )
    # Also check correlation_analysis.json output exists
    check(
        "correlation_analysis.json exists",
        os.path.exists(os.path.join(agent_workspace, "correlation_analysis.json")),
        "missing correlation analysis output",
    )

    # Database checks
    try:
        conn = get_conn()
        cur = conn.cursor()
        # Exact title match for the strategy review event
        cur.execute(
            """
            SELECT summary, start_datetime, end_datetime, description
            FROM gcal.events
            WHERE LOWER(summary) = LOWER('Market-Sales Strategy Review')
            """
        )
        event_row = cur.fetchone()
        check(
            "Calendar event 'Market-Sales Strategy Review' exists",
            event_row is not None,
            "no event with exact title found",
        )
        if event_row:
            summary, start_dt, end_dt, description = event_row
            # Validate start datetime: 2026-03-16 11:00 UTC
            import datetime as _dt

            def _to_utc_naive(d):
                if d is None:
                    return None
                if hasattr(d, "tzinfo") and d.tzinfo is not None:
                    return d.astimezone(_dt.timezone.utc).replace(tzinfo=None)
                return d

            start_utc = _to_utc_naive(start_dt)
            end_utc = _to_utc_naive(end_dt)
            expected_start = _dt.datetime(2026, 3, 16, 11, 0, 0)
            expected_end = _dt.datetime(2026, 3, 16, 12, 30, 0)
            check(
                "Calendar event start matches 2026-03-16 11:00 UTC",
                start_utc is not None and abs((start_utc - expected_start).total_seconds()) <= 60,
                f"start_datetime: {start_utc}",
            )
            check(
                "Calendar event end matches 2026-03-16 12:30 UTC",
                end_utc is not None and abs((end_utc - expected_end).total_seconds()) <= 60,
                f"end_datetime: {end_utc}",
            )
            # Description should mention correlation findings
            desc_lower = (description or "").lower()
            check(
                "Calendar event description mentions correlation findings",
                "correlation" in desc_lower,
                f"description: {desc_lower[:100]}",
            )
        # Reverse verification: noise events should not match task keyword
        cur.execute("SELECT COUNT(*) FROM gcal.events WHERE summary ILIKE '%standup%' OR summary ILIKE '%lunch%'")
        noise_events = cur.fetchone()[0]
        check("Noise events exist (not deleted by agent)", noise_events >= 1, f"noise events: {noise_events}")
        conn.close()
    except Exception as e:
        check("DB checks", False, str(e))

    return FAIL_COUNT == 0, f"Passed {PASS_COUNT}/{PASS_COUNT + FAIL_COUNT} checks"

def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--agent_workspace", required=False, default=".")
    parser.add_argument("--groundtruth_workspace", required=False, default=".")
    parser.add_argument("--launch_time", required=False, default="2026-03-07 10:00:00")
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()

    success, message = run_evaluation(
        args.agent_workspace, args.groundtruth_workspace,
        args.launch_time, args.res_log_file
    )
    print(message)
    sys.exit(0 if success else 1)

if __name__ == "__main__":
    main()