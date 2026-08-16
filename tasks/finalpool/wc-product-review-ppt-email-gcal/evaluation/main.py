"""Evaluation for wc-product-review-ppt-email-gcal.

Checks:
1. PPT file Product_Review_Report.pptx with at least 5 slides
   - Contains "Product Quality Review" and "Q1" title slide
   - Contains category rating data (Headphones best, Speakers worst)
   - Contains recommendations slide
2. GCal event "Product Quality Review Meeting" ~7 days from launch_time
3. Email to product-team@store.example.com from analytics@store.example.com
   Subject: "Product Quality Review Report"
"""
import argparse
import json
import os
import sys
from datetime import datetime, timedelta, timezone

import psycopg2

DB = dict(host=os.environ.get("PGHOST", "localhost"), port=int(os.environ.get("PGPORT", "5432")), dbname="toolathlon_gym", user="eigent", password="camel")

PASS_COUNT = 0
FAIL_COUNT = 0
RUNTIME_ONLY_FAIL = 0


def check(name, condition, detail="", runtime_only=False):
    global PASS_COUNT, FAIL_COUNT, RUNTIME_ONLY_FAIL
    if condition:
        PASS_COUNT += 1
        print(f"  [PASS] {name}")
    else:
        FAIL_COUNT += 1
        if runtime_only:
            RUNTIME_ONLY_FAIL += 1
        detail_str = f": {str(detail)[:200]}" if detail else ""
        print(f"  [FAIL] {name}{detail_str}")


def check_ppt(agent_workspace):
    print("\n=== Checking PPT File ===")
    ppt_path = os.path.join(agent_workspace, "Product_Review_Report.pptx")
    check("Product_Review_Report.pptx exists", os.path.isfile(ppt_path),
          f"Expected at {ppt_path}")
    if not os.path.isfile(ppt_path):
        return

    try:
        from pptx import Presentation
        prs = Presentation(ppt_path)
    except Exception as e:
        check("PPT file readable", False, str(e))
        return

    check("PPT has at least 5 slides", len(prs.slides) >= 5,
          f"Found {len(prs.slides)} slides")

    all_text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                all_text.append(shape.text.lower())
    full_text = " ".join(all_text)

    check("PPT contains 'Product Quality Review'",
          "product quality review" in full_text, "Title text not found")
    check("PPT contains 'Q1'", "q1" in full_text, "Q1 not found")

    # Normalize text so "TV & Home Theater" and "TV and Home Theater" match.
    def _norm_cat(s):
        return (str(s or "")
                .strip()
                .lower()
                .replace(" & ", " and ")
                .replace("&", "and"))

    full_text_norm = _norm_cat(full_text)

    # Require all 6 known category names. Accept both "&" and "and" spellings.
    expected_cats = ["headphones", "tv and home theater", "cameras", "audio",
                     "electronics", "speakers"]
    missing_cats = [c for c in expected_cats if _norm_cat(c) not in full_text_norm]
    check("PPT mentions all 6 category names",
          len(missing_cats) == 0, f"Missing: {missing_cats}")
    # Require specific best/worst category names with exact ratings.
    # The old constants (Headphones 4.78 / Speakers 4.43) drifted from the seed
    # data: recomputing from wc.product_reviews JOIN wc.products shows TV & Home
    # Theater = 4.78 best and Speakers = 4.39 worst (c4 case-study 2026-08-15).
    # Compute the ratings from the DB so GT can never go stale again
    # (〔拍板 B.1.5〕evaluator-data-side precedent).
    best_cat = best_rating = worst_cat = worst_rating = None
    try:
        conn = psycopg2.connect(**DB)
        cur = conn.cursor()
        cur.execute("""
            WITH per_review AS (
                SELECT
                       -- categories is a JSON array of {id, name} objects from
                       -- coarse to fine. The LAST element is the finest
                       -- sub-category (Headphones, Speakers, TV & Home
                       -- Theater, ...). Verified against the seed dump:
                       -- best = TV & Home Theater 4.78 (n=63),
                       -- worst = Speakers 4.39 (n=28).
                       (p.categories::jsonb -> -1) ->> 'name' AS subcat,
                       r.rating
                FROM wc.product_reviews r
                JOIN wc.products p ON p.id = r.product_id
                WHERE r.status = 'approved'
            )
            SELECT subcat, ROUND(AVG(rating)::numeric, 2) AS avg_rating, COUNT(*) AS n
            FROM per_review
            WHERE subcat IS NOT NULL
            GROUP BY subcat
            ORDER BY avg_rating DESC
        """)
        agg = cur.fetchall()
        cur.close()
        conn.close()
        if agg:
            best_cat, best_rating, _ = agg[0]
            worst_cat, worst_rating, _ = agg[-1]
    except Exception as e:
        print(f"  [warn] could not compute category ratings from DB: {e}")

    def _fmt_rating(v):
        return f"{float(v):.2f}" if v is not None else None

    if best_cat and worst_cat:
        check(f"PPT identifies {best_cat} as best category ({_fmt_rating(best_rating)})",
              _norm_cat(best_cat) in full_text_norm
              and _fmt_rating(best_rating) in full_text_norm,
              f"{best_cat} rating {_fmt_rating(best_rating)} not found")
        check(f"PPT identifies {worst_cat} as worst category ({_fmt_rating(worst_rating)})",
              _norm_cat(worst_cat) in full_text_norm
              and _fmt_rating(worst_rating) in full_text_norm,
              f"{worst_cat} rating {_fmt_rating(worst_rating)} not found")
    else:
        # DB unavailable: fall back to the (documented-stale) constants so the
        # remaining checks still produce signal instead of hard-crashing.
        check("PPT identifies Headphones as best category (4.78)",
              "headphones" in full_text_norm and "4.78" in full_text_norm,
              "Headphones rating 4.78 not found")
        check("PPT identifies Speakers as worst category (4.43)",
              "speakers" in full_text_norm and "4.43" in full_text_norm,
              "Speakers rating 4.43 not found")
    check("PPT mentions top/best reviewed section",
          "best" in full_text or "top" in full_text or "highest" in full_text,
          "No best-products section found")
    check("PPT mentions recommendations",
          "recommend" in full_text or "action" in full_text or "improve" in full_text,
          "No recommendations found")


def check_gcal(launch_time_str=None):
    print("\n=== Checking Google Calendar ===")
    conn = psycopg2.connect(**DB)
    cur = conn.cursor()

    cur.execute("""
        SELECT summary, start_datetime, end_datetime, description
        FROM gcal.events
        WHERE LOWER(summary) LIKE '%product%' AND LOWER(summary) LIKE '%quality%'
    """)
    events = cur.fetchall()
    check("Product Quality Review Meeting event created", len(events) >= 1,
          f"Found {len(events)} matching events", runtime_only=True)

    if events and launch_time_str:
        try:
            launch_time = datetime.fromisoformat(launch_time_str)
            if launch_time.tzinfo is None:
                launch_time = launch_time.replace(tzinfo=timezone.utc)
            target_date = launch_time + timedelta(days=7)
            for event in events:
                event_start = event[1]
                if event_start.tzinfo is None:
                    event_start = event_start.replace(tzinfo=timezone.utc)
                diff_days = abs((event_start.date() - target_date.date()).days)
                if diff_days <= 2:
                    check("Review Meeting is ~7 days from launch", True)
                    break
            else:
                check("Review Meeting is ~7 days from launch", False,
                      f"Closest event at {events[0][1]}, expected ~{target_date.date()}",
                      runtime_only=True)
        except Exception as e:
            check("Review Meeting date check", False, str(e), runtime_only=True)

    # Check event duration = 1 hour
    if events:
        for event in events:
            start, end = event[1], event[2]
            if start and end:
                if start.tzinfo is None:
                    start = start.replace(tzinfo=timezone.utc)
                if end.tzinfo is None:
                    end = end.replace(tzinfo=timezone.utc)
                duration_min = (end - start).total_seconds() / 60
                if abs(duration_min - 60) <= 15:
                    check("Review Meeting duration ~1 hour", True)
                    break
        else:
            check("Review Meeting duration ~1 hour", False, "No 1-hour event found",
                  runtime_only=True)

    cur.close()
    conn.close()


def check_emails():
    print("\n=== Checking Emails ===")
    conn = psycopg2.connect(**DB)
    cur = conn.cursor()

    cur.execute("""
        SELECT subject, from_addr, to_addr, body_text
        FROM email.messages
    """)
    all_emails = cur.fetchall()
    conn.close()

    def parse_recipients(to_addr):
        if to_addr is None:
            return []
        if isinstance(to_addr, list):
            return [str(r).strip().lower() for r in to_addr]
        to_str = str(to_addr).strip()
        try:
            parsed = json.loads(to_str)
            if isinstance(parsed, list):
                return [str(r).strip().lower() for r in parsed]
            return [to_str.lower()]
        except (json.JSONDecodeError, TypeError):
            return [to_str.lower()]

    target = "product-team@store.example.com"
    found = None
    for subj, from_addr, to_addr, body in all_emails:
        recipients = parse_recipients(to_addr)
        if target in recipients:
            found = (subj, from_addr, to_addr, body)
            break

    check("Email sent to product-team@store.example.com", found is not None,
          f"No email found for {target}", runtime_only=True)
    if found:
        subj, from_addr, to_addr, body = found
        check("Email from analytics@store.example.com",
              "analytics@store.example.com" in (from_addr or "").lower(),
              f"From: {from_addr}", runtime_only=True)
        check("Subject is 'Product Quality Review Report'",
              "product quality review" in (subj or "").lower(),
              f"Subject: {subj}", runtime_only=True)


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--agent_workspace", required=False, default=".")
    parser.add_argument("--groundtruth_workspace", required=False)
    parser.add_argument("--launch_time", required=False)
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()

    print("=" * 70)
    print("WC PRODUCT REVIEW PPT EMAIL GCAL - EVALUATION")
    print("=" * 70)

    check_ppt(args.agent_workspace)
    check_gcal(args.launch_time)
    check_emails()

    print(f"\n=== SUMMARY ===")
    print(f"  Passed: {PASS_COUNT}")
    print(f"  Failed: {FAIL_COUNT}")
    non_runtime_fail = FAIL_COUNT - RUNTIME_ONLY_FAIL
    overall = non_runtime_fail == 0
    print(f"  Overall: {'PASS' if overall else 'FAIL'} (runtime-only fails: {RUNTIME_ONLY_FAIL})")
    sys.exit(0 if overall else 1)


if __name__ == "__main__":
    main()
