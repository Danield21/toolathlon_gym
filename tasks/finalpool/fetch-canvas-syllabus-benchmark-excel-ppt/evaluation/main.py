"""
Evaluation script for fetch-canvas-syllabus-benchmark-excel-ppt task.

Checks:
1. Course_Benchmark_Analysis.xlsx with 3 sheets and correct data
2. Academic_Benchmark_Presentation.pptx with required slides
3. Email sent with benchmark report
"""

import argparse
import json
import os
import sys

import openpyxl
import psycopg2

DB_CONFIG = {
    "host": os.environ.get("PGHOST", "localhost"),
    "port": int(os.environ.get("PGPORT", "5432")),
    "dbname": os.environ.get("PGDATABASE", "toolathlon_gym"),
    "user": "eigent",
    "password": "camel",
}

PASS_COUNT = 0
FAIL_COUNT = 0


def record(name, passed, detail=""):
    global PASS_COUNT, FAIL_COUNT
    if passed:
        PASS_COUNT += 1
        print(f"  [PASS] {name}")
    else:
        FAIL_COUNT += 1
        msg = f": {detail[:300]}" if detail else ""
        print(f"  [FAIL] {name}{msg}")


def num_close(a, b, tol=50.0):
    try:
        return abs(float(a) - float(b)) <= tol
    except (TypeError, ValueError):
        return False


def str_contains(haystack, needle):
    if haystack is None or needle is None:
        return False
    return needle.strip().lower() in str(haystack).strip().lower()


# Expected our averages from Canvas data (validated against GT file).
# Per course type (lower-cased key): enrollment, assignments, quizzes.
OUR_AVG = {
    "applied analytics & algorithms":     {"enr": 374,  "asg": 6,    "quz": 0},
    "biochemistry & bioinformatics":      {"enr": 1977, "asg": 10.5, "quz": 3.8},
    "creative computing & culture":       {"enr": 2217, "asg": 10,   "quz": 4},
    "data-driven design":                 {"enr": 1568, "asg": 8.8,  "quz": 1.8},
    "environmental economics & ethics":   {"enr": 978,  "asg": 5,    "quz": 0},
    "foundations of finance":             {"enr": 1941, "asg": 13,   "quz": 7},
    "global governance & geopolitics":    {"enr": 845,  "asg": 10,   "quz": 6},
}

NATIONAL_AVG = {
    "applied analytics & algorithms":     {"enr": 250,  "asg": 8,  "quz": 4},
    "biochemistry & bioinformatics":      {"enr": 1500, "asg": 12, "quz": 6},
    "creative computing & culture":       {"enr": 1800, "asg": 8,  "quz": 3},
    "data-driven design":                 {"enr": 1200, "asg": 10, "quz": 5},
    "environmental economics & ethics":   {"enr": 800,  "asg": 7,  "quz": 3},
    "foundations of finance":             {"enr": 1600, "asg": 11, "quz": 5},
    "global governance & geopolitics":    {"enr": 600,  "asg": 9,  "quz": 4},
}

ALL_COURSE_TYPES = set(OUR_AVG.keys())


def check_excel(agent_workspace):
    """Check Course_Benchmark_Analysis.xlsx."""
    print("\n=== Checking Excel Output ===")

    fpath = os.path.join(agent_workspace, "Course_Benchmark_Analysis.xlsx")
    if not os.path.isfile(fpath):
        record("Excel file exists", False, f"Not found: {fpath}")
        return False

    record("Excel file exists", True)

    try:
        wb = openpyxl.load_workbook(fpath, data_only=True)
    except Exception as e:
        record("Excel file readable", False, str(e))
        return False

    all_ok = True

    # --- Sheet 1: National Benchmarks ---
    bench_sheet = None
    for name in wb.sheetnames:
        if "benchmark" in name.lower() or "national" in name.lower():
            bench_sheet = name
            break
    if not bench_sheet:
        record("National Benchmarks sheet exists", False, f"Sheets: {wb.sheetnames}")
        all_ok = False
    else:
        record("National Benchmarks sheet exists", True)
        ws = wb[bench_sheet]
        rows = list(ws.iter_rows(values_only=True))
        data_rows = rows[1:] if len(rows) > 1 else []
        ok = len(data_rows) == 7
        record("Benchmarks has 7 rows", ok, f"Found {len(data_rows)}")
        if not ok:
            all_ok = False

    # --- Sheet 2: Our Courses ---
    our_sheet = None
    for name in wb.sheetnames:
        if "our" in name.lower() or "course" in name.lower():
            if "benchmark" not in name.lower() and "comparison" not in name.lower():
                our_sheet = name
                break
    if not our_sheet:
        record("Our Courses sheet exists", False, f"Sheets: {wb.sheetnames}")
        all_ok = False
    else:
        record("Our Courses sheet exists", True)
        ws = wb[our_sheet]
        rows = list(ws.iter_rows(values_only=True))
        data_rows = rows[1:] if len(rows) > 1 else []
        ok = len(data_rows) == 22
        record("Our Courses has 22 rows", ok, f"Found {len(data_rows)}")
        if not ok:
            all_ok = False

    # --- Sheet 3: Comparison ---
    comp_sheet = None
    for name in wb.sheetnames:
        if "comparison" in name.lower() or "compare" in name.lower():
            comp_sheet = name
            break
    if not comp_sheet:
        record("Comparison sheet exists", False, f"Sheets: {wb.sheetnames}")
        all_ok = False
    else:
        record("Comparison sheet exists", True)
        ws = wb[comp_sheet]
        rows = list(ws.iter_rows(values_only=True))
        data_rows = rows[1:] if len(rows) > 1 else []
        ok = len(data_rows) == 7
        record("Comparison has 7 rows", ok, f"Found {len(data_rows)}")
        if not ok:
            all_ok = False

        # Verify ALL 7 course types present with correct our_avg and national_avg
        # Comparison cols layout (header):
        #   0: Course_Type
        #   1: Our_Avg_Enrollment
        #   2: National_Avg_Enrollment
        #   3: Enrollment_Diff
        #   4: Our_Avg_Assignments
        #   5: National_Avg_Assignments
        #   6: Assignment_Diff
        #   7: Our_Avg_Quizzes
        #   8: National_Avg_Quizzes
        #   9: Quiz_Diff
        types_seen = set()
        for row in data_rows:
            if not row or not row[0]:
                continue
            ct = str(row[0]).strip().lower()
            if ct not in OUR_AVG:
                continue
            types_seen.add(ct)
            our = OUR_AVG[ct]
            nat = NATIONAL_AVG[ct]
            # Enrollment: tolerance 50
            e_our_ok = num_close(row[1] if len(row) > 1 else None, our["enr"], tol=50)
            e_nat_ok = num_close(row[2] if len(row) > 2 else None, nat["enr"], tol=50)
            record(f"Comparison[{ct[:30]}] Our_Avg_Enrollment~={our['enr']}", e_our_ok,
                   f"row: {row[:4]}")
            record(f"Comparison[{ct[:30]}] National_Avg_Enrollment~={nat['enr']}", e_nat_ok,
                   f"row: {row[:4]}")
            if not (e_our_ok and e_nat_ok):
                all_ok = False
            # Assignments: tolerance 0.5 because avg can be fractional
            a_our_ok = num_close(row[4] if len(row) > 4 else None, our["asg"], tol=0.5)
            a_nat_ok = num_close(row[5] if len(row) > 5 else None, nat["asg"], tol=0.5)
            record(f"Comparison[{ct[:30]}] Our_Avg_Assignments~={our['asg']}", a_our_ok,
                   f"got {row[4] if len(row) > 4 else None}")
            record(f"Comparison[{ct[:30]}] National_Avg_Assignments~={nat['asg']}", a_nat_ok,
                   f"got {row[5] if len(row) > 5 else None}")
            if not (a_our_ok and a_nat_ok):
                all_ok = False
            # Quizzes: tolerance 0.5
            q_our_ok = num_close(row[7] if len(row) > 7 else None, our["quz"], tol=0.5)
            q_nat_ok = num_close(row[8] if len(row) > 8 else None, nat["quz"], tol=0.5)
            record(f"Comparison[{ct[:30]}] Our_Avg_Quizzes~={our['quz']}", q_our_ok,
                   f"got {row[7] if len(row) > 7 else None}")
            record(f"Comparison[{ct[:30]}] National_Avg_Quizzes~={nat['quz']}", q_nat_ok,
                   f"got {row[8] if len(row) > 8 else None}")
            if not (q_our_ok and q_nat_ok):
                all_ok = False
        missing = ALL_COURSE_TYPES - types_seen
        record("Comparison covers all 7 course types", not missing,
               f"Missing: {missing}")
        if missing:
            all_ok = False

    wb.close()
    return all_ok


def check_pptx(agent_workspace):
    """Check Academic_Benchmark_Presentation.pptx."""
    print("\n=== Checking PowerPoint Output ===")

    fpath = os.path.join(agent_workspace, "Academic_Benchmark_Presentation.pptx")
    if not os.path.isfile(fpath):
        record("PPT file exists", False, f"Not found: {fpath}")
        return False

    record("PPT file exists", True)

    try:
        from pptx import Presentation
        prs = Presentation(fpath)
    except Exception as e:
        record("PPT file readable", False, str(e))
        return False

    all_ok = True

    # Should have title + 7 course types + summary = 9 slides minimum
    slide_count = len(prs.slides)
    ok = slide_count >= 9
    record(f"PPT has >= 9 slides", ok, f"Found {slide_count}")
    if not ok:
        all_ok = False

    # Check title slide
    first_slide = prs.slides[0]
    title_text = ""
    for shape in first_slide.shapes:
        if shape.has_text_frame:
            title_text += shape.text_frame.text.lower()
    ok = "benchmark" in title_text or "course" in title_text
    record("Title slide has benchmark/course", ok, f"Title: {title_text[:100]}")
    if not ok:
        all_ok = False

    # Check that at least 3 course types appear in slide titles
    course_types_found = set()
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.lower()
                for ct in ["finance", "biochemistry", "analytics", "data-driven",
                           "environmental", "computing", "governance"]:
                    if ct in text:
                        course_types_found.add(ct)

    # Task requires all 7 course-type slides
    ok = len(course_types_found) == 7
    record(f"PPT covers all 7 course types", ok,
           f"Found {len(course_types_found)}: {course_types_found}")
    if not ok:
        all_ok = False

    return all_ok


def check_email():
    """Check email was sent with benchmark report."""
    print("\n=== Checking Email ===")

    try:
        conn = psycopg2.connect(**DB_CONFIG)
        cur = conn.cursor()
        cur.execute("SELECT subject, from_addr, to_addr, body_text FROM email.messages")
        emails = cur.fetchall()
        cur.close()
        conn.close()
    except Exception as e:
        record("Email DB accessible", False, str(e))
        return False

    all_ok = True
    found_email = False

    for subject, from_addr, to_addr, body_text in emails:
        subj_lower = (subject or "").lower().strip()
        # Exact required subject per task.md: "Annual Course Benchmark Report"
        if "annual course benchmark report" in subj_lower or (
            "annual" in subj_lower and "benchmark" in subj_lower and ("course" in subj_lower or "report" in subj_lower)):
            found_email = True
            record("Benchmark report email exists with correct subject", True,
                   f"Subject: {subject}")

            # Sender address is fixed by the email environment and not checked
            to_str = str(to_addr).lower()
            to_ok = "dean" in to_str
            record("Email to dean", to_ok, f"To: {to_addr}")
            if not to_ok:
                all_ok = False

            body_lower = (body_text or "").lower()
            body_ok = ("enrollment" in body_lower or "benchmark" in body_lower or
                       "above" in body_lower or "below" in body_lower)
            record("Email body discusses benchmarks", body_ok,
                   f"Body preview: {(body_text or '')[:200]}")
            if not body_ok:
                all_ok = False
            break

    if not found_email:
        record("Benchmark report email exists", False,
               f"Found {len(emails)} emails, none with benchmark/course in subject")
        all_ok = False

    return all_ok


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--agent_workspace", required=False, default=".")
    parser.add_argument("--groundtruth_workspace", required=False, default=".")
    parser.add_argument("--launch_time", required=False)
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()

    excel_ok = check_excel(args.agent_workspace)
    pptx_ok = check_pptx(args.agent_workspace)
    email_ok = check_email()

    print(f"\n=== SUMMARY ===")
    print(f"  Excel:  {'PASS' if excel_ok else 'FAIL'}")
    print(f"  PPT:    {'PASS' if pptx_ok else 'FAIL'}")
    print(f"  Email:  {'PASS' if email_ok else 'FAIL'}")
    print(f"  Passed: {PASS_COUNT}, Failed: {FAIL_COUNT}")

    overall = excel_ok and pptx_ok and email_ok
    print(f"  Overall: {'PASS' if overall else 'FAIL'}")

    sys.exit(0 if overall else 1)


if __name__ == "__main__":
    main()
