"""Evaluation for sf-sales-segment-ppt-gform."""
import argparse
import os
import re
import sys

import psycopg2

DB_CONFIG = {
    "host": os.environ.get("PGHOST", "localhost"),
    "port": int(os.environ.get("PGPORT", "5432")),
    "dbname": "toolathlon_gym",
    "user": "eigent",
    "password": "camel",
}


def _fetch_expected_segments():
    """Query DB to get all segment revenue/count."""
    try:
        conn = psycopg2.connect(**DB_CONFIG)
        cur = conn.cursor()
        cur.execute("""
            SELECT "SEGMENT", COUNT(*), ROUND(SUM("TOTAL_AMOUNT")::numeric, 2)
            FROM sf_data."SALES_DW__PUBLIC__ORDERS" o
            JOIN sf_data."SALES_DW__PUBLIC__CUSTOMERS" c ON o."CUSTOMER_ID" = c."CUSTOMER_ID"
            GROUP BY "SEGMENT"
            ORDER BY SUM("TOTAL_AMOUNT") DESC
        """)
        rows = cur.fetchall()
        cur.close()
        conn.close()
        # [(segment, order_count, revenue), ...] sorted by revenue desc
        return [(r[0], int(r[1]), float(r[2])) for r in rows]
    except Exception as e:
        print(f"[WARN] segment query: {e}")
        return None


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--agent_workspace", required=False)
    parser.add_argument("--groundtruth_workspace", required=False)
    parser.add_argument("--launch_time", required=False)
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()

    task_root = os.path.abspath(os.path.join(os.path.dirname(__file__), ".."))
    agent_ws = args.agent_workspace or task_root

    all_errors = []

    # Fetch expected segments dynamically
    expected_segments = _fetch_expected_segments()
    expected_segment_names = [s[0] for s in expected_segments] if expected_segments else []
    top_segment = expected_segments[0] if expected_segments else None
    most_orders_segment = max(expected_segments, key=lambda s: s[1])[0] if expected_segments else None
    total_revenue = round(sum(s[2] for s in expected_segments), 2) if expected_segments else 0.0

    # --- Check 1: PowerPoint ---
    print("Checking PowerPoint presentation...")
    pptx_path = os.path.join(agent_ws, "Segment_Performance.pptx")
    if not os.path.exists(pptx_path):
        all_errors.append("Segment_Performance.pptx not found in agent workspace")
    else:
        from pptx import Presentation
        prs = Presentation(pptx_path)

        if len(prs.slides) < 4:
            all_errors.append(f"PPT has only {len(prs.slides)} slides, expected at least 4")

        # Per-slide text
        slide_texts = []
        for slide in prs.slides:
            stext = ""
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    stext += " " + shape.text
            slide_texts.append(stext)
        all_text = " ".join(slide_texts)

        # Validate ALL 4 segment names from DB
        if expected_segment_names:
            for seg in expected_segment_names:
                if seg not in all_text:
                    all_errors.append(f"PPT missing segment '{seg}'")
        else:
            # Fallback hardcoded
            for seg in ["Consumer", "Enterprise", "Government", "SMB"]:
                if seg not in all_text:
                    all_errors.append(f"PPT missing segment '{seg}'")

        # Validate top revenue figure (dynamic)
        if top_segment:
            rev_int = int(top_segment[2])
            rev_str = f"{rev_int:,}"  # "839,609"
            rev_str_no_comma = str(rev_int)
            rev_decimal = f"{top_segment[2]:.2f}"
            if not any(x in all_text for x in [rev_str, rev_str_no_comma, rev_decimal]):
                all_errors.append(f"PPT missing top-segment revenue figure ({rev_decimal})")

        # Slide titles
        title_slide_text = slide_texts[0] if len(slide_texts) > 0 else ""
        if "Customer Segment" not in title_slide_text and "Segment Performance" not in title_slide_text:
            all_errors.append("First slide missing 'Customer Segment Performance Report' title")

        # Slide 2: Revenue by Segment
        if len(slide_texts) > 1:
            if "Revenue by Segment" not in slide_texts[1]:
                all_errors.append("Slide 2 missing title 'Revenue by Segment'")

        # Slide 3: Key Insights
        if len(slide_texts) > 2:
            if "Key Insights" not in slide_texts[2] and "Insights" not in slide_texts[2]:
                all_errors.append("Slide 3 missing title 'Key Insights'")

        # Slide 4: Strategy Recommendations (require literal phrase, not just 'Recommendations')
        if len(slide_texts) > 3:
            if "Strategy Recommendations" not in slide_texts[3]:
                all_errors.append("Slide 4 missing literal title 'Strategy Recommendations'")

        print(f"    PPT checks done ({len(prs.slides)} slides)")

    # --- Check 2: GForm ---
    print("Checking Google Form...")
    try:
        conn = psycopg2.connect(**DB_CONFIG)
        cur = conn.cursor()
        cur.execute("SELECT id, title FROM gform.forms WHERE LOWER(title) LIKE '%sales strategy%'")
        forms = cur.fetchall()
        if not forms:
            all_errors.append("Google Form 'Sales Strategy Feedback' not found in gform.forms")
        else:
            form_id = forms[0][0]
            # Find questions table - check both gform.questions and gform.items
            try:
                cur.execute("""
                    SELECT title, type, options
                    FROM gform.items
                    WHERE form_id = %s
                """, (form_id,))
                questions = cur.fetchall()
            except Exception:
                conn.rollback()
                try:
                    cur.execute("""
                        SELECT title, type, options
                        FROM gform.questions
                        WHERE form_id = %s
                    """, (form_id,))
                    questions = cur.fetchall()
                except Exception:
                    questions = []
            if len(questions) < 3:
                all_errors.append(f"GForm has only {len(questions)} questions, expected >=3")
            # Verify investment question lists 4 segment options
            invest_q = None
            for q in questions:
                qt = (q[0] or "").lower()
                if "investment" in qt or "invest" in qt or "priority" in qt:
                    invest_q = q
                    break
            if invest_q and expected_segment_names:
                opts_str = str(invest_q[2] or "").lower()
                missing = [s for s in expected_segment_names if s.lower() not in opts_str]
                if len(missing) > 1:
                    all_errors.append(f"GForm investment question missing segments: {missing}")
            print(f"    GForm found with {len(questions)} questions")
        cur.close()
        conn.close()
    except Exception as e:
        all_errors.append(f"Error checking GForm: {e}")

    # --- Check 3: Email ---
    print("Checking email...")
    try:
        conn = psycopg2.connect(**DB_CONFIG)
        cur = conn.cursor()
        cur.execute("""
            SELECT subject, body_text, to_addr
            FROM email.messages
            WHERE to_addr::text ILIKE '%sales_director@company.com%'
        """)
        emails = cur.fetchall()
        if not emails:
            all_errors.append("No email sent to sales_director@company.com")
        else:
            # Find email matching segment-related subject
            target = None
            for subj, body, to in emails:
                s = (subj or "").lower()
                if "segment" in s or "sales" in s:
                    target = (subj, body, to)
                    break
            if target is None:
                target = (emails[0][0], emails[0][1], emails[0][2])
            subj, body, to = target
            if not subj or "segment" not in (subj or "").lower():
                all_errors.append(f"Email subject does not reference segment performance: '{subj}'")
            body_lower = (body or "").lower()
            # Body mentions top revenue segment as a word-bounded token
            if top_segment:
                seg_pat = re.compile(r"\b" + re.escape(top_segment[0].lower()) + r"\b")
                if not seg_pat.search(body_lower):
                    all_errors.append(f"Email body missing top revenue segment '{top_segment[0]}' as standalone token")
            # Body mentions form created
            if "form" not in body_lower:
                all_errors.append("Email body missing mention of feedback form")
            print(f"    Email found ({len(emails)} messages)")
        cur.close()
        conn.close()
    except Exception as e:
        all_errors.append(f"Error checking email: {e}")

    # --- Final result ---
    if all_errors:
        print(f"\n=== RESULT: FAIL ({len(all_errors)} errors) ===")
        for e in all_errors[:10]:
            print(f"  {e}")
        sys.exit(1)
    else:
        print("\n=== RESULT: PASS ===")
        sys.exit(0)


if __name__ == "__main__":
    main()
