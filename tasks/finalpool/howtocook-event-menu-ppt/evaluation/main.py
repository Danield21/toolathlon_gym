"""
Evaluation script for howtocook-event-menu-ppt task.

Checks:
1. PPT exists with >=5 slides, title slide has "dinner" or "menu", course types present
2. Excel exists with correct sheet structure
Since HowToCook data is dynamic, we evaluate structure not specific values.
"""

import argparse
import json
import os
import sys


def check_pptx(workspace):
    """Check PowerPoint file structure and content."""
    from pptx import Presentation

    errors = []
    pptx_path = os.path.join(workspace, "Event_Menu_Presentation.pptx")
    if not os.path.exists(pptx_path):
        return False, "Event_Menu_Presentation.pptx not found"

    prs = Presentation(pptx_path)
    slides = list(prs.slides)

    if len(slides) < 5:
        return False, f"PPT has {len(slides)} slides, expected at least 5"

    # Per-slide text
    slide_texts = []
    for slide in slides:
        t = ""
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    t += paragraph.text + " "
            elif hasattr(shape, "text"):
                t += shape.text + " "
        slide_texts.append(t)
    slide_lower = [t.lower() for t in slide_texts]
    all_text = " ".join(slide_lower)

    # Title slide must contain exact "Company Team Dinner Menu"
    if "company team dinner menu" not in slide_lower[0]:
        errors.append(f"Title slide missing exact phrase 'Company Team Dinner Menu': got {slide_texts[0][:100]}")

    # Each course type appears in a dedicated slide (not just somewhere)
    for course in ["appetizer", "main", "dessert"]:
        slide_with_course = [i for i, t in enumerate(slide_lower[1:-1], 1) if course in t]
        if not slide_with_course:
            errors.append(f"No course slide for '{course}'")

    # Each course slide should have ingredients content
    has_ingredients = any("ingredient" in t for t in slide_lower[1:-1])
    if not has_ingredients:
        errors.append("No course slide mentions ingredients")

    # Summary slide (last) should mention total ingredients count
    summary_text = slide_lower[-1]
    import re
    # Look for a number in the summary slide
    if not re.search(r"\b\d+\b", summary_text):
        errors.append("Summary slide has no numeric total")
    if not any(kw in summary_text for kw in ["total", "summary", "menu summary"]):
        errors.append("Summary slide missing 'total' or 'summary' keyword")

    if errors:
        return False, "; ".join(errors)
    return True, "PPT check passed"


def check_excel(workspace):
    """Check Excel file structure."""
    from openpyxl import load_workbook

    errors = []
    xlsx_path = os.path.join(workspace, "Menu_Budget.xlsx")
    if not os.path.exists(xlsx_path):
        return False, "Menu_Budget.xlsx not found"

    wb = load_workbook(xlsx_path)
    sheet_names = [s.lower() for s in wb.sheetnames]

    if "menu items" not in sheet_names:
        return False, f"Missing 'Menu Items' sheet. Found: {wb.sheetnames}"

    if "summary" not in sheet_names:
        return False, f"Missing 'Summary' sheet. Found: {wb.sheetnames}"

    # Check Menu Items sheet has data
    menu_sheet = wb[wb.sheetnames[sheet_names.index("menu items")]]
    headers = [str(cell.value).lower() if cell.value else "" for cell in menu_sheet[1]]

    required_headers = ["course", "recipe_name", "ingredients_count", "difficulty"]
    for rh in required_headers:
        if not any(rh.replace("_", " ") in h or rh in h for h in headers):
            errors.append(f"Menu Items sheet missing header: {rh}")

    # Check exactly 3 data rows (one per course)
    data_rows = [row for row in menu_sheet.iter_rows(min_row=2, values_only=True) if row[0] is not None]
    if len(data_rows) != 3:
        errors.append(f"Menu Items sheet has {len(data_rows)} data rows, expected exactly 3")

    # Check all 3 courses appear in Course column (col 0)
    courses_found = set()
    for row in data_rows:
        if row[0]:
            c = str(row[0]).strip().lower()
            for course in ["appetizer", "main", "dessert"]:
                if course in c:
                    courses_found.add(course)
    missing_courses = {"appetizer", "main", "dessert"} - courses_found
    if missing_courses:
        errors.append(f"Menu Items missing course(s): {missing_courses}")

    # Difficulty must be 1/2/3
    for row in data_rows:
        if len(row) > 3 and row[3] is not None:
            try:
                d = int(row[3])
                if d not in (1, 2, 3):
                    errors.append(f"Difficulty {d} not in {{1,2,3}} for {row[1]}")
            except (TypeError, ValueError):
                errors.append(f"Difficulty not numeric for {row[1]}: {row[3]}")

    # Compute Total_Ingredients from Menu Items (col 2)
    total_ingredients = 0
    for row in data_rows:
        if len(row) > 2 and row[2] is not None:
            try:
                total_ingredients += int(row[2])
            except (TypeError, ValueError):
                pass

    # Check Summary sheet
    summary_sheet = wb[wb.sheetnames[sheet_names.index("summary")]]
    summary_rows = list(summary_sheet.iter_rows(values_only=True))
    summary_map = {}
    for row in summary_rows[1:]:
        if row[0] is not None:
            summary_map[str(row[0]).strip().lower()] = row[1] if len(row) > 1 else None

    # Total_Recipes == 3 (exact)
    tr = summary_map.get("total_recipes") or summary_map.get("total recipes")
    try:
        if int(tr) != 3:
            errors.append(f"Total_Recipes={tr}, expected 3")
    except (TypeError, ValueError):
        errors.append(f"Total_Recipes missing or non-numeric: {tr}")

    # Total_Ingredients == sum of Menu Items column
    ti = summary_map.get("total_ingredients") or summary_map.get("total ingredients")
    try:
        if int(ti) != total_ingredients:
            errors.append(f"Total_Ingredients={ti} does not equal sum of Menu Items = {total_ingredients}")
    except (TypeError, ValueError):
        errors.append(f"Total_Ingredients missing or non-numeric: {ti}")

    # Avg_Difficulty present + numeric
    ad = summary_map.get("avg_difficulty") or summary_map.get("average_difficulty")
    if ad is None:
        errors.append("Avg_Difficulty missing from Summary")
    else:
        try:
            float(ad)
        except (TypeError, ValueError):
            errors.append(f"Avg_Difficulty not numeric: {ad}")

    if errors:
        return False, "; ".join(errors)
    return True, "Excel check passed"


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--agent_workspace", required=False, default=".")
    parser.add_argument("--groundtruth_workspace", required=False, default=".")
    parser.add_argument("--launch_time", required=False)
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()

    errors = []

    # Check PPT
    print("\n=== Checking PowerPoint ===")
    pptx_ok, pptx_msg = check_pptx(args.agent_workspace)
    if pptx_ok:
        print(f"  [PASS] {pptx_msg}")
    else:
        print(f"  [FAIL] {pptx_msg}")
        errors.append(pptx_msg)

    # Check Excel
    print("\n=== Checking Excel ===")
    excel_ok, excel_msg = check_excel(args.agent_workspace)
    if excel_ok:
        print(f"  [PASS] {excel_msg}")
    else:
        print(f"  [FAIL] {excel_msg}")
        errors.append(excel_msg)

    # Summary
    print(f"\n=== SUMMARY ===")
    if errors:
        for e in errors:
            print(f"  [ERROR] {e}")
        print("  Overall: FAIL")
    else:
        print("  Overall: PASS")

    if args.res_log_file:
        result = {"errors": errors, "success": len(errors) == 0}
        with open(args.res_log_file, "w") as f:
            json.dump(result, f, indent=2)

    sys.exit(1 if errors else 0)


if __name__ == "__main__":
    main()
