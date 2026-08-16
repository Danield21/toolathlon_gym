"""Evaluation for yf-gold-vs-stocks-ppt."""
import argparse
import json
import os
import sys

import openpyxl
import psycopg2

DB = dict(host=os.environ.get("PGHOST", "localhost"), port=int(os.environ.get("PGPORT", "5432")), dbname="toolathlon_gym", user="eigent", password="camel")
PASS_COUNT = 0
FAIL_COUNT = 0


def record(name, passed, detail=""):
    global PASS_COUNT, FAIL_COUNT
    if passed:
        PASS_COUNT += 1; print(f"  [PASS] {name}")
    else:
        FAIL_COUNT += 1; print(f"  [FAIL] {name}: {str(detail)[:300]}")


def num_close(a, b, tol=1.0):
    try: return abs(float(a) - float(b)) <= tol
    except: return False


def get_expected():
    conn = psycopg2.connect(**DB)
    cur = conn.cursor()
    monthly = {}
    for symbol in ['GC=F', '^DJI']:
        cur.execute("""SELECT date, close FROM yf.stock_prices
            WHERE symbol=%s AND date>='2025-03-06' AND date<='2026-03-05' ORDER BY date""", (symbol,))
        by_month = {}
        for d, c in cur.fetchall():
            mk = d.strftime("%Y-%m")
            by_month[mk] = float(c)
        monthly[symbol] = by_month
    conn.close()
    months = sorted(set(list(monthly['GC=F'].keys()) + list(monthly['^DJI'].keys())))
    prices = []
    for m in months:
        prices.append({"month": m, "gold": monthly['GC=F'].get(m), "dji": monthly['^DJI'].get(m)})
    returns = []
    for i in range(1, len(prices)):
        gr = dr = None
        if prices[i-1]["gold"] and prices[i]["gold"] and prices[i-1]["gold"] != 0:
            gr = round((prices[i]["gold"] - prices[i-1]["gold"]) / prices[i-1]["gold"] * 100, 2)
        if prices[i-1]["dji"] and prices[i]["dji"] and prices[i-1]["dji"] != 0:
            dr = round((prices[i]["dji"] - prices[i-1]["dji"]) / prices[i-1]["dji"] * 100, 2)
        returns.append({"month": prices[i]["month"], "gold_ret": gr, "dji_ret": dr})
    return {"prices": prices, "returns": returns, "months": months}


def sheet_dicts(wb, name):
    for sn in wb.sheetnames:
        if sn.strip().lower() == name.strip().lower():
            ws = wb[sn]
            rows = list(ws.iter_rows(values_only=True))
            if len(rows) < 2: return []
            hdrs = [str(h).strip() if h else "" for h in rows[0]]
            return [{hdrs[i]: row[i] for i in range(len(hdrs))} for row in rows[1:] if not all(v is None for v in row)]
    return None


def check_excel(ws_path, exp):
    print("\n=== Checking Excel ===")
    p = os.path.join(ws_path, "Gold_vs_DJI.xlsx")
    if not os.path.isfile(p):
        record("Excel file exists", False, p); return
    record("Excel file exists", True)
    wb = openpyxl.load_workbook(p, data_only=True)

    # Monthly Prices
    d = sheet_dicts(wb, "Monthly Prices")
    if d is None:
        record("Sheet Monthly Prices", False, str(wb.sheetnames))
    else:
        record("Sheet Monthly Prices", True)
        # Task: 2025-03 through 2026-03 = 13 months
        record(f"Monthly Prices row count == {len(exp['prices'])}",
               len(d) == len(exp["prices"]),
               f"Got {len(d)}, expected {len(exp['prices'])}")
        # Validate ALL months
        for ep in exp["prices"]:
            m = next((r for r in d if str(r.get("Month","")).strip() == ep["month"]), None)
            if not m:
                record(f"Month {ep['month']} present", False, "Missing"); continue
            record(f"Month {ep['month']} present", True)
            if ep["gold"]:
                record(f"Month {ep['month']} gold close",
                       num_close(m.get("Gold_Close"), ep["gold"], 20.0),
                       f"{m.get('Gold_Close')} vs {ep['gold']}")
            if ep["dji"]:
                record(f"Month {ep['month']} DJI close",
                       num_close(m.get("DJI_Close"), ep["dji"], 200.0),
                       f"{m.get('DJI_Close')} vs {ep['dji']}")

    # Returns
    d = sheet_dicts(wb, "Returns")
    if d is None:
        record("Sheet Returns", False, str(wb.sheetnames))
    else:
        record("Sheet Returns", True)
        # n-1 returns
        expected_ret_count = len(exp["returns"])
        record(f"Returns row count == {expected_ret_count}",
               len(d) == expected_ret_count,
               f"Got {len(d)}")
        # Validate ALL returns; tightened tol from 2.0 to 0.1 (rounded to 2 decimals)
        for er in exp["returns"]:
            m = next((r for r in d if str(r.get("Month","")).strip() == er["month"]), None)
            if not m:
                record(f"Return {er['month']} present", False, "Missing"); continue
            if er["gold_ret"] is not None:
                record(f"Return {er['month']} gold",
                       num_close(m.get("Gold_Return_Pct"), er["gold_ret"], 0.1),
                       f"{m.get('Gold_Return_Pct')} vs {er['gold_ret']}")
            if er["dji_ret"] is not None:
                record(f"Return {er['month']} DJI",
                       num_close(m.get("DJI_Return_Pct"), er["dji_ret"], 0.1),
                       f"{m.get('DJI_Return_Pct')} vs {er['dji_ret']}")
    wb.close()


def check_pptx(ws_path, exp):
    print("\n=== Checking PPTX ===")
    p = os.path.join(ws_path, "Gold_vs_Stocks.pptx")
    if not os.path.isfile(p):
        record("PPTX file exists", False, p); return
    record("PPTX file exists", True)
    try:
        from pptx import Presentation
        prs = Presentation(p)
        slides = list(prs.slides)
        # Task: exactly 3 slides
        record("Slide count == 3", len(slides) == 3, f"Got {len(slides)}")

        # ----- Slide 1: title + subtitle -----
        if len(slides) >= 1:
            slide_text = " ".join(sh.text for sh in slides[0].shapes if sh.has_text_frame).lower()
            title = slides[0].shapes.title.text.lower() if slides[0].shapes.title else ""
            record(
                "Slide 1 title 'Gold vs Stock Market Performance'",
                "gold vs stock market performance" in title or "gold vs stock market performance" in slide_text,
                f"title text: {slide_text[:200]}",
            )
            record(
                "Slide 1 subtitle mentions March 2025 to March 2026",
                ("march 2025" in slide_text and "march 2026" in slide_text)
                or ("2025-03" in slide_text and "2026-03" in slide_text),
                f"slide text: {slide_text[:200]}",
            )

        # ----- Slide 2: 'Monthly Price Comparison' table content -----
        if len(slides) >= 2:
            slide2_title = slides[1].shapes.title.text.lower() if slides[1].shapes.title else ""
            slide2_text = " ".join(sh.text for sh in slides[1].shapes if sh.has_text_frame).lower()
            record(
                "Slide 2 title 'Monthly Price Comparison'",
                "monthly price comparison" in slide2_title or "monthly price comparison" in slide2_text,
                f"title text: {slide2_text[:200]}",
            )
            # Slide 2 should mention at least the first month and the last month
            if exp["prices"]:
                first_m = exp["prices"][0]["month"]
                last_m = exp["prices"][-1]["month"]
                record(
                    f"Slide 2 mentions first month {first_m}",
                    first_m in slide2_text,
                    f"slide text: {slide2_text[:300]}",
                )
                record(
                    f"Slide 2 mentions last month {last_m}",
                    last_m in slide2_text,
                    f"slide text: {slide2_text[:300]}",
                )

        # ----- Slide 3: 'Conclusion' with which asset performed better -----
        if len(slides) >= 3:
            slide3_title = slides[2].shapes.title.text.lower() if slides[2].shapes.title else ""
            slide3_text = " ".join(sh.text for sh in slides[2].shapes if sh.has_text_frame).lower()
            record(
                "Slide 3 title 'Conclusion'",
                "conclusion" in slide3_title or "conclusion" in slide3_text,
                f"title: {slide3_title}",
            )
            # Determine which asset performed better
            if exp["prices"] and exp["prices"][0]["gold"] and exp["prices"][-1]["gold"]:
                gold_overall = (exp["prices"][-1]["gold"] - exp["prices"][0]["gold"]) / exp["prices"][0]["gold"] * 100
                dji_overall = (exp["prices"][-1]["dji"] - exp["prices"][0]["dji"]) / exp["prices"][0]["dji"] * 100
                if gold_overall > dji_overall:
                    expected_winner = "gold"
                else:
                    expected_winner = "dji"  # also check 'dow', 'jones', 'stock'
                if expected_winner == "gold":
                    record(
                        "Slide 3 conclusion says gold performed better",
                        "gold" in slide3_text,
                        f"slide3 text: {slide3_text[:300]}",
                    )
                else:
                    record(
                        "Slide 3 conclusion says DJI/Dow Jones/stock performed better",
                        any(k in slide3_text for k in ["dji", "dow", "jones", "stock"]),
                        f"slide3 text: {slide3_text[:300]}",
                    )
    except ImportError:
        record("python-pptx available", False, "Cannot import pptx")
    except Exception as e:
        record("PPTX readable", False, str(e))


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--agent_workspace", default=".")
    parser.add_argument("--groundtruth_workspace", default=".")
    parser.add_argument("--launch_time", required=False)
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()
    exp = get_expected()
    check_excel(args.agent_workspace, exp)
    check_pptx(args.agent_workspace, exp)
    print(f"\n=== SUMMARY: {PASS_COUNT} passed, {FAIL_COUNT} failed ===")
    if args.res_log_file:
        with open(args.res_log_file, "w") as f:
            json.dump({"passed": PASS_COUNT, "failed": FAIL_COUNT, "success": FAIL_COUNT == 0}, f)
    sys.exit(0 if FAIL_COUNT == 0 else 1)

if __name__ == "__main__":
    main()
