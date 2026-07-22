"""
Evaluation for train-yf-roadshow-excel-ppt-email task.

Tightened verifier: FAIL_COUNT==0 zero-tolerance,
DB-derived expected values, exact column ordering,
PPT slide-by-slide content, exact email subject + body content.
"""
import json
import os
import re
import sys
from argparse import ArgumentParser

import psycopg2
import openpyxl


DB_CONFIG = {
    "host": os.environ.get("PGHOST", "localhost"),
    "port": int(os.environ.get("PGPORT", "5432")),
    "dbname": "toolathlon_gym",
    "user": "eigent",
    "password": "camel",
}

PASS_COUNT = 0
FAIL_COUNT = 0


def record(name, passed, detail=""):
    global PASS_COUNT, FAIL_COUNT
    if passed:
        PASS_COUNT += 1
        print(f"  [PASS] {name}")
    else:
        FAIL_COUNT += 1
        msg = f": {detail[:300]}" if detail else ""
        print(f"  [FAIL] {name}{msg}")


def num_close(a, b, rel_tol=0.05, abs_tol=0.5):
    try:
        return abs(float(a) - float(b)) <= max(abs_tol, abs(float(b)) * rel_tol)
    except Exception:
        return False


def norm(s):
    return str(s or "").strip().lower().replace("_", " ")


# ---------- Build expected ----------

def build_expected():
    """Return dict with G1 train info and AMZN stock info."""
    fallback = {
        "g1": {
            "train_code": "G1",
            "from_name_keywords": ["beijing"],
            "to_name_keywords": ["shanghai"],
            "depart_time": "09:00",
            "arrive_time": "13:28",
            "duration_min": 268,
            "business_price": 1748.5,
        },
        "amzn": {
            "ticker": "AMZN",
            "name_keywords": ["amazon"],
            "sector": "Consumer Cyclical",
            "latest_close": 218.94,
            "high5": 218.94,
            "low5": 208.39,
            "avg5": 212.58,
        }
    }
    try:
        conn = psycopg2.connect(**DB_CONFIG)
        cur = conn.cursor()
        # G1 train info (Beijing South VNP -> Shanghai Hongqiao SHH on 2026-03-10)
        cur.execute("""
            SELECT t.id, t.station_train_code, t.from_station_telecode, t.to_station_telecode,
                   t.start_time, t.arrive_time
            FROM train.trains t
            WHERE t.station_train_code = 'G1' AND t.depart_date = '2026-03-10'
            LIMIT 1
        """)
        train = cur.fetchone()
        if train:
            tid, code, fcode, tcode, depart, arrive = train
            # business class price: seat_short = 'swz' (商务座)
            cur.execute("SELECT price FROM train.train_seats WHERE train_id = %s AND seat_short = 'swz'", (tid,))
            row = cur.fetchone()
            biz_price = float(row[0]) if row else fallback["g1"]["business_price"]

            try:
                d_h, d_m = map(int, depart.split(":"))
                a_h, a_m = map(int, arrive.split(":"))
                duration = (a_h * 60 + a_m) - (d_h * 60 + d_m)
                if duration < 0:
                    duration += 24 * 60
            except Exception:
                duration = fallback["g1"]["duration_min"]

            fallback["g1"]["depart_time"] = depart
            fallback["g1"]["arrive_time"] = arrive
            fallback["g1"]["duration_min"] = duration
            fallback["g1"]["business_price"] = biz_price

        # AMZN stock data
        cur.execute("""SELECT data->>'longName', data->>'sector', data->>'currentPrice'
                       FROM yf.stock_info WHERE symbol='AMZN'""")
        sinfo = cur.fetchone()
        if sinfo:
            lname, sector, _cp = sinfo
            if lname and "amazon" in lname.lower():
                fallback["amzn"]["name_keywords"] = ["amazon"]
            if sector:
                fallback["amzn"]["sector"] = sector
        cur.execute("""SELECT date, close FROM yf.stock_prices
                       WHERE symbol='AMZN' ORDER BY date DESC LIMIT 5""")
        prices = [float(r[1]) for r in cur.fetchall()]
        if len(prices) == 5:
            fallback["amzn"]["latest_close"] = prices[0]
            fallback["amzn"]["high5"] = max(prices)
            fallback["amzn"]["low5"] = min(prices)
            fallback["amzn"]["avg5"] = round(sum(prices) / 5, 2)
        cur.close()
        conn.close()
    except Exception as e:
        print(f"[expected] DB fetch failed: {e}; using fallback")
    return fallback


EXPECTED = build_expected()


# ---------- Excel ----------

EXPECTED_TRAVEL_COLS = ["Train_Code", "From", "To", "Depart", "Arrive",
                       "Duration", "Seat_Class", "Price_CNY"]
EXPECTED_STOCK_LABELS = ["Ticker", "Company_Name", "Sector", "Latest_Close_Price",
                        "5_Day_High", "5_Day_Low", "5_Day_Avg"]


def check_excel(agent_workspace):
    print("\n=== Check 1: Roadshow_Prep.xlsx ===")

    xlsx_path = os.path.join(agent_workspace, "Roadshow_Prep.xlsx")
    if not os.path.exists(xlsx_path):
        record("Roadshow_Prep.xlsx exists", False, xlsx_path)
        for k in ["Excel readable", "Travel_Details sheet", "Travel column order",
                  "Travel G1 row data correct", "Stock_Summary sheet", "Stock label order",
                  "Stock Ticker AMZN", "Stock Sector", "Stock Latest_Close",
                  "Stock 5_Day_Avg", "Stock 5_Day_High", "Stock 5_Day_Low"]:
            record(k, False, "skipped")
        return
    record("Roadshow_Prep.xlsx exists", True)
    try:
        wb = openpyxl.load_workbook(xlsx_path)
    except Exception as e:
        record("Excel readable", False, str(e))
        return
    record("Excel readable", True)

    sheet_names_lower = [s.lower() for s in wb.sheetnames]

    # Travel_Details
    if "travel_details" not in sheet_names_lower:
        record("Travel_Details sheet", False, f"Sheets: {wb.sheetnames}")
        for k in ["Travel column order", "Travel G1 row data correct"]:
            record(k, False, "skipped")
    else:
        record("Travel_Details sheet", True)
        ws = wb[wb.sheetnames[sheet_names_lower.index("travel_details")]]
        rows = list(ws.iter_rows(values_only=True))
        if not rows:
            record("Travel column order", False, "empty")
            record("Travel G1 row data correct", False, "empty")
        else:
            header = [norm(c) for c in rows[0]]
            expected_norm = [norm(c) for c in EXPECTED_TRAVEL_COLS]
            ok_cols = header[:len(expected_norm)] == expected_norm
            record("Travel column order", ok_cols, f"got {rows[0]}")

            data_rows = [r for r in rows[1:] if any(c is not None and str(c).strip() != "" for c in r)]
            if not data_rows:
                record("Travel G1 row data correct", False, "no data row")
            else:
                row = data_rows[0]
                # Build dict
                d = {}
                for i, col in enumerate(EXPECTED_TRAVEL_COLS):
                    if i < len(row):
                        d[col] = row[i]
                ok = True
                detail = []
                if str(d.get("Train_Code") or "").strip().upper() != "G1":
                    ok = False
                    detail.append(f"Train_Code={d.get('Train_Code')}")
                # From - must mention 'beijing'
                if not any(k in str(d.get("From") or "").lower() for k in EXPECTED["g1"]["from_name_keywords"]):
                    ok = False
                    detail.append(f"From={d.get('From')}")
                # To - must mention 'shanghai'
                if not any(k in str(d.get("To") or "").lower() for k in EXPECTED["g1"]["to_name_keywords"]):
                    ok = False
                    detail.append(f"To={d.get('To')}")
                # Depart and Arrive - normalize HH:MM
                def time_str(v):
                    if v is None:
                        return ""
                    s = str(v).strip()
                    if hasattr(v, "strftime"):
                        try:
                            s = v.strftime("%H:%M")
                        except Exception:
                            pass
                    m = re.match(r"^(\d{1,2}):(\d{2})", s)
                    if m:
                        return f"{int(m.group(1)):02d}:{int(m.group(2)):02d}"
                    return s
                if time_str(d.get("Depart")) != EXPECTED["g1"]["depart_time"]:
                    ok = False
                    detail.append(f"Depart={d.get('Depart')} expected {EXPECTED['g1']['depart_time']}")
                if time_str(d.get("Arrive")) != EXPECTED["g1"]["arrive_time"]:
                    ok = False
                    detail.append(f"Arrive={d.get('Arrive')} expected {EXPECTED['g1']['arrive_time']}")
                # Seat_Class
                if "business" not in str(d.get("Seat_Class") or "").lower():
                    ok = False
                    detail.append(f"Seat_Class={d.get('Seat_Class')}")
                # Price_CNY
                price = d.get("Price_CNY")
                try:
                    price_f = float(price)
                except Exception:
                    # try strip non-numeric
                    s = re.search(r"-?\d+(?:\.\d+)?", str(price or ""))
                    price_f = float(s.group(0)) if s else None
                if price_f is None or not num_close(price_f, EXPECTED["g1"]["business_price"]):
                    ok = False
                    detail.append(f"Price={price} expected {EXPECTED['g1']['business_price']}")
                record("Travel G1 row data correct", ok, "; ".join(detail)[:280])

    # Stock_Summary
    if "stock_summary" not in sheet_names_lower:
        record("Stock_Summary sheet", False, f"Sheets: {wb.sheetnames}")
        for k in ["Stock label order", "Stock Ticker AMZN", "Stock Sector",
                  "Stock Latest_Close", "Stock 5_Day_Avg",
                  "Stock 5_Day_High", "Stock 5_Day_Low"]:
            record(k, False, "skipped")
    else:
        record("Stock_Summary sheet", True)
        ws2 = wb[wb.sheetnames[sheet_names_lower.index("stock_summary")]]
        rows2 = list(ws2.iter_rows(values_only=True))
        # First row is header (Metric/Value); subsequent rows are label/value pairs
        # Build pairs in order
        pairs_in_order = []
        for r in rows2:
            if not r or len(r) < 2:
                continue
            lab = norm(r[0])
            if not lab or lab == "metric":
                continue
            pairs_in_order.append((lab, r[1]))
        # Check labels appear in expected order
        expected_labels_norm = [norm(l) for l in EXPECTED_STOCK_LABELS]
        actual_labels = [p[0] for p in pairs_in_order]
        ok_label_order = actual_labels[:len(expected_labels_norm)] == expected_labels_norm
        record("Stock label order", ok_label_order,
               f"got {actual_labels}; expected {expected_labels_norm}")

        pairs = dict(pairs_in_order)
        # Ticker AMZN
        record("Stock Ticker AMZN",
               str(pairs.get(norm("Ticker"), "")).strip().upper() == "AMZN",
               f"got {pairs.get(norm('Ticker'))}")
        # Sector
        sec = str(pairs.get(norm("Sector"), "")).lower()
        record("Stock Sector", "consumer cyclical" in sec, f"got {pairs.get(norm('Sector'))}")
        # Latest close
        lc = pairs.get(norm("Latest_Close_Price"))
        try:
            lc_f = float(re.search(r"-?\d+(?:\.\d+)?", str(lc or "")).group(0))
        except Exception:
            lc_f = None
        record("Stock Latest_Close",
               lc_f is not None and num_close(lc_f, EXPECTED["amzn"]["latest_close"], abs_tol=0.5),
               f"got {lc}")
        # 5-day avg
        avg = pairs.get(norm("5_Day_Avg"))
        try:
            avg_f = float(re.search(r"-?\d+(?:\.\d+)?", str(avg or "")).group(0))
        except Exception:
            avg_f = None
        record("Stock 5_Day_Avg",
               avg_f is not None and num_close(avg_f, EXPECTED["amzn"]["avg5"], abs_tol=0.5),
               f"got {avg}")
        hi = pairs.get(norm("5_Day_High"))
        try:
            hi_f = float(re.search(r"-?\d+(?:\.\d+)?", str(hi or "")).group(0))
        except Exception:
            hi_f = None
        record("Stock 5_Day_High",
               hi_f is not None and num_close(hi_f, EXPECTED["amzn"]["high5"], abs_tol=0.5),
               f"got {hi}")
        lo = pairs.get(norm("5_Day_Low"))
        try:
            lo_f = float(re.search(r"-?\d+(?:\.\d+)?", str(lo or "")).group(0))
        except Exception:
            lo_f = None
        record("Stock 5_Day_Low",
               lo_f is not None and num_close(lo_f, EXPECTED["amzn"]["low5"], abs_tol=0.5),
               f"got {lo}")


# ---------- PPT ----------

def check_ppt(agent_workspace):
    print("\n=== Check 2: Roadshow_Agenda.pptx ===")

    pptx_path = os.path.join(agent_workspace, "Roadshow_Agenda.pptx")
    if not os.path.exists(pptx_path):
        record("Roadshow_Agenda.pptx exists", False, pptx_path)
        for k in ["PPT readable", "PPT exactly 4 slides",
                  "Slide 1 title 'Investor Roadshow - March 2026'",
                  "Slide 1 content has 'March 10, 2026'",
                  "Slide 2 title 'Journey Details'",
                  "Slide 2 content has G1 + departure time + arrival time",
                  "Slide 3 title 'Portfolio Overview'",
                  "Slide 3 content has AMZN + sector + latest close",
                  "Slide 4 title 'Meeting Schedule'",
                  "Slide 4 content has 'TBD'"]:
            record(k, False, "skipped")
        return
    record("Roadshow_Agenda.pptx exists", True)
    try:
        from pptx import Presentation
        prs = Presentation(pptx_path)
    except Exception as e:
        record("PPT readable", False, str(e))
        return
    record("PPT readable", True)

    record("PPT exactly 4 slides", len(prs.slides) == 4, f"found {len(prs.slides)}")

    # Build slide data
    slide_titles = []
    slide_texts = []
    for slide in prs.slides:
        title_shape = slide.shapes.title
        title = title_shape.text.strip() if title_shape else ""
        slide_titles.append(title)
        all_text = ""
        for shape in slide.shapes:
            if shape.has_text_frame:
                all_text += shape.text_frame.text + " "
        slide_texts.append(all_text)

    # Slide 1
    if len(prs.slides) >= 1:
        t1 = slide_titles[0].lower()
        record("Slide 1 title 'Investor Roadshow - March 2026'",
               "investor roadshow" in t1 and ("march" in t1 or "2026" in t1),
               f"title={slide_titles[0]}")
        s1 = slide_texts[0].lower()
        ok_date = ("march 10" in s1 or "2026-03-10" in s1 or "03/10/2026" in s1) and "2026" in s1
        record("Slide 1 content has 'March 10, 2026'", ok_date, f"text={s1[:200]}")
    else:
        record("Slide 1 title 'Investor Roadshow - March 2026'", False, "missing slide")
        record("Slide 1 content has 'March 10, 2026'", False, "missing slide")

    # Slide 2
    if len(prs.slides) >= 2:
        t2 = slide_titles[1].lower()
        record("Slide 2 title 'Journey Details'", "journey" in t2 and "detail" in t2,
               f"title={slide_titles[1]}")
        s2 = slide_texts[1].lower()
        # Word boundary G1
        has_g1 = bool(re.search(r"\bg1\b", s2))
        has_depart = EXPECTED["g1"]["depart_time"] in s2
        has_arrive = EXPECTED["g1"]["arrive_time"] in s2
        record("Slide 2 content has G1 + departure time + arrival time",
               has_g1 and has_depart and has_arrive,
               f"g1={has_g1} depart={has_depart} arrive={has_arrive}")
        # Slide 2 must also mention duration (in mins or HH:MM format) and seat class
        dur_min = EXPECTED["g1"]["duration_min"]
        # accept either "Xh Ym" "X minutes" or "Xh"
        dur_h = dur_min // 60
        dur_m = dur_min % 60
        has_duration = (
            (str(dur_min) in s2 and "min" in s2)
            or (f"{dur_h}h" in s2 or f"{dur_h} h" in s2)
            or (f"{dur_h}:{dur_m:02d}" in s2)
        )
        record(
            "Slide 2 content includes journey duration",
            has_duration,
            f"target {dur_min} min ({dur_h}h{dur_m}m); text snippet={s2[:200]}",
        )
        record(
            "Slide 2 content includes seat class (Business)",
            "business" in s2,
            f"text snippet={s2[:200]}",
        )
    else:
        record("Slide 2 title 'Journey Details'", False, "missing")
        record("Slide 2 content has G1 + departure time + arrival time", False, "missing")
        record("Slide 2 content includes journey duration", False, "missing")
        record("Slide 2 content includes seat class (Business)", False, "missing")

    # Slide 3
    if len(prs.slides) >= 3:
        t3 = slide_titles[2].lower()
        record("Slide 3 title 'Portfolio Overview'",
               "portfolio" in t3 and "overview" in t3, f"title={slide_titles[2]}")
        s3 = slide_texts[2].lower()
        has_amzn = "amzn" in s3 or "amazon" in s3
        has_sector = "consumer cyclical" in s3
        # latest close ~ 218.94
        nums = [float(m.group(0)) for m in re.finditer(r"-?\d+(?:\.\d+)?", s3)]
        has_close = any(num_close(n, EXPECTED["amzn"]["latest_close"], abs_tol=0.5) for n in nums)
        record("Slide 3 content has AMZN + sector + latest close",
               has_amzn and has_sector and has_close,
               f"amzn={has_amzn} sector={has_sector} close={has_close}")
    else:
        record("Slide 3 title 'Portfolio Overview'", False, "missing")
        record("Slide 3 content has AMZN + sector + latest close", False, "missing")

    # Slide 4
    if len(prs.slides) >= 4:
        t4 = slide_titles[3].lower()
        record("Slide 4 title 'Meeting Schedule'",
               "meeting" in t4 and "schedule" in t4, f"title={slide_titles[3]}")
        s4 = slide_texts[3].lower()
        record("Slide 4 content has 'TBD'", "tbd" in s4, f"text={s4[:200]}")
    else:
        record("Slide 4 title 'Meeting Schedule'", False, "missing")
        record("Slide 4 content has 'TBD'", False, "missing")


# ---------- Email ----------

def check_email():
    print("\n=== Check 3: Email to roadshow@bank.com ===")
    try:
        conn = psycopg2.connect(**DB_CONFIG)
        cur = conn.cursor()
        cur.execute("SELECT to_addr, subject, body_text, body_html FROM email.messages")
        rows = cur.fetchall()
        cur.close()
        conn.close()
    except Exception as e:
        record("Email DB query", False, str(e))
        for k in ["Email to roadshow@bank.com",
                  "Email subject 'Roadshow Travel and Prep Confirmed'",
                  "Email body has G1 + departure time",
                  "Email body has business class price",
                  "Email body has AMZN latest close"]:
            record(k, False, "skipped")
        return

    matched = []
    for to_addr, subj, btxt, bhtml in rows:
        to_str = str(to_addr or "").lower()
        if "roadshow@bank.com" in to_str:
            matched.append((subj or "", (btxt or "") + "\n" + (bhtml or "")))
    record("Email to roadshow@bank.com", len(matched) >= 1, f"count={len(matched)}")
    if not matched:
        for k in ["Email subject 'Roadshow Travel and Prep Confirmed'",
                  "Email body has G1 + departure time",
                  "Email body has business class price",
                  "Email body has AMZN latest close"]:
            record(k, False, "skipped")
        return

    subj, body = matched[0]
    subj_l = subj.strip().lower()
    expected_subj = "roadshow travel and prep confirmed"
    # Strict: subject must match exactly (case-insensitive)
    ok_subj = subj_l == expected_subj
    record("Email subject 'Roadshow Travel and Prep Confirmed'", ok_subj, f"subj={subj!r}")

    body_l = body.lower()
    has_g1 = bool(re.search(r"\bg1\b", body_l))
    has_depart = EXPECTED["g1"]["depart_time"] in body_l
    record("Email body has G1 + departure time",
           has_g1 and has_depart, f"g1={has_g1} depart={has_depart}")

    nums = [float(m.group(0)) for m in re.finditer(r"-?\d+(?:\.\d+)?", body)]
    has_price = any(num_close(n, EXPECTED["g1"]["business_price"], abs_tol=1.0) for n in nums)
    record("Email body has business class price", has_price,
           f"target={EXPECTED['g1']['business_price']} nums sample={sorted(nums)[:10]}")

    has_close = any(num_close(n, EXPECTED["amzn"]["latest_close"], abs_tol=0.5) for n in nums)
    record("Email body has AMZN latest close", has_close,
           f"target={EXPECTED['amzn']['latest_close']}")


def main():
    parser = ArgumentParser()
    parser.add_argument("--agent_workspace", required=False, default=".")
    parser.add_argument("--groundtruth_workspace", required=False, default=".")
    parser.add_argument("--launch_time", required=False)
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()

    print(f"[expected G1] {EXPECTED['g1']}")
    print(f"[expected AMZN] {EXPECTED['amzn']}")

    check_excel(args.agent_workspace)
    check_ppt(args.agent_workspace)
    check_email()

    total = PASS_COUNT + FAIL_COUNT
    print(f"\nOverall: {PASS_COUNT}/{total} checks passed [FAIL_COUNT={FAIL_COUNT}]")
    if FAIL_COUNT == 0:
        print("PASS")
        sys.exit(0)
    else:
        print("FAIL")
        sys.exit(1)


if __name__ == "__main__":
    main()
