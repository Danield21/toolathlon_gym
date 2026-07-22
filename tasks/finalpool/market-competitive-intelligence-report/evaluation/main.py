#!/usr/bin/env python3
"""Evaluation script for market-competitive-intelligence-report."""

import argparse
import os
import sys
from pathlib import Path

PASS_COUNT = 0
FAIL_COUNT = 0


def check(name, condition, detail=""):
    global PASS_COUNT, FAIL_COUNT
    if condition:
        PASS_COUNT += 1
        print(f"  [PASS] {name}")
    else:
        FAIL_COUNT += 1
        detail_str = str(detail)[:200] if detail else ""
        print(f"  [FAIL] {name}: {detail_str}")


def get_conn():
    import psycopg2
    return psycopg2.connect(
        host=os.environ.get("PGHOST", "localhost"), port=5432,
        dbname=os.environ.get("PGDATABASE", "toolathlon_gym"),
        user="eigent", password="camel",
    )


def run_evaluation(agent_workspace, groundtruth_workspace, launch_time, res_log_file):
    global PASS_COUNT, FAIL_COUNT
    PASS_COUNT = 0
    FAIL_COUNT = 0

    if not agent_workspace:
        return False, "No agent workspace provided"

    agent_ws = Path(agent_workspace)
    if not agent_ws.exists():
        return False, f"Agent workspace not found: {agent_workspace}"

    # Detect GT self-test mode (V1 parity test).
    try:
        gt_canon = os.path.realpath(groundtruth_workspace) if groundtruth_workspace else ""
        ag_canon = os.path.realpath(agent_workspace) if agent_workspace else ""
        is_gt_self_test = bool(gt_canon) and (gt_canon == ag_canon)
    except Exception:
        is_gt_self_test = False

    # Phase 4: Excel competitive analysis
    xlsx_files = [p for p in agent_ws.glob("*.xlsx")
                  if 'data.csv' not in p.name.lower()]
    check("At least one xlsx competitive analysis report exists",
          len(xlsx_files) >= 1, f"found {len(xlsx_files)} xlsx files")
    if xlsx_files:
        try:
            import openpyxl
            audit_path = max(xlsx_files, key=lambda p: p.stat().st_size)
            wb = openpyxl.load_workbook(str(audit_path), data_only=True)
            sheets = [s.lower() for s in wb.sheetnames]
            has_comp_sheet = any(
                kw in s
                for s in sheets
                for kw in ['competitor', 'pricing', 'feature', 'comparison',
                          'matrix', 'analysis', 'summary', 'opportunity']
            )
            check(f"Competitive xlsx ({audit_path.name}) has competitor-related sheet",
                  has_comp_sheet, f"sheets={wb.sheetnames}")
            # Phase 1: at least 5 competitors mentioned somewhere
            max_rows = max((ws.max_row for ws in wb.worksheets), default=0)
            # Tightened: >=6 rows (header + 5 competitors)
            check("Competitive xlsx has >=6 rows of data (header + >=5 competitors)",
                  max_rows >= 6, f"max rows={max_rows}")

            # Tightened (per QA): count distinct candidate competitor entries
            # in column A across all sheets (data rows only, skipping headers).
            # We accept any string >=3 chars that is not a known header token.
            # Uppercase requirement is relaxed (per QA) since some brand names
            # are stylized lowercase. Restricting to column A still constrains
            # the check to per-row entity entries.
            proper_nouns = set()
            HEADER_TOKENS = {
                'competitor', 'feature', 'opportunity', 'priority',
                'tier', 'price', 'impact', 'description', 'metric',
                'value', 'action', 'item', 'notes', 'category', 'name',
                'rank', 'segment', 'channel', 'feature matrix',
                'competitor pricing', 'opportunities', 'summary', 'total',
            }
            for ws in wb.worksheets:
                # Only column A (index 0), skip header (min_row=2)
                for row in ws.iter_rows(min_row=2, values_only=True):
                    if not row:
                        continue
                    v = row[0]
                    if not isinstance(v, str):
                        continue
                    s = v.strip()
                    if len(s) < 3:
                        continue
                    if s.lower() in HEADER_TOKENS:
                        continue
                    proper_nouns.add(s.lower())
            check("Competitive xlsx column A has >=5 distinct candidate entries",
                  len(proper_nouns) >= 5,
                  f"distinct candidates count={len(proper_nouns)}, sample={list(proper_nouns)[:5]}")

            # Value-level GT comparison: check headers + numeric coverage
            gt_xlsx_path = Path(groundtruth_workspace) / "Competitive_Analysis.xlsx"
            if gt_xlsx_path.exists():
                try:
                    gt_wb = openpyxl.load_workbook(str(gt_xlsx_path), data_only=True)
                    agent_headers = set()
                    for ws in wb.worksheets:
                        first = next(ws.iter_rows(min_row=1, max_row=1, values_only=True), ())
                        for h in first:
                            if h:
                                agent_headers.add(str(h).strip().lower().replace(' ', '_'))
                    expected_kw = ['competitor', 'price', 'tier', 'feature',
                                   'opportunity', 'priority', 'impact']
                    matched_kw = sum(
                        1 for kw in expected_kw
                        if any(kw in h for h in agent_headers)
                    )
                    check("Competitive xlsx column-coverage (>=2 of expected keywords)",
                          matched_kw >= 2,
                          f"matched {matched_kw}/{len(expected_kw)} headers={list(agent_headers)[:10]}")

                    # Numeric coverage: pricing/feature counts
                    agent_numerics = []
                    for ws in wb.worksheets:
                        for row in ws.iter_rows(min_row=2, values_only=True):
                            for v in row:
                                if isinstance(v, (int, float)):
                                    agent_numerics.append(float(v))
                    check("Competitive xlsx has >=3 numeric values",
                          len(agent_numerics) >= 3,
                          f"got {len(agent_numerics)} numerics")
                    gt_wb.close()
                except Exception as _e:
                    check("Competitive xlsx GT comparison", False, str(_e))
        except Exception as e:
            check("Competitive xlsx parse", False, str(e))

    # Phase 5: Word executive summary
    docx_files = list(agent_ws.glob("*.docx"))
    check("At least one docx executive summary exists",
          len(docx_files) >= 1, f"found {len(docx_files)} docx files")
    if docx_files:
        try:
            from docx import Document
            audit_docx = max(docx_files, key=lambda p: p.stat().st_size)
            doc = Document(str(audit_docx))
            text = " ".join(p.text for p in doc.paragraphs).lower()
            check("Word doc has substantive content (>=300 chars)",
                  len(text) >= 300, f"len={len(text)}")
            has_topic = any(kw in text for kw in
                            ['competitor', 'competitive', 'pricing',
                             'market', 'feature', 'differentiation'])
            has_recs = any(kw in text for kw in
                           ['recommendation', 'opportunity', 'finding',
                            'strategy', 'risk'])
            check("Word doc mentions competitor/pricing/market topic",
                  has_topic, f"text head: {text[:200]}")
            check("Word doc mentions recommendations/strategies/findings",
                  has_recs, f"text head: {text[:200]}")
        except Exception as e:
            check("Word doc parse", False, str(e))

    # Phase 6: PowerPoint deck
    pptx_files = list(agent_ws.glob("*.pptx"))
    check("At least one pptx executive deck exists",
          len(pptx_files) >= 1, f"found {len(pptx_files)} pptx files")
    if pptx_files:
        try:
            from pptx import Presentation
            audit_ppt = max(pptx_files, key=lambda p: p.stat().st_size)
            prs = Presentation(str(audit_ppt))
            n_slides = len(prs.slides)
            check("PPTX has at least 5 slides", n_slides >= 5,
                  f"n_slides={n_slides}")
            all_text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        all_text += " " + shape.text.lower()
            has_topic = any(kw in all_text for kw in
                            ['competitor', 'pricing', 'market', 'feature',
                             'opportunity', 'recommendation'])
            check("PPTX mentions competitor/pricing/market topic",
                  has_topic, f"text head: {all_text[:200]}")
        except Exception as e:
            check("PPTX parse", False, str(e))

    # Phase 6: Email to executives
    try:
        conn = get_conn()
        cur = conn.cursor()
        cur.execute(
            """SELECT id, subject FROM email.messages
               WHERE subject ILIKE %s OR subject ILIKE %s OR subject ILIKE %s
                  OR body_text ILIKE %s""",
            ('%competitive%', '%competitor%', '%market%', '%intelligence%'))
        rows = cur.fetchall()
        if is_gt_self_test and len(rows) == 0:
            print("  [WARN] Executive briefing email sent: 0 found (GT self-test, non-blocking)")
        else:
            check("Executive briefing email sent",
                  len(rows) >= 1, f"found {len(rows)} matching emails")
        conn.close()
    except Exception as e:
        check("Email check", False, str(e))

    # Phase 6: Strategy briefing meeting
    try:
        conn = get_conn()
        cur = conn.cursor()
        cur.execute(
            """SELECT summary FROM gcal.events
               WHERE summary ILIKE %s OR summary ILIKE %s OR summary ILIKE %s""",
            ('%briefing%', '%strategy%', '%competitive%'))
        rows = cur.fetchall()
        if is_gt_self_test and len(rows) == 0:
            print("  [WARN] Strategy briefing meeting scheduled: 0 found (GT self-test, non-blocking)")
        else:
            check("Strategy briefing meeting scheduled",
                  len(rows) >= 1, f"found {len(rows)} matching events")
        conn.close()
    except Exception as e:
        check("Calendar check", False, str(e))

    return FAIL_COUNT == 0, f"Passed {PASS_COUNT}/{PASS_COUNT + FAIL_COUNT} checks"


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--agent_workspace", required=False, default=".")
    parser.add_argument("--groundtruth_workspace", required=False, default=".")
    parser.add_argument("--launch_time", required=False, default="2026-03-07 10:00:00")
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()

    success, message = run_evaluation(
        args.agent_workspace, args.groundtruth_workspace,
        args.launch_time, args.res_log_file
    )
    print(message)
    sys.exit(0 if success else 1)


if __name__ == "__main__":
    main()
