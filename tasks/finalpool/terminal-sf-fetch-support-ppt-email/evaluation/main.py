"""Evaluation for terminal-sf-fetch-support-ppt-email."""
import argparse
import json
import os
import sys

import psycopg2

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

DB = dict(host=os.environ.get("PGHOST", "localhost"), port=5432,
          dbname=os.environ.get("PGDATABASE", "toolathlon_gym"),
          user="eigent", password="camel")

PASS_COUNT = 0
FAIL_COUNT = 0


def check(name, condition, detail=""):
    global PASS_COUNT, FAIL_COUNT
    if condition:
        PASS_COUNT += 1
        print(f"  [PASS] {name}")
    else:
        FAIL_COUNT += 1
        d = f": {str(detail)[:200]}" if detail else ""
        print(f"  [FAIL] {name}{d}")


def _get_live_ticket_counts():
    """Fetch live per-priority ticket counts from Snowflake mirror."""
    try:
        conn = psycopg2.connect(**DB)
        cur = conn.cursor()
        cur.execute(
            'SELECT "PRIORITY", COUNT(*) '
            'FROM sf_data."SUPPORT_CENTER__PUBLIC__TICKETS" '
            'GROUP BY "PRIORITY"'
        )
        rows = cur.fetchall()
        cur.close()
        conn.close()
        out = {str(r[0]).strip().lower(): int(r[1]) for r in rows}
        out["total"] = sum(out.values())
        return out
    except Exception:
        return {}


def _has_number_token(text, val):
    """True if number `val` appears as a standalone token in text. Accepts
    comma-grouped form (e.g., 6,466) and bare digits (6466)."""
    import re
    val = int(round(float(val)))
    bare = str(val)
    grouped = f"{val:,}"
    pat = (
        r"(?<!\d)" + re.escape(bare) + r"(?!\d|\.\d)"
        + r"|" + r"(?<!\d)" + re.escape(grouped) + r"(?!\d)"
    )
    return bool(re.search(pat, text))


def _has_priority_word(text, level):
    """True if the priority-level word appears in a context that maps to
    priority semantics, not as a substring of unrelated words. Accept:
    '<level> priority', 'priority <level>', '<level>-priority',
    or the level word followed by a quantitative tag (count/tickets)."""
    import re
    level = level.lower()
    patterns = [
        rf"\b{level}[- ]priority\b",
        rf"\bpriority[: ]+{level}\b",
        rf"\b{level}\b\s*\((\d|tickets)",
        rf"\b{level}\b[ ]+(tickets|priority|csat|satisfaction|sla|compliant|gap)",
    ]
    for p in patterns:
        if re.search(p, text):
            return True
    # Fallback: appears as standalone word (still much stricter than 'in text')
    return bool(re.search(rf"\b{level}\b", text))


def check_pptx(ws_path):
    """Check Support_Performance_Review.pptx."""
    print("\n=== Checking PowerPoint ===")
    path = os.path.join(ws_path, "Support_Performance_Review.pptx")
    if not os.path.isfile(path):
        check("PPT file exists", False, f"Not found: {path}")
        return
    check("PPT file exists", True)

    if Presentation is None:
        check("python-pptx available", False, "Cannot import pptx")
        return

    try:
        prs = Presentation(path)
        slides = prs.slides
        check("PPT has >= 7 slides", len(slides) >= 7, f"Found {len(slides)} slides")

        all_text = ""
        for slide in slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    all_text += shape.text_frame.text.lower() + " "

        check("PPT mentions benchmark or industry",
              "benchmark" in all_text or "industry" in all_text)
        # Use word-boundary checks for priority levels to avoid 'high' matching 'highlight'.
        check("PPT mentions all 3 priority levels (word-bounded)",
              _has_priority_word(all_text, "high")
              and _has_priority_word(all_text, "medium")
              and _has_priority_word(all_text, "low"))
        check("PPT mentions satisfaction or CSAT",
              "satisfaction" in all_text or "csat" in all_text)
        check("PPT mentions compliance or SLA",
              "compliance" in all_text or "sla" in all_text or "compliant" in all_text)
        check("PPT mentions agents and performance",
              "agent" in all_text and "performance" in all_text)
        check("PPT mentions recommendations or takeaways",
              "recommend" in all_text or "takeaway" in all_text or "improvement" in all_text)

        # === Numeric checks against GT values ===
        # Industry benchmark CSAT targets (3.50, 3.40, 3.20). Accept either
        # the 2-decimal form or the 1-decimal token surrounded by non-digit
        # boundaries (so '3.5' won't match inside '13.5' or '3.55').
        import re as _re
        def _csat_match(value_str_2dp, value_str_1dp):
            return bool(
                value_str_2dp in all_text
                or _re.search(r"(?<!\d)" + _re.escape(value_str_1dp) + r"(?!\d|\.\d)", all_text)
            )
        check("PPT contains industry CSAT target 3.50 (High)",
              _csat_match("3.50", "3.5"))
        check("PPT contains industry CSAT target 3.40 (Medium)",
              _csat_match("3.40", "3.4"))
        check("PPT contains industry CSAT target 3.20 (Low)",
              _csat_match("3.20", "3.2"))

        # Company actual ticket counts: pull live from Snowflake to avoid
        # hardcoded GT drift.
        live = _get_live_ticket_counts()
        if live:
            for level in ("high", "medium", "low"):
                expected = live.get(level)
                if expected is None:
                    continue
                check(f"PPT contains live ticket count {expected:,} ({level})",
                      _has_number_token(all_text, expected))
            total = live.get("total")
            if total:
                check(f"PPT contains total ticket count {total:,}",
                      _has_number_token(all_text, total))

        # SLA Compliance status presence
        check("PPT mentions Non-Compliant",
              "non-compliant" in all_text or "non compliant" in all_text or "noncompliant" in all_text)
        check("PPT mentions At Risk status",
              "at risk" in all_text or "at-risk" in all_text)
    except Exception as e:
        check("PPT readable", False, str(e))


def check_email():
    """Check email sent to managers."""
    print("\n=== Checking Email ===")
    conn = psycopg2.connect(**DB)
    cur = conn.cursor()

    cur.execute("SELECT subject, from_addr, to_addr, body_text FROM email.messages")
    all_emails = cur.fetchall()

    target_email = None
    for subj, from_addr, to_addr, body in all_emails:
        if to_addr:
            recipients = []
            if isinstance(to_addr, list):
                recipients = [str(r).strip().lower() for r in to_addr]
            elif isinstance(to_addr, str):
                try:
                    parsed = json.loads(to_addr)
                    if isinstance(parsed, list):
                        recipients = [str(r).strip().lower() for r in parsed]
                    else:
                        recipients = [str(to_addr).strip().lower()]
                except (json.JSONDecodeError, TypeError):
                    recipients = [str(to_addr).strip().lower()]
            if "managers@support-team.example.com" in recipients:
                target_email = (subj, from_addr, to_addr, body)
                break

    check("Email sent to managers@support-team.example.com", target_email is not None,
          f"Total emails: {len(all_emails)}")

    if target_email:
        subj, from_addr, to_addr, body = target_email
        subj_l = (subj or "").lower()
        check("Email subject contains 'Support Performance Benchmark Analysis'",
              "support performance benchmark analysis" in subj_l,
              f"Subject: {subj}")
        check("Email from analytics@support-team.example.com",
              "analytics@support-team.example.com" in (from_addr or "").lower(),
              f"From: {from_addr}")
        body_lower = (body or "").lower()
        check("Email body mentions satisfaction or CSAT",
              "satisfaction" in body_lower or "csat" in body_lower,
              "Expected satisfaction/CSAT in body")
        check("Email body mentions compliance or priority",
              "compliant" in body_lower or "compliance" in body_lower or "priority" in body_lower,
              "Expected compliance/priority in body")
        # Body should mention top performers and improvement areas
        check("Email body mentions agents (top performers)",
              "agent" in body_lower or "top" in body_lower or "performer" in body_lower,
              "Expected agent/top/performer in body")
        check("Email body mentions improvement or area",
              "improv" in body_lower or "area" in body_lower or "need" in body_lower,
              "Expected improvement/area/need in body")

    conn.close()


def check_reverse_validation(workspace):
    """Verify things that should NOT exist in output."""
    print("\n=== Reverse Validation ===")

    # PPT: should not have empty slides
    path = os.path.join(workspace, "Support_Performance_Review.pptx")
    if os.path.isfile(path) and Presentation is not None:
        try:
            prs = Presentation(path)
            empty_slides = 0
            for slide in prs.slides:
                slide_text = ""
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        slide_text += shape.text_frame.text.strip()
                if not slide_text:
                    empty_slides += 1
            check("No empty slides in PPT", empty_slides == 0,
                  f"Found {empty_slides} empty slides")
        except Exception:
            pass

    # Email: no emails sent to wrong recipients about support performance
    try:
        conn = psycopg2.connect(**DB)
        cur = conn.cursor()
        cur.execute("""
            SELECT COUNT(*) FROM email.messages
            WHERE (lower(subject) LIKE '%%benchmark%%' OR lower(subject) LIKE '%%performance%%')
              AND to_addr::text ILIKE '%%competitor%%'
        """)
        bad_count = cur.fetchone()[0]
        check("No benchmark emails to competitor addresses", bad_count == 0,
              f"Found {bad_count}")
        cur.close()
        conn.close()
    except Exception:
        pass


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--agent_workspace", required=False, default=".")
    parser.add_argument("--groundtruth_workspace", required=False)
    parser.add_argument("--launch_time", required=False)
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()

    print("=" * 70)
    print("TERMINAL-SF-FETCH-SUPPORT-PPT-EMAIL - EVALUATION")
    print("=" * 70)

    check_pptx(args.agent_workspace)
    check_email()
    check_reverse_validation(args.agent_workspace)

    print(f"\n=== SUMMARY ===")
    print(f"  Passed: {PASS_COUNT}")
    print(f"  Failed: {FAIL_COUNT}")
    overall = FAIL_COUNT == 0
    print(f"  Overall: {'PASS' if overall else 'FAIL'}")
    sys.exit(0 if overall else 1)


if __name__ == "__main__":
    main()
