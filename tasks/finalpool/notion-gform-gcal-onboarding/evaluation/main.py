"""
Evaluation script for notion-gform-gcal-onboarding task.

Checks (strict):
1. PPT file with title slide containing 'Welcome to the Team!' + March 16, 2026,
   plus checklist slides (>=6) and a 'New Team Members' slide listing 3 hires
   (names, departments, emails).
2. Calendar Orientation Session (9-12) and Team Lunch (12-13) on 2026-03-16,
   Orientation includes the three hires as attendees.
3. 3 individual welcome emails to each new hire with exact subject and body
   content (name, department, dates).
4. Notion page updated with a 'March 2026 New Hires' section listing all
   three new hires by exact full name.
"""

import argparse
import json
import os
import re
import sys

import psycopg2
from pptx import Presentation

# ──────────────────────────────────────────────────────────────────────────
# EVALUATION GROUND TRUTH SPEC (gcal tz root-fix v3, case-study 2026-08-13)
# gcal.events.start_datetime is TIMESTAMPTZ; bare start_dt.hour / .strftime
# silently compares wrong in non-UTC PG sessions. Use gcal_helpers.
# ──────────────────────────────────────────────────────────────────────────
# task.md line 5: "9:00 AM to 12:00 PM" / "12:00 PM to 1:00 PM" (corporate
# onboarding, ET business default)
EXPECTED_TIMEZONE = "America/New_York"

from utils.evaluation.gcal_helpers import get_zone_components  # noqa: E402

DB_CONFIG = {
    "host": os.environ.get("PGHOST", "localhost"),
    "port": int(os.environ.get("PGPORT", "5432")),
    "dbname": "toolathlon_gym",
    "user": "eigent",
    "password": "camel",
}

PASS_COUNT = 0
FAIL_COUNT = 0

NEW_HIRES = [
    {"name": "Sarah Park", "email": "sarah.park@company.com", "dept": "Engineering"},
    {"name": "Mike Chen", "email": "mike.chen@company.com", "dept": "Sales"},
    {"name": "Amy Rodriguez", "email": "amy.rodriguez@company.com", "dept": "Marketing"},
]


def record(name, passed, detail=""):
    global PASS_COUNT, FAIL_COUNT
    if passed:
        PASS_COUNT += 1
        print(f"  [PASS] {name}")
    else:
        FAIL_COUNT += 1
        msg = f": {str(detail)[:300]}" if detail else ""
        print(f"  [FAIL] {name}{msg}")


def str_contains(haystack, needle):
    if haystack is None or needle is None:
        return False
    return needle.lower() in str(haystack).lower()


# ============================================================================
# Check 1: PowerPoint file
# ============================================================================

def check_pptx(agent_workspace):
    print("\n=== Checking PowerPoint ===")

    pptx_path = os.path.join(agent_workspace, "Onboarding_Presentation.pptx")
    if not os.path.isfile(pptx_path):
        record("PPT file exists", False, f"Not found: {pptx_path}")
        return
    record("PPT file exists", True)

    try:
        prs = Presentation(pptx_path)
    except Exception as e:
        record("PPT file readable", False, str(e))
        return
    record("PPT file readable", True)

    slide_count = len(prs.slides)
    # Title + 6 checklist + New Team Members = 8
    record("PPT has at least 8 slides (title + 6 checklist + members)", slide_count >= 8,
           f"Found {slide_count} slides")

    # Title slide
    if slide_count > 0:
        first_slide = prs.slides[0]
        slide_text = ""
        for shape in first_slide.shapes:
            if shape.has_text_frame:
                slide_text += " " + shape.text_frame.text
        slide_lower = slide_text.lower()
        record("Title slide contains 'Welcome to the Team!'",
               "welcome to the team" in slide_lower,
               f"Slide text: {slide_text[:200]}")
        # Subtitle date - accept various forms but must contain March 16, 2026 (or 2026-03-16)
        date_ok = (
            ("march 16" in slide_lower and "2026" in slide_lower)
            or "2026-03-16" in slide_lower
            or "03/16/2026" in slide_lower
        )
        record("Title slide subtitle mentions March 16, 2026", date_ok,
               f"Slide text: {slide_text[:200]}")

    all_slide_text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                all_slide_text += " " + shape.text_frame.text
    all_slide_lower = all_slide_text.lower()

    # Checklist content check - require all 6 items
    checklist_items = [
        ("hr paperwork", "complete hr paperwork"),
        ("workstation", "set up workstation"),
        ("team members", "meet team members"),
        ("orientation", "orientation"),
        ("handbook", "company handbook"),
        ("compliance", "compliance training"),
    ]
    found = sum(1 for kw, _ in checklist_items if kw in all_slide_lower)
    record("PPT covers all 6 checklist items", found >= 6,
           f"Found {found}/6 checklist keywords")

    # New Team Members slide present
    record("PPT has 'New Team Members' slide", "new team members" in all_slide_lower,
           "Phrase 'New Team Members' missing")

    # Each hire: name + email + department
    for h in NEW_HIRES:
        nm = h["name"].lower(); em = h["email"].lower(); dp = h["dept"].lower()
        ok = (nm in all_slide_lower) and (em in all_slide_lower) and (dp in all_slide_lower)
        record(f"PPT lists {h['name']} (name+email+department)", ok,
               f"name={nm in all_slide_lower}, email={em in all_slide_lower}, dept={dp in all_slide_lower}")


# ============================================================================
# Check 2: Google Calendar events
# ============================================================================

def check_gcal():
    print("\n=== Checking Google Calendar ===")

    conn = psycopg2.connect(**DB_CONFIG)
    cur = conn.cursor()

    cur.execute("""
        SELECT summary, description, start_datetime, end_datetime, attendees
        FROM gcal.events
        ORDER BY start_datetime
    """)
    events = cur.fetchall()
    cur.close()
    conn.close()

    print(f"[check_gcal] Found {len(events)} calendar events.")
    record("At least 2 calendar events created", len(events) >= 2, f"Found {len(events)}")

    # Orientation Session
    orientation = None
    for ev in events:
        s = (ev[0] or "").strip().lower()
        if "orientation" in s:
            orientation = ev
            break
    record("'Orientation Session' event exists", orientation is not None,
           f"Events: {[e[0] for e in events]}")
    if orientation is not None:
        summary, desc, start_dt, end_dt, attendees = orientation
        # gcal.events.start_datetime is TIMESTAMPTZ; use session-tz-independent
        # helper to extract ET date/hour.
        ev_date, ev_sh, _ = get_zone_components(start_dt, EXPECTED_TIMEZONE)
        _ev_ed_date, ev_eh, _ = get_zone_components(end_dt, EXPECTED_TIMEZONE)
        start_date_str = ev_date.strftime("%Y-%m-%d") if ev_date else ""
        record("Orientation on 2026-03-16", start_date_str == "2026-03-16",
               f"Start date: {start_date_str}")
        if start_dt and end_dt and ev_sh is not None and ev_eh is not None:
            record("Orientation 9 AM - 12 PM",
                   ev_sh == 9 and ev_eh == 12,
                   f"start={start_dt}, end={end_dt}")
        # Attendees include 3 emails
        attendees_str = json.dumps(attendees).lower() if attendees is not None else ""
        for h in NEW_HIRES:
            record(f"Orientation attendee includes {h['email']}",
                   h["email"].lower() in attendees_str,
                   f"attendees={attendees_str[:200]}")

    # Team Lunch
    lunch = None
    for ev in events:
        s = (ev[0] or "").strip().lower()
        if "team lunch" in s or s == "team lunch":
            lunch = ev
            break
    record("'Team Lunch' event exists", lunch is not None,
           f"Events: {[e[0] for e in events]}")
    if lunch is not None:
        summary, desc, start_dt, end_dt, attendees = lunch
        # gcal.events.start_datetime is TIMESTAMPTZ; use helper for ET date.
        ev_date, ev_sh, _ = get_zone_components(start_dt, EXPECTED_TIMEZONE)
        _ed, ev_eh, _ = get_zone_components(end_dt, EXPECTED_TIMEZONE)
        start_date_str = ev_date.strftime("%Y-%m-%d") if ev_date else ""
        record("Team Lunch on 2026-03-16", start_date_str == "2026-03-16",
               f"Start date: {start_date_str}")
        if start_dt and end_dt and ev_sh is not None and ev_eh is not None:
            record("Team Lunch 12 PM - 1 PM",
                   ev_sh == 12 and ev_eh == 13,
                   f"start={start_dt}, end={end_dt}")


# ============================================================================
# Check 3: Emails
# ============================================================================

def check_emails():
    print("\n=== Checking Emails ===")

    conn = psycopg2.connect(**DB_CONFIG)
    cur = conn.cursor()

    # Try Sent folder first; fallback to all messages
    cur.execute("""
        SELECT subject, from_addr, to_addr, body_text
        FROM email.messages
        WHERE folder_id = (SELECT id FROM email.folders WHERE name = 'Sent' LIMIT 1)
    """)
    sent_emails = cur.fetchall()
    if not sent_emails:
        cur.execute("SELECT subject, from_addr, to_addr, body_text FROM email.messages")
        sent_emails = cur.fetchall()
    cur.close()
    conn.close()

    print(f"[check_emails] Found {len(sent_emails)} emails.")
    record("At least 3 emails sent", len(sent_emails) >= 3, f"Found {len(sent_emails)}")

    expected_subject = "Welcome to Our Company - Onboarding Information"
    for h in NEW_HIRES:
        match = None
        for subj, from_addr, to_addr, body in sent_emails:
            to_str = json.dumps(to_addr).lower() if isinstance(to_addr, (list, dict)) else str(to_addr or "").lower()
            if h["email"].lower() in to_str:
                match = (subj, from_addr, to_addr, body)
                break
        record(f"Email to {h['email']} present", match is not None,
               f"to_addrs: {[e[2] for e in sent_emails]}")
        if match is None:
            continue
        subj, _, _, body = match
        record(f"Email to {h['email']} has exact subject",
               (subj or "").strip().lower() == expected_subject.lower(),
               f"subject={subj}")
        body_lower = (body or "").lower()
        record(f"Email to {h['email']} body greets by name",
               h["name"].lower() in body_lower,
               f"body snippet: {body_lower[:150]}")
        record(f"Email to {h['email']} body mentions department {h['dept']}",
               h["dept"].lower() in body_lower,
               f"body snippet: {body_lower[:150]}")
        date_ok = (
            ("march 16" in body_lower and "2026" in body_lower)
            or "2026-03-16" in body_lower or "03/16/2026" in body_lower
        )
        record(f"Email to {h['email']} body mentions start date Mar 16, 2026",
               date_ok, f"body snippet: {body_lower[:200]}")
        # Check both 'orientation' and 'lunch' are mentioned
        record(f"Email to {h['email']} body mentions orientation",
               "orientation" in body_lower, "")
        record(f"Email to {h['email']} body mentions team lunch",
               "lunch" in body_lower, "")


# ============================================================================
# Check 4: Notion page updated
# ============================================================================

def check_notion():
    print("\n=== Checking Notion ===")

    conn = psycopg2.connect(**DB_CONFIG)
    cur = conn.cursor()

    cur.execute("SELECT id, properties FROM notion.pages WHERE archived=false")
    pages = cur.fetchall()
    print(f"[check_notion] Found {len(pages)} Notion pages.")

    def _extract_notion_title(props):
        """Extract title plain_text from properties JSON."""
        if not isinstance(props, dict):
            return ""
        for v in props.values():
            if isinstance(v, dict) and (v.get("type") == "title" or "title" in v):
                arr = v.get("title")
                if isinstance(arr, list):
                    out = ""
                    for piece in arr:
                        if isinstance(piece, dict):
                            out += piece.get("plain_text", "")
                            inner = piece.get("text") or {}
                            if isinstance(inner, dict):
                                out += inner.get("content", "")
                    return out
        return ""

    checklist_page_id = None
    for page_id, props in pages:
        title_text = _extract_notion_title(props or {}).strip().lower()
        if "new employee onboarding checklist" in title_text:
            checklist_page_id = page_id
            break
    record("Onboarding checklist page exists", checklist_page_id is not None,
           f"Pages: {len(pages)}")
    if not checklist_page_id:
        cur.close()
        conn.close()
        return

    cur.execute("""SELECT block_data::text FROM notion.blocks WHERE parent_id=%s""", (checklist_page_id,))
    blocks = cur.fetchall()
    all_block_text = " ".join(str(b[0]) for b in blocks)
    all_block_lower = all_block_text.lower()

    record("Notion page contains 'March 2026 New Hires' section",
           "march 2026 new hires" in all_block_lower,
           f"Section heading missing in {len(blocks)} blocks")

    # Each hire's full name must be present (post-update)
    for h in NEW_HIRES:
        record(f"Notion section lists {h['name']}",
               h["name"].lower() in all_block_lower,
               f"Looking for '{h['name']}' in blocks")

    cur.close()
    conn.close()


# ============================================================================
# Main
# ============================================================================

def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--agent_workspace", required=False, default=".")
    parser.add_argument("--groundtruth_workspace", required=False, default=".")
    parser.add_argument("--launch_time", required=False)
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()

    check_pptx(args.agent_workspace)
    check_gcal()
    check_emails()
    check_notion()

    all_passed = (FAIL_COUNT == 0)

    print(f"\n=== SUMMARY ===")
    print(f"  Passed: {PASS_COUNT}")
    print(f"  Failed: {FAIL_COUNT}")
    print(f"  Overall: {'PASS' if all_passed else 'FAIL'}")

    if args.res_log_file:
        result = {
            "passed": PASS_COUNT,
            "failed": FAIL_COUNT,
            "success": all_passed,
        }
        with open(args.res_log_file, "w") as f:
            json.dump(result, f, indent=2)

    sys.exit(0 if all_passed else 1)


if __name__ == "__main__":
    main()
