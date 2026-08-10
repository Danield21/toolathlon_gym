"""
Evaluation for arxiv-conference-prep task.

Checks:
1. PPTX file exists and has at least 7 slides
2. Title slide mentions RLHF
3. Paper slides contain the 5 target RLHF paper keywords
4. Summary/conclusion slide present
5. Calendar event with RLHF on April 10, 2026
6. Email sent to collaborators@rlhf-lab.org
"""

import os
import sys
import json
from argparse import ArgumentParser
from zoneinfo import ZoneInfo, ZoneInfoNotFoundError

import psycopg2
from pptx import Presentation


DB_CONFIG = {
    "host": os.environ.get("PGHOST", "localhost"),
    "port": int(os.environ.get("PGPORT", "5432")),
    "dbname": os.environ.get("PGDATABASE", "toolathlon_gym"),
    "user": os.environ.get("PGUSER", "eigent"),
    "password": os.environ.get("PGPASSWORD", "camel"),
}

UTC = ZoneInfo("UTC")
# The conference is at the San Francisco Convention Center; the conference-local
# (America/Los_Angeles) reading of an event instant is always a valid view.
LA = ZoneInfo("America/Los_Angeles")

TARGET_PAPER_KEYWORDS = [
    "instructgpt",
    "follow instructions",
    "summarize from human feedback",
    "constitutional ai",
    "direct preference optimization",
    "proximal policy optimization",
]

# The task asks to EXCLUDE non-RLHF papers such as self-play / board-game papers.
# We only reject a deck when a slide is *dedicated* to such a noise paper
# (see _is_dedicated_noise_slide); mentioning AlphaZero/self-play in an overview,
# comparison, or exclusion context is legitimate and must not be penalized.
NOISE_KEYWORDS = [
    "alphazero",
    "alpha zero",
    "self-play",
    "self play",
    "board game",
    "chess",
    "shogi",
    "go game",
]

# A slide containing any of these is treated as a context/structural slide
# (overview, summary, comparison, exclusion, an RLHF-related slide, a
# references/bibliography slide, a related-work slide, or an appendix/scope
# slide), not as a dedicated noise-paper slide. The references/related-work/
# appendix/scope entries protect legitimate slides that merely *mention* an
# excluded paper (e.g. a References slide listing "AlphaZero: Mastering Chess
# and Shogi by Self-Play", or a "Related Work" slide) without naming one of the
# RLHF target papers or the generic exclusion words below.
CONTEXT_KEYWORDS = [
    "rlhf",
    "reinforcement learning from human feedback",
    "human feedback",
    "overview",
    "introduction",
    "agenda",
    "table of contents",
    "summary",
    "conclusion",
    "synthesis",
    "themes",
    "key takeaways",
    "comparison",
    "contrast",
    "exclud",
    "filter",
    "omitt",
    "not include",
    "not relevant",
    "reference",
    "references",
    "bibliography",
    "citation",
    "citations",
    "cited",
    "further reading",
    "works cited",
    "appendix",
    "supplement",
    "supplementary",
    "related work",
    "related works",
    "prior work",
    "out of scope",
    "beyond the scope",
    "beyond scope",
    "not discussed",
    "not covered",
    "did not cover",
    "we do not cover",
]

PASS_COUNT = 0
FAIL_COUNT = 0
RUNTIME_ONLY_FAIL = 0


def record(name, passed, detail="", runtime_only=False):
    global PASS_COUNT, FAIL_COUNT, RUNTIME_ONLY_FAIL
    if passed:
        PASS_COUNT += 1
        print(f"  [PASS] {name}")
    else:
        FAIL_COUNT += 1
        if runtime_only:
            RUNTIME_ONLY_FAIL += 1
        msg = f": {detail[:300]}" if detail else ""
        print(f"  [FAIL] {name}{msg}")


def slide_text(slide):
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                texts.append(para.text)
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    texts.append(cell.text)
    return " ".join(texts)


def _is_dedicated_noise_slide(text):
    """True only when a slide is dedicated to a non-RLHF noise paper.

    A slide is flagged as a dedicated noise-paper slide only if it contains a
    noise signature (AlphaZero / self-play / board game / chess / shogi / Go)
    AND does NOT discuss RLHF, any target RLHF paper, or serve as an
    overview/summary/comparison/exclusion context. This keeps the anti-check
    aligned with the task intent ("exclude self-play board-game papers as one
    of the RLHF papers") while not penalizing legitimate contrast mentions.
    """
    low = text.lower()
    if not any(kw in low for kw in NOISE_KEYWORDS):
        return False
    if any(kw in low for kw in CONTEXT_KEYWORDS):
        return False
    if any(kw in low for kw in TARGET_PAPER_KEYWORDS):
        return False
    return True


def slide_title(slide):
    if slide.shapes.title and slide.shapes.title.has_text_frame:
        return slide.shapes.title.text
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text_frame.text.strip():
            return shape.text_frame.text.strip()
    return ""


def check_pptx(agent_workspace):
    """Check the PowerPoint presentation."""
    print("\n=== Check 1: PowerPoint Presentation ===")

    pptx_path = os.path.join(agent_workspace, "RLHF_Conference_Report.pptx")
    if not os.path.exists(pptx_path):
        record("PPTX file exists", False, f"Not found at {pptx_path}")
        return
    record("PPTX file exists", True)

    try:
        prs = Presentation(pptx_path)
    except Exception as e:
        record("PPTX file readable", False, str(e))
        return
    record("PPTX file readable", True)

    slides = list(prs.slides)
    record("At least 7 slides", len(slides) >= 7, f"Found {len(slides)} slides")

    # Check title slide mentions RLHF
    if slides:
        first_text = slide_text(slides[0]).lower()
        has_rlhf = "rlhf" in first_text or "reinforcement learning from human feedback" in first_text
        record("Title slide mentions RLHF", has_rlhf, f"First slide text: {first_text[:100]}")

    # Check a summary/conclusion slide exists anywhere in the deck. The task
    # asks for a summary slide; it does not forbid extra trailing slides such as
    # references / thank-you / Q&A, so we do not require it to be the last slide.
    if slides:
        summary_keywords = ["summary", "conclusion", "themes", "synthesis"]
        has_summary = any(
            any(kw in slide_text(s).lower() for kw in summary_keywords)
            for s in slides
        )
        record("Summary/conclusion slide present", has_summary,
               "No slide contains summary/conclusion/themes/synthesis")

    # Check paper keywords appear across all slides (require 4+, tighter than prior 3+)
    all_text = " ".join(slide_text(s) for s in slides).lower()
    papers_found = sum(1 for kw in TARGET_PAPER_KEYWORDS if kw.lower() in all_text)
    record(
        "At least 4 RLHF paper keywords in slides",
        papers_found >= 4,
        f"Found {papers_found}/{len(TARGET_PAPER_KEYWORDS)} keywords",
    )

    # Check for specific paper titles - require all 5 (per task: 5 specific papers)
    has_instructgpt = "instruct" in all_text and "human feedback" in all_text
    has_summarize = "summarize" in all_text and "human feedback" in all_text
    has_constitutional = "constitutional" in all_text
    has_dpo = "direct preference" in all_text or "dpo" in all_text
    has_ppo = "proximal policy" in all_text or "ppo" in all_text

    papers_present = sum([has_instructgpt, has_summarize, has_constitutional, has_dpo, has_ppo])
    record(
        "All 5 of 5 target papers discussed",
        papers_present == 5,
        f"Found {papers_present}/5 papers",
    )

    # Anti-check: reject slides *dedicated* to a non-RLHF paper (AlphaZero /
    # self-play / board games) as required by the task. The check is per-slide
    # and scoped so that overview / summary / comparison mentions of AlphaZero
    # in otherwise-RLHF context are not penalized.
    noise_slides = [s for s in slides if _is_dedicated_noise_slide(slide_text(s))]
    record("No slide dedicated to AlphaZero / self-play / board-game paper",
           len(noise_slides) == 0,
           f"{len(noise_slides)} slide(s) dedicated to a non-RLHF paper")


def check_calendar():
    """Verify the RLHF Summit 2026 event was created."""
    print("\n=== Check 2: Google Calendar Event ===")

    conn = psycopg2.connect(**DB_CONFIG)
    cur = conn.cursor()
    cur.execute("""
        SELECT summary, description, start_datetime, start_timezone, end_datetime, location
        FROM gcal.events
        ORDER BY start_datetime
    """)
    events = cur.fetchall()
    cur.close()
    conn.close()

    # Look for RLHF event, prefer one matching 2026-04-10
    rlhf_candidates = []
    for summary, description, start_dt, start_tz, end_dt, location in events:
        summary_lower = (summary or "").lower()
        desc_lower = (description or "").lower() if description else ""
        if "rlhf" in summary_lower or "reinforcement learning" in summary_lower or "rlhf" in desc_lower:
            rlhf_candidates.append((summary, description, start_dt, start_tz, end_dt, location))

    # Prefer a candidate that is on the target date; fall back to the first
    rlhf_event = None
    for cand in rlhf_candidates:
        _, _, sdt, _, _, _ = cand
        if sdt is not None and sdt.astimezone(UTC).strftime("%Y-%m-%d") == "2026-04-10":
            rlhf_event = cand
            break
    if rlhf_event is None and rlhf_candidates:
        rlhf_event = rlhf_candidates[0]

    record("Calendar event with RLHF exists", rlhf_event is not None,
           "No event found with RLHF in summary/description",
           runtime_only=True)

    # If a candidate event exists, become blocking (wrong-date/wrong-location shouldn't pass).
    # Full absence is still runtime_only (GT self-test).
    if rlhf_event:
        summary, description, start_dt, start_tz, end_dt, location = rlhf_event
        if start_dt is not None:
            # Guard against naive datetimes (should not happen for timestamptz,
            # but be defensive): treat a missing tzinfo as UTC.
            if start_dt.tzinfo is None:
                start_dt = start_dt.replace(tzinfo=UTC)

            # The conference is in San Francisco. When the agent recorded a
            # timezone on the event, respect it; otherwise assume the conference
            # local timezone.
            ref_tz = UTC
            if start_tz:
                try:
                    ref_tz = ZoneInfo(start_tz)
                except (ZoneInfoNotFoundError, ValueError):
                    ref_tz = UTC

            # Four interpretations of the same instant, so the evaluation is
            # robust to (i) agents that correctly write SF-local time with an
            # RFC3339 offset but omit the optional timeZone field (in which case
            # start_timezone is NULL and we fall back to UTC), (ii) agents that
            # write naive or 'Z' datetimes (which PostgreSQL parses in the
            # session timezone), (iii) the session timezone of the judge's own
            # postgres connection, and (iv) the conference-local (San
            # Francisco) reading, which is always a valid interpretation since
            # the event's location is the San Francisco Convention Center.
            local_dt = start_dt.astimezone(ref_tz)   # event/declared timezone
            utc_dt = start_dt.astimezone(UTC)        # UTC reading
            sess_dt = start_dt                        # session wall clock
            sf_dt = start_dt.astimezone(LA)          # conference-local reading

            def _is_target(dt):
                return dt.strftime("%Y-%m-%d") == "2026-04-10" and dt.hour == 9

            date_ok = any(
                dt.strftime("%Y-%m-%d") == "2026-04-10"
                for dt in (sess_dt, local_dt, utc_dt, sf_dt)
            )
            record("Calendar event on 2026-04-10", date_ok,
                   f"Session {sess_dt.strftime('%Y-%m-%d %H:%M')}, "
                   f"SF {sf_dt.strftime('%Y-%m-%d %H:%M')}, "
                   f"local {local_dt.strftime('%Y-%m-%d %H:%M')}, "
                   f"UTC {utc_dt.strftime('%Y-%m-%d %H:%M')}")

            # Duration 8 hours (9-17); duration is timezone-invariant.
            if end_dt:
                duration_hours = (end_dt - start_dt).total_seconds() / 3600
                record("Event duration exactly 8 hours", abs(duration_hours - 8.0) <= 0.25,
                       f"Got {duration_hours} hours")

            # Start time 9:00 (task says 9am-5pm). Accept 9:00 on 2026-04-10 in
            # ANY of the four interpretations above.
            record("Event starts at 9:00",
                   any(_is_target(dt) for dt in (sess_dt, local_dt, utc_dt, sf_dt)),
                   f"Session {sess_dt.strftime('%Y-%m-%d %H:%M')}, "
                   f"SF {sf_dt.strftime('%Y-%m-%d %H:%M')}, "
                   f"local {local_dt.strftime('%Y-%m-%d %H:%M')}, "
                   f"UTC {utc_dt.strftime('%Y-%m-%d %H:%M')}")
        else:
            record("Calendar event on 2026-04-10", False, "start_datetime is NULL")

        # Check location
        loc_lower = (location or "").lower()
        record("Event location is 'San Francisco Convention Center'",
               "san francisco convention center" in loc_lower,
               f"Got location: {location}")


def check_email():
    """Verify the email was sent to collaborators."""
    print("\n=== Check 3: Email to Collaborators ===")

    conn = psycopg2.connect(**DB_CONFIG)
    cur = conn.cursor()
    cur.execute("""
        SELECT subject, from_addr, to_addr, body_text
        FROM email.messages
    """)
    all_messages = cur.fetchall()
    cur.close()
    conn.close()

    # Find email with RLHF in subject
    matching_email = None
    for subject, from_addr, to_addr, body_text in all_messages:
        subject_lower = (subject or "").lower()
        if "rlhf" in subject_lower or "reinforcement learning" in subject_lower:
            matching_email = (subject, from_addr, to_addr, body_text)
            break

    record("Email with RLHF in subject exists", matching_email is not None,
           "No email found with RLHF in subject",
           runtime_only=True)

    # When email exists, content checks become blocking (not runtime_only).
    if matching_email:
        subject, from_addr, to_addr, body_text = matching_email
        # Check recipient
        to_str = ""
        if isinstance(to_addr, list):
            to_str = " ".join(str(r).lower() for r in to_addr)
        elif isinstance(to_addr, str):
            try:
                parsed = json.loads(to_addr)
                if isinstance(parsed, list):
                    to_str = " ".join(str(r).lower() for r in parsed)
                else:
                    to_str = str(to_addr).lower()
            except (json.JSONDecodeError, TypeError):
                to_str = str(to_addr).lower()

        record("Email to collaborators@rlhf-lab.org",
               "collaborators@rlhf-lab.org" in to_str,
               f"Recipient: {to_addr}")

        # Stricter: primary (first) recipient must be the collaborators address
        primary_recipient = ""
        if isinstance(to_addr, list) and to_addr:
            primary_recipient = str(to_addr[0]).lower()
        elif isinstance(to_addr, str):
            try:
                parsed = json.loads(to_addr)
                if isinstance(parsed, list) and parsed:
                    primary_recipient = str(parsed[0]).lower()
                else:
                    primary_recipient = str(to_addr).lower()
            except (json.JSONDecodeError, TypeError):
                primary_recipient = str(to_addr).lower()
        record("Primary recipient is collaborators@rlhf-lab.org",
               "collaborators@rlhf-lab.org" in primary_recipient,
               f"Primary: {primary_recipient}")

        # Check body mentions conference date. Accept the task's own phrasing
        # ("April 10, 2026") plus common compact/alternate formats a correct
        # agent might use ("April 10th", "Apr 10", "2026-04-10", "4/10/2026",
        # "10 April").
        body_lower = (body_text or "").lower()
        date_forms = [
            "april 10",
            "2026-04-10",
            "10 april",
            "apr 10",
            "4/10/2026",
            "4-10-2026",
            "04/10/2026",
        ]
        has_date = any(form in body_lower for form in date_forms)
        record("Email body mentions conference date", has_date,
               "Date not found in email body")
        # Check body mentions location (per task)
        has_loc = "san francisco" in body_lower or "convention center" in body_lower
        record("Email body mentions location", has_loc,
               "Expected mention of San Francisco / Convention Center")


def main():
    parser = ArgumentParser()
    parser.add_argument("--agent_workspace", required=False, default=".")
    parser.add_argument("--groundtruth_workspace", required=False, default=".")
    parser.add_argument("--launch_time", required=False)
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()

    check_pptx(args.agent_workspace)
    check_calendar()
    check_email()

    total = PASS_COUNT + FAIL_COUNT
    if total == 0:
        print("\nFAIL: No checks were performed.")
        sys.exit(1)

    accuracy = PASS_COUNT / total * 100
    print(f"\nOverall: {PASS_COUNT}/{total} checks passed ({accuracy:.1f}%)")

    result = {
        "total_passed": PASS_COUNT,
        "total_checks": total,
        "accuracy": accuracy,
    }

    if args.res_log_file:
        with open(args.res_log_file, "w") as f:
            json.dump(result, f, indent=2)

    # Gate on non-runtime-only FAIL_COUNT == 0 (was 80% accuracy)
    non_runtime_fail = FAIL_COUNT - RUNTIME_ONLY_FAIL
    if non_runtime_fail == 0:
        print(f"PASS (runtime-only fails: {RUNTIME_ONLY_FAIL})")
        sys.exit(0)
    else:
        print(f"FAIL (non-runtime fails: {non_runtime_fail}, runtime-only fails: {RUNTIME_ONLY_FAIL})")
        sys.exit(1)


if __name__ == "__main__":
    main()
