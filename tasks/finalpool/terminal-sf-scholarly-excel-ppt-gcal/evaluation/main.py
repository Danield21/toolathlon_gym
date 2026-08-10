"""Evaluation for terminal-sf-scholarly-excel-ppt-gcal.
Checks:
1. Sales_Strategy_Analysis.xlsx (4 sheets with correct data)
2. Sales_Strategy_Presentation.pptx (6 slides)
3. Google Calendar event "Q1 Sales Strategy Review"
4. Python scripts exist (analyze_sales_gaps.py, match_recommendations.py, generate_summary.py)
5. Output files exist (sales_gaps.json, research_recommendations.json, executive_summary.txt)
"""
import argparse
import json
import os
import re
import sys
from datetime import date, timezone

import openpyxl
import psycopg2

# R1: every DB connection must read environment variables (PGHOST/PGPORT/PGDATABASE/
# PGUSER/PGPASSWORD) with the same defaults as preprocess/main.py. Never hard-code
# the port (the Harbor runner injects a per-case PGPORT for worker-DB isolation).
DB_CONFIG = {
    "host": os.environ.get("PGHOST", "localhost"),
    "port": int(os.environ.get("PGPORT", "5432")),
    "dbname": os.environ.get("PGDATABASE", "toolathlon_gym"),
    "user": os.environ.get("PGUSER", "eigent"),
    "password": os.environ.get("PGPASSWORD", "camel"),
}

PASS_COUNT = 0
FAIL_COUNT = 0


def get_expected_from_db():
    """Query sf_data schema dynamically for regional revenue values.

    The task brief is 'Q1 2026', so a reasonable model may aggregate regional
    revenue either over the full sales history (the ground-truth interpretation)
    or over the current calendar quarter. To avoid false FAILs for a correct
    model that picks the quarterly reading, this returns BOTH aggregates:
      * *_rev        -> full-period total (all orders)
      * *_rev_q      -> current-calendar-quarter total (quarter of the newest
                        order date, e.g. Q1 2026 for the seeded data)
    Checks accept either value within tolerance. (Fix2 issue 1)
    """
    defaults = {
        "europe_rev": 648798.0,
        "latam_rev": 549129.0,
        "total_rev": 3048998.0,
        "europe_rev_q": 58788.81,
        "latam_rev_q": 43356.91,
        "total_rev_q": 267149.36,
    }

    def _fill(rows, key_europe, key_latam):
        for region, rev in rows:
            rev_f = float(rev)
            if region and region.strip().lower() == "europe":
                defaults[key_europe] = rev_f
            elif region and "latin" in region.strip().lower():
                defaults[key_latam] = rev_f

    try:
        conn = psycopg2.connect(**DB_CONFIG)
        cur = conn.cursor()
        cur.execute("""
            SELECT c."REGION", SUM(o."TOTAL_AMOUNT") as rev
            FROM sf_data."SALES_DW__PUBLIC__ORDERS" o
            JOIN sf_data."SALES_DW__PUBLIC__CUSTOMERS" c ON o."CUSTOMER_ID" = c."CUSTOMER_ID"
            GROUP BY c."REGION"
        """)
        rows = cur.fetchall()
        defaults["total_rev"] = sum(float(r) for _, r in rows)
        _fill(rows, "europe_rev", "latam_rev")

        # Current-calendar-quarter window (the quarter containing the newest order).
        cur.execute("SELECT MAX(o.\"ORDER_DATE\") FROM sf_data.\"SALES_DW__PUBLIC__ORDERS\" o")
        max_date = cur.fetchone()[0]
        if max_date is not None:
            q_start = date(max_date.year, ((max_date.month - 1) // 3) * 3 + 1, 1)
            q_end_month = q_start.month + 3
            q_end_year = q_start.year
            if q_end_month > 12:
                q_end_month = 1
                q_end_year += 1
            q_end = date(q_end_year, q_end_month, 1)
            cur.execute("""
                SELECT c."REGION", SUM(o."TOTAL_AMOUNT") as rev
                FROM sf_data."SALES_DW__PUBLIC__ORDERS" o
                JOIN sf_data."SALES_DW__PUBLIC__CUSTOMERS" c ON o."CUSTOMER_ID" = c."CUSTOMER_ID"
                WHERE o."ORDER_DATE" >= %s AND o."ORDER_DATE" < %s
                GROUP BY c."REGION"
            """, (q_start, q_end))
            qrows = cur.fetchall()
            defaults["total_rev_q"] = sum(float(r) for _, r in qrows)
            _fill(qrows, "europe_rev_q", "latam_rev_q")
        cur.close()
        conn.close()
    except Exception as e:
        print(f"  [WARN] DB query for expected values failed, using defaults: {e}")
    return defaults


EXPECTED = get_expected_from_db()


def check(name, condition, detail=""):
    global PASS_COUNT, FAIL_COUNT
    if condition:
        PASS_COUNT += 1
        print(f"  [PASS] {name}")
    else:
        FAIL_COUNT += 1
        print(f"  [FAIL] {name}: {str(detail)[:200]}")


def _to_float(v):
    """Parse a numeric value from int/float/str/None.

    Strips thousand-separator commas, currency symbols ($/¥/€/£) and a trailing
    '%' (e.g. "90%" -> 90.0). Returns None when the value is absent or cannot
    be parsed. (R3)
    """
    if v is None or isinstance(v, bool):
        return None
    if isinstance(v, (int, float)):
        return float(v)
    if isinstance(v, str):
        s = v.strip().replace(",", "").replace("$", "").replace("¥", "").replace("€", "").replace("£", "").strip()
        if s.endswith("%"):
            s = s[:-1].strip()
        if not s:
            return None
        try:
            return float(s)
        except ValueError:
            return None
    return None


def num_close(a, b, tol=2.0):
    """Robust numeric comparison. (R3)
    If both sides parse to float, compare with tolerance. Only when one side is
    non-numeric/absent does it fall back to a case-insensitive string comparison.
    """
    fa, fb = _to_float(a), _to_float(b)
    if fa is not None and fb is not None:
        return abs(fa - fb) <= tol
    if a is None or b is None:
        return a is None and b is None
    return str(a).strip().lower() == str(b).strip().lower()


_MONEY_TOKEN = re.compile(
    r"[\$€£¥]?\s?\d[\d,]*(?:\.\d+)?\s?(?:m|k|million|thousand)?", re.IGNORECASE)


def _parse_amount(token):
    """Parse a money token like '$3.05M', '3,048,998', '3.05 million', '3,049K'.
    Returns a float, or None when unparsable.
    """
    s = token.strip().replace(",", "").replace("$", "").replace("€", "")
    s = s.replace("£", "").replace("¥", "").strip()
    if not s:
        return None
    low = s.lower()
    mult = 1.0
    for suffix, m in (("million", 1e6), ("thousand", 1e3)):
        if suffix in low:
            mult = m
            s = low.replace(suffix, "").strip()
            break
    if s.endswith(("m", "k")):
        mult = 1e6 if s[-1] == "m" else 1e3
        s = s[:-1].strip()
    try:
        return float(s) * mult
    except ValueError:
        return None


def _near_candidate_total(content, candidates):
    """True if any money amount in `content` is within tolerance of a candidate
    total. Tolerates rounding like '$3.05M' or '3.0 million' against the exact
    DB aggregate. (Fix2 issue 5)
    """
    for m in _MONEY_TOKEN.finditer(content):
        val = _parse_amount(m.group(0))
        if val is None:
            continue
        for cand in candidates:
            if cand and abs(val - cand) <= max(200000.0, cand * 0.05):
                return True
    return False


def _load_workbook_pair(path):
    """Load a workbook twice: data_only=False (preserve formulas/structure) and
    data_only=True (cached formula values). Returns (wb, wb_data) or (None,None). (R2)
    """
    try:
        wb = openpyxl.load_workbook(path, data_only=False)
    except Exception as e:
        return None, None, e
    try:
        wb_data = openpyxl.load_workbook(path, data_only=True)
    except Exception:
        wb_data = None
    return wb, wb_data, None


def _sheet_rows(ws):
    """Yield (row_idx, values) for non-blank rows."""
    for r_i in range(1, ws.max_row + 1):
        vals = [ws.cell(row=r_i, column=c).value for c in range(1, ws.max_column + 1)]
        if any(v is not None and str(v).strip() for v in vals):
            yield r_i, vals


def _find_header_row(ws, header_keywords, min_hits=2):
    """Probe for the header row by requiring >= min_hits keyword hits."""
    for r_i, vals in _sheet_rows(ws):
        text = " ".join(str(c).lower() for c in vals if c is not None)
        hits = sum(1 for k in header_keywords if k in text)
        if hits >= min_hits:
            return r_i
    return None


def _data_rows(ws, wb_data, header_keywords, min_hits=2):
    """Return data rows below the header row.

    Probes and skips any title/header row, drops blank rows, and resolves formula
    cells to cached values from the data_only workbook when available. (R2/R10)
    """
    hr = _find_header_row(ws, header_keywords, min_hits)
    ws_v = wb_data[ws.title] if wb_data is not None and ws.title in wb_data.sheetnames else None
    rows = []
    for r_i, vals in _sheet_rows(ws):
        if hr is not None and r_i <= hr:
            continue
        resolved = []
        for c_i, v in enumerate(vals, start=1):
            if isinstance(v, str) and v.lstrip().startswith("="):
                # Formula cell: prefer the cached value (data_only workbook).
                # When there is no cached value and the groundtruth is a literal,
                # the check is skipped (R2) rather than failing on the formula text.
                cached = None
                if ws_v is not None:
                    try:
                        cached = ws_v.cell(row=r_i, column=c_i).value
                    except Exception:
                        cached = None
                resolved.append(cached if cached is not None else None)
            else:
                resolved.append(v)
        rows.append(tuple(resolved))
    return rows


def check_excel(workspace):
    print("\n=== Check 1: Sales_Strategy_Analysis.xlsx ===")
    path = os.path.join(workspace, "Sales_Strategy_Analysis.xlsx")
    if not os.path.exists(path):
        check("Excel file exists", False, f"Not found at {path}")
        return
    check("Excel file exists", True)

    wb, wb_data, err = _load_workbook_pair(path)
    if wb is None:
        check("Excel readable", False, str(err))
        return
    check("Excel readable", True)

    sheets = wb.sheetnames
    check("Has 4 sheets", len(sheets) >= 4, f"Found {len(sheets)}: {sheets}")

    sheets_lower = [s.lower().replace(" ", "_") for s in sheets]

    # Sheet 1: Regional_Performance
    rp_idx = next((i for i, s in enumerate(sheets_lower) if "regional" in s or "region" in s), 0)
    ws1 = wb[sheets[rp_idx]]
    data1 = _data_rows(ws1, wb_data, ["region", "revenue", "order", "priority"])
    check("Regional_Performance has 5 region rows", len(data1) >= 5, f"Found {len(data1)}")

    all_text1 = " ".join(str(c) for r in data1 for c in r if c).lower()
    check("Contains Europe region", "europe" in all_text1)
    check("Contains Latin America region", "latin" in all_text1)
    check("Contains Priority Focus tag", "priority" in all_text1 and "focus" in all_text1,
          f"Text: {all_text1[:150]}")
    check("Contains On Track tag", "on track" in all_text1 or "on_track" in all_text1 or "ontrack" in all_text1,
          f"Text: {all_text1[:150]}")

    # Check revenue values are reasonable (within tolerance of actual DB values).
    # Accepts either the full-period aggregate or the current-calendar-quarter
    # aggregate (Fix2 issue 1). Formula-aware: an uncached formula cell is
    # skipped (R2 courtesy — the model computed, just as a formula), but a truly
    # empty revenue cell FAILs so the arithmetic check cannot be bypassed.
    def _rev_for(region_sub):
        """Return (state, value) for the region's revenue cell.
        state in {'number', 'empty', 'formula_no_cache', 'not_found'}.
        """
        hr = _find_header_row(ws1, ["region", "revenue", "order", "priority"])
        if hr is None:
            return "not_found", None
        rev_col = 2
        for c in range(1, ws1.max_column + 1):
            h = ws1.cell(row=hr, column=c).value
            if h and "revenue" in str(h).lower():
                rev_col = c
                break
        for r_i in range(hr + 1, ws1.max_row + 1):
            first = ws1.cell(row=r_i, column=1).value
            if first is None or region_sub not in str(first).lower():
                continue
            raw = ws1.cell(row=r_i, column=rev_col).value
            if isinstance(raw, str) and raw.lstrip().startswith("="):
                cached = None
                if wb_data is not None and ws1.title in wb_data.sheetnames:
                    cached = wb_data[ws1.title].cell(row=r_i, column=rev_col).value
                if cached is not None:
                    return "number", cached
                return "formula_no_cache", None
            if raw is None or (isinstance(raw, str) and not raw.strip()):
                return "empty", None
            return "number", raw
        return "not_found", None

    for label, sub, k_full, k_q in [
            ("Europe", "europe", "europe_rev", "europe_rev_q"),
            ("Latin America", "latin", "latam_rev", "latam_rev_q")]:
        state, val = _rev_for(sub)
        if state == "number":
            ok = num_close(val, EXPECTED[k_full], tol=5000) or num_close(val, EXPECTED[k_q], tol=5000)
            check(f"{label} revenue reasonable", ok,
                  f"Got {val}, expected ~{EXPECTED[k_full]:.0f} (full) or ~{EXPECTED[k_q]:.0f} (Q)")
        elif state == "empty":
            check(f"{label} revenue reasonable", False, f"{label} revenue cell is empty")
        elif state == "formula_no_cache":
            check(f"{label} revenue reasonable", True, "uncached formula (skipped)")
        else:
            check(f"{label} revenue reasonable", False, f"{label} region row not found")

    # Check priority tags: Middle East and Latin America should be Priority Focus.
    # Detect the tag by scanning the whole row text (robust to column ordering /
    # extra trailing columns), and region by the first cell.
    priority_regions = []
    for row in data1:
        rowtext = " ".join(str(c).lower() for c in row if c is not None)
        if "priority" in rowtext and "focus" in rowtext and row[0] is not None:
            priority_regions.append(str(row[0]).lower())
    check("Middle East is Priority Focus", any("middle" in r for r in priority_regions),
          f"Priority regions: {priority_regions}")
    check("Latin America is Priority Focus", any("latin" in r for r in priority_regions),
          f"Priority regions: {priority_regions}")

    # Sheet 2: Segment_Analysis
    sa_idx = next((i for i, s in enumerate(sheets_lower) if "segment" in s), 1)
    if sa_idx < len(sheets):
        ws2 = wb[sheets[sa_idx]]
        data2 = _data_rows(ws2, wb_data, ["segment", "customer", "spend", "ltv"])
        check("Segment_Analysis has 4 segment rows", len(data2) >= 4, f"Found {len(data2)}")
        all_text2 = " ".join(str(c) for r in data2 for c in r if c).lower()
        check("Contains Consumer segment", "consumer" in all_text2)
        check("Contains Enterprise segment", "enterprise" in all_text2)

    # Sheet 3: Research_Insights
    ri_idx = next((i for i, s in enumerate(sheets_lower) if "research" in s or "insight" in s), 2)
    if ri_idx < len(sheets):
        ws3 = wb[sheets[ri_idx]]
        data3 = _data_rows(ws3, wb_data, ["paper", "key_finding", "applicable", "impact"])
        check("Research_Insights has at least 2 rows", len(data3) >= 2, f"Found {len(data3)}")
        all_text3 = " ".join(str(c) for r in data3 for c in r if c).lower()
        # Any of the four relevant papers' signature terms (paper titles /
        # key findings) is accepted, so a solution built on a single relevant
        # paper still passes. (Fix2 issue 4)
        check("Research references territory or segmentation",
              "territory" in all_text3 or "segmentation" in all_text3 or "optimization" in all_text3
              or "specialization" in all_text3 or "pricing" in all_text3,
              f"Text: {all_text3[:150]}")

    # Sheet 4: Action_Items
    ai_idx = next((i for i, s in enumerate(sheets_lower) if "action" in s), 3)
    if ai_idx < len(sheets):
        ws4 = wb[sheets[ai_idx]]
        data4 = _data_rows(ws4, wb_data, ["action", "region", "owner", "research", "timeline"])
        check("Action_Items has at least 2 rows", len(data4) >= 2, f"Found {len(data4)}")
        all_text4 = " ".join(str(c) for r in data4 for c in r if c).lower()
        check("Action items reference research papers",
              "territory" in all_text4 or "segmentation" in all_text4 or "specialization" in all_text4 or "pricing" in all_text4,
              f"Text: {all_text4[:150]}")
        check("Action items include owner names",
              "ahmed" in all_text4 or "carlos" in all_text4 or "hassan" in all_text4 or "rivera" in all_text4,
              f"Text: {all_text4[:150]}")


def check_pptx(workspace):
    print("\n=== Check 2: Sales_Strategy_Presentation.pptx ===")
    path = os.path.join(workspace, "Sales_Strategy_Presentation.pptx")
    if not os.path.exists(path):
        check("PPTX file exists", False, f"Not found at {path}")
        return
    check("PPTX file exists", True)

    try:
        from pptx import Presentation
        prs = Presentation(path)
        slides = list(prs.slides)
        check("Has 6 slides", len(slides) >= 6, f"Found {len(slides)} slides")

        # Check slide content
        all_text = ""
        for slide in slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    all_text += shape.text_frame.text + " "

        all_lower = all_text.lower()
        check("Title slide mentions Q1 or strategy", "q1" in all_lower or "strategy" in all_lower,
              f"Text: {all_lower[:200]}")
        check("Contains regional data", "europe" in all_lower or "region" in all_lower,
              f"Text snippet: {all_lower[:200]}")
        check("Contains segment data", "consumer" in all_lower or "segment" in all_lower,
              f"Text snippet: {all_lower[:200]}")
        check("Contains research findings", "research" in all_lower or "finding" in all_lower or "study" in all_lower,
              f"Text snippet: {all_lower[:200]}")
        check("Contains action plan", "action" in all_lower or "plan" in all_lower or "initiative" in all_lower,
              f"Text snippet: {all_lower[:200]}")
        check("Contains timeline", "timeline" in all_lower or "milestone" in all_lower or "week" in all_lower,
              f"Text snippet: {all_lower[:200]}")
    except ImportError:
        check("python-pptx available", False, "Cannot import pptx module")


def check_gcal():
    print("\n=== Check 3: Calendar Event ===")
    try:
        conn = psycopg2.connect(**DB_CONFIG)
        cur = conn.cursor()
    except Exception as e:
        # A DB failure must degrade to a failed check, never crash the evaluator.
        check("Q1 Sales Strategy Review event exists", False, f"DB connect failed: {e}")
        return

    try:
        cur.execute("""
            SELECT summary, start_datetime, end_datetime, description
            FROM gcal.events
            WHERE lower(summary) LIKE '%%q1%%strategy%%'
               OR lower(summary) LIKE '%%sales%%strategy%%review%%'
            ORDER BY start_datetime
        """)
        events = cur.fetchall()
        check("Q1 Sales Strategy Review event exists", len(events) >= 1,
              f"Found {len(events)} matching events")

        if events:
            ev = events[0]
            summary, start, end, desc = ev
            # Check it's 90 minutes
            if start and end:
                duration = (end - start).total_seconds() / 60
                check("Event is 90 minutes", num_close(duration, 90, tol=15),
                      f"Duration: {duration} minutes")

            # Check during business hours (9-5; a 90-min meeting that starts
            # by 16:00 still ends before 17:30). psycopg2 returns timestamptz in
            # the DB session timezone, but the google_calendar MCP may store an
            # event with a different offset (e.g. UTC). To avoid a false FAIL for
            # a correctly scheduled business-hours meeting, accept the hour in the
            # DB session timezone OR in UTC. (Fix2 issue 2)
            if start:
                hour = start.hour
                utc_hour = start.astimezone(timezone.utc).hour
                check("Event during business hours",
                      9 <= hour <= 16 or 9 <= utc_hour <= 16,
                      f"Start hour: {hour} (session) / {utc_hour} (UTC)")

            # Check weekday (accept session-timezone OR UTC view for the same
            # offset-robustness as the business-hours check).
            if start:
                utc = start.astimezone(timezone.utc)
                check("Event on a weekday", start.weekday() < 5 or utc.weekday() < 5,
                      f"Day: {start.strftime('%A')} / {utc.strftime('%A')} (UTC)")

            # Check no conflict with existing events. Exclude ALL strategy-review
            # events so concurrent swarm agents each creating the same intended
            # event are not counted as conflicting with one another.
            if start and end:
                cur.execute("""
                    SELECT COUNT(*) FROM gcal.events
                    WHERE NOT (lower(summary) LIKE '%%q1%%strategy%%'
                           OR lower(summary) LIKE '%%sales%%strategy%%review%%')
                      AND start_datetime < %s AND end_datetime > %s
                """, (end, start))
                conflicts = cur.fetchone()[0]
                check("No calendar conflicts", conflicts == 0,
                      f"Found {conflicts} conflicting events")

            # Check description mentions review
            if desc:
                check("Description mentions review or action",
                      "review" in desc.lower() or "action" in desc.lower() or "regional" in desc.lower(),
                      f"Description: {str(desc)[:100]}")
    finally:
        cur.close()
        conn.close()


def check_reverse_validation(workspace):
    print("\n=== Reverse Validation ===")
    try:
        conn = psycopg2.connect(**DB_CONFIG)
        cur = conn.cursor()
    except Exception as e:
        check("No strategy review events on weekends", False, f"DB connect failed: {e}")
        return
    try:
        # Check no gcal events on weekends
        cur.execute("""
            SELECT summary, start_datetime
            FROM gcal.events
            WHERE (lower(summary) LIKE '%%q1%%strategy%%'
               OR lower(summary) LIKE '%%sales%%strategy%%review%%')
              AND EXTRACT(DOW FROM start_datetime) IN (0, 6)
        """)
        weekend_events = cur.fetchall()
        check("No strategy review events on weekends", len(weekend_events) == 0,
              f"Found {len(weekend_events)} weekend events: {weekend_events}")

        # Duplicate tolerance: this is a swarm task (SWARM_PARALLEL=1) where
        # multiple concurrent agents may each create the same named event. Accept
        # a small number of copies; only pathological duplicate spam fails.
        cur.execute("""
            SELECT COUNT(*) FROM gcal.events
            WHERE lower(summary) LIKE '%%q1%%strategy%%'
               OR lower(summary) LIKE '%%sales%%strategy%%review%%'
        """)
        event_count = cur.fetchone()[0]
        check("Strategy review events limited (no duplicate spam)", event_count <= 3,
              f"Found {event_count} strategy review events")

        # Check Research_Insights sheet does not contain noise papers
        # (healthcare, machine learning theory papers should be excluded)
        xlsx_path = os.path.join(workspace, "Sales_Strategy_Analysis.xlsx")
        if os.path.exists(xlsx_path):
            wb, wb_data, _ = _load_workbook_pair(xlsx_path)
            if wb is not None:
                sheets_lower = [s.lower().replace(" ", "_") for s in wb.sheetnames]
                ri_idx = next((i for i, s in enumerate(sheets_lower) if "research" in s or "insight" in s), None)
                if ri_idx is not None:
                    ws = wb[wb.sheetnames[ri_idx]]
                    data = _data_rows(ws, wb_data, ["paper", "key_finding", "applicable", "impact"])
                    all_text = " ".join(str(c) for r in data for c in r if c).lower()
                    # Expanded noise topics: healthcare, ML theory
                    noise_topics = ["healthcare", "medical", "clinical trial", "genomic",
                                    "federated learning", "convergence analysis",
                                    "non-iid", "fedavg", "fedprox", "oncology",
                                    "cardiology", "neurology", "patient recruitment"]
                    for topic in noise_topics:
                        check(f"Research_Insights does not contain noise topic '{topic}'",
                              topic not in all_text,
                              f"Found '{topic}' in Research_Insights")

        # Reverse: noise papers preserved in scholarly DB
        cur.execute("""
            SELECT COUNT(*) FROM scholarly.scholar_papers
            WHERE lower(title) LIKE '%%clinical trial%%'
               OR lower(title) LIKE '%%federated learning%%'
        """)
        noise_paper_count = cur.fetchone()[0]
        check("Reverse: noise scholarly papers preserved",
              noise_paper_count >= 2,
              f"Only {noise_paper_count}/2 noise papers remain")
    except Exception as e:
        check("Reverse validation", False, str(e))
    finally:
        cur.close()
        conn.close()


def check_scripts(workspace):
    print("\n=== Check 4: Python Scripts ===")
    check("analyze_sales_gaps.py exists",
          os.path.exists(os.path.join(workspace, "analyze_sales_gaps.py")))
    check("match_recommendations.py exists",
          os.path.exists(os.path.join(workspace, "match_recommendations.py")))
    check("generate_summary.py exists",
          os.path.exists(os.path.join(workspace, "generate_summary.py")))


def check_outputs(workspace):
    print("\n=== Check 5: Output Files ===")

    # sales_gaps.json
    sg_path = os.path.join(workspace, "sales_gaps.json")
    if os.path.exists(sg_path):
        check("sales_gaps.json exists", True)
        try:
            with open(sg_path) as f:
                data = json.load(f)
            check("sales_gaps.json is valid JSON", True)
            text = json.dumps(data).lower()
            check("sales_gaps contains priority tags",
                  "priority" in text or "focus" in text,
                  f"Keys: {list(data.keys()) if isinstance(data, dict) else 'list'}")
        except Exception as e:
            check("sales_gaps.json is valid JSON", False, str(e))
    else:
        check("sales_gaps.json exists", False)

    # research_recommendations.json
    rr_path = os.path.join(workspace, "research_recommendations.json")
    if os.path.exists(rr_path):
        check("research_recommendations.json exists", True)
        try:
            with open(rr_path) as f:
                data = json.load(f)
            check("research_recommendations.json is valid JSON", True)
            text = json.dumps(data).lower()
            # Any of the four relevant papers' signature terms, or the paper_title
            # field itself, proves the recommendations reference research. A
            # solution built on a single relevant paper (e.g. Sales Team
            # Specialization) still passes. (Fix2 issue 4)
            check("Recommendations reference papers",
                  "territory" in text or "segmentation" in text or "optimization" in text
                  or "specialization" in text or "pricing" in text
                  or "paper" in text or "title" in text,
                  f"Content: {text[:150]}")
        except Exception as e:
            check("research_recommendations.json is valid JSON", False, str(e))
    else:
        check("research_recommendations.json exists", False)

    # executive_summary.txt
    es_path = os.path.join(workspace, "executive_summary.txt")
    if os.path.exists(es_path):
        check("executive_summary.txt exists", True)
        with open(es_path) as f:
            content = f.read().lower()
        total_rev_str = f"{EXPECTED['total_rev']:.0f}"
        total_rev_comma = f"{EXPECTED['total_rev']:,.0f}"
        # Accept the word 'revenue', the exact total digits, or a rounded money
        # expression (e.g. '$3.05M', '3.0 million', '3,049K') near either the
        # full-period or the current-quarter total. (Fix2 issue 5 / issue 1)
        near_total = _near_candidate_total(
            content, [EXPECTED['total_rev'], EXPECTED.get('total_rev_q')])
        check("Summary mentions total revenue",
              "revenue" in content or total_rev_str[:4] in content
              or total_rev_comma[:5] in content or near_total,
              f"Content: {content[:150]}")
        check("Summary mentions priority regions",
              "priority" in content or "focus" in content or "underperform" in content,
              f"Content: {content[:150]}")
    else:
        check("executive_summary.txt exists", False)


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--agent_workspace", required=False, default=".")
    parser.add_argument("--groundtruth_workspace", required=False, default=".")
    parser.add_argument("--launch_time", required=False)
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()

    check_excel(args.agent_workspace)
    check_pptx(args.agent_workspace)
    check_gcal()
    check_scripts(args.agent_workspace)
    check_outputs(args.agent_workspace)
    check_reverse_validation(args.agent_workspace)

    total = PASS_COUNT + FAIL_COUNT
    if total == 0:
        print("\nFAIL: No checks performed.")
        sys.exit(1)

    accuracy = PASS_COUNT / total * 100
    print(f"\nOverall: {PASS_COUNT}/{total} checks passed ({accuracy:.1f}%)")

    result = {"total_passed": PASS_COUNT, "total_checks": total, "accuracy": accuracy}
    if args.res_log_file:
        with open(args.res_log_file, "w") as f:
            json.dump(result, f, indent=2)

    # Tightened: require ALL checks to pass (previously >= 70%).
    if FAIL_COUNT == 0:
        print("PASS")
        sys.exit(0)
    else:
        print("FAIL")
        sys.exit(1)


if __name__ == "__main__":
    main()
