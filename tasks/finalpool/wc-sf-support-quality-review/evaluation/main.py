"""Evaluation for wc-sf-support-quality-review."""
import argparse
import os
import sys

def num_close(a, b, rel_tol=0.15, abs_tol=0.5):
    return abs(float(a) - float(b)) <= max(abs_tol, abs(float(b)) * rel_tol)



def nums_close(a, b, abs_tol=1.0, rel_tol=0.05):
    try:
        a, b = float(a), float(b)
    except (TypeError, ValueError):
        return False
    if abs(a - b) <= abs_tol:
        return True
    if b != 0 and abs(a - b) / abs(b) <= rel_tol:
        return True
    return False


def load_sheet_rows(wb, sheet_name):
    for name in wb.sheetnames:
        if name.strip().lower() == sheet_name.strip().lower():
            return [[cell.value for cell in row] for row in wb[name].iter_rows()]
    return None


def check_excel(agent_workspace, gt_workspace):
    errors = []
    import openpyxl

    path = os.path.join(agent_workspace, "CX_Quality_Review.xlsx")
    gt_path = os.path.join(gt_workspace, "CX_Quality_Review.xlsx")
    if not os.path.exists(path):
        return ["CX_Quality_Review.xlsx not found"]
    if not os.path.exists(gt_path):
        return ["GT CX_Quality_Review.xlsx not found"]
    try:
        wb = openpyxl.load_workbook(path, data_only=True)
        gt_wb = openpyxl.load_workbook(gt_path, data_only=True)

        # Sheet 1: Order Issues - validate per-category rows
        rows = load_sheet_rows(wb, "Order Issues")
        gt_rows = load_sheet_rows(gt_wb, "Order Issues")
        if rows is None:
            errors.append("Sheet 'Order Issues' not found")
        elif gt_rows is None:
            errors.append("GT Sheet 'Order Issues' not found")
        else:
            data_rows = [r for r in rows[1:] if r and r[0] is not None]
            gt_data = [r for r in gt_rows[1:] if r and r[0] is not None]
            if len(data_rows) != len(gt_data):
                errors.append(f"Order Issues has {len(data_rows)} rows, expected {len(gt_data)}")
            a_lookup = {str(r[0]).strip().lower(): r for r in data_rows}
            for g in gt_data:
                key = str(g[0]).strip().lower()
                a = a_lookup.get(key)
                if a is None:
                    errors.append(f"Order Issues: missing category '{g[0]}'")
                    continue
                # Total_Orders (col 1) - exact integer
                if a[1] is None or int(a[1]) != int(g[1]):
                    errors.append(f"Order Issues {key} Total_Orders={a[1]}, expected {g[1]}")
                # Issue_Orders (col 2) - exact integer
                if a[2] is None or int(a[2]) != int(g[2]):
                    errors.append(f"Order Issues {key} Issue_Orders={a[2]}, expected {g[2]}")
                # Issue_Rate_Pct (col 3) - tol 0.5
                if a[3] is None or abs(float(a[3]) - float(g[3])) > 0.5:
                    errors.append(f"Order Issues {key} Issue_Rate_Pct={a[3]}, expected {g[3]}")

        # Sheet 2: Support Metrics - validate all 3 priority rows
        rows2 = load_sheet_rows(wb, "Support Metrics")
        gt_rows2 = load_sheet_rows(gt_wb, "Support Metrics")
        if rows2 is None:
            errors.append("Sheet 'Support Metrics' not found")
        elif gt_rows2 is None:
            errors.append("GT Sheet 'Support Metrics' not found")
        else:
            data_rows2 = [r for r in rows2[1:] if r and r[0] is not None]
            gt_data2 = [r for r in gt_rows2[1:] if r and r[0] is not None]
            if len(data_rows2) != 3:
                errors.append(f"Support Metrics has {len(data_rows2)} rows, expected exactly 3")
            a_lookup2 = {str(r[0]).strip().lower(): r for r in data_rows2}
            for g in gt_data2:
                key = str(g[0]).strip().lower()
                a = a_lookup2.get(key)
                if a is None:
                    errors.append(f"Support Metrics: missing priority '{g[0]}'")
                    continue
                # Ticket_Count exact (allow tol 50 since DB might have noise)
                if a[1] is None or abs(int(a[1]) - int(g[1])) > 50:
                    errors.append(f"Support Metrics {key} Ticket_Count={a[1]}, expected ~{g[1]}")
                # Avg_Response_Time_Hours tol 0.5
                if a[2] is None or abs(float(a[2]) - float(g[2])) > 0.5:
                    errors.append(f"Support Metrics {key} Avg_Response_Time={a[2]}, expected ~{g[2]}")
                # Avg_Satisfaction tol 0.1
                if a[3] is None or abs(float(a[3]) - float(g[3])) > 0.1:
                    errors.append(f"Support Metrics {key} Avg_Satisfaction={a[3]}, expected ~{g[3]}")
                # Response_Target_Hours exact (must come from portal)
                if len(a) > 5 and a[5] is not None:
                    try:
                        if abs(float(a[5]) - float(g[5])) > 0.1:
                            errors.append(f"Support Metrics {key} Response_Target={a[5]}, expected {g[5]}")
                    except (ValueError, TypeError):
                        errors.append(f"Support Metrics {key} Response_Target not numeric: {a[5]}")
                # Meets_Response_Target Yes/No
                if len(a) > 6 and a[6] is not None and len(g) > 6 and g[6] is not None:
                    if str(a[6]).strip().lower() not in ("yes", "no"):
                        errors.append(f"Support Metrics {key} Meets_Response_Target={a[6]} not yes/no")
                    elif str(a[6]).strip().lower() != str(g[6]).strip().lower():
                        errors.append(f"Support Metrics {key} Meets_Response_Target={a[6]}, expected {g[6]}")

        # Sheet 3: Cross Reference
        rows3 = load_sheet_rows(wb, "Cross Reference")
        gt_rows3 = load_sheet_rows(gt_wb, "Cross Reference")
        if rows3 is None:
            errors.append("Sheet 'Cross Reference' not found")
        elif gt_rows3 is None:
            errors.append("GT Sheet 'Cross Reference' not found")
        else:
            data_rows3 = [r for r in rows3[1:] if r and r[0] is not None]
            gt_data3 = [r for r in gt_rows3[1:] if r and r[0] is not None]
            if len(data_rows3) != len(gt_data3):
                errors.append(f"Cross Reference has {len(data_rows3)} rows, expected {len(gt_data3)}")
            a_lookup3 = {str(r[0]).strip().lower(): r for r in data_rows3}
            for g in gt_data3:
                key = str(g[0]).strip().lower()
                a = a_lookup3.get(key)
                if a is None:
                    errors.append(f"Cross Reference: missing issue type '{g[0]}'")
                    continue
                # Ticket_Count tol 50
                if a[1] is None or abs(int(a[1]) - int(g[1])) > 50:
                    errors.append(f"Cross Reference {key} Ticket_Count={a[1]}, expected ~{g[1]}")
                # Avg_Response_Time tol 0.5
                if a[2] is None or abs(float(a[2]) - float(g[2])) > 0.5:
                    errors.append(f"Cross Reference {key} Avg_Response_Time={a[2]}, expected ~{g[2]}")
                # Avg_Satisfaction tol 0.15
                if a[3] is None or abs(float(a[3]) - float(g[3])) > 0.15:
                    errors.append(f"Cross Reference {key} Avg_Satisfaction={a[3]}, expected ~{g[3]}")

        # Sheet 4: Executive Summary
        rows4 = load_sheet_rows(wb, "Executive Summary")
        gt_rows4 = load_sheet_rows(gt_wb, "Executive Summary")
        if rows4 is None:
            errors.append("Sheet 'Executive Summary' not found")
        elif gt_rows4 is None:
            errors.append("GT Sheet 'Executive Summary' not found")
        else:
            data_rows4 = [r for r in rows4[1:] if r and r[0] is not None]
            gt_data4 = [r for r in gt_rows4[1:] if r and r[0] is not None]
            if len(data_rows4) < 5:
                errors.append(f"Executive Summary has {len(data_rows4)} rows, expected >= 5")
            a_lookup4 = {str(r[0]).strip().lower(): r for r in data_rows4}
            for g in gt_data4:
                key = str(g[0]).strip().lower()
                a = a_lookup4.get(key)
                if a is None:
                    errors.append(f"Executive Summary: missing metric '{g[0]}'")
                    continue
                if isinstance(g[1], str):
                    if str(a[1] or "").strip().lower() != g[1].strip().lower():
                        errors.append(f"Executive Summary {key}={a[1]}, expected {g[1]}")
                else:
                    try:
                        gv = float(g[1])
                        # Tol: integer counts: 50 leniency for tickets; rates 0.5; ratings 0.1
                        if "rate" in key or "satisfaction" in key or "rating" in key:
                            tol = max(0.1, abs(gv) * 0.05)
                        elif "ticket" in key:
                            tol = 50
                        elif "order" in key:
                            tol = 5
                        elif "time" in key:
                            tol = 0.5
                        else:
                            tol = 0.1
                        if a[1] is None or abs(float(a[1]) - gv) > tol:
                            errors.append(f"Executive Summary {key}={a[1]}, expected ~{g[1]} (tol={tol})")
                    except (ValueError, TypeError):
                        errors.append(f"Executive Summary {key}={a[1]} not numeric")

    except Exception as e:
        errors.append(f"Error reading Excel: {e}")
    return errors


def check_pptx(agent_workspace):
    errors = []
    from pptx import Presentation

    path = os.path.join(agent_workspace, "Quality_Review.pptx")
    if not os.path.exists(path):
        return ["Quality_Review.pptx not found"]
    try:
        prs = Presentation(path)
        slide_count = len(prs.slides)
        if slide_count < 6:
            errors.append(f"PowerPoint has {slide_count} slides, expected >= 6")

        # Check content across slides
        all_text = ""
        slide_titles = []
        for i, slide in enumerate(prs.slides):
            slide_text = ""
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    all_text += shape.text.lower() + " "
                    slide_text += shape.text.lower() + " "
            slide_titles.append(slide_text)

        # Required content topics (each MUST be present)
        if "benchmark" not in all_text and "target" not in all_text:
            errors.append("PowerPoint does not mention benchmarks or targets")
        if not ("issue" in all_text or "order" in all_text):
            errors.append("PowerPoint does not discuss order issues")
        if not ("support" in all_text or "ticket" in all_text):
            errors.append("PowerPoint does not discuss support metrics")
        if "cross" not in all_text and "reference" not in all_text:
            errors.append("PowerPoint does not include cross-reference findings")
        if "recommend" not in all_text:
            errors.append("PowerPoint does not include recommendations")
        # Validate that targets from portal are mentioned (CSAT 4.2 or NPS, etc.)
        if not any(s in all_text for s in ["4.2", "csat", "satisfaction"]):
            errors.append("PowerPoint does not mention CSAT/satisfaction target metrics")

    except Exception as e:
        errors.append(f"Error reading PowerPoint: {e}")
    return errors


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--agent_workspace", required=False)
    parser.add_argument("--groundtruth_workspace", required=False)
    parser.add_argument("--launch_time", required=False)
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()
    task_root = os.path.abspath(os.path.join(os.path.dirname(__file__), ".."))
    agent_ws = args.agent_workspace or os.path.join(task_root, "groundtruth_workspace")
    gt_ws = args.groundtruth_workspace or os.path.join(task_root, "groundtruth_workspace")

    all_errors = []

    print("  Checking Excel file...")
    errs = check_excel(agent_ws, gt_ws)
    if errs:
        all_errors.extend(errs)
        for e in errs[:5]:
            print(f"    ERROR: {e}")
    else:
        print("    PASS")

    print("  Checking PowerPoint...")
    errs = check_pptx(agent_ws)
    if errs:
        all_errors.extend(errs)
        for e in errs[:3]:
            print(f"    ERROR: {e}")
    else:
        print("    PASS")

    if all_errors:
        print(f"\n=== RESULT: FAIL ({len(all_errors)} errors) ===")
        for e in all_errors[:10]:
            print(f"  {e}")
        sys.exit(1)
    else:
        print("\n=== RESULT: PASS ===")
        sys.exit(0)


if __name__ == "__main__":
    main()
