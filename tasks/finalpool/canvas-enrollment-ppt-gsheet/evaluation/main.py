"""Evaluation for canvas-enrollment-ppt-gsheet.

Blocking checks: Enrollment_Overview.xlsx and Enrollment_Overview.pptx.
GSheet 'Enrollment Dashboard' check: runtime_only by default; upgraded to
blocking if the agent populated relevant cells (catches partial-deliverable
attacks).
"""
import argparse
import os
import re
import sys
import openpyxl
from pptx import Presentation


PASS_COUNT = 0
FAIL_COUNT = 0
RUNTIME_ONLY_FAIL = 0


def record(name, passed, detail="", runtime_only=False):
    global PASS_COUNT, FAIL_COUNT, RUNTIME_ONLY_FAIL
    if passed:
        PASS_COUNT += 1
        print(f"  [PASS] {name}")
    else:
        FAIL_COUNT += 1
        if runtime_only:
            RUNTIME_ONLY_FAIL += 1
        msg = f": {detail[:300]}" if detail else ""
        suffix = " (runtime-only)" if runtime_only else ""
        print(f"  [FAIL] {name}{suffix}{msg}")


def num_close(a, b, tol=1.0):
    if a is None and b is None:
        return True
    if a is None or b is None:
        return False
    try:
        return abs(float(a) - float(b)) <= tol
    except (TypeError, ValueError):
        return str(a).strip().lower() == str(b).strip().lower()


def str_match(a, b):
    if a is None or b is None:
        return a is None and b is None
    return str(a).strip().lower() == str(b).strip().lower()


def load_sheet_rows(wb, sheet_name):
    for name in wb.sheetnames:
        if name.strip().lower() == sheet_name.strip().lower():
            return [[cell.value for cell in row] for row in wb[name].iter_rows()]
    return None


def check_excel(agent_workspace, gt_dir):
    print("\n=== Checking Excel Output ===")

    agent_excel = os.path.join(agent_workspace, "Enrollment_Overview.xlsx")
    gt_excel = os.path.join(gt_dir, "Enrollment_Overview.xlsx")

    if not os.path.exists(agent_excel):
        record("Excel file exists", False, f"Not found: {agent_excel}")
        return False
    record("Excel file exists", True)

    if not os.path.exists(gt_excel):
        record("Groundtruth Excel exists", False)
        return False

    try:
        agent_wb = openpyxl.load_workbook(agent_excel, data_only=True)
        gt_wb = openpyxl.load_workbook(gt_excel, data_only=True)
    except Exception as e:
        record("Excel readable", False, str(e))
        return False

    # ----- Enrollment Details -----
    a_rows = load_sheet_rows(agent_wb, "Enrollment Details")
    g_rows = load_sheet_rows(gt_wb, "Enrollment Details")
    if a_rows is None:
        record("Sheet 'Enrollment Details' exists", False)
    elif g_rows is None:
        record("Groundtruth has 'Enrollment Details'", False)
    else:
        record("Sheet 'Enrollment Details' exists", True)
        a_data = a_rows[1:] if len(a_rows) > 1 else []
        g_data = g_rows[1:] if len(g_rows) > 1 else []
        record("Enrollment Details row count",
               abs(len(a_data) - len(g_data)) <= 2,
               f"agent={len(a_data)}, expected={len(g_data)}")

        # Build lookup by course code (col 1)
        a_lookup = {}
        for row in a_data:
            if row and len(row) > 1 and row[1] is not None:
                a_lookup[str(row[1]).strip().lower()] = row
        for g_row in g_data:
            if not g_row or len(g_row) < 2 or g_row[1] is None:
                continue
            key = str(g_row[1]).strip().lower()
            a_row = a_lookup.get(key)
            if a_row is None:
                record(f"Course '{g_row[1]}' present", False, "Missing")
                continue
            # Col 2: Total_Enrollments — INTEGER count, tight tolerance
            if len(a_row) > 2 and len(g_row) > 2:
                # Tolerance tightened: 10 -> 2 (was overly loose for integer count)
                ok = num_close(a_row[2], g_row[2], 2)
                record(f"{key}.Total_Enrollments",
                       ok,
                       f"agent={a_row[2]}, expected={g_row[2]} (tol=2)")
            # Col 3: Students
            if len(a_row) > 3 and len(g_row) > 3:
                ok = num_close(a_row[3], g_row[3], 2)
                record(f"{key}.Students", ok,
                       f"agent={a_row[3]}, expected={g_row[3]} (tol=2)")

            # Cols 4-7: Teachers, TAs, Active, Completed (N1: previously
            # unchecked — only Total_Enrollments and Students were validated).
            extra_cols = {4: "Teachers", 5: "TAs", 6: "Active", 7: "Completed"}
            for ci, cname in extra_cols.items():
                if len(a_row) > ci and len(g_row) > ci and g_row[ci] is not None:
                    ok = num_close(a_row[ci], g_row[ci], 2)
                    record(f"{key}.{cname}", ok,
                           f"agent={a_row[ci]}, expected={g_row[ci]} (tol=2)")

    # ----- Summary -----
    a_rows = load_sheet_rows(agent_wb, "Summary")
    g_rows = load_sheet_rows(gt_wb, "Summary")
    if a_rows is None:
        record("Sheet 'Summary' exists", False)
    elif g_rows is None:
        record("Groundtruth has 'Summary'", False)
    else:
        record("Sheet 'Summary' exists", True)
        a_data = a_rows[1:] if len(a_rows) > 1 else []
        g_data = g_rows[1:] if len(g_rows) > 1 else []

        a_lookup = {}
        for row in a_data:
            if row and row[0] is not None:
                a_lookup[str(row[0]).strip().lower()] = row
        for g_row in g_data:
            if not g_row or g_row[0] is None:
                continue
            key = str(g_row[0]).strip().lower()
            a_row = a_lookup.get(key)
            if a_row is None:
                record(f"Summary '{g_row[0]}' present", False)
                continue
            if len(a_row) > 1 and len(g_row) > 1:
                try:
                    float(g_row[1])
                    # Tighter tolerances. Integer counts like Total_Enrollments
                    # need tight tol (2). Averages allow 0.05.
                    if "avg" in key or "average" in key:
                        tol = 0.05
                    else:
                        tol = 2
                    ok = num_close(a_row[1], g_row[1], tol)
                    record(f"Summary.{key}", ok,
                           f"agent={a_row[1]}, expected={g_row[1]} (tol={tol})")
                except (TypeError, ValueError):
                    ok = str_match(a_row[1], g_row[1])
                    record(f"Summary.{key}", ok,
                           f"agent={a_row[1]}, expected={g_row[1]}")

    return True


def check_pptx(agent_workspace):
    print("\n=== Checking PowerPoint ===")
    agent_ppt = os.path.join(agent_workspace, "Enrollment_Overview.pptx")
    if not os.path.exists(agent_ppt):
        record("PPT file exists", False, f"Not found: {agent_ppt}")
        return False
    record("PPT file exists", True)

    prs = Presentation(agent_ppt)
    slides = list(prs.slides)
    record("PPT has >= 4 slides", len(slides) >= 4, f"Got {len(slides)}")
    if len(slides) < 4:
        return False

    title_text = ""
    for shape in slides[0].shapes:
        if shape.has_text_frame:
            title_text += shape.text_frame.text.lower() + " "
    record("Title slide has 'enrollment'", "enrollment" in title_text,
           f"Found: {title_text[:100]}")
    # Accept any YYYY-MM-DD date in the title-slide subtitle. The task prompt
    # asks for today's date, so the check must not be pinned to the DB dump date.
    date_match = re.search(r"\d{4}-\d{2}-\d{2}", title_text)
    record("Title slide subtitle has a YYYY-MM-DD date",
           date_match is not None,
           f"Found: {title_text[:120]}")

    all_ppt_text = ""
    for slide in slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                all_ppt_text += shape.text_frame.text.lower() + " "

    record("PPT has 'Top 5' slide",
           "top 5" in all_ppt_text or "top five" in all_ppt_text)
    record("PPT has 'Distribution' slide",
           "distribution" in all_ppt_text)
    record("PPT mentions Creative Computing course",
           "creative computing" in all_ppt_text)
    return True


def check_gsheet():
    """Check the 'Enrollment Dashboard' shared spreadsheet.

    Pattern: by default runtime_only (will FAIL on V1 GT-only test). BUT if
    the agent populated cells in the gsheet schema, the missing/incorrect
    sheet IS a real failure (blocking).
    """
    print("\n=== Checking Google Sheet (Enrollment Dashboard) ===")
    try:
        import psycopg2
        conn = psycopg2.connect(
            host=os.environ.get("PGHOST", "localhost"),
            port=5432, dbname="toolathlon_gym",
            user="eigent", password="camel",
        )
        cur = conn.cursor()
        # Look for a spreadsheet titled like "Enrollment Dashboard"
        cur.execute("""
            SELECT id, title FROM gsheet.spreadsheets
            WHERE LOWER(title) LIKE '%enrollment%' AND LOWER(title) LIKE '%dashboard%'
        """)
        rows = cur.fetchall()

        # Detect agent population specifically: a relevant spreadsheet exists,
        # OR an enrollment-related spreadsheet has been created (NOT counting
        # noise from other tasks).
        cur.execute("""
            SELECT COUNT(*) FROM gsheet.spreadsheets
            WHERE LOWER(title) LIKE '%enroll%' OR LOWER(title) LIKE '%course%'
        """)
        relevant_ss_count = cur.fetchone()[0]
        agent_populated = (len(rows) > 0) or (relevant_ss_count > 0)
        is_runtime_only = not agent_populated

        # Block 1: spreadsheet 'Enrollment Dashboard' exists
        record("Enrollment Dashboard spreadsheet exists",
               len(rows) > 0,
               f"Relevant spreadsheets: {relevant_ss_count}",
               runtime_only=is_runtime_only)

        if rows:
            spreadsheet_id = rows[0][0]
            # Look for sheet 'Course Data' (LIKE pattern as parameter to avoid % escaping issue)
            cur.execute("""
                SELECT id, title FROM gsheet.sheets
                WHERE spreadsheet_id = %s AND LOWER(title) LIKE %s
            """, (spreadsheet_id, "%course%data%"))
            sheet_rows = cur.fetchall()
            # Once spreadsheet exists, expecting Course Data is now blocking
            record("Sheet 'Course Data' exists in Enrollment Dashboard",
                   len(sheet_rows) > 0,
                   runtime_only=False)
            if sheet_rows:
                sheet_id = sheet_rows[0][0]
                # Count populated cells (>= 22 courses x 8 columns + header = ~184)
                cur.execute("""
                    SELECT COUNT(*) FROM gsheet.cells
                    WHERE spreadsheet_id = %s AND sheet_id = %s
                """, (spreadsheet_id, sheet_id))
                cell_count = cur.fetchone()[0]
                record("Course Data has reasonable cell count (>= 50)",
                       cell_count >= 50,
                       f"Got {cell_count}",
                       runtime_only=False)

        cur.close()
        conn.close()
        return True

    except Exception as e:
        record("GSheet DB accessible", False, str(e), runtime_only=True)
        return False


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--agent_workspace", required=False, default=".")
    parser.add_argument("--groundtruth_workspace", required=False)
    parser.add_argument("--launch_time", required=False)
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()

    task_root = os.path.abspath(os.path.join(os.path.dirname(__file__), ".."))
    gt_dir = args.groundtruth_workspace or os.path.join(task_root, "groundtruth_workspace")

    check_excel(args.agent_workspace, gt_dir)
    check_pptx(args.agent_workspace)
    check_gsheet()

    blocking_fail = FAIL_COUNT - RUNTIME_ONLY_FAIL
    print(f"\n=== SUMMARY ===")
    print(f"  Passed: {PASS_COUNT}")
    print(f"  Failed: {FAIL_COUNT} (runtime-only fails: {RUNTIME_ONLY_FAIL})")
    print(f"  Blocking failures: {blocking_fail}")
    overall = blocking_fail == 0
    print(f"  Overall: {'PASS' if overall else 'FAIL'}")
    sys.exit(0 if overall else 1)


if __name__ == "__main__":
    main()
