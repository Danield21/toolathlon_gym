"""Evaluation for sf-hr-dept-budget-ppt-email."""
import argparse
import json
import os
import sys

import psycopg2

DB = {"host": os.environ.get("PGHOST", "localhost"), "port": int(os.environ.get("PGPORT", "5432")), "dbname": "toolathlon_gym", "user": "eigent", "password": "camel"}


def num_close(a, b, tol=1.0):
    if a is None and b is None:
        return True
    if a is None or b is None:
        return False
    try:
        return abs(float(a) - float(b)) <= tol
    except (TypeError, ValueError):
        return str(a).strip().lower() == str(b).strip().lower()


def str_match(a, b):
    if a is None or b is None:
        return a is None and b is None
    return str(a).strip().lower() == str(b).strip().lower()


def load_sheet_rows(wb, sheet_name):
    for name in wb.sheetnames:
        if name.strip().lower() == sheet_name.strip().lower():
            return [[cell.value for cell in row] for row in wb[name].iter_rows()]
    return None


def _fetch_dept_expected():
    """Fetch expected dept budgets/headcount from Snowflake mirror with fallback."""
    try:
        conn = psycopg2.connect(**DB, connect_timeout=5)
        cur = conn.cursor()
        cur.execute('SELECT "DEPARTMENT_NAME", "HEADCOUNT"::integer, "BUDGET"::numeric '
                    'FROM sf_data."HR_ANALYTICS__PUBLIC__DEPARTMENTS"')
        data = {}
        for name, hc, bud in cur.fetchall():
            data[str(name).strip().lower()] = {"headcount": int(hc), "budget": float(bud)}
        cur.close(); conn.close()
        if data:
            return data
    except Exception as e:
        print(f"[WARN] DB fetch failed: {e}")
    return None


def check_excel(agent_workspace, gt_dir):
    errors = []
    # Fetch dynamic expected values
    dept_expected = _fetch_dept_expected()
    try:
        import openpyxl
    except ImportError:
        errors.append("openpyxl not installed")
        return errors

    agent_file = os.path.join(agent_workspace, "Department_Budget_Analysis.xlsx")
    gt_file = os.path.join(gt_dir, "Department_Budget_Analysis.xlsx")

    if not os.path.exists(agent_file):
        errors.append("Department_Budget_Analysis.xlsx not found in agent workspace")
        return errors
    if not os.path.exists(gt_file):
        errors.append("Groundtruth Department_Budget_Analysis.xlsx not found")
        return errors

    agent_wb = openpyxl.load_workbook(agent_file, data_only=True)
    gt_wb = openpyxl.load_workbook(gt_file, data_only=True)

    # Check Department Summary sheet
    a_rows = load_sheet_rows(agent_wb, "Department Summary")
    g_rows = load_sheet_rows(gt_wb, "Department Summary")
    if a_rows is None:
        errors.append("Sheet 'Department Summary' not found in agent output")
    else:
        a_data = [r for r in (a_rows[1:] if len(a_rows) > 1 else []) if r and r[0] is not None]
        g_data = [r for r in (g_rows[1:] if g_rows and len(g_rows) > 1 else []) if r and r[0] is not None]

        # Header check (order-independent): all expected columns must be present.
        expected_headers = [
            "Department", "Headcount", "Budget", "Location", "Budget_Per_Employee",
        ]
        actual_headers = [str(c or "").strip() for c in (a_rows[0] if a_rows else [])]
        header_lower = [h.lower() for h in actual_headers]
        missing_headers = [h for h in expected_headers if h.lower() not in header_lower]
        if missing_headers:
            errors.append(f"Department Summary missing headers: {missing_headers}; found {actual_headers}")
        # Map expected columns to the agent's column indices by header name so
        # the column order in the output does not matter.
        a_col = {h.lower(): i for i, h in enumerate(actual_headers)}

        if len(a_data) < 7:
            errors.append(f"Department Summary: expected 7 data rows, got {len(a_data)}")
        else:
            a_lookup = {str(r[0]).strip().lower(): r for r in a_data if r[0]}
            for g_row in g_data:
                key = str(g_row[0]).strip().lower()
                a_row = a_lookup.get(key)
                if a_row is None:
                    errors.append(f"Missing department row: {g_row[0]}")
                    continue
                # Use DB-fetched expected when available (more robust), else GT
                exp_headcount = (dept_expected or {}).get(key, {}).get("headcount", g_row[1]) if dept_expected else g_row[1]
                exp_budget = (dept_expected or {}).get(key, {}).get("budget", g_row[2]) if dept_expected else g_row[2]
                # Headcount exact int
                hc_i = a_col.get("headcount")
                if hc_i is not None and hc_i < len(a_row) and not num_close(a_row[hc_i], exp_headcount, 0):
                    errors.append(f"{g_row[0]} Headcount: got {a_row[hc_i]}, expected {exp_headcount}")
                # Budget tol=1 (rounded to 2 decimals in source)
                bud_i = a_col.get("budget")
                if bud_i is not None and bud_i < len(a_row) and not num_close(a_row[bud_i], exp_budget, 1.0):
                    errors.append(f"{g_row[0]} Budget: got {a_row[bud_i]}, expected {exp_budget}")
                # Budget_Per_Employee (GT col 4) - tol=1
                bpe_i = a_col.get("budget_per_employee")
                if bpe_i is not None and bpe_i < len(a_row) and 4 < len(g_row) and not num_close(a_row[bpe_i], g_row[4], 1.0):
                    errors.append(f"{g_row[0]} Budget_Per_Employee: got {a_row[bpe_i]}, expected {g_row[4]}")

    # Check Summary sheet
    a_sum = load_sheet_rows(agent_wb, "Summary")
    g_sum = load_sheet_rows(gt_wb, "Summary")
    if a_sum is None:
        errors.append("Sheet 'Summary' not found in agent output")
    else:
        a_sum_data = {str(r[0]).strip().lower(): r[1] for r in (a_sum[1:] if len(a_sum) > 1 else []) if r and r[0]}
        g_sum_data = {str(r[0]).strip().lower(): r[1] for r in (g_sum[1:] if g_sum and len(g_sum) > 1 else []) if r and r[0]}

        # Compute dynamic expected totals from DB when available
        if dept_expected:
            exp_tot_hc = sum(v["headcount"] for v in dept_expected.values())
            exp_tot_bud = sum(v["budget"] for v in dept_expected.values())
            exp_top_dept = max(dept_expected.items(), key=lambda kv: kv[1]["budget"])[0]
        else:
            exp_tot_hc, exp_tot_bud, exp_top_dept = 50000, 4564347744.31, "r&d"

        # Total_Headcount - exact int
        th = a_sum_data.get("total_headcount")
        if th is None:
            errors.append("Summary missing Total_Headcount")
        elif not num_close(th, exp_tot_hc, 0):
            errors.append(f"Total_Headcount: got {th}, expected {exp_tot_hc}")

        # Total_Budget - tol=1
        tb = a_sum_data.get("total_budget")
        if tb is None:
            errors.append("Summary missing Total_Budget")
        elif not num_close(tb, exp_tot_bud, 1.0):
            errors.append(f"Total_Budget: got {tb}, expected {exp_tot_bud}")

        # Highest_Budget_Dept - exact match (case insensitive)
        hbd = a_sum_data.get("highest_budget_dept")
        if hbd is None:
            errors.append("Summary missing Highest_Budget_Dept")
        elif str(hbd).strip().lower() != exp_top_dept:
            errors.append(f"Highest_Budget_Dept: got '{hbd}', expected '{exp_top_dept}'")

        # Avg_Budget_Per_Employee validation (if present)
        # tol=25: task.md now defines this as Total_Budget / Total_Headcount, but a
        # model may legitimately compute the average of the per-department
        # Budget_Per_Employee values instead (differs by ~13 on this dataset).
        abpe = a_sum_data.get("avg_budget_per_employee")
        if abpe is not None and exp_tot_hc > 0:
            expected_avg = exp_tot_bud / exp_tot_hc
            if not num_close(abpe, expected_avg, 25.0):
                errors.append(f"Avg_Budget_Per_Employee: got {abpe}, expected ~{expected_avg:.2f}")

    return errors


def check_pptx(agent_workspace):
    errors = []
    pptx_path = os.path.join(agent_workspace, "HR_Department_Overview.pptx")
    if not os.path.exists(pptx_path):
        errors.append("HR_Department_Overview.pptx not found")
        return errors
    try:
        from pptx import Presentation
        prs = Presentation(pptx_path)
        if len(prs.slides) < 7:
            errors.append(f"HR_Department_Overview.pptx has {len(prs.slides)} slides, expected at least 7")
        # Check for department names in slides
        all_text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    all_text += " " + shape.text_frame.text
        all_text_lower = all_text.lower()
        # Department names must match the data source verbatim (e.g. "HR", "R&D");
        # pull the actual names from the warehouse so rephrased/expanded names fail.
        dept_names = list((_fetch_dept_expected() or {}).keys())
        if not dept_names:
            dept_names = ["engineering", "finance", "hr", "operations", "r&d", "sales", "support"]
        for dept in dept_names:
            if dept not in all_text_lower:
                errors.append(f"PPTX missing department: {dept}")
    except ImportError:
        if os.path.getsize(pptx_path) < 1000:
            errors.append("HR_Department_Overview.pptx too small")
    except Exception as e:
        errors.append(f"Error reading HR_Department_Overview.pptx: {e}")
    return errors


def check_gsheet():
    errors = []
    try:
        conn = psycopg2.connect(**DB)
        cur = conn.cursor()
        cur.execute("""
            SELECT s.id, s.title
            FROM gsheet.spreadsheets s
            WHERE LOWER(s.title) LIKE '%hr department%' OR LOWER(s.title) LIKE '%department tracker%'
        """)
        sheets = cur.fetchall()
        if not sheets:
            errors.append("No Google Sheet named 'HR Department Tracker' found")
        else:
            ss_id = sheets[0][0]
            cur.execute("""
                SELECT COUNT(DISTINCT c.row_index)
                FROM gsheet.cells c
                WHERE c.spreadsheet_id = %s AND c.row_index > 0
            """, (ss_id,))
            row_count = cur.fetchone()[0]
            if row_count < 7:
                errors.append(f"HR Department Tracker has only {row_count} data rows, expected at least 7")
            else:
                # Verify the sheet actually contains the department data (not just
                # a title + empty rows). Pull every cell value and check that the
                # expected department names + headcounts (derived from the live DB)
                # appear among the cells.
                cur.execute("""
                    SELECT c.value FROM gsheet.cells c
                    WHERE c.spreadsheet_id = %s AND c.row_index > 0
                """, (ss_id,))
                cell_values = [str(r[0]) for r in cur.fetchall() if r[0] is not None]
                cell_blob = " ".join(cell_values).lower()
                cur.execute('''
                    SELECT "DEPARTMENT", COUNT(*)
                    FROM sf_data."HR_ANALYTICS__PUBLIC__EMPLOYEES"
                    GROUP BY "DEPARTMENT"
                ''')
                dept_rows = cur.fetchall()
                # Require every department name to be present; headcount check is
                # best-effort (at least half of the headcount values must appear).
                missing_depts = [
                    d for (d, _) in dept_rows
                    if str(d).strip().lower() not in cell_blob
                ]
                if missing_depts:
                    errors.append(
                        f"HR Department Tracker missing department names in cells: {missing_depts}"
                    )
                headcount_hits = 0
                for _, hc in dept_rows:
                    if str(int(hc)) in cell_values or f"{int(hc):,}" in cell_values:
                        headcount_hits += 1
                if dept_rows and headcount_hits < (len(dept_rows) + 1) // 2:
                    errors.append(
                        f"HR Department Tracker cells do not contain enough department "
                        f"headcount values ({headcount_hits}/{len(dept_rows)} found)"
                    )
        cur.close()
        conn.close()
    except Exception as e:
        errors.append(f"GSheet DB check error: {e}")
    return errors


def check_email():
    errors = []
    try:
        conn = psycopg2.connect(**DB)
        cur = conn.cursor()
        cur.execute("""
            SELECT subject, to_addr
            FROM email.messages
            WHERE LOWER(subject) LIKE '%hr department budget%' OR LOWER(subject) LIKE '%budget summary%'
        """)
        emails = cur.fetchall()
        cur.close()
        conn.close()
        if not emails:
            errors.append("No email with 'hr department budget' or 'budget summary' in subject found")
        else:
            found_to = False
            for em in emails:
                if "hr.directors" in str(em[1]).lower():
                    found_to = True
                    break
            if not found_to:
                errors.append("No email sent to hr.directors@company.com")
    except Exception as e:
        errors.append(f"Email DB check error: {e}")
    return errors


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--agent_workspace", required=False, default=".")
    parser.add_argument("--groundtruth_workspace", required=False)
    parser.add_argument("--launch_time", required=False)
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()

    task_root = os.path.abspath(os.path.join(os.path.dirname(__file__), ".."))
    gt_dir = args.groundtruth_workspace or os.path.join(task_root, "groundtruth_workspace")

    all_errors = []

    print("\n=== Checking Excel ===")
    excel_errors = check_excel(args.agent_workspace, gt_dir)
    if excel_errors:
        for e in excel_errors:
            print(f"  [FAIL] {e}")
        all_errors.extend(excel_errors)
    else:
        print("  [PASS] Excel check passed")

    print("\n=== Checking PowerPoint ===")
    pptx_errors = check_pptx(args.agent_workspace)
    if pptx_errors:
        for e in pptx_errors:
            print(f"  [FAIL] {e}")
        all_errors.extend(pptx_errors)
    else:
        print("  [PASS] PPTX check passed")

    print("\n=== Checking Google Sheet ===")
    gsheet_errors = check_gsheet()
    if gsheet_errors:
        for e in gsheet_errors:
            print(f"  [FAIL] {e}")
        all_errors.extend(gsheet_errors)
    else:
        print("  [PASS] GSheet check passed")

    print("\n=== Checking Email ===")
    email_errors = check_email()
    if email_errors:
        for e in email_errors:
            print(f"  [FAIL] {e}")
        all_errors.extend(email_errors)
    else:
        print("  [PASS] Email check passed")

    if args.res_log_file:
        with open(args.res_log_file, "w") as f:
            json.dump({"errors": all_errors, "success": len(all_errors) == 0}, f, indent=2)

    if all_errors:
        print(f"\n=== RESULT: FAIL ({len(all_errors)} errors) ===")
        sys.exit(1)
    else:
        print("\n=== RESULT: PASS ===")
        sys.exit(0)


if __name__ == "__main__":
    main()
