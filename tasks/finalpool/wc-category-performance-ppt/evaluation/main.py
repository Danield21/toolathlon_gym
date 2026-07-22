"""
Evaluation for wc-category-performance-ppt.
Checks:
1. Excel file Category_Data.xlsx with correct category metrics
2. PPT file Category_Review.pptx with correct slide count and content
"""
import argparse
import json
import os
import sys
from collections import defaultdict

import openpyxl
import psycopg2

DB_CONFIG = {
    "host": os.environ.get("PGHOST", "localhost"),
    "port": 5432,
    "dbname": "toolathlon_gym",
    "user": "eigent",
    "password": "camel",
}

PASS_COUNT = 0
FAIL_COUNT = 0


def record(name, passed, detail=""):
    global PASS_COUNT, FAIL_COUNT
    if passed:
        PASS_COUNT += 1
        print(f"  [PASS] {name}")
    else:
        FAIL_COUNT += 1
        msg = f": {detail[:300]}" if detail else ""
        print(f"  [FAIL] {name}{msg}")


def num_close(a, b, tol=1.0):
    try:
        return abs(float(a) - float(b)) <= tol
    except (TypeError, ValueError):
        return False


def str_match(a, b):
    if a is None or b is None:
        return a is None and b is None
    return str(a).strip().lower() == str(b).strip().lower()


def load_sheet_rows(wb, sheet_name):
    for name in wb.sheetnames:
        if name.strip().lower() == sheet_name.strip().lower():
            return [[cell.value for cell in row] for row in wb[name].iter_rows()]
    return None


def get_expected_data():
    """Query WC DB for expected category performance."""
    conn = psycopg2.connect(**DB_CONFIG)
    cur = conn.cursor()

    cur.execute("SELECT id, name, categories, total_sales, price, average_rating FROM wc.products")
    rows = cur.fetchall()

    cat_data = defaultdict(lambda: {'count': 0, 'prices': [], 'total_units_sold': 0, 'ratings': []})

    for r in rows:
        cats = r[2]
        if isinstance(cats, str):
            cats = json.loads(cats)
        if isinstance(cats, list):
            for c in cats:
                cat_name = c.get('name', '') if isinstance(c, dict) else str(c)
                cat_data[cat_name]['count'] += 1
                try:
                    cat_data[cat_name]['prices'].append(float(r[4]) if r[4] else 0)
                except (TypeError, ValueError):
                    pass
                cat_data[cat_name]['total_units_sold'] += int(r[3]) if r[3] else 0
                try:
                    if r[5] and float(r[5]) > 0:
                        cat_data[cat_name]['ratings'].append(float(r[5]))
                except (TypeError, ValueError):
                    pass

    results = []
    for cat in sorted(cat_data.keys()):
        d = cat_data[cat]
        avg_price = round(sum(d['prices']) / len(d['prices']), 2) if d['prices'] else 0
        avg_rating = round(sum(d['ratings']) / len(d['ratings']), 2) if d['ratings'] else 0
        results.append({
            'category': cat,
            'product_count': d['count'],
            'avg_price': avg_price,
            'total_units_sold': d['total_units_sold'],
            'avg_rating': avg_rating,
        })

    cur.close()
    conn.close()
    return results


def check_excel(agent_workspace, gt_workspace, expected):
    """Check Category_Data.xlsx."""
    print("\n=== Checking Excel ===")
    agent_file = os.path.join(agent_workspace, "Category_Data.xlsx")
    gt_file = os.path.join(gt_workspace, "Category_Data.xlsx")

    if not os.path.exists(agent_file):
        record("Excel file exists", False, f"Not found: {agent_file}")
        return False
    record("Excel file exists", True)

    agent_wb = openpyxl.load_workbook(agent_file, data_only=True)
    gt_wb = openpyxl.load_workbook(gt_file, data_only=True)

    a_rows = load_sheet_rows(agent_wb, "Category Performance")
    g_rows = load_sheet_rows(gt_wb, "Category Performance")

    if a_rows is None:
        record("Sheet 'Category Performance' exists", False, f"Sheets: {agent_wb.sheetnames}")
        return False
    record("Sheet 'Category Performance' exists", True)

    a_data = a_rows[1:] if len(a_rows) > 1 else []
    g_data = g_rows[1:] if len(g_rows) > 1 else []

    record("Category row count", len(a_data) == len(g_data),
           f"Expected {len(g_data)}, got {len(a_data)}")

    a_lookup = {}
    for row in a_data:
        if row and row[0]:
            a_lookup[str(row[0]).strip().lower()] = row

    for g_row in g_data:
        if not g_row or not g_row[0]:
            continue
        cat = str(g_row[0]).strip()
        key = cat.lower()
        a_row = a_lookup.get(key)
        if a_row is None:
            record(f"Category '{cat}' found", False, "Missing")
            continue
        record(f"Category '{cat}' found", True)

        if len(a_row) > 1 and len(g_row) > 1:
            record(f"  {cat} Product_Count",
                   num_close(a_row[1], g_row[1], 1),
                   f"Agent={a_row[1]}, GT={g_row[1]}")
        if len(a_row) > 2 and len(g_row) > 2:
            # Tightened tol from 5.0 to 1.0 for Avg_Price
            record(f"  {cat} Avg_Price",
                   num_close(a_row[2], g_row[2], 1.0),
                   f"Agent={a_row[2]}, GT={g_row[2]}")
        if len(a_row) > 3 and len(g_row) > 3:
            # Tightened tol from 5 to 2 for Total_Units_Sold
            record(f"  {cat} Total_Units_Sold",
                   num_close(a_row[3], g_row[3], 2),
                   f"Agent={a_row[3]}, GT={g_row[3]}")
        if len(a_row) > 4 and len(g_row) > 4:
            record(f"  {cat} Avg_Rating",
                   num_close(a_row[4], g_row[4], 0.1),
                   f"Agent={a_row[4]}, GT={g_row[4]}")

    # Sort order: alphabetical by category
    cat_names = [str(r[0]).strip() for r in a_data if r and r[0]]
    sorted_cats = sorted(cat_names, key=lambda s: s.lower())
    record(
        "Categories sorted alphabetically",
        cat_names == sorted_cats,
        f"got {cat_names}",
    )

    return True


def check_pptx(agent_workspace, expected):
    """Check Category_Review.pptx."""
    print("\n=== Checking PowerPoint ===")
    pptx_path = os.path.join(agent_workspace, "Category_Review.pptx")

    if not os.path.exists(pptx_path):
        record("PPTX file exists", False, f"Not found: {pptx_path}")
        return False
    record("PPTX file exists", True)

    try:
        from pptx import Presentation
        prs = Presentation(pptx_path)
        num_slides = len(prs.slides)

        expected_slides = len(expected) + 2  # title + per-category + summary
        # Tighten: must be exactly title + N + summary (no ±2)
        record(f"Slide count exactly {expected_slides}",
               num_slides == expected_slides,
               f"Got {num_slides}")

        # Check for title slide (slide 1)
        first_slide = prs.slides[0]
        first_title = first_slide.shapes.title.text.strip().lower() if first_slide.shapes.title else ""
        first_text = " ".join(
            shape.text for shape in first_slide.shapes if shape.has_text_frame
        ).lower()
        record(
            "Title slide title 'Product Category Performance Review'",
            "product category performance review" in first_title or "product category performance review" in first_text,
            f"Title text: {first_text[:120]}",
        )

        # Per-category slides (slides 2 .. N+1): each should have category name as title
        slide_titles = []
        for slide in prs.slides:
            try:
                title_text = (slide.shapes.title.text if slide.shapes.title else "").strip()
            except Exception:
                title_text = ""
            slide_titles.append(title_text)

        cat_names = [c['category'] for c in expected]
        # Slides 2..N+1 (index 1..N): each title should equal a category name
        if num_slides >= len(cat_names) + 2:
            cat_slide_titles = [t.lower() for t in slide_titles[1: 1 + len(cat_names)]]
            cat_names_lower = sorted(c.lower() for c in cat_names)
            actual_cat_titles_sorted = sorted(cat_slide_titles)
            record(
                "Each category has its own slide with category name as title",
                actual_cat_titles_sorted == cat_names_lower,
                f"actual={actual_cat_titles_sorted} expected={cat_names_lower}",
            )

        # Each per-category slide must include 4 metric values
        # We accept loose check: slide content must have product count, avg price, units, rating
        for i, cat_info in enumerate(expected):
            slide_idx = i + 1  # slide 2 onward
            if slide_idx >= num_slides - 1:  # exclude summary slide
                continue
            slide = prs.slides[slide_idx]
            slide_text = " ".join(
                shape.text for shape in slide.shapes if shape.has_text_frame
            ).lower()
            # Check the 4 metrics by value (text-based contains)
            pc = str(cat_info['product_count'])
            av = str(cat_info['avg_price'])
            tu = str(cat_info['total_units_sold'])
            ar = str(cat_info['avg_rating'])
            metrics_found = sum(
                1 for v in [pc, av, tu, ar] if v in slide_text
            )
            record(
                f"Slide for '{cat_info['category']}' has at least 3 of 4 metric values",
                metrics_found >= 3,
                f"found {metrics_found}/4 in: {slide_text[:160]}",
            )

        # Summary slide (last)
        last_slide = prs.slides[-1]
        last_title = (last_slide.shapes.title.text if last_slide.shapes.title else "").strip().lower()
        last_text = " ".join(
            shape.text for shape in last_slide.shapes if shape.has_text_frame
        ).lower()
        record("Last slide title 'Overall Summary'",
               "overall summary" in last_title or "overall summary" in last_text,
               f"Last slide title: {last_title}")
        # Determine expected highlights from `expected` (handle ties)
        if expected:
            # Most products: tied
            mp_count = max(c['product_count'] for c in expected)
            mp_cats = [c['category'].lower() for c in expected if c['product_count'] == mp_count]
            # Highest rating: tied
            hr_val = max(c['avg_rating'] for c in expected)
            hr_cats = [c['category'].lower() for c in expected if c['avg_rating'] == hr_val]
            # Most units sold: tied
            mu_count = max(c['total_units_sold'] for c in expected)
            mu_cats = [c['category'].lower() for c in expected if c['total_units_sold'] == mu_count]

            record(
                f"Summary mentions most-products category (one of {mp_cats})",
                any(c in last_text for c in mp_cats),
                f"text: {last_text[:200]}",
            )
            record(
                f"Summary mentions highest-rating category (one of {hr_cats})",
                any(c in last_text for c in hr_cats),
                f"text: {last_text[:200]}",
            )
            record(
                f"Summary mentions most-units category (one of {mu_cats})",
                any(c in last_text for c in mu_cats),
                f"text: {last_text[:200]}",
            )

    except ImportError:
        record("python-pptx available", False, "Cannot import pptx module")
        return False
    except Exception as e:
        record("PPTX readable", False, str(e))
        return False

    return True


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--agent_workspace", required=False, default=".")
    parser.add_argument("--groundtruth_workspace", required=False)
    parser.add_argument("--launch_time", required=False)
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()

    task_root = os.path.abspath(os.path.join(os.path.dirname(__file__), ".."))
    gt_dir = args.groundtruth_workspace or os.path.join(task_root, "groundtruth_workspace")

    expected = get_expected_data()
    print(f"[eval] {len(expected)} categories found")

    excel_ok = check_excel(args.agent_workspace, gt_dir, expected)
    pptx_ok = check_pptx(args.agent_workspace, expected)

    print(f"\n=== SUMMARY ===")
    print(f"  Passed: {PASS_COUNT}")
    print(f"  Failed: {FAIL_COUNT}")
    overall = PASS_COUNT > 0 and FAIL_COUNT == 0
    print(f"  Overall: {'PASS' if overall else 'FAIL'}")

    sys.exit(0 if overall else 1)


if __name__ == "__main__":
    main()
