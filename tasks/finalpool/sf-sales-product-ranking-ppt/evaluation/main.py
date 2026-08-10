"""
Evaluation script for sf-sales-product-ranking-ppt task.

Checks:
1. Excel file (Product_Rankings.xlsx) - Top Products and Category Summary sheets
2. PowerPoint (Product_Performance_Review.pptx) - >=4 slides with required content
3. Google Sheet 'Product Rankings Dashboard' with 'Top 20' worksheet
"""

import argparse
import json
import os
import sys

import openpyxl
import psycopg2

DB_CONFIG = {
    "host": os.environ.get("PGHOST", "localhost"),
    "port": 5432,
    "dbname": "toolathlon_gym",
    "user": "eigent",
    "password": "camel",
}

PASS_COUNT = 0
FAIL_COUNT = 0


def check(name, condition, detail=""):
    global PASS_COUNT, FAIL_COUNT
    if condition:
        PASS_COUNT += 1
        print(f"  [PASS] {name}")
    else:
        FAIL_COUNT += 1
        detail_str = f": {str(detail)[:300]}" if detail else ""
        print(f"  [FAIL] {name}{detail_str}")


def num_close(a, b, rel_tol=0.05, abs_tol=0.5):
    try:
        af = float(a); bf = float(b)
        return abs(af - bf) <= max(abs_tol, abs(bf) * rel_tol)
    except (TypeError, ValueError):
        return False


def norm_name(s):
    """Lowercase and collapse whitespace."""
    return " ".join(str(s or "").lower().split()).rstrip(".")


def base_name(s):
    """Strip variable promotional suffixes (e.g. 'with No Cost EMI ...')."""
    n = norm_name(s)
    for sep in (" with ", " - "):
        idx = n.find(sep)
        if idx >= 10:
            n = n[:idx]
    return n.strip()


def product_name_match(agent_name, exp_name):
    """Relaxed name match: GT names may carry promotional suffixes that the
    data source omits, so compare base names / prefixes instead of exact text."""
    a = norm_name(agent_name)
    e = norm_name(exp_name)
    if a == e:
        return True
    ab, eb = base_name(a), base_name(e)
    if len(eb) >= 10 and ab == eb:
        return True
    shorter, longer = (a, e) if len(a) <= len(e) else (e, a)
    return len(shorter) >= 10 and longer.startswith(shorter)


def get_expected_top20():
    conn = psycopg2.connect(**DB_CONFIG)
    cur = conn.cursor()
    cur.execute('''
        SELECT p."PRODUCT_NAME", p."CATEGORY",
               SUM(o."TOTAL_AMOUNT")::float as revenue,
               SUM(o."QUANTITY")::int as units,
               COUNT(o."ORDER_ID")::int as order_count
        FROM sf_data."SALES_DW__PUBLIC__ORDERS" o
        JOIN sf_data."SALES_DW__PUBLIC__PRODUCTS" p ON o."PRODUCT_ID" = p."PRODUCT_ID"
        GROUP BY p."PRODUCT_NAME", p."CATEGORY"
        ORDER BY revenue DESC
        LIMIT 20
    ''')
    rows = cur.fetchall()

    cur.execute('''
        SELECT p."CATEGORY", COUNT(DISTINCT p."PRODUCT_NAME")::int,
               SUM(o."TOTAL_AMOUNT")::float
        FROM sf_data."SALES_DW__PUBLIC__ORDERS" o
        JOIN sf_data."SALES_DW__PUBLIC__PRODUCTS" p ON o."PRODUCT_ID" = p."PRODUCT_ID"
        GROUP BY p."CATEGORY"
        ORDER BY 3 DESC
    ''')
    cat_rows = cur.fetchall()
    cur.close()
    conn.close()
    return rows, cat_rows


def check_excel(agent_workspace, expected_top20, expected_cat):
    print("\n=== Checking Excel Output ===")
    excel_path = os.path.join(agent_workspace, "Product_Rankings.xlsx")
    check("Excel file exists", os.path.isfile(excel_path), f"Expected {excel_path}")
    if not os.path.isfile(excel_path):
        return

    try:
        wb = openpyxl.load_workbook(excel_path, data_only=True)
    except Exception as e:
        check("Excel file readable", False, str(e))
        return

    # === Top Products sheet ===
    ws = None
    for s in wb.sheetnames:
        sl = s.lower().strip()
        if sl == "top products":
            ws = wb[s]; break
    if ws is None:
        for s in wb.sheetnames:
            if "top" in s.lower() and "product" in s.lower():
                ws = wb[s]; break

    check("Sheet 'Top Products' exists", ws is not None, f"Sheets: {wb.sheetnames}")
    if ws is None:
        return

    rows = list(ws.iter_rows(min_row=2, values_only=True))
    rows = [r for r in rows if r and any(v is not None for v in r)]
    check("Top Products has exactly 20 data rows",
          len(rows) == 20,
          f"Got {len(rows)}")

    # Header check
    header = [str(c.value or "").strip().lower().replace(" ", "_") for c in ws[1]]
    expected_headers = ["rank", "product_name", "category", "total_revenue", "units_sold", "avg_order_value"]
    for i, eh in enumerate(expected_headers):
        ah = header[i] if i < len(header) else "MISSING"
        check(f"Top Products header[{i}] = '{eh}'",
              ah == eh,
              f"Got '{ah}'")

    # Validate Rank column 1..20
    if rows:
        ranks = []
        for r in rows:
            try:
                ranks.append(int(r[0]))
            except (TypeError, ValueError):
                ranks.append(None)
        check("Rank column = [1..20]",
              ranks == list(range(1, len(rows) + 1)),
              f"Got: {ranks}")

    # Sort: descending by revenue
    if rows:
        try:
            revs = [float(r[3]) for r in rows]
            check("Top Products sorted descending by revenue",
                  revs == sorted(revs, reverse=True),
                  f"Revenues: {revs[:5]}...")
        except (TypeError, ValueError):
            check("Top Products sorted descending by revenue", False, "non-numeric revenue")

    # Per-row product + revenue match
    if expected_top20 and rows:
        for i, exp in enumerate(expected_top20):
            if i >= len(rows):
                break
            agent_row = rows[i]
            exp_name = exp[0]
            exp_rev = float(exp[2])
            exp_units = int(exp[3])
            exp_aov = exp_rev / exp[4] if exp[4] else 0  # avg_order_value

            agent_name = str(agent_row[1] or "").strip()
            check(f"Rank {i+1} product_name matches",
                  product_name_match(agent_name, exp_name),
                  f"Got '{agent_name}' vs '{exp_name}'")
            try:
                ar = float(agent_row[3])
                check(f"Rank {i+1} revenue ~ {exp_rev:.2f}",
                      num_close(ar, exp_rev, rel_tol=0.02, abs_tol=10),
                      f"Got {ar} vs {exp_rev}")
            except (TypeError, ValueError):
                check(f"Rank {i+1} revenue numeric", False, f"Got {agent_row[3]}")
            try:
                au = int(float(agent_row[4]))
                check(f"Rank {i+1} units_sold matches",
                      abs(au - exp_units) <= 0,
                      f"Got {au} vs {exp_units}")
            except (TypeError, ValueError):
                check(f"Rank {i+1} units_sold numeric", False, f"Got {agent_row[4]}")

    # === Category Summary sheet ===
    ws2 = None
    for s in wb.sheetnames:
        if s.strip().lower() == "category summary":
            ws2 = wb[s]; break
    check("Sheet 'Category Summary' exists", ws2 is not None, f"Sheets: {wb.sheetnames}")
    if ws2 is not None and expected_cat:
        cat_rows = [r for r in ws2.iter_rows(min_row=2, values_only=True) if r and any(v is not None for v in r)]
        check(f"Category Summary has {len(expected_cat)} row(s)",
              len(cat_rows) == len(expected_cat),
              f"Got {len(cat_rows)}")
        if cat_rows:
            # Validate category total revenue matches
            cat_lookup = {str(r[0] or "").strip().lower(): r for r in cat_rows}
            for ec in expected_cat:
                cat_name = str(ec[0] or "").strip().lower()
                ar = cat_lookup.get(cat_name)
                if not ar:
                    check(f"Category '{ec[0]}' present", False, "Missing")
                    continue
                try:
                    rev = float(ar[2])
                    check(f"Category '{ec[0]}' Total_Revenue ~ {ec[2]:.2f}",
                          num_close(rev, ec[2], rel_tol=0.02, abs_tol=10),
                          f"Got {rev} vs {ec[2]}")
                except (TypeError, ValueError):
                    check(f"Category '{ec[0]}' revenue numeric", False, f"Got {ar[2]}")
                try:
                    pc = int(float(ar[1]))
                    check(f"Category '{ec[0]}' Product_Count = {ec[1]}",
                          pc == ec[1],
                          f"Got {pc}")
                except (TypeError, ValueError):
                    pass


def check_pptx(agent_workspace, expected_top20):
    print("\n=== Checking PowerPoint Output ===")
    pptx_path = os.path.join(agent_workspace, "Product_Performance_Review.pptx")
    check("PPTX file exists", os.path.isfile(pptx_path), f"Expected {pptx_path}")
    if not os.path.isfile(pptx_path):
        return

    try:
        from pptx import Presentation
        prs = Presentation(pptx_path)
    except Exception as e:
        check("PPTX file readable", False, str(e))
        return

    slide_count = len(prs.slides)
    check("PPTX has >= 4 slides", slide_count >= 4, f"Got {slide_count} slides")

    slide_texts = []
    for slide in prs.slides:
        text = ""
        for shape in slide.shapes:
            if shape.has_text_frame:
                text += shape.text_frame.text + " "
        slide_texts.append(text.strip())

    # Slide 1 must contain Q1 2026 Product Performance Review
    if slide_texts:
        s1 = slide_texts[0].lower()
        check("Slide 1 contains 'Q1 2026 Product Performance Review'",
              "q1 2026" in s1 and "product" in s1 and ("performance" in s1 or "review" in s1),
              f"Slide 1 text: {slide_texts[0][:200]}")

    # Slide 2/3/4 should have content (not empty) AND topical semantic markers
    for i in range(1, min(4, len(slide_texts))):
        check(f"Slide {i+1} has content (>=20 chars)",
              len(slide_texts[i]) >= 20,
              f"Slide {i+1}: {slide_texts[i][:100]}")

    # Tighten: enforce slide-2 mentions top/rank, slide-3 mentions category,
    # slide-4 mentions total revenue or top product (semantic anchor).
    if len(slide_texts) >= 2:
        s2 = slide_texts[1].lower()
        check("Slide 2 mentions 'top' or 'rank' (top-products theme)",
              ("top" in s2) or ("rank" in s2),
              f"Slide 2: {slide_texts[1][:120]}")
    if len(slide_texts) >= 3:
        s3 = slide_texts[2].lower()
        check("Slide 3 mentions 'category' (category breakdown theme)",
              "categor" in s3,
              f"Slide 3: {slide_texts[2][:120]}")
    if len(slide_texts) >= 4:
        s4 = slide_texts[3].lower()
        check("Slide 4 mentions 'total revenue' or 'top product' (summary theme)",
              ("total revenue" in s4) or ("top product" in s4),
              f"Slide 4: {slide_texts[3][:120]}")

    # Top product appears somewhere in PPTX
    if expected_top20:
        all_text_norm = norm_name(" ".join(slide_texts))
        # GT name may carry a promotional suffix; anchor on the base name
        anchor = base_name(expected_top20[0][0])
        check("PPTX mentions top-1 product (by base name)",
              anchor in all_text_norm,
              f"Looking for '{anchor}'")


def check_gsheet(expected_top20):
    print("\n=== Checking Google Sheet ===")
    conn = psycopg2.connect(**DB_CONFIG)
    cur = conn.cursor()
    cur.execute("SELECT id, title FROM gsheet.spreadsheets")
    sheets = cur.fetchall()

    target_title = "product rankings dashboard"
    found_ss = None
    for sid, t in sheets:
        if t and target_title in t.strip().lower():
            found_ss = sid
            break
    check(f"Google Sheet 'Product Rankings Dashboard' exists",
          found_ss is not None,
          f"Titles: {[t for _, t in sheets]}")
    if not found_ss:
        cur.close(); conn.close()
        return

    # Find 'Top 20' sub-sheet
    cur.execute("SELECT id, title FROM gsheet.sheets WHERE spreadsheet_id = %s", (found_ss,))
    subs = cur.fetchall()
    top20_id = None
    for sid, t in subs:
        if t and t.strip().lower() in ("top 20", "top20"):
            top20_id = sid
            break
    check("Worksheet 'Top 20' exists", top20_id is not None, f"Worksheets: {[t for _, t in subs]}")
    if top20_id is None:
        cur.close(); conn.close()
        return

    cur.execute("""SELECT row_index, col_index, COALESCE(formatted_value, value)
                   FROM gsheet.cells WHERE spreadsheet_id = %s AND sheet_id = %s
                   ORDER BY row_index, col_index""",
                (found_ss, top20_id))
    cells = cur.fetchall()
    cur.close()
    conn.close()

    # Build grid
    grid = {}
    for r, c, v in cells:
        grid.setdefault(r, {})[c] = v

    # 21 rows: header + 20 data
    data_rows = sorted(r for r in grid if r > 0)
    check("'Top 20' has exactly 20 data rows",
          len(data_rows) == 20,
          f"Got {len(data_rows)}")

    # Validate first data row product name
    if expected_top20 and data_rows:
        first_row = grid[data_rows[0]]
        # product name is likely in col 1 (after Rank)
        candidates = [first_row.get(1), first_row.get(0), first_row.get(2)]
        top_name = expected_top20[0][0]
        anchor = base_name(top_name)
        match = any(c and (product_name_match(str(c), top_name)
                           or anchor in norm_name(c)) for c in candidates)
        check("'Top 20' rank-1 row has top product",
              match,
              f"Got: {candidates}")


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--agent_workspace", required=False, default=".")
    parser.add_argument("--groundtruth_workspace", required=False, default=".")
    parser.add_argument("--launch_time", required=False)
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()

    expected_top20, expected_cat = get_expected_top20()
    print(f"Expected top product: {expected_top20[0][0] if expected_top20 else 'N/A'}")

    check_excel(args.agent_workspace, expected_top20, expected_cat)
    check_pptx(args.agent_workspace, expected_top20)
    check_gsheet(expected_top20)

    all_ok = FAIL_COUNT == 0
    print(f"\n=== SUMMARY ===")
    print(f"  Total checks - Passed: {PASS_COUNT}, Failed: {FAIL_COUNT}")
    print(f"  Overall: {'PASS' if all_ok else 'FAIL'}")

    if args.res_log_file:
        with open(args.res_log_file, "w") as f:
            json.dump({"passed": PASS_COUNT, "failed": FAIL_COUNT, "success": all_ok}, f, indent=2)

    sys.exit(0 if all_ok else 1)


if __name__ == "__main__":
    main()
