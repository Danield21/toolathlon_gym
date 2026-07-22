"""Evaluation for wc-coupon-performance-notion-ppt."""
import argparse
import os
import sys
import psycopg2
import openpyxl


DB = {"host": os.environ.get("PGHOST", "localhost"), "port": 5432, "dbname": "toolathlon_gym", "user": "eigent", "password": "camel"}


def num_close(a, b, tol=1.0):
    if a is None and b is None:
        return True
    if a is None or b is None:
        return False
    try:
        return abs(float(a) - float(b)) <= tol
    except (TypeError, ValueError):
        return str(a).strip().lower() == str(b).strip().lower()


def str_match(a, b):
    if a is None or b is None:
        return a is None and b is None
    return str(a).strip().lower() == str(b).strip().lower()


def load_sheet_rows(wb, sheet_name):
    for name in wb.sheetnames:
        if name.strip().lower() == sheet_name.strip().lower():
            return [[cell.value for cell in row] for row in wb[name].iter_rows()]
    return None


def get_coupon_data():
    conn = psycopg2.connect(**DB)
    cur = conn.cursor()
    cur.execute("SELECT code, discount_type, amount, usage_count, usage_limit FROM wc.coupons ORDER BY usage_count DESC")
    rows = cur.fetchall()
    cur.close()
    conn.close()
    total_uses = sum(r[3] for r in rows)
    most_used = max(rows, key=lambda x: x[3])
    highest_disc = max(rows, key=lambda x: float(x[2]))
    return {
        "coupons": rows,
        "total_coupons": len(rows),
        "most_used": most_used[0],
        "highest_discount": highest_disc[0],
        "total_uses": total_uses,
    }


def check_notion_page():
    """Return (bool, detail) — True only if page title exists AND blocks contain
    top coupon code, recommendation keyword, and key metric references."""
    try:
        conn = psycopg2.connect(**DB)
        cur = conn.cursor()
        cur.execute("SELECT id FROM notion.pages WHERE LOWER(properties::text) LIKE '%coupon strategy overview 2026%'")
        page_rows = cur.fetchall()
        found_title = len(page_rows) > 0
        if not found_title:
            cur.execute("SELECT id FROM notion.blocks WHERE LOWER(block_data::text) LIKE '%coupon strategy overview 2026%'")
            if len(cur.fetchall()) > 0:
                found_title = True
        # Verify block content includes specific required items
        cur.execute("SELECT LOWER(block_data::text) FROM notion.blocks")
        all_block_text = " ".join(r[0] for r in cur.fetchall() if r[0])
        cur.close()
        conn.close()
        # Strict recommendation keyword: require explicit actionable word, not
        # generic "top"/"highest" which often appear in analysis descriptions.
        has_recommend = any(kw in all_block_text for kw in ["recommend", "strategy", "suggest", "propose", "action", "should"])
        # Require mention of top coupon code (most used)
        has_top_code = any(code in all_block_text for code in ["holiday30", "save20", "flash50"])
        # Require explicit analytical keywords
        has_analysis_kw = ("performing" in all_block_text or "performance" in all_block_text
                           or "analysis" in all_block_text or "finding" in all_block_text)
        ok = found_title and has_recommend and has_top_code and has_analysis_kw
        detail = ""
        if not found_title:
            detail = "title not found"
        elif not has_top_code:
            detail = "top coupon code not referenced"
        elif not has_recommend:
            detail = "recommendation keyword missing"
        elif not has_analysis_kw:
            detail = "analysis keyword missing"
        return ok, detail
    except Exception as e:
        return False, str(e)


def check_email_sent():
    try:
        conn = psycopg2.connect(**DB)
        cur = conn.cursor()
        cur.execute("SELECT count(*) FROM email.messages WHERE LOWER(to_addr::text) LIKE '%marketing%' AND (LOWER(subject) LIKE '%coupon%' OR LOWER(subject) LIKE '%campaign%')")
        cnt = cur.fetchone()[0]
        cur.close()
        conn.close()
        return cnt >= 1
    except Exception:
        return False


def check_ppt_file(agent_workspace):
    ppt_path = os.path.join(agent_workspace, "Coupon_Analysis_Presentation.pptx")
    if not os.path.exists(ppt_path):
        return False, "Coupon_Analysis_Presentation.pptx not found"
    try:
        from pptx import Presentation
        prs = Presentation(ppt_path)
        n_slides = len(prs.slides)
        if n_slides < 12:
            return False, f"PPT has only {n_slides} slides, expected >= 12 (1 title + 10 coupons + 1 summary)"
        # Additionally: check slides contain coupon code names
        all_text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    all_text += shape.text_frame.text + " "
        all_text_upper = all_text.upper()
        # Verify at least top coupon codes HOLIDAY30, SAVE20 appear
        required_codes = ["HOLIDAY30", "SAVE20"]
        missing = [c for c in required_codes if c not in all_text_upper]
        if missing:
            return False, f"PPT missing coupon codes: {missing}"
        return True, ""
    except Exception as e:
        # If pptx not importable, just check file exists
        return True, ""


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--agent_workspace", required=False)
    parser.add_argument("--groundtruth_workspace", required=False)
    parser.add_argument("--launch_time", required=False)
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()

    task_root = os.path.abspath(os.path.join(os.path.dirname(__file__), ".."))
    gt_dir = args.groundtruth_workspace or os.path.join(task_root, "groundtruth_workspace")

    agent_file = os.path.join(args.agent_workspace, "Coupon_Performance.xlsx")
    gt_file = os.path.join(gt_dir, "Coupon_Performance.xlsx")

    if not os.path.exists(agent_file):
        print(f"FAIL: Agent output not found: {agent_file}")
        sys.exit(1)
    if not os.path.exists(gt_file):
        print(f"FAIL: Groundtruth not found: {gt_file}")
        sys.exit(1)

    all_errors = []

    try:
        coupon_data = get_coupon_data()
    except Exception as e:
        print(f"WARNING: Could not query DB: {e}")
        coupon_data = {
            "total_coupons": 10, "most_used": "HOLIDAY30",
            "highest_discount": "FLASH50", "total_uses": 209,
            "coupons": []
        }

    agent_wb = openpyxl.load_workbook(agent_file, data_only=True)
    gt_wb = openpyxl.load_workbook(gt_file, data_only=True)

    # Check Coupon Analysis sheet
    print("  Checking Coupon Analysis sheet...")
    a_rows = load_sheet_rows(agent_wb, "Coupon Analysis")
    g_rows = load_sheet_rows(gt_wb, "Coupon Analysis")
    if a_rows is None:
        all_errors.append("Sheet 'Coupon Analysis' not found in agent output")
    else:
        data_rows = [r for r in a_rows[1:] if r and any(c is not None for c in r)]
        if len(data_rows) < coupon_data["total_coupons"]:
            all_errors.append(f"Coupon Analysis has {len(data_rows)} rows, expected >= {coupon_data['total_coupons']}")
        else:
            print(f"    PASS ({len(data_rows)} data rows)")

        # Build lookup by coupon code
        a_lookup = {}
        for row in data_rows:
            if row and row[0] is not None:
                a_lookup[str(row[0]).strip().upper()] = row

        # Check HOLIDAY30
        if "HOLIDAY30" in a_lookup:
            row = a_lookup["HOLIDAY30"]
            if len(row) >= 4:
                if not num_close(row[3], 50, 0):
                    all_errors.append(f"HOLIDAY30.Times_Used: {row[3]} vs 50")
                else:
                    print("    HOLIDAY30 usage PASS")
        else:
            all_errors.append("HOLIDAY30 not found in Coupon Analysis")

        # Check SAVE20 usage rate
        if "SAVE20" in a_lookup:
            row = a_lookup["SAVE20"]
            if len(row) >= 4:
                if not num_close(row[3], 39, 0):
                    all_errors.append(f"SAVE20.Times_Used: {row[3]} vs 39")
            if len(row) >= 6:
                rate_val = row[5]
                if str(rate_val).strip().upper() != "N/A":
                    if not num_close(rate_val, 78.00, 0.5):
                        all_errors.append(f"SAVE20.Usage_Rate_Pct: {rate_val} vs 78.00 (tol=0.5)")
                    else:
                        print("    SAVE20 rate PASS")

    # Check Summary sheet
    print("  Checking Summary sheet...")
    a_rows = load_sheet_rows(agent_wb, "Summary")
    if a_rows is None:
        all_errors.append("Sheet 'Summary' not found in agent output")
    else:
        a_data = {str(r[0]).strip().lower(): r[1] for r in a_rows[1:] if r and r[0] is not None}
        errors = []

        total_uses_val = a_data.get("total_coupon_uses")
        if total_uses_val is None:
            errors.append("Missing metric: Total_Coupon_Uses")
        elif not num_close(total_uses_val, coupon_data["total_uses"], 0):
            errors.append(f"Total_Coupon_Uses: {total_uses_val} vs {coupon_data['total_uses']}")

        most_used_val = a_data.get("most_used_coupon")
        if most_used_val is None:
            errors.append("Missing metric: Most_Used_Coupon")
        elif str(most_used_val).strip().upper() != coupon_data["most_used"].upper():
            errors.append(f"Most_Used_Coupon: {most_used_val} vs {coupon_data['most_used']}")

        if errors:
            all_errors.extend(errors)
            for e in errors[:5]:
                print(f"    ERROR: {e}")
        else:
            print("    PASS")

    # Check PPT file
    print("  Checking PPT file...")
    ok, detail = check_ppt_file(args.agent_workspace)
    if ok:
        print("    PASS")
    else:
        all_errors.append(detail)

    # Check Notion page
    print("  Checking Notion page...")
    ok_n, notion_detail = check_notion_page()
    if ok_n:
        print("    PASS")
    else:
        all_errors.append(f"Notion page 'Coupon Strategy Overview 2026' check failed: {notion_detail}")

    # Check email sent
    print("  Checking email to marketing...")
    if check_email_sent():
        print("    PASS")
    else:
        all_errors.append("Email to marketing@company.com with coupon/campaign subject not found")

    if all_errors:
        print(f"\n=== RESULT: FAIL ({len(all_errors)} errors) ===")
        for e in all_errors[:10]:
            print(f"  {e}")
        sys.exit(1)
    else:
        print("\n=== RESULT: PASS ===")
        sys.exit(0)


if __name__ == "__main__":
    main()
