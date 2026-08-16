"""
Evaluation script for fetch-sf-sales-forecast-ppt-gcal task.

Checks:
1. Sales_Forecast_Data.xlsx with Q1_Actuals, Q2_Forecast, Segment_Mix sheets
2. Q2_Sales_Forecast.pptx with forecast content
3. Calendar event for board presentation
"""

import argparse
import json
import os
import sys

import openpyxl
import psycopg2

# ──────────────────────────────────────────────────────────────────────────
# EVALUATION GROUND TRUTH SPEC (gcal tz root-fix v3, case-study 2026-08-13)
# gcal.events.start_datetime is TIMESTAMPTZ; bare end_dt.hour / start_dt
# silently compares wrong in non-UTC PG sessions. Use gcal_helpers.
# ──────────────────────────────────────────────────────────────────────────
# task.md: SF company board presentation → America/Los_Angeles (PT)
EXPECTED_TIMEZONE = "America/Los_Angeles"

from utils.evaluation.gcal_helpers import get_zone_components  # noqa: E402

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

DB_CONFIG = {
    "host": os.environ.get("PGHOST", "localhost"),
    "port": int(os.environ.get("PGPORT", "5432")),
    "dbname": os.environ.get("PGDATABASE", "toolathlon_gym"),
    "user": "eigent",
    "password": "camel",
}

PASS_COUNT = 0
FAIL_COUNT = 0
BLOCKING_FAIL_COUNT = 0


def record(name, passed, detail="", runtime_only=False):
    global PASS_COUNT, FAIL_COUNT, BLOCKING_FAIL_COUNT
    if passed:
        PASS_COUNT += 1
        print(f"  [PASS] {name}")
    else:
        FAIL_COUNT += 1
        if not runtime_only:
            BLOCKING_FAIL_COUNT += 1
        msg = f": {detail[:300]}" if detail else ""
        suffix = " (runtime-only)" if runtime_only else ""
        print(f"  [FAIL] {name}{suffix}{msg}")


def num_close(a, b, tol=5000.0):
    try:
        return abs(float(a) - float(b)) <= tol
    except (TypeError, ValueError):
        return False


def str_contains(haystack, needle):
    if haystack is None or needle is None:
        return False
    return needle.strip().lower() in str(haystack).strip().lower()


def check_excel(agent_workspace):
    """Check Sales_Forecast_Data.xlsx."""
    print("\n=== Checking Excel Output ===")

    agent_file = os.path.join(agent_workspace, "Sales_Forecast_Data.xlsx")
    if not os.path.isfile(agent_file):
        record("Excel file exists", False, f"Not found: {agent_file}")
        return False

    record("Excel file exists", True)

    try:
        wb = openpyxl.load_workbook(agent_file, data_only=True)
    except Exception as e:
        record("Excel file readable", False, str(e))
        return False

    all_ok = True

    # Check Q1_Actuals sheet (exact match preferred)
    q1_sheet = None
    for name in wb.sheetnames:
        if name.strip().lower() == "q1_actuals":
            q1_sheet = name; break
    if not q1_sheet:
        for name in wb.sheetnames:
            if "q1" in name.lower() or "actual" in name.lower():
                q1_sheet = name
                break

    if not q1_sheet:
        record("Q1_Actuals sheet exists", False, f"Sheets: {wb.sheetnames}")
        all_ok = False
    else:
        record("Q1_Actuals sheet exists (exact)",
               q1_sheet.strip().lower() == "q1_actuals",
               f"Got '{q1_sheet}'")
        record("Q1_Actuals sheet exists", True)
        ws = wb[q1_sheet]
        rows = list(ws.iter_rows(values_only=True))
        data_rows = [r for r in rows[1:] if r and r[0]] if len(rows) > 1 else []
        record(
            "Q1_Actuals has 15 rows (5 regions x 3 months)",
            len(data_rows) >= 15,
            f"Found {len(data_rows)} data rows",
        )
        if len(data_rows) < 15:
            all_ok = False

        # Spot check: Asia Pacific Jan revenue ~57557
        for r in data_rows:
            if r[0] and "asia" in str(r[0]).lower() and r[1] in (1, "1", "January"):
                ok = num_close(r[3] if len(r) > 3 else 0, 57557.80, tol=5000)
                record("Asia Pacific Jan revenue ~57558", ok, f"Got {r[3] if len(r) > 3 else 'N/A'}")
                if not ok:
                    all_ok = False
                break

        # More Q1 spot checks (one per region in a different month)
        EXPECTED_Q1 = {
            ("asia pacific", 2): 46608.06,
            ("europe", 1): 48744.82,
            ("europe", 3): None,  # just check row exists
            ("latin america", 2): None,
            ("north america", 3): None,
        }
        region_month_found = set()
        for r in data_rows:
            if not (r and r[0] and len(r) > 1):
                continue
            region = str(r[0]).strip().lower()
            try:
                month = int(r[1]) if not isinstance(r[1], int) else r[1]
            except Exception:
                continue
            key = (region, month)
            if key in EXPECTED_Q1:
                region_month_found.add(key)
                expected = EXPECTED_Q1[key]
                if expected is not None and len(r) > 3:
                    ok = num_close(r[3], expected, tol=max(expected * 0.05, 100))
                    record(f"Q1 {region} M{month} revenue ~{expected:.2f}", ok,
                           f"Got {r[3]}")
        # Ensure we saw all five region/month combos (every region appears)
        regions_seen = {k[0] for k in region_month_found}
        expected_regions = {"asia pacific", "europe", "latin america", "north america"}
        missing = expected_regions - regions_seen
        record("Q1_Actuals covers all major regions",
               len(missing) == 0, f"Missing: {missing}")

    # Check Q2_Forecast sheet (exact match preferred)
    q2_sheet = None
    for name in wb.sheetnames:
        if name.strip().lower() == "q2_forecast":
            q2_sheet = name; break
    if not q2_sheet:
        for name in wb.sheetnames:
            if "q2" in name.lower() or "forecast" in name.lower():
                q2_sheet = name
                break

    if not q2_sheet:
        record("Q2_Forecast sheet exists", False, f"Sheets: {wb.sheetnames}")
        all_ok = False
    else:
        record("Q2_Forecast sheet exists (exact)",
               q2_sheet.strip().lower() == "q2_forecast",
               f"Got '{q2_sheet}'")
        record("Q2_Forecast sheet exists", True)
        ws2 = wb[q2_sheet]
        rows2 = list(ws2.iter_rows(values_only=True))
        data_rows2 = [r for r in rows2[1:] if r and r[0]] if len(rows2) > 1 else []
        record(
            "Q2_Forecast has exactly 5 region rows",
            len(data_rows2) == 5,
            f"Found {len(data_rows2)} data rows",
        )
        if len(data_rows2) != 5:
            all_ok = False

        # Exact expected growth rates per region (from GT)
        EXPECTED_GROWTH = {
            "asia pacific": 8.5,
            "europe": 4.2,
            "latin america": 6.8,
            "middle east": 7.3,
            "north america": 3.5,
        }
        growth_errors = 0
        for r in data_rows2:
            if r and len(r) >= 3 and r[0]:
                region = str(r[0]).strip().lower()
                exp = EXPECTED_GROWTH.get(region)
                if exp is not None:
                    try:
                        if abs(float(r[2]) - exp) > 0.5:
                            growth_errors += 1
                    except (TypeError, ValueError):
                        growth_errors += 1
        record("All 5 region growth rates match expected", growth_errors == 0,
               f"{growth_errors} mismatches")

    # Check Segment_Mix sheet (exact match preferred)
    seg_sheet = None
    for name in wb.sheetnames:
        if name.strip().lower() == "segment_mix":
            seg_sheet = name; break
    if not seg_sheet:
        for name in wb.sheetnames:
            if "segment" in name.lower() or "mix" in name.lower():
                seg_sheet = name
                break

    if not seg_sheet:
        record("Segment_Mix sheet exists", False, f"Sheets: {wb.sheetnames}")
        all_ok = False
    else:
        record("Segment_Mix sheet exists (exact)",
               seg_sheet.strip().lower() == "segment_mix",
               f"Got '{seg_sheet}'")
        record("Segment_Mix sheet exists", True)
        ws3 = wb[seg_sheet]
        rows3 = list(ws3.iter_rows(values_only=True))
        data_rows3 = [r for r in rows3[1:] if r and r[0]] if len(rows3) > 1 else []
        record(
            "Segment_Mix has 20 rows (5 regions x 4 segments)",
            len(data_rows3) == 20,
            f"Found {len(data_rows3)} data rows",
        )
        if len(data_rows3) != 20:
            all_ok = False
        # Revenue_Share_Pct per region should sum to ~100
        region_sums = {}
        for r in data_rows3:
            if r and len(r) >= 4 and r[0]:
                region = str(r[0]).strip().lower()
                try:
                    region_sums[region] = region_sums.get(region, 0) + float(r[3])
                except Exception:
                    pass
        bad_regions = [r for r, s in region_sums.items() if abs(s - 100) > 2]
        record("Revenue_Share_Pct sums to ~100 per region",
               len(bad_regions) == 0, f"Bad regions: {bad_regions}")

    wb.close()
    return all_ok


def check_pptx(agent_workspace):
    """Check Q2_Sales_Forecast.pptx."""
    print("\n=== Checking PowerPoint ===")

    pptx_file = os.path.join(agent_workspace, "Q2_Sales_Forecast.pptx")
    if not os.path.isfile(pptx_file):
        record("PowerPoint file exists", False, f"Not found: {pptx_file}")
        return False

    record("PowerPoint file exists", True)

    if Presentation is None:
        record("python-pptx available", False, "Cannot import pptx")
        return True  # File exists, can't verify content

    try:
        prs = Presentation(pptx_file)
        slides = prs.slides

        record(
            "PPT has >= 4 slides",
            len(slides) >= 4,
            f"Found {len(slides)} slides",
        )

        # Check content across all slides
        all_text = ""
        for slide in slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    all_text += shape.text_frame.text.lower() + " "

        record("PPT mentions Q2 forecast", "q2" in all_text and "forecast" in all_text)
        record(
            "PPT mentions regions",
            any(r in all_text for r in ["asia", "europe", "north america"]),
        )
        record(
            "PPT mentions growth",
            "growth" in all_text or "%" in all_text or "projection" in all_text,
        )

        return True
    except Exception as e:
        record("PPT readable", False, str(e))
        return False


def check_calendar():
    """Check calendar event for board presentation."""
    print("\n=== Checking Google Calendar ===")

    try:
        conn = psycopg2.connect(**DB_CONFIG)
        cur = conn.cursor()
        cur.execute(
            "SELECT summary, description, start_datetime, end_datetime FROM gcal.events"
        )
        events = cur.fetchall()
        cur.close()
        conn.close()
    except Exception as e:
        record("Calendar DB accessible", False, str(e))
        return False

    found = False
    for summary, description, start_dt, end_dt in events:
        summary_lower = (summary or "").lower()
        if ("board" in summary_lower or "forecast" in summary_lower or "sales" in summary_lower) and (
            "presentation" in summary_lower or "meeting" in summary_lower or "q2" in summary_lower
        ):
            found = True
            record("Board presentation event exists", True)

            # Check date is March 28, 2026 (blocking once event exists)
            # gcal.events.start_datetime is TIMESTAMPTZ; use session-tz-
            # independent helper to extract PT date/hour.
            ev_date, ev_hour, ev_minute = get_zone_components(start_dt, EXPECTED_TIMEZONE)
            date_str = ev_date.strftime("%Y-%m-%d") if ev_date else "None"
            record(
                "Event on March 28, 2026",
                date_str == "2026-03-28",
                f"Start date (PT): {date_str}",
            )

            # Check end_datetime is 11:30 per task (PT)
            if end_dt is not None:
                _ed, end_hour, end_min = get_zone_components(end_dt, EXPECTED_TIMEZONE)
                record("Event end_datetime is 11:30",
                       end_hour == 11 and end_min == 30,
                       f"End: {end_dt.isoformat()}")

            # Check description has forecast info
            desc_lower = (description or "").lower()
            has_info = any(
                kw in desc_lower for kw in ["revenue", "forecast", "growth", "region"]
            )
            record("Event description has forecast info", has_info)
            break

    if not found:
        record(
            "Board presentation event exists",
            False,
            f"Found {len(events)} events but none for board/forecast/sales presentation",
            runtime_only=True,
        )

    return found


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--agent_workspace", required=False, default=".")
    parser.add_argument("--groundtruth_workspace", required=False, default=".")
    parser.add_argument("--launch_time", required=False)
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()

    excel_ok = check_excel(args.agent_workspace)
    pptx_ok = check_pptx(args.agent_workspace)
    _ = check_calendar()  # returns bool but we gate on BLOCKING_FAIL_COUNT instead

    print(f"\n=== SUMMARY ===")
    print(f"  Excel:    {'PASS' if excel_ok else 'FAIL'}")
    print(f"  PPT:      {'PASS' if pptx_ok else 'FAIL'}")
    print(f"  Passed: {PASS_COUNT}, Failed: {FAIL_COUNT} (blocking_fail={BLOCKING_FAIL_COUNT})")

    # Blocking-fail gate: calendar absence is runtime_only, but
    # wrong-value calendar/excel/pptx checks are blocking.
    overall = excel_ok and pptx_ok and BLOCKING_FAIL_COUNT == 0
    print(f"  Overall:  {'PASS' if overall else 'FAIL'}")

    sys.exit(0 if overall else 1)


if __name__ == "__main__":
    main()
