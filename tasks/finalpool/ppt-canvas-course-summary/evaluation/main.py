"""
Evaluation script for ppt-canvas-course-summary task.

Dynamically queries PostgreSQL Canvas schema to compute expected values,
then checks agent output files for correctness.
"""

from argparse import ArgumentParser
import sys
import os
from pathlib import Path

def num_close(a, b, rel_tol=0.15, abs_tol=0.5):
    return abs(float(a) - float(b)) <= max(abs_tol, abs(float(b)) * rel_tol)



def get_expected_data():
    """Query PostgreSQL to compute expected course data from Canvas schema."""
    import psycopg2

    conn = psycopg2.connect(
        host=os.environ.get("PGHOST", "localhost"), port=5432, dbname="toolathlon_gym",
        user="eigent", password="camel"
    )
    cur = conn.cursor()

    # Dynamically resolve course_id by matching the course name and Fall 2013 term
    cur.execute("""
        SELECT id FROM canvas.courses
        WHERE (LOWER(name) LIKE %s OR LOWER(name) LIKE %s)
           AND (LOWER(COALESCE(name,'')) LIKE %s
                OR CAST(id AS TEXT) = '1')
        ORDER BY id ASC
        LIMIT 1
    """, ('%applied analytics%algorithms%', '%analytics & algorithms%', '%fall 2013%'))
    row = cur.fetchone()
    if row:
        course_id = row[0]
    else:
        # Fallback: first try any 'applied analytics' course, else id=1
        cur.execute("""SELECT id FROM canvas.courses
                       WHERE LOWER(name) LIKE %s
                       ORDER BY id ASC LIMIT 1""",
                    ('%applied analytics%',))
        r2 = cur.fetchone()
        course_id = r2[0] if r2 else 1

    # Enrollment counts
    cur.execute("""
        SELECT type, COUNT(*)
        FROM canvas.enrollments
        WHERE course_id = %s
        GROUP BY type ORDER BY type
    """, (course_id,))
    enrollment_raw = dict(cur.fetchall())
    enrollments = {
        "student": enrollment_raw.get("StudentEnrollment", 0),
        "teacher": enrollment_raw.get("TeacherEnrollment", 0),
        "ta": enrollment_raw.get("TaEnrollment", 0),
    }

    # Assignment performance
    cur.execute("""
        SELECT a.name,
               ROUND(AVG(s.score)::numeric, 2)::float as avg_score,
               MAX(s.score)::float as max_score,
               MIN(s.score)::float as min_score,
               COUNT(*) as sub_count
        FROM canvas.assignments a
        JOIN canvas.submissions s ON a.id = s.assignment_id
        WHERE a.course_id = %s AND s.score IS NOT NULL
        GROUP BY a.name
        ORDER BY a.name
    """, (course_id,))
    assignments = cur.fetchall()

    # Grade distribution
    cur.execute("""
        SELECT
          CASE
            WHEN avg_score >= 90 THEN 'A (90-100)'
            WHEN avg_score >= 80 THEN 'B (80-89)'
            WHEN avg_score >= 70 THEN 'C (70-79)'
            WHEN avg_score >= 60 THEN 'D (60-69)'
            ELSE 'F (<60)'
          END as grade_range,
          COUNT(*) as student_count
        FROM (
          SELECT s.user_id, AVG(s.score)::float as avg_score
          FROM canvas.submissions s
          JOIN canvas.assignments a ON s.assignment_id = a.id
          WHERE a.course_id = %s AND s.score IS NOT NULL
          GROUP BY s.user_id
        ) sub
        GROUP BY grade_range
        ORDER BY grade_range
    """, (course_id,))
    grades = cur.fetchall()
    total_graded_students = sum(g[1] for g in grades)

    conn.close()
    return enrollments, assignments, grades, total_graded_students


def check_excel(workspace, enrollments, assignments, grades, total_graded_students):
    """Check Course_Summary_AAA_F13.xlsx for correctness."""
    import openpyxl

    xlsx_path = Path(workspace) / "Course_Summary_AAA_F13.xlsx"
    if not xlsx_path.exists():
        return False, f"Course_Summary_AAA_F13.xlsx not found in {workspace}"

    wb = openpyxl.load_workbook(xlsx_path, data_only=True)

    # --- Check sheet names ---
    required_sheets = ["Enrollment Stats", "Assignment Performance", "Grade Distribution"]
    for sheet_name in required_sheets:
        # Case-insensitive sheet name check
        found = False
        for sn in wb.sheetnames:
            if sn.strip().lower() == sheet_name.lower():
                found = True
                break
        if not found:
            return False, f"Missing sheet '{sheet_name}'. Found: {wb.sheetnames}"

    # Helper to find sheet case-insensitively
    def get_sheet(name):
        for sn in wb.sheetnames:
            if sn.strip().lower() == name.lower():
                return wb[sn]
        return None

    # --- Check Enrollment Stats ---
    ws1 = get_sheet("Enrollment Stats")
    rows1 = list(ws1.iter_rows(values_only=True))
    if len(rows1) < 2:
        return False, "Enrollment Stats sheet has no data rows"

    header1 = [str(h).strip().lower() if h else "" for h in rows1[0]]
    for col in ["student_count", "teacher_count", "ta_count"]:
        if col.lower() not in header1:
            return False, f"Enrollment Stats missing column '{col}'. Found: {[str(h) for h in rows1[0]]}"

    data_row = rows1[1]
    idx_student = header1.index("student_count")
    idx_teacher = header1.index("teacher_count")
    idx_ta = header1.index("ta_count")

    if int(data_row[idx_student]) != enrollments["student"]:
        return False, f"Student count: expected {enrollments['student']}, got {data_row[idx_student]}"
    if int(data_row[idx_teacher]) != enrollments["teacher"]:
        return False, f"Teacher count: expected {enrollments['teacher']}, got {data_row[idx_teacher]}"
    if int(data_row[idx_ta]) != enrollments["ta"]:
        return False, f"TA count: expected {enrollments['ta']}, got {data_row[idx_ta]}"
    print("  [PASS] Enrollment Stats correct")

    # --- Check Assignment Performance ---
    ws2 = get_sheet("Assignment Performance")
    rows2 = list(ws2.iter_rows(values_only=True))
    if len(rows2) < 2:
        return False, "Assignment Performance sheet has no data rows"

    header2 = [str(h).strip().lower() if h else "" for h in rows2[0]]
    for col in ["assignment_name", "avg_score", "max_score", "min_score", "submission_count"]:
        if col.lower() not in header2:
            return False, f"Assignment Performance missing column '{col}'. Found: {[str(h) for h in rows2[0]]}"

    idx_map = {col: header2.index(col) for col in ["assignment_name", "avg_score", "max_score", "min_score", "submission_count"]}

    data_rows2 = rows2[1:]
    if len(data_rows2) != len(assignments):
        return False, f"Assignment Performance: expected {len(assignments)} rows, got {len(data_rows2)}"

    # Build lookup by assignment_name (case-insensitive) for order-independent comparison
    agent_lookup = {}
    for row in data_rows2:
        nv = row[idx_map["assignment_name"]]
        if nv is not None:
            agent_lookup[str(nv).strip().lower()] = row

    for (exp_name, exp_avg, exp_max, exp_min, exp_count) in assignments:
        key = exp_name.lower()
        if key not in agent_lookup:
            return False, f"Assignment '{exp_name}' not found in Assignment Performance"
        row = agent_lookup[key]

        avg_val = float(row[idx_map["avg_score"]]) if row[idx_map["avg_score"]] is not None else None
        if avg_val is None or abs(avg_val - exp_avg) > 0.5:
            return False, f"Assignment '{exp_name}' avg: expected {exp_avg}, got {avg_val}"

        max_val = float(row[idx_map["max_score"]]) if row[idx_map["max_score"]] is not None else None
        if max_val is None or abs(max_val - exp_max) > 0.5:
            return False, f"Assignment '{exp_name}' max: expected {exp_max}, got {max_val}"

        min_val = float(row[idx_map["min_score"]]) if row[idx_map["min_score"]] is not None else None
        if min_val is None or abs(min_val - exp_min) > 0.5:
            return False, f"Assignment '{exp_name}' min: expected {exp_min}, got {min_val}"

        count_val = int(row[idx_map["submission_count"]]) if row[idx_map["submission_count"]] is not None else None
        if count_val is None or count_val != exp_count:
            return False, f"Assignment '{exp_name}' count: expected {exp_count}, got {count_val}"

    print("  [PASS] Assignment Performance correct")

    # --- Check Grade Distribution ---
    ws3 = get_sheet("Grade Distribution")
    rows3 = list(ws3.iter_rows(values_only=True))
    if len(rows3) < 2:
        return False, "Grade Distribution sheet has no data rows"

    header3 = [str(h).strip().lower() if h else "" for h in rows3[0]]
    for col in ["grade_range", "student_count", "percentage"]:
        if col.lower() not in header3:
            return False, f"Grade Distribution missing column '{col}'. Found: {[str(h) for h in rows3[0]]}"

    idx_grade = {col: header3.index(col) for col in ["grade_range", "student_count", "percentage"]}
    data_rows3 = rows3[1:]

    if len(data_rows3) != len(grades):
        return False, f"Grade Distribution: expected {len(grades)} rows, got {len(data_rows3)}"

    for i, (exp_range, exp_count) in enumerate(grades):
        row = data_rows3[i]
        range_val = str(row[idx_grade["grade_range"]]).strip() if row[idx_grade["grade_range"]] else ""
        if range_val.lower() != exp_range.lower():
            return False, f"Grade row {i+1}: expected '{exp_range}', got '{range_val}'"

        count_val = int(row[idx_grade["student_count"]]) if row[idx_grade["student_count"]] is not None else None
        if count_val is None or count_val != exp_count:
            return False, f"Grade '{exp_range}' count: expected {exp_count}, got {count_val}"

        exp_pct = round(exp_count / total_graded_students * 100, 1)
        pct_val = float(row[idx_grade["percentage"]]) if row[idx_grade["percentage"]] is not None else None
        if pct_val is None or abs(pct_val - exp_pct) > 0.2:
            return False, f"Grade '{exp_range}' pct: expected {exp_pct}, got {pct_val}"

    print("  [PASS] Grade Distribution correct")

    wb.close()
    return True, "Excel file checks passed"


def check_pptx(workspace, enrollments, assignments):
    """Check Course_Summary_AAA_F13.pptx for correctness."""
    from pptx import Presentation

    pptx_path = Path(workspace) / "Course_Summary_AAA_F13.pptx"
    if not pptx_path.exists():
        return False, f"Course_Summary_AAA_F13.pptx not found in {workspace}"

    prs = Presentation(str(pptx_path))
    slides = list(prs.slides)

    # At least 5 slides
    if len(slides) < 5:
        return False, f"Expected at least 5 slides, got {len(slides)}"
    print(f"  Slide count: {len(slides)}")

    # Collect all text from all slides
    all_text = []
    for slide in slides:
        slide_texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    slide_texts.append(paragraph.text)
        all_text.append("\n".join(slide_texts))

    full_text = "\n".join(all_text).lower()

    # Check title slide contains course name
    first_slide_text = all_text[0].lower()
    if "applied analytics" not in first_slide_text or "algorithms" not in first_slide_text:
        return False, f"Title slide does not contain course name. Text: {all_text[0][:200]}"
    if "fall 2013" not in first_slide_text and "2013" not in first_slide_text:
        return False, f"Title slide does not contain 'Fall 2013'. Text: {all_text[0][:200]}"
    print("  [PASS] Title slide contains course name")

    # Check enrollment numbers appear somewhere with appropriate label context
    import re as _re
    student_str = str(enrollments["student"])
    teacher_str = str(enrollments["teacher"])
    ta_str = str(enrollments["ta"])

    # Student count: require explicit context (labels with number on same slide/paragraph)
    student_pattern = _re.compile(
        r"student[^\n]{0,20}[:\s\-]+" + _re.escape(student_str) + r"\b"
        r"|" + _re.escape(student_str) + r"\s*students?",
        _re.IGNORECASE,
    )
    if not student_pattern.search(full_text):
        return False, (f"Student count ({student_str}) not found with label "
                       f"in presentation")
    print("  [PASS] Student count found with label context")

    # Teacher count
    teacher_pattern = _re.compile(
        r"teacher[^\n]{0,20}[:\s\-]+" + _re.escape(teacher_str) + r"\b"
        r"|" + _re.escape(teacher_str) + r"\s*teachers?",
        _re.IGNORECASE,
    )
    if not teacher_pattern.search(full_text):
        return False, (f"Teacher count ({teacher_str}) not found with label "
                       f"in presentation")
    print("  [PASS] Teacher count found with label context")

    # TA count
    ta_pattern = _re.compile(
        r"\bta[^\n]{0,20}[:\s\-]+" + _re.escape(ta_str) + r"\b"
        r"|" + _re.escape(ta_str) + r"\s*tas?\b"
        r"|teaching\s+assistant[^\n]{0,20}[:\s\-]+" + _re.escape(ta_str) + r"\b",
        _re.IGNORECASE,
    )
    if not ta_pattern.search(full_text):
        return False, (f"TA count ({ta_str}) not found with label "
                       f"in presentation")
    print("  [PASS] TA count found with label context")

    # Check all assignment names appear
    for a in assignments:
        if a[0].lower() not in full_text:
            return False, f"Assignment '{a[0]}' not found in presentation"
    print("  [PASS] All assignment names found in presentation")

    # Check grade-related keywords
    grade_keywords = ["grade", "distribution", "a (90", "b (80", "c (70", "d (60", "f (<60"]
    grade_found = sum(1 for kw in grade_keywords if kw in full_text)
    if grade_found < 3:
        return False, f"Grade distribution content insufficient (found {grade_found}/7 keywords)"
    print("  [PASS] Grade distribution content found")

    # Check for findings/recommendations slide
    findings_keywords = ["finding", "recommendation", "key", "summary", "overall", "average"]
    findings_found = sum(1 for kw in findings_keywords if kw in full_text)
    if findings_found < 2:
        return False, f"Key findings slide content insufficient (found {findings_found}/6 keywords)"
    print("  [PASS] Key findings content found")

    # Check last slide mentions highest-performing assignment (max avg_score)
    if assignments:
        top_assign = max(assignments, key=lambda a: a[1])  # a[1]=exp_avg
        last_slide = all_text[-1].lower() if all_text else ""
        # The highest performing assignment name should appear in the last slide OR in the overall text
        if top_assign[0].lower() not in last_slide and top_assign[0].lower() not in full_text:
            return False, f"Highest performing assignment '{top_assign[0]}' not found in presentation"
        print(f"  [PASS] Highest-performing assignment '{top_assign[0]}' referenced")

    return True, "PPTX file checks passed"


def check_gsheet(enrollments=None, assignments=None):
    """Check that AAA F13 Course Dashboard was created in gsheet schema."""
    import psycopg2

    conn = psycopg2.connect(
        host=os.environ.get("PGHOST", "localhost"), port=5432, dbname="toolathlon_gym",
        user="eigent", password="camel"
    )
    cur = conn.cursor()

    cur.execute("""
        SELECT id, title FROM gsheet.spreadsheets
        WHERE LOWER(title) LIKE '%aaa%' AND LOWER(title) LIKE '%f13%'
    """)
    rows = cur.fetchall()

    if not rows:
        # Try broader search
        cur.execute("""
            SELECT id, title FROM gsheet.spreadsheets
            WHERE LOWER(title) LIKE '%course%dashboard%'
               OR LOWER(title) LIKE '%aaa%dashboard%'
               OR LOWER(title) LIKE '%analytics%algorithms%'
        """)
        rows = cur.fetchall()

    conn.close()

    if not rows:
        return False, "No 'AAA F13 Course Dashboard' spreadsheet found in online spreadsheet platform"

    spreadsheet_id = rows[0][0]
    spreadsheet_title = rows[0][1]
    print(f"  Found spreadsheet: '{spreadsheet_title}' (id={spreadsheet_id})")

    # Check that it has some content
    conn = psycopg2.connect(
        host=os.environ.get("PGHOST", "localhost"), port=5432, dbname="toolathlon_gym",
        user="eigent", password="camel"
    )
    cur = conn.cursor()
    cur.execute("""
        SELECT COUNT(*) FROM gsheet.cells
        WHERE spreadsheet_id = %s
    """, (spreadsheet_id,))
    cell_count = cur.fetchone()[0]
    # Fetch all cell values to verify content
    cur.execute("""
        SELECT value FROM gsheet.cells
        WHERE spreadsheet_id = %s
    """, (spreadsheet_id,))
    all_cells = cur.fetchall()
    conn.close()

    if cell_count < 4:
        return False, f"Dashboard spreadsheet has only {cell_count} cells, expected at least 4"

    all_text = " ".join(str(c[0]).lower() for c in all_cells if c[0] is not None)
    # Must contain course name reference and some enrollment number
    has_course = "applied analytics" in all_text or "aaa" in all_text or "algorithms" in all_text
    if not has_course:
        return False, f"Dashboard does not mention course name. Cells: {all_text[:300]}"

    # Key-field check: student count + assignment count must appear verbatim
    if enrollments is not None:
        student_str = str(enrollments["student"])
        if student_str not in all_text:
            return False, (f"Dashboard missing student count '{student_str}'. "
                           f"Cells excerpt: {all_text[:300]}")
        print(f"  [PASS] Dashboard contains student count {student_str}")
    if assignments is not None:
        assign_count = str(len(assignments))
        if assign_count not in all_text:
            return False, (f"Dashboard missing assignment count '{assign_count}'. "
                           f"Cells excerpt: {all_text[:300]}")
        print(f"  [PASS] Dashboard contains assignment count {assign_count}")

    print(f"  [PASS] Dashboard spreadsheet found with {cell_count} cells, contains course name")
    return True, "Google Sheet check passed"


if __name__ == "__main__":
    parser = ArgumentParser()
    parser.add_argument("--agent_workspace", required=False)
    parser.add_argument("--groundtruth_workspace", required=False)
    parser.add_argument("--res_log_file", required=False)
    parser.add_argument("--launch_time", required=False, help="Launch time")
    args = parser.parse_args()

    workspace = args.agent_workspace
    if not workspace:
        print("Error: --agent_workspace is required")
        exit(1)

    print("Fetching expected data from database...")
    try:
        enrollments, assignments, grades, total_graded_students = get_expected_data()
        print(f"  Enrollments: S={enrollments['student']}, T={enrollments['teacher']}, TA={enrollments['ta']}")
        print(f"  Assignments: {len(assignments)}")
        print(f"  Grade ranges: {len(grades)}, total graded students: {total_graded_students}")
    except Exception as e:
        print(f"Error querying database: {e}")
        import traceback
        traceback.print_exc()
        exit(1)

    all_passed = True

    # Check Excel
    print("\n--- Check 1: Excel File ---")
    try:
        ok, msg = check_excel(workspace, enrollments, assignments, grades, total_graded_students)
        if not ok:
            print(f"  [FAIL] {msg}")
            all_passed = False
        else:
            print(f"  {msg}")
    except Exception as e:
        print(f"  [FAIL] Excel check error: {e}")
        import traceback
        traceback.print_exc()
        all_passed = False

    # Check PPTX
    print("\n--- Check 2: PowerPoint File ---")
    try:
        ok, msg = check_pptx(workspace, enrollments, assignments)
        if not ok:
            print(f"  [FAIL] {msg}")
            all_passed = False
        else:
            print(f"  {msg}")
    except Exception as e:
        print(f"  [FAIL] PPTX check error: {e}")
        import traceback
        traceback.print_exc()
        all_passed = False

    # Check Google Sheet
    print("\n--- Check 3: Online Spreadsheet Dashboard ---")
    try:
        ok, msg = check_gsheet(enrollments, assignments)
        if not ok:
            print(f"  [FAIL] {msg}")
            all_passed = False
        else:
            print(f"  {msg}")
    except Exception as e:
        print(f"  [FAIL] Google Sheet check error: {e}")
        import traceback
        traceback.print_exc()
        all_passed = False

    if all_passed:
        print("\nPass all tests!")
        sys.exit(0)
    else:
        print("\nSome checks failed.")
        sys.exit(1)
