"""
Evaluation for sf-hr-performance-ppt task.

Dynamically queries PostgreSQL to compute expected values,
then checks agent output files for correctness.
"""
from argparse import ArgumentParser
import sys
import os
from pathlib import Path


def get_expected_data():
    """Query PostgreSQL to compute expected department performance values."""
    import psycopg2

    conn = psycopg2.connect(
        host=os.environ.get("PGHOST", "localhost"), port=5432, dbname="toolathlon_gym",
        user="eigent", password="camel"
    )
    cur = conn.cursor()

    cur.execute('''
        SELECT "DEPARTMENT",
               COUNT(*) as total_emp,
               ROUND(AVG("PERFORMANCE_RATING")::numeric, 2) as avg_rating,
               SUM(CASE WHEN "PERFORMANCE_RATING" >= 4 THEN 1 ELSE 0 END) as top_performers,
               ROUND(AVG("SALARY")::numeric, 2) as avg_salary
        FROM sf_data."HR_ANALYTICS__PUBLIC__EMPLOYEES"
        GROUP BY "DEPARTMENT"
        ORDER BY "DEPARTMENT"
    ''')
    departments = [(r[0], int(r[1]), float(r[2]), int(r[3]), float(r[4])) for r in cur.fetchall()]

    cur.execute('''
        SELECT COUNT(*) as total,
               ROUND(AVG("PERFORMANCE_RATING")::numeric, 2),
               SUM(CASE WHEN "PERFORMANCE_RATING" >= 4 THEN 1 ELSE 0 END),
               ROUND(AVG("SALARY")::numeric, 2)
        FROM sf_data."HR_ANALYTICS__PUBLIC__EMPLOYEES"
    ''')
    row = cur.fetchone()
    summary = (int(row[0]), float(row[1]), int(row[2]), float(row[3]))

    conn.close()
    return departments, summary


def check_excel(workspace, departments, summary):
    """Check Performance_Review.xlsx for correctness."""
    import openpyxl

    xlsx_path = Path(workspace) / "Performance_Review.xlsx"
    if not xlsx_path.exists():
        return False, f"Performance_Review.xlsx not found in {workspace}"

    wb = openpyxl.load_workbook(xlsx_path, data_only=True)

    if "By Department" not in wb.sheetnames:
        return False, f"Missing 'By Department' sheet. Found: {wb.sheetnames}"
    if "Summary" not in wb.sheetnames:
        return False, f"Missing 'Summary' sheet. Found: {wb.sheetnames}"

    # Check By Department sheet
    ws1 = wb["By Department"]
    rows1 = list(ws1.iter_rows(values_only=True))
    if len(rows1) < 2:
        return False, "By Department sheet has no data rows"

    header1 = [str(h).strip() if h else "" for h in rows1[0]]
    expected_cols = ["Department", "Total_Employees", "Avg_Rating", "Top_Performers_Count", "Avg_Salary"]
    for col in expected_cols:
        if col not in header1:
            return False, f"By Department missing column '{col}'. Found: {header1}"

    idx = {col: header1.index(col) for col in expected_cols}
    data_rows = rows1[1:]
    if len(data_rows) != len(departments):
        return False, f"By Department: expected {len(departments)} data rows, got {len(data_rows)}"

    for i, (exp_dept, exp_total, exp_rating, exp_top, exp_salary) in enumerate(departments):
        row = data_rows[i]
        dept_val = str(row[idx["Department"]]).strip() if row[idx["Department"]] else ""
        total_val = row[idx["Total_Employees"]]
        rating_val = row[idx["Avg_Rating"]]
        top_val = row[idx["Top_Performers_Count"]]
        salary_val = row[idx["Avg_Salary"]]

        if dept_val != exp_dept:
            return False, f"By Department row {i+1}: expected '{exp_dept}', got '{dept_val}'"
        if total_val is None or int(total_val) != exp_total:
            return False, f"'{exp_dept}' Total_Employees: expected {exp_total}, got {total_val}"
        if rating_val is None or abs(float(rating_val) - exp_rating) > 0.02:
            return False, f"'{exp_dept}' Avg_Rating: expected {exp_rating}, got {rating_val}"
        if top_val is None or int(top_val) != exp_top:
            return False, f"'{exp_dept}' Top_Performers_Count: expected {exp_top}, got {top_val}"
        if salary_val is None or abs(float(salary_val) - exp_salary) > 1.0:
            return False, f"'{exp_dept}' Avg_Salary: expected {exp_salary}, got {salary_val}"

    print("  [PASS] By Department data correct")

    # Check Summary sheet
    ws2 = wb["Summary"]
    rows2 = list(ws2.iter_rows(values_only=True))
    if len(rows2) < 2:
        return False, "Summary sheet has no data rows"

    summary_map = {}
    for row in rows2[1:]:
        if row[0]:
            summary_map[str(row[0]).strip()] = row[1]

    exp_metrics = {
        "Total_Employees": summary[0],
        "Overall_Avg_Rating": summary[1],
        "Total_Top_Performers": summary[2],
        "Overall_Avg_Salary": summary[3],
    }

    for metric, exp_val in exp_metrics.items():
        if metric not in summary_map:
            return False, f"Summary missing metric '{metric}'"
        actual = summary_map[metric]
        if isinstance(exp_val, int):
            if actual is None or int(actual) != exp_val:
                return False, f"Summary '{metric}': expected {exp_val}, got {actual}"
        else:
            if actual is None or abs(float(actual) - exp_val) > 1.0:
                return False, f"Summary '{metric}': expected {exp_val}, got {actual}"

    print("  [PASS] Summary data correct")
    wb.close()
    return True, "Excel file checks passed"


def check_pptx(workspace, departments, summary):
    """Check Performance_Review.pptx for correctness."""
    from pptx import Presentation

    pptx_path = Path(workspace) / "Performance_Review.pptx"
    if not pptx_path.exists():
        return False, f"Performance_Review.pptx not found in {workspace}"

    prs = Presentation(str(pptx_path))
    slides = list(prs.slides)

    # 1 title + 7 departments + 1 summary = 9 minimum
    if len(slides) < 9:
        return False, f"Expected at least 9 slides, got {len(slides)}"
    print(f"  Slide count: {len(slides)}")

    all_text = []
    for slide in slides:
        slide_texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    slide_texts.append(paragraph.text)
        all_text.append("\n".join(slide_texts))

    full_text = "\n".join(all_text).lower()

    # Check title slide
    first_text = all_text[0].lower()
    if "performance review" not in first_text:
        return False, f"Title slide does not contain 'Performance Review'. Text: {all_text[0][:200]}"
    print("  [PASS] Title slide correct")

    # Check all departments appear
    for dept, _, _, _, _ in departments:
        if dept.lower() not in full_text:
            return False, f"Department '{dept}' not found in presentation"
    print("  [PASS] All departments present")

    # Check summary values appear somewhere (total_employees, top_performers, avg_salary, avg_rating)
    summary_slides_text = "\n".join(all_text[-3:])
    total_emp_str = str(summary[0])
    top_str = str(summary[2])
    avg_sal_int = str(int(summary[3]))
    summary_matches = sum([
        total_emp_str in summary_slides_text,
        top_str in summary_slides_text,
        avg_sal_int in summary_slides_text,
    ])
    if summary_matches < 2:
        return False, f"Summary slides missing values: total={total_emp_str}, top={top_str}, avg_sal={avg_sal_int}"
    print(f"  [PASS] Summary values present ({summary_matches}/3)")

    # Each department slide must mention the dept AND at least the avg_rating AND avg_salary
    # (matching task.md "total employees, average rating, top performers count, and average salary")
    missing_per_dept = []
    for dept, total, rating, top, avg_sal in departments:
        dept_slide_texts = [s for s in all_text if dept.lower() in s.lower()]
        if not dept_slide_texts:
            missing_per_dept.append(f"{dept}: no slide mentions department")
            continue
        combined = " ".join(dept_slide_texts)
        combined_lower = combined.lower()
        # Check all four metrics appear on the slide containing this dept
        rating_str = f"{rating:.2f}"
        salary_int_str = str(int(avg_sal))
        salary_2d_str = f"{avg_sal:.2f}"
        has_total = str(total) in combined
        has_rating = rating_str in combined or str(rating) in combined
        has_top = str(top) in combined
        has_salary = salary_int_str in combined or salary_2d_str in combined
        missing = []
        if not has_total: missing.append(f"total={total}")
        if not has_rating: missing.append(f"avg_rating={rating_str}")
        if not has_top: missing.append(f"top={top}")
        if not has_salary: missing.append(f"avg_salary={salary_int_str}")
        if missing:
            missing_per_dept.append(f"{dept}: missing {missing}")
    if missing_per_dept:
        return False, "Dept slides incomplete: " + "; ".join(missing_per_dept[:3])
    print(f"  [PASS] All {len(departments)} department slides include total/rating/top/salary")

    # Verify summary slide title contains 'Overall Summary'
    summary_slide_found = False
    for s_text in all_text:
        # title is typically at top; check any occurrence of 'overall summary' (case-insensitive)
        if "overall summary" in s_text.lower():
            summary_slide_found = True
            break
    if not summary_slide_found:
        return False, "No slide contains 'Overall Summary' title"
    print("  [PASS] 'Overall Summary' slide present")

    return True, "PPTX file checks passed"


if __name__ == "__main__":
    parser = ArgumentParser()
    parser.add_argument("--agent_workspace", required=False)
    parser.add_argument("--groundtruth_workspace", required=False)
    parser.add_argument("--res_log_file", required=False)
    parser.add_argument("--launch_time", required=False, help="Launch time")
    args = parser.parse_args()

    workspace = args.agent_workspace
    if not workspace:
        print("Error: --agent_workspace is required")
        sys.exit(1)

    print("Fetching expected data from database...")
    try:
        departments, summary = get_expected_data()
        print(f"  Departments: {len(departments)}")
        print(f"  Summary: {summary}")
    except Exception as e:
        print(f"Error querying database: {e}")
        sys.exit(1)

    all_passed = True

    print("\n--- Check 1: Excel File ---")
    try:
        ok, msg = check_excel(workspace, departments, summary)
        if not ok:
            print(f"  [FAIL] {msg}")
            all_passed = False
        else:
            print(f"  {msg}")
    except Exception as e:
        print(f"  [FAIL] Excel check error: {e}")
        all_passed = False

    print("\n--- Check 2: PowerPoint File ---")
    try:
        ok, msg = check_pptx(workspace, departments, summary)
        if not ok:
            print(f"  [FAIL] {msg}")
            all_passed = False
        else:
            print(f"  {msg}")
    except Exception as e:
        print(f"  [FAIL] PPTX check error: {e}")
        all_passed = False

    if all_passed:
        print("\nPass all tests!")
        sys.exit(0)
    else:
        print("\nSome checks failed.")
        sys.exit(1)
