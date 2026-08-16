"""Evaluation script for pw-yf-risk-assessment-excel-ppt."""
import os
import argparse, json, os, sys
import openpyxl

# --- verify_v2 smart primitives ---
_EVAL_DIR = os.path.dirname(os.path.abspath(__file__))
_PACK_ROOT = os.path.abspath(os.path.join(_EVAL_DIR, "..", "..", "..", ".."))
if _PACK_ROOT not in sys.path:
    sys.path.insert(0, _PACK_ROOT)
try:
    from utils.verify_v2 import smart_column_exists
    from utils.verify_v2.eval_helpers import get_sheet_rows_as_dicts, get_gt_column_values
    _HAS_VERIFY_V2 = True
except Exception:
    _HAS_VERIFY_V2 = False

TASK_NAME = "pw-yf-risk-assessment-excel-ppt"


DB_CONFIG = {
    "host": os.environ.get("PGHOST", "localhost"), "port": int(os.environ.get("PGPORT", "5432")),
    "dbname": os.environ.get("PGDATABASE", "toolathlon_gym"),
    "user": "eigent", "password": "camel"
}

PASS_COUNT = 0
FAIL_COUNT = 0

def check(name, condition, detail=""):
    global PASS_COUNT, FAIL_COUNT
    if condition:
        PASS_COUNT += 1
        print(f"  [PASS] {name}")
    else:
        FAIL_COUNT += 1
        detail_str = str(detail)[:200] if detail else ""
        print(f"  [FAIL] {name}: {detail_str}")

def safe_float(val, default=None):
    try:
        if val is None:
            return default
        return float(str(val).replace(',', '').replace('%', '').replace('$', '').strip())
    except (ValueError, TypeError):
        return default

def get_conn():
    import psycopg2
    return psycopg2.connect(**DB_CONFIG)

def run_evaluation(agent_workspace, groundtruth_workspace, launch_time, res_log_file):
    global PASS_COUNT, FAIL_COUNT
    PASS_COUNT = 0
    FAIL_COUNT = 0

    
    excel_path = os.path.join(agent_workspace, "Risk_Assessment_Report.xlsx")
    check("Risk_Assessment_Report.xlsx exists", os.path.exists(excel_path))
    if os.path.exists(excel_path):
        wb = openpyxl.load_workbook(excel_path)
        gt_path = os.path.join(groundtruth_workspace, "Risk_Assessment_Report.xlsx")
        gt_wb = openpyxl.load_workbook(gt_path) if os.path.exists(gt_path) else None

        # Deterministic synonym table — pre-LLM fallback for column matching.
        # Bridges generic task.md vocabulary ("internal/external/gap") to the
        # YF-specific GT schema (Symbol/Current_Price/Target_Price/Upside).
        COLUMN_SYNONYMS = {
            'Symbol': ['symbol', 'ticker', 'stock', 'code', 'name',
                       'company', 'identifier', 'dimension', 'primary_dimension'],
            'Current_Price': ['current_price', 'current price', 'price',
                              'market_price', 'last_price', 'close', 'internal',
                              'internal_metric', 'internal_value', 'actual'],
            'Target_Price': ['target_price', 'target price', 'analyst_target',
                             'fair_value', 'benchmark', 'external',
                             'external_benchmark', 'expected', 'forecast',
                             'projected_price'],
            'Upside': ['upside', 'upside_pct', 'upside_percent', 'gap',
                       'difference', 'delta', 'variance', 'gap_pct',
                       'percent_change', 'pct_diff', 'risk', 'risk_score'],
            'Metric': ['metric', 'name', 'key', 'kpi', 'measure'],
            'Value': ['value', 'amount', 'number', 'val'],
            'Priority': ['priority', 'rank', 'order', 'level'],
            'Action': ['action', 'recommendation', 'recommended_action',
                       'next_step', 'item'],
        }

        def _header_match_synonym(expected, headers_lower):
            for syn in COLUMN_SYNONYMS.get(expected, []):
                if syn.lower() in headers_lower:
                    return syn
            for syn in COLUMN_SYNONYMS.get(expected, []):
                for h in headers_lower:
                    if syn.lower() in h or h in syn.lower():
                        return h
            return None

        def check_columns(sheet_name, expected_cols, min_rows):
            """Verify sheet exists, >= min_rows, contains required columns.
            Lookup order: exact -> deterministic synonym -> LLM mapper."""
            check(f"{sheet_name} sheet exists", sheet_name in wb.sheetnames)
            if sheet_name not in wb.sheetnames:
                return
            _ws = wb[sheet_name]
            _data_rows = list(_ws.iter_rows(min_row=2, values_only=True))
            check(f"{sheet_name} has >= {min_rows} rows",
                  len(_data_rows) >= min_rows, f"got {len(_data_rows)}")
            _raw_headers = [str(c.value).strip() if c.value else "" for c in _ws[1]]
            _headers_lower = [h.lower() for h in _raw_headers]
            for _exp in expected_cols:
                if _exp.lower() in _headers_lower:
                    check(f"{sheet_name} has {_exp} column", True, "exact match")
                    continue
                _syn = _header_match_synonym(_exp, _headers_lower)
                if _syn:
                    check(f"{sheet_name} has {_exp} column", True,
                          f"matched via synonym {_syn!r}")
                    continue
                if _HAS_VERIFY_V2 and gt_wb is not None:
                    _raw_h_v2, _agent_rows = get_sheet_rows_as_dicts(wb, sheet_name)
                    _gt_vals = get_gt_column_values(gt_wb, sheet_name, _exp)
                    _ok, _matched, _reason = smart_column_exists(
                        expected_col=_exp, agent_headers=_raw_h_v2,
                        gt_samples=_gt_vals[:3], agent_rows=_agent_rows,
                        task_name=TASK_NAME,
                    )
                    _detail = _reason
                    if _ok and _matched and _matched.lower() != _exp.lower():
                        _detail = f"LLM-mapped to {_matched!r}"
                    check(f"{sheet_name} has {_exp} column", _ok, _detail)
                else:
                    check(f"{sheet_name} has {_exp} column",
                          False, f"headers: {_raw_headers[:8]}")

        check("Data_Analysis sheet exists", "Data_Analysis" in wb.sheetnames)
        if "Data_Analysis" in wb.sheetnames:
            ws = wb["Data_Analysis"]
            data_rows = list(ws.iter_rows(min_row=2, values_only=True))
            check("Data_Analysis has >= 5 rows", len(data_rows) >= 5, f"got {len(data_rows)}")

            check_columns('Data_Analysis', ['Symbol', 'Current_Price', 'Target_Price', 'Upside'], 5)
        check("Metrics sheet exists", "Metrics" in wb.sheetnames)
        if "Metrics" in wb.sheetnames:
            ws = wb["Metrics"]
            data_rows = list(ws.iter_rows(min_row=2, values_only=True))
            check("Metrics has >= 3 rows", len(data_rows) >= 3, f"got {len(data_rows)}")

            check_columns('Metrics', ['Metric', 'Value'], 3)
        check("Recommendations sheet exists", "Recommendations" in wb.sheetnames)
        if "Recommendations" in wb.sheetnames:
            ws = wb["Recommendations"]
            data_rows = list(ws.iter_rows(min_row=2, values_only=True))
            check("Recommendations has >= 2 rows", len(data_rows) >= 2, f"got {len(data_rows)}")

            check_columns('Recommendations', ['Priority', 'Action'], 2)

    # PPT/processor checks moved out of Excel-existence block to ensure they
    # are always enumerated even if Excel is missing.
    import glob as globmod
    pptx_path = os.path.join(agent_workspace, "Risk_Assessment_Presentation.pptx")
    if os.path.exists(pptx_path):
        check("Risk_Assessment_Presentation.pptx exists", True)
    else:
        pptx_files = globmod.glob(os.path.join(agent_workspace, "*.pptx"))
        check("Risk_Assessment_Presentation.pptx exists", len(pptx_files) >= 1, f"found {len(pptx_files)} pptx files")
        if pptx_files:
            pptx_path = pptx_files[0]
    if os.path.exists(pptx_path):
        try:
            from pptx import Presentation
            prs = Presentation(pptx_path)
            slide_count = len(prs.slides)
            check("Presentation has >= 4 slides", slide_count >= 4, f"got {slide_count}")
            # Aggregate text across slides for content validation
            all_text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        all_text += " " + shape.text
            all_text_low = all_text.lower()
            # Risk/recommendation content keywords expected per task topic
            expected_terms = ["risk", "recommend"]
            missing = [t for t in expected_terms if t not in all_text_low]
            check("Presentation covers risk/recommendation content",
                  len(missing) == 0, f"missing terms: {missing}")
        except ImportError:
            check("python-pptx available", False, "python-pptx not installed; cannot validate slide count")

    check("yf_risk_processor.py exists", os.path.exists(os.path.join(agent_workspace, "yf_risk_processor.py")))


    return FAIL_COUNT == 0, f"Passed {PASS_COUNT}/{PASS_COUNT + FAIL_COUNT} checks"

def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--agent_workspace", required=False, default=".")
    parser.add_argument("--groundtruth_workspace", required=False, default=".")
    parser.add_argument("--launch_time", required=False, default="2026-03-07 10:00:00")
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()

    success, message = run_evaluation(
        args.agent_workspace, args.groundtruth_workspace,
        args.launch_time, args.res_log_file
    )
    print(message)
    sys.exit(0 if success else 1)

if __name__ == "__main__":
    main()
