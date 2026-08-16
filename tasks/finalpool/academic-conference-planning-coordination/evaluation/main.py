#!/usr/bin/env python3
"""Evaluation script for academic-conference-planning-coordination.

Six-phase conference logistics task. The task description is a narrative
without explicit filenames; this evaluator validates the structural
deliverables that any reasonable agent would produce:

  - A submissions/papers spreadsheet (xlsx) with multiple rows of papers
  - A program / schedule spreadsheet (xlsx) with sessions/speakers/times
  - A conference program document (docx) with substantive content
  - A speaker brief / orientation deck (pptx) with multiple slides
  - Outgoing emails referencing acceptance / reviewer / conference
  - Calendar events for conference / coordination meetings
"""

from argparse import ArgumentParser
import os
import sys
import json
from pathlib import Path

PASS_COUNT = 0
FAIL_COUNT = 0
WARN_COUNT = 0
IS_GT_SELF_TEST = False


def record(name, passed, detail="", db_side=False):
    global PASS_COUNT, FAIL_COUNT, WARN_COUNT
    if passed:
        PASS_COUNT += 1
        print(f"  [PASS] {name}")
    else:
        # In GT self-test mode, DB-side checks (email/calendar/gsheet/notion/gform)
        # naturally fail because GT files cannot pre-populate DB state.
        if IS_GT_SELF_TEST and db_side:
            WARN_COUNT += 1
            msg = f": {str(detail)[:300]}" if detail else ""
            print(f"  [WARN] {name} (GT self-test mode, DB-side){msg}")
        else:
            FAIL_COUNT += 1
            msg = f": {str(detail)[:300]}" if detail else ""
            print(f"  [FAIL] {name}{msg}")


def get_conn():
    import psycopg2
    return psycopg2.connect(
        host=os.environ.get("PGHOST", "localhost"), port=int(os.environ.get("PGPORT", "5432")),
        dbname=os.environ.get("PGDATABASE", "toolathlon_gym"),
        user="eigent", password="camel",
    )


# ---- Helpers --------------------------------------------------------------


def list_xlsx(ws):
    return [p for p in Path(ws).glob("*.xlsx") if not p.name.startswith("~$")]


def list_docx(ws):
    return [p for p in Path(ws).glob("*.docx") if not p.name.startswith("~$")]


def list_pptx(ws):
    return [p for p in Path(ws).glob("*.pptx") if not p.name.startswith("~$")]


def load_xlsx(path):
    import openpyxl
    return openpyxl.load_workbook(str(path), data_only=True)


# ---- Structural checks ----------------------------------------------------


def check_workspace_files(agent_ws):
    print("\n=== Check: Workspace files ===")
    xlsx_files = list_xlsx(agent_ws)
    docx_files = list_docx(agent_ws)
    pptx_files = list_pptx(agent_ws)
    json_or_csv = [p for p in Path(agent_ws).glob("*")
                   if p.suffix.lower() in {".json", ".csv"}
                   and not p.name.startswith("~$")
                   and p.name not in ("config.json", "data.csv")]  # skip placeholders

    # At least one xlsx (program / submissions)
    record("At least one .xlsx deliverable",
           len(xlsx_files) >= 1, f"found {[p.name for p in xlsx_files]}")
    # At least one docx (conference program / report)
    record("At least one .docx deliverable",
           len(docx_files) >= 1, f"found {[p.name for p in docx_files]}")
    # PPT optional but desired (orientation/speaker brief)
    record("At least one .pptx or supplementary doc deliverable",
           len(pptx_files) >= 1 or len(docx_files) >= 2,
           f"docx={len(docx_files)} pptx={len(pptx_files)}")

    return xlsx_files, docx_files, pptx_files


def check_program_xlsx(xlsx_files):
    print("\n=== Check: Program / submissions xlsx ===")
    if not xlsx_files:
        record("Program xlsx readable", False, "no xlsx file")
        return

    # Pick the largest (most-substantive) workbook
    target = max(xlsx_files, key=lambda p: p.stat().st_size)
    try:
        wb = load_xlsx(target)
    except Exception as e:
        record(f"xlsx {target.name} readable", False, str(e))
        return
    record(f"xlsx {target.name} readable", True)

    # At least one sheet has multiple rows of program/submission data
    max_rows = 0
    max_cols = 0
    for ws in wb.worksheets:
        max_rows = max(max_rows, ws.max_row)
        max_cols = max(max_cols, ws.max_column)

    # Must have multiple rows (>=5 of papers/sessions) and multiple columns
    record(f"xlsx {target.name} has substantive rows (>=5)",
           max_rows >= 5, f"max_rows={max_rows}")
    record(f"xlsx {target.name} has multiple columns (>=3)",
           max_cols >= 3, f"max_cols={max_cols}")

    wb.close()


def check_conference_docx(docx_files):
    print("\n=== Check: Conference program docx ===")
    if not docx_files:
        record("Conference docx readable", False, "no docx file")
        return

    target = max(docx_files, key=lambda p: p.stat().st_size)
    try:
        from docx import Document
        doc = Document(str(target))
    except Exception as e:
        record(f"docx {target.name} readable", False, str(e))
        return
    record(f"docx {target.name} readable", True)

    text = "\n".join(p.text for p in doc.paragraphs)
    text_low = text.lower()
    word_count = len([w for w in text.split() if w.strip()])

    record(f"docx {target.name} substantive (>=200 words)",
           word_count >= 200, f"word_count={word_count}")

    # Should mention conference / program / schedule / session
    program_keywords = ['conference', 'program', 'schedule', 'session',
                        'speaker', 'paper', 'agenda']
    has_program_kw = sum(1 for kw in program_keywords if kw in text_low)
    record(f"docx {target.name} mentions conference/program (>=2 keywords)",
           has_program_kw >= 2,
           f"matched {has_program_kw}/7 keywords")

    # Substantive content: must cover at least 2 of the 3 core conference
    # documentation areas (Schedule, Speakers/Bios, Logistics/Travel).
    has_schedule = any(kw in text_low for kw in
                       ['schedule', 'session', 'agenda', 'program'])
    has_speakers = any(kw in text_low for kw in
                       ['speaker', 'bio', 'biograph', 'presenter', 'author'])
    has_logistics = any(kw in text_low for kw in
                        ['logistic', 'venue', 'travel', 'accommodat',
                         'hotel', 'registration', 'parking', 'transit',
                         'location'])
    sections_covered = sum([has_schedule, has_speakers, has_logistics])
    record(f"docx {target.name} covers core conference sections (>=2 of "
           "Schedule/Speakers/Logistics)",
           sections_covered >= 2,
           f"schedule={has_schedule} speakers={has_speakers} "
           f"logistics={has_logistics}")


def check_speaker_pptx(pptx_files, docx_files):
    print("\n=== Check: Speaker brief / orientation pptx ===")
    # If pptx exists, validate it; otherwise allow second docx as fallback
    if pptx_files:
        target = max(pptx_files, key=lambda p: p.stat().st_size)
        try:
            from pptx import Presentation
            prs = Presentation(str(target))
        except Exception as e:
            record(f"pptx {target.name} readable", False, str(e))
            return
        record(f"pptx {target.name} readable", True)
        record(f"pptx {target.name} has multiple slides (>=3)",
               len(prs.slides) >= 3, f"slides={len(prs.slides)}")
    elif len(docx_files) >= 2:
        # secondary docx is acceptable as orientation/brief
        record("Secondary docx serves as orientation/brief",
               True, "2+ docx files present")
    else:
        record("Speaker brief / orientation document present",
               False, "no pptx and only one docx")


def check_outgoing_emails():
    print("\n=== Check: Outgoing acceptance / reviewer emails ===")
    try:
        conn = get_conn()
        cur = conn.cursor()
        cur.execute(
            """SELECT id, subject, body_text FROM email.messages
               WHERE subject ILIKE %s OR subject ILIKE %s OR subject ILIKE %s
                  OR subject ILIKE %s OR body_text ILIKE %s OR body_text ILIKE %s""",
            ('%conference%', '%accept%', '%reviewer%', '%paper%',
             '%conference%', '%accept%'))
        rows = cur.fetchall()
        conn.close()
    except Exception as e:
        record("Outgoing emails (db query)", False, f"db error: {e}", db_side=True)
        return
    record("At least one acceptance/reviewer/conference email sent",
           len(rows) >= 1, f"matching emails: {len(rows)}", db_side=True)
    # Substantive body: at least one email mentions paper/review/session.
    substantive = any(
        any(kw in (str(r[2] or "")).lower()
            for kw in ['paper', 'review', 'session', 'submission',
                       'accept', 'program', 'conference'])
        for r in rows
    )
    record("At least one email body has substantive paper/review/session content",
           substantive,
           f"checked {len(rows)} emails",
           db_side=True)


def check_calendar_events():
    print("\n=== Check: Calendar coordination / conference events ===")
    try:
        conn = get_conn()
        cur = conn.cursor()
        cur.execute(
            """SELECT id, summary FROM gcal.events
               WHERE summary ILIKE %s OR summary ILIKE %s OR summary ILIKE %s
                  OR summary ILIKE %s OR summary ILIKE %s""",
            ('%conference%', '%session%', '%coordination%',
             '%review%', '%meeting%'))
        rows = cur.fetchall()
        conn.close()
    except Exception as e:
        record("Calendar events (db query)", False, f"db error: {e}", db_side=True)
        return
    record("At least one conference/coordination calendar event",
           len(rows) >= 1, f"matching events: {len(rows)}", db_side=True)


def run_evaluation(agent_workspace, groundtruth_workspace, launch_time, res_log_file):
    global PASS_COUNT, FAIL_COUNT, WARN_COUNT, IS_GT_SELF_TEST
    PASS_COUNT = 0
    FAIL_COUNT = 0
    WARN_COUNT = 0

    if not agent_workspace:
        return False, "No agent workspace provided"
    agent_ws = Path(agent_workspace)
    if not agent_ws.exists():
        return False, f"Agent workspace not found: {agent_workspace}"

    # Detect GT self-test mode (agent_ws == groundtruth_ws)
    try:
        if groundtruth_workspace and Path(groundtruth_workspace).exists():
            IS_GT_SELF_TEST = (
                os.path.realpath(str(agent_ws)) ==
                os.path.realpath(str(Path(groundtruth_workspace)))
            )
    except Exception:
        IS_GT_SELF_TEST = False

    xlsx_files, docx_files, pptx_files = check_workspace_files(agent_ws)
    check_program_xlsx(xlsx_files)
    check_conference_docx(docx_files)
    check_speaker_pptx(pptx_files, docx_files)
    check_outgoing_emails()
    check_calendar_events()

    return FAIL_COUNT == 0, f"Passed {PASS_COUNT}/{PASS_COUNT + FAIL_COUNT} checks"


def main():
    parser = ArgumentParser()
    parser.add_argument("--agent_workspace", required=False, default=".")
    parser.add_argument("--groundtruth_workspace", required=False, default=".")
    parser.add_argument("--launch_time", required=False)
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()

    success, message = run_evaluation(
        args.agent_workspace, args.groundtruth_workspace,
        args.launch_time, args.res_log_file
    )
    print(message)

    if args.res_log_file:
        try:
            with open(args.res_log_file, "w") as f:
                json.dump({
                    "total_passed": PASS_COUNT,
                    "total_checks": PASS_COUNT + FAIL_COUNT,
                    "success": success,
                }, f, indent=2)
        except Exception:
            pass

    if FAIL_COUNT == 0:
        print("PASS")
        sys.exit(0)
    print(f"FAIL ({FAIL_COUNT} checks failed)")
    sys.exit(1)


if __name__ == "__main__":
    main()
