"""
Evaluation for academic-presentation-builder task.

Checks:
1. LLM_Reasoning_Review.docx: headings, summary table, required keywords
2. LLM_Reasoning_Slides.pptx: slide count, title, key terms
3. word_count.txt: exists with two lines containing word counts
"""
import os
import sys
import json
import re
from argparse import ArgumentParser
from datetime import datetime

from docx import Document
from pptx import Presentation


def _norm(s):
    """Lowercase and strip spaces/hyphens/underscores for fuzzy matching."""
    return re.sub(r"[\s\-_]+", "", s.lower())


def text_contains_mention(text, mention):
    """Case-insensitive presence check that tolerates common writing variants.

    In addition to a plain substring check, this accepts:
      - hyphen/space variants ("chain-of-thought" vs "chain of thought"),
      - the standard "ToT" abbreviation for "tree of thoughts" (word-bounded so
        "total"/"totally" do not falsely match).
    """
    if mention.lower() in text:
        return True
    nm = _norm(mention)
    if nm and nm in _norm(text):
        return True
    if mention.lower() == "tree of thoughts" and re.search(r"\btot\b", text):
        return True
    return False


def check_word_doc(agent_workspace, gt_data):
    """Check the Word document for required structure and content.

    Returns (passed, total, critical_failed) — critical_failed is a count of
    critical checks that must all pass for the overall eval to succeed.
    """
    passed = 0
    total = 0
    critical_failed = 0
    filename = gt_data["review_doc"]["filename"]
    doc_path = os.path.join(agent_workspace, filename)

    # Check file exists (critical)
    total += 1
    if not os.path.exists(doc_path):
        print(f"  FAIL: {filename} not found at {doc_path}")
        critical_failed += 1
        return passed, total, critical_failed

    passed += 1
    print(f"  PASS: {filename} exists")

    doc = Document(doc_path)

    # Extract text from paragraphs (tracking heading-style paragraphs and
    # short/bold candidate headings for robustness) AND from all table cells.
    all_text = []
    headings_found = []   # paragraphs that use a Heading paragraph style
    short_paras = []      # short standalone paragraphs (heading fallback candidates)
    bold_paras = []       # paragraphs that are entirely bold (heading fallback candidates)
    for para in doc.paragraphs:
        txt = para.text.strip()
        all_text.append(para.text)
        if not txt:
            continue
        if para.style and para.style.name and "Heading" in para.style.name:
            headings_found.append(txt)
        else:
            is_bold = bool(para.runs) and all((r.bold or not r.text.strip()) for r in para.runs)
            if is_bold:
                bold_paras.append(txt)
            if len(txt) <= 60:
                short_paras.append(txt)

    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                all_text.append(cell.text)

    full_text = " ".join(all_text).lower()

    # Check required headings. A heading counts as present if:
    #   1) a Heading-styled paragraph contains it (case-insensitive), OR
    #   2) a short standalone paragraph contains it (case-insensitive), OR
    #   3) a bold paragraph contains it (case-insensitive).
    # This accepts python-docx Heading styles, bold-plain-paragraph headings,
    # and plain-text headings that carry a numbering prefix ("1. Introduction")
    # or a slightly extended title ("Introduction to LLM Reasoning").
    required_headings = gt_data["review_doc"]["required_headings"]
    for heading in required_headings:
        total += 1
        hl = heading.lower()
        found = any(hl in h.lower() for h in headings_found)
        if not found:
            found = any(hl in h.lower() for h in short_paras)
        if not found:
            found = any(hl in h.lower() for h in bold_paras)
        if found:
            passed += 1
            print(f"  PASS: Heading '{heading}' found")
        else:
            print(f"  FAIL: Heading '{heading}' not found. Styled headings: {headings_found}")

    # Locate the summary table by matching its header row against the required
    # columns (with synonym support), instead of assuming it is tables[0]. This
    # is robust to other tables appearing earlier in the document.
    tables = doc.tables
    expected_cols = [c.lower() for c in gt_data["review_doc"]["required_table_columns"]]
    synonyms = {
        "title": ["paper", "name", "title"],
        "method": ["approach", "technique", "method"],
        "contribution": ["contribution", "key", "main"],
        "finding": ["result", "outcome", "finding"],
        "author": ["author"],
        "year": ["year", "date"],
    }

    def col_matches(expected_col, header_cells):
        if any(expected_col in hc for hc in header_cells):
            return True
        for keyword, syns in synonyms.items():
            if keyword in expected_col:
                if any(any(s in hc for s in syns) for hc in header_cells):
                    return True
        return False

    def count_non_empty(rows):
        n = 0
        for row in rows:
            cell_text = " ".join(cell.text.strip() for cell in row.cells)
            if cell_text.strip():
                n += 1
        return n

    def looks_like_header(row):
        """True if the row is a column-header row rather than a data row.

        A header row matches at least 2 required columns (via synonyms) AND
        does not itself look like paper data (no 4-digit year cell, no long
        paper-title cell). This keeps a no-header table whose first data row
        happens to contain column-like words (e.g. a method cell "Method: CoT",
        a title cell "Paper: X") from being misread as a header row, which
        would otherwise undercount its data rows.
        """
        cells = [cell.text.strip().lower() for cell in row.cells]
        score = sum(1 for ec in expected_cols if col_matches(ec, cells))
        if score < 2:
            return False
        if any(re.match(r'^\s*(19|20)\d{2}\s*$', c) for c in cells):
            return False
        if any(len(c) > 40 for c in cells):
            return False
        return True

    # Locate the summary table. For each table decide whether its first row is
    # a header or already paper data; derive the data-row count under that
    # interpretation. Prefer the table whose data-row count equals the expected
    # count; among ties, prefer more column matches, then header tables. This
    # accepts both header-row tables and no-header tables, and tolerates other
    # tables appearing elsewhere in the document.
    best_table = None
    best_data_rows = None
    best_col_score = -1
    best_has_header = False
    expected_data_rows = gt_data["review_doc"]["required_table_rows"]
    for table in tables:
        if len(table.rows) == 0:
            continue
        has_header = looks_like_header(table.rows[0])
        if has_header:
            header_cells = [cell.text.strip().lower() for cell in table.rows[0].cells]
            data_rows = count_non_empty(table.rows[1:])
            col_score = sum(1 for ec in expected_cols if col_matches(ec, header_cells))
        else:
            data_rows = count_non_empty(table.rows)
            col_score = 0
        key = (data_rows == expected_data_rows, col_score, has_header)
        if best_table is None or key > (best_data_rows == expected_data_rows, best_col_score, best_has_header):
            best_table = table
            best_data_rows = data_rows
            best_col_score = col_score
            best_has_header = has_header

    # Check summary table (CRITICAL: exact data row count and required width).
    total += 1
    if best_table is None:
        print("  FAIL: No table found in document [CRITICAL]")
        critical_failed += 1
    elif len(best_table.columns) < len(expected_cols):
        print(f"  FAIL: Summary table has {len(best_table.columns)} columns (expected >= {len(expected_cols)}) [CRITICAL]")
        critical_failed += 1
    elif best_data_rows == expected_data_rows:
        passed += 1
        print(f"  PASS: Summary table has exactly {expected_data_rows} data rows")
    else:
        print(f"  FAIL: Summary table has {best_data_rows} data rows (expected == {expected_data_rows}) [CRITICAL]")
        critical_failed += 1

    # Check table columns against the located table. For a header table, verify
    # the header text matches the required columns (with synonyms). For a
    # no-header table, verify the table is at least as wide as the required
    # columns.
    total += 1
    if best_table is not None:
        if best_has_header:
            header_cells = [cell.text.strip().lower() for cell in best_table.rows[0].cells]
            cols_found = sum(1 for ec in expected_cols if col_matches(ec, header_cells))
            if cols_found >= len(expected_cols):
                passed += 1
                print(f"  PASS: All {len(expected_cols)} required table columns found")
            else:
                print(f"  FAIL: Only {cols_found}/{len(expected_cols)} required columns found. Headers: {header_cells}")
        else:
            if len(best_table.columns) >= len(expected_cols):
                passed += 1
                print(f"  PASS: No-header table has {len(best_table.columns)} columns (>= {len(expected_cols)})")
            else:
                print(f"  FAIL: No-header table has {len(best_table.columns)} columns (expected >= {len(expected_cols)})")
    else:
        print("  FAIL: No tables to check columns")

    # Check required mentions (keywords) against paragraph + table cell text.
    required_mentions = gt_data["review_doc"]["required_mentions"]
    for mention in required_mentions:
        total += 1
        if text_contains_mention(full_text, mention):
            passed += 1
            print(f"  PASS: Keyword '{mention}' found in document")
        else:
            print(f"  FAIL: Keyword '{mention}' not found in document text")

    return passed, total, critical_failed


def check_pptx(agent_workspace, gt_data):
    """Check the PowerPoint presentation for required structure and content."""
    passed = 0
    total = 0
    filename = gt_data["slides"]["filename"]
    pptx_path = os.path.join(agent_workspace, filename)

    # Check file exists
    total += 1
    if not os.path.exists(pptx_path):
        print(f"  FAIL: {filename} not found at {pptx_path}")
        return passed, total

    passed += 1
    print(f"  PASS: {filename} exists")

    prs = Presentation(pptx_path)

    # Check slide count
    total += 1
    slide_count = len(prs.slides)
    required_count = gt_data["slides"]["required_slide_count"]
    if slide_count >= required_count:
        passed += 1
        print(f"  PASS: Presentation has {slide_count} slides (required >= {required_count})")
    else:
        print(f"  FAIL: Presentation has {slide_count} slides (required >= {required_count})")

    # Extract all text from slides
    all_slide_text = []
    for slide in prs.slides:
        slide_texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    slide_texts.append(para.text)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        slide_texts.append(cell.text)
        all_slide_text.append(" ".join(slide_texts))

    full_slides_text = " ".join(all_slide_text).lower()

    # Check title slide contains required title
    total += 1
    required_title = gt_data["slides"]["required_title"].lower()
    first_slide_text = all_slide_text[0].lower() if all_slide_text else ""
    # Also check across all slides in case title is on a different slide
    if required_title in first_slide_text or required_title in full_slides_text:
        passed += 1
        print(f"  PASS: Required title '{gt_data['slides']['required_title']}' found")
    else:
        # Try partial match
        title_words = required_title.split()
        match_count = sum(1 for w in title_words if w in full_slides_text)
        if match_count >= len(title_words) - 1:
            passed += 1
            print(f"  PASS: Required title approximately matched ({match_count}/{len(title_words)} words)")
        else:
            print(f"  FAIL: Required title '{gt_data['slides']['required_title']}' not found. First slide: {first_slide_text[:100]}")

    # Check required mentions across slides
    required_mentions = gt_data["slides"]["required_mentions"]
    for mention in required_mentions:
        total += 1
        if text_contains_mention(full_slides_text, mention):
            passed += 1
            print(f"  PASS: Keyword '{mention}' found in slides")
        else:
            print(f"  FAIL: Keyword '{mention}' not found in slides")

    return passed, total


def check_word_count(agent_workspace):
    """Check word_count.txt exists and has valid content."""
    passed = 0
    total = 0

    wc_path = os.path.join(agent_workspace, "word_count.txt")

    # Check file exists
    total += 1
    if not os.path.exists(wc_path):
        print(f"  FAIL: word_count.txt not found at {wc_path}")
        return passed, total

    passed += 1
    print("  PASS: word_count.txt exists")

    with open(wc_path, "r") as f:
        content = f.read().strip()

    lines = [line.strip() for line in content.split("\n") if line.strip()]

    # Check at least 2 lines
    total += 1
    if len(lines) >= 2:
        passed += 1
        print(f"  PASS: word_count.txt has {len(lines)} lines (expected >= 2)")
    else:
        print(f"  FAIL: word_count.txt has {len(lines)} lines (expected >= 2)")

    # Check each line has a word count > 0
    for line in lines:
        total += 1
        # Try to extract a number from the line
        numbers = re.findall(r'\d+', line)
        if numbers and any(int(n) > 0 for n in numbers):
            passed += 1
            print(f"  PASS: Line contains valid word count: {line}")
        else:
            print(f"  FAIL: Line does not contain valid word count: {line}")

    # Check both filenames are mentioned
    total += 1
    if "LLM_Reasoning_Review" in content and "LLM_Reasoning_Slides" in content:
        passed += 1
        print("  PASS: Both filenames referenced in word_count.txt")
    else:
        print("  FAIL: Not both filenames found in word_count.txt")

    return passed, total


def main(args):
    gt_path = os.path.join(args.groundtruth_workspace, "expected_results.json")
    if not os.path.exists(gt_path):
        print(f"FAIL: expected_results.json not found at {gt_path}")
        sys.exit(1)

    with open(gt_path, "r") as f:
        gt_data = json.load(f)

    total_passed = 0
    total_checks = 0

    total_critical_failed = 0

    # Check 1: Word document
    print("--- Check 1: Word Document (LLM_Reasoning_Review.docx) ---")
    p, t, cf = check_word_doc(args.agent_workspace, gt_data)
    print(f"  Word Doc: {p}/{t} checks passed")
    total_passed += p
    total_checks += t
    total_critical_failed += cf

    # Check 2: PowerPoint
    print("\n--- Check 2: PowerPoint (LLM_Reasoning_Slides.pptx) ---")
    p, t = check_pptx(args.agent_workspace, gt_data)
    print(f"  PowerPoint: {p}/{t} checks passed")
    total_passed += p
    total_checks += t

    # Check 3: Word count file
    print("\n--- Check 3: Word Count File (word_count.txt) ---")
    p, t = check_word_count(args.agent_workspace)
    print(f"  Word Count: {p}/{t} checks passed")
    total_passed += p
    total_checks += t

    # Overall
    if total_checks == 0:
        print("\nFAIL: No checks were performed.")
        accuracy = 0.0
    else:
        accuracy = total_passed / total_checks * 100
        print(f"\nOverall: {total_passed}/{total_checks} checks passed ({accuracy:.1f}%)")

    result = {
        "total_passed": total_passed,
        "total_checks": total_checks,
        "accuracy": accuracy,
        "timestamp": datetime.now().isoformat(),
    }

    if args.res_log_file:
        with open(args.res_log_file, "w") as f:
            json.dump(result, f, indent=2)
        print(f"Report saved to {args.res_log_file}")

    if total_critical_failed > 0:
        print(f"FAIL (critical checks failed: {total_critical_failed})")
        sys.exit(1)
    if accuracy >= 80:
        print("PASS")
        sys.exit(0)
    else:
        print("FAIL")
        sys.exit(1)


if __name__ == "__main__":
    parser = ArgumentParser()
    parser.add_argument("--agent_workspace", required=False, default=".")
    parser.add_argument("--groundtruth_workspace", required=False, default=".")
    parser.add_argument("--launch_time", required=False)
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()
    main(args)
