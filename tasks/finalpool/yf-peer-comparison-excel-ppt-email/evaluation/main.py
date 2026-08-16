"""Evaluation for yf-peer-comparison-excel-ppt-email.

Checks:
1. Excel file (Peer_Comparison.xlsx) - 3 sheets with correct data
2. PowerPoint (Investor_Presentation.pptx) - >=6 slides
3. Emails - 3 emails to correct recipients
"""

import argparse
import json
import os
import sys

import openpyxl
import psycopg2

DB_CONFIG = {
    "host": os.environ.get("PGHOST", "localhost"),
    "port": int(os.environ.get("PGPORT", "5432")),
    "dbname": "toolathlon_gym",
    "user": "eigent",
    "password": "camel",
}

FILE_PASS = 0
FILE_FAIL = 0
DB_PASS = 0
DB_FAIL = 0


def check(name, condition, detail="", db=False):
    global FILE_PASS, FILE_FAIL, DB_PASS, DB_FAIL
    if condition:
        if db:
            DB_PASS += 1
        else:
            FILE_PASS += 1
        print(f"  [PASS] {name}")
    else:
        if db:
            DB_FAIL += 1
        else:
            FILE_FAIL += 1
        detail_str = f": {detail[:300]}" if detail else ""
        print(f"  [FAIL] {name}{detail_str}")


def num_close(a, b, tol=1.0):
    try:
        return abs(float(a) - float(b)) <= tol
    except (TypeError, ValueError):
        return str(a).strip().lower() == str(b).strip().lower()


def load_gt_sheet(gt_ws, sheet_name):
    """Load groundtruth sheet rows as list of tuples."""
    gt_path = os.path.join(gt_ws, "Peer_Comparison.xlsx")
    if not os.path.exists(gt_path):
        return None
    wb = openpyxl.load_workbook(gt_path, data_only=True)
    for name in wb.sheetnames:
        if name.strip().lower() == sheet_name.strip().lower():
            return [[cell.value for cell in row] for row in wb[name].iter_rows()]
    return None


def check_excel(agent_workspace, groundtruth_workspace):
    print("\n=== Checking Excel Output ===")
    excel_path = os.path.join(agent_workspace, "Peer_Comparison.xlsx")
    check("Excel file exists", os.path.isfile(excel_path), f"Expected {excel_path}")
    if not os.path.isfile(excel_path):
        return

    try:
        wb = openpyxl.load_workbook(excel_path, data_only=True)
    except Exception as e:
        check("Excel file readable", False, str(e))
        return

    # --- Sheet 1: Company Profiles ---
    ws1 = None
    for s in wb.sheetnames:
        if "company" in s.lower() and "profile" in s.lower():
            ws1 = wb[s]
            break
    if ws1 is None:
        for s in wb.sheetnames:
            if "company" in s.lower() or "profile" in s.lower():
                ws1 = wb[s]
                break

    check("Sheet 'Company Profiles' exists", ws1 is not None, f"Sheets: {wb.sheetnames}")
    if ws1 is not None:
        rows = list(ws1.iter_rows(min_row=2, values_only=True))
        data_rows = [r for r in rows if r and r[0] is not None]
        check("Company Profiles has 5 rows", len(data_rows) == 5, f"Got {len(data_rows)}")

        # Check symbols present
        symbols_found = {str(r[0]).strip().upper() for r in data_rows if r[0]}
        for sym in ["AMZN", "GOOGL", "JNJ", "JPM", "XOM"]:
            check(f"{sym} in Company Profiles", sym in symbols_found,
                  f"Found: {symbols_found}")

        # Check alphabetical order
        sym_list = [str(r[0]).strip().upper() for r in data_rows if r[0]]
        check("Symbols sorted alphabetically", sym_list == sorted(sym_list),
              f"Order: {sym_list}")

        # Compare with groundtruth values
        gt_rows = load_gt_sheet(groundtruth_workspace, "Company Profiles")
        if gt_rows:
            gt_data = {}
            for r in gt_rows[1:]:
                if r and r[0]:
                    gt_data[str(r[0]).strip().upper()] = r

            agent_data = {}
            for r in data_rows:
                if r and r[0]:
                    agent_data[str(r[0]).strip().upper()] = r

            for sym in ["AMZN", "GOOGL", "JNJ", "JPM", "XOM"]:
                if sym in agent_data and sym in gt_data:
                    ar = agent_data[sym]
                    gr = gt_data[sym]
                    # Market Cap (col 3, tol 1M)
                    if len(ar) > 3 and len(gr) > 3 and ar[3] is not None:
                        check(f"{sym} Market_Cap", num_close(ar[3], gr[3], 1e6),
                              f"Agent={ar[3]}, GT={gr[3]}")
                    # Trailing PE (col 4, tol 0.5)
                    if len(ar) > 4 and len(gr) > 4 and ar[4] is not None:
                        check(f"{sym} Trailing_PE", num_close(ar[4], gr[4], 0.5),
                              f"Agent={ar[4]}, GT={gr[4]}")
                    # YTD Return (col 10, tol 0.2)
                    if len(ar) > 10 and len(gr) > 10 and ar[10] is not None:
                        check(f"{sym} YTD_Return_Pct", num_close(ar[10], gr[10], 0.2),
                              f"Agent={ar[10]}, GT={gr[10]}")
                    # Beta (col 6) and Latest_Close (col 9)
                    if len(ar) > 6 and len(gr) > 6 and ar[6] is not None:
                        check(f"{sym} Beta", num_close(ar[6], gr[6], 0.05),
                              f"Agent={ar[6]}, GT={gr[6]}")
                    if len(ar) > 9 and len(gr) > 9 and ar[9] is not None:
                        check(f"{sym} Latest_Close", num_close(ar[9], gr[9], 0.5),
                              f"Agent={ar[9]}, GT={gr[9]}")
                    # Sector exact match
                    if len(ar) > 2 and len(gr) > 2 and ar[2] is not None:
                        check(f"{sym} Sector", str(ar[2]).strip().lower() == str(gr[2]).strip().lower(),
                              f"Agent={ar[2]}, GT={gr[2]}")

    # --- Sheet 2: Financial Comparison ---
    ws2 = None
    for s in wb.sheetnames:
        if "financial" in s.lower() and "comparison" in s.lower():
            ws2 = wb[s]
            break
    if ws2 is None:
        for s in wb.sheetnames:
            if "financial" in s.lower():
                ws2 = wb[s]
                break

    check("Sheet 'Financial Comparison' exists", ws2 is not None, f"Sheets: {wb.sheetnames}")
    if ws2 is not None:
        rows2 = list(ws2.iter_rows(min_row=2, values_only=True))
        data_rows2 = [r for r in rows2 if r and r[0] is not None]
        check("Financial Comparison has 5 rows", len(data_rows2) == 5, f"Got {len(data_rows2)}")

        # Sorted by Symbol
        sym_list2 = [str(r[0]).strip().upper() for r in data_rows2 if r and r[0]]
        check("Financial Comparison sorted alphabetically",
              sym_list2 == sorted(sym_list2), f"Order: {sym_list2}")

        # Verify ALL 4 financial metrics for each of 5 symbols
        gt_rows2 = load_gt_sheet(groundtruth_workspace, "Financial Comparison")
        if gt_rows2:
            gt_fin = {str(r[0]).strip().upper(): r for r in gt_rows2[1:] if r and r[0]}
            agent_fin = {str(r[0]).strip().upper(): r for r in data_rows2 if r and r[0]}
            cols = [(1, "Revenue"), (2, "Net_Income"), (3, "Total_Assets"), (4, "Free_Cash_Flow")]
            for sym in ["AMZN", "GOOGL", "JNJ", "JPM", "XOM"]:
                if sym in agent_fin and sym in gt_fin:
                    ar = agent_fin[sym]
                    gr = gt_fin[sym]
                    for col, name in cols:
                        if len(ar) > col and len(gr) > col and gr[col] is not None:
                            try:
                                exp = float(gr[col])
                                act = float(ar[col]) if ar[col] is not None else None
                            except (TypeError, ValueError):
                                check(f"{sym} {name} numeric", False,
                                      f"Agent={ar[col]}, GT={gr[col]}")
                                continue
                            if act is None:
                                check(f"{sym} {name} present", False, "missing")
                                continue
                            # Use 5% relative tolerance for large numbers
                            tol = max(abs(exp) * 0.05, 1e6)
                            ok = abs(act - exp) <= tol
                            check(f"{sym} {name} within 5%", ok,
                                  f"Agent={act}, GT={exp}")
                else:
                    check(f"{sym} present in Financial Comparison",
                          sym in agent_fin, f"Missing")

    # --- Sheet 3: Scoring ---
    ws3 = None
    for s in wb.sheetnames:
        if "scor" in s.lower():
            ws3 = wb[s]
            break

    check("Sheet 'Scoring' exists", ws3 is not None, f"Sheets: {wb.sheetnames}")
    if ws3 is not None:
        rows3 = list(ws3.iter_rows(min_row=2, values_only=True))
        data_rows3 = [r for r in rows3 if r and r[0] is not None]
        check("Scoring has 5 rows", len(data_rows3) == 5, f"Got {len(data_rows3)}")

        # Sorted by Symbol
        sym_list3 = [str(r[0]).strip().upper() for r in data_rows3 if r and r[0]]
        check("Scoring sorted alphabetically",
              sym_list3 == sorted(sym_list3), f"Order: {sym_list3}")

        gt_rows3 = load_gt_sheet(groundtruth_workspace, "Scoring")
        if gt_rows3:
            gt_score = {str(r[0]).strip().upper(): r for r in gt_rows3[1:] if r and r[0]}
            agent_score = {str(r[0]).strip().upper(): r for r in data_rows3 if r and r[0]}

            rank_cols = [(1, "Valuation_Rank"), (2, "Growth_Rank"), (3, "Income_Rank"),
                         (4, "Risk_Rank"), (5, "Momentum_Rank")]
            for sym in ["AMZN", "GOOGL", "JNJ", "JPM", "XOM"]:
                if sym in agent_score and sym in gt_score:
                    ar = agent_score[sym]
                    gr = gt_score[sym]
                    # Per-dimension ranks: integer, exact match
                    for col, name in rank_cols:
                        if len(gr) > col and gr[col] is not None:
                            exp_rank = gr[col]
                            act_rank = ar[col] if len(ar) > col else None
                            check(f"{sym} {name}",
                                  num_close(act_rank, exp_rank, 0),
                                  f"Agent={act_rank}, GT={exp_rank}")
                    # Weighted Score (col 6, tighter tol 0.1)
                    if len(ar) > 6 and len(gr) > 6 and ar[6] is not None:
                        check(f"{sym} Weighted_Score", num_close(ar[6], gr[6], 0.1),
                              f"Agent={ar[6]}, GT={gr[6]}")
                    # Overall Rating (col 7, exact match)
                    if len(ar) > 7 and len(gr) > 7 and ar[7] is not None:
                        agent_rating = str(ar[7]).strip().lower()
                        gt_rating = str(gr[7]).strip().lower()
                        check(f"{sym} Overall_Rating", agent_rating == gt_rating,
                              f"Agent='{ar[7]}', GT='{gr[7]}'")
                else:
                    check(f"{sym} present in Scoring",
                          sym in agent_score, "Missing")


def check_pptx(agent_workspace):
    print("\n=== Checking PowerPoint Output ===")
    pptx_path = os.path.join(agent_workspace, "Investor_Presentation.pptx")
    check("PPTX file exists", os.path.isfile(pptx_path), f"Expected {pptx_path}")
    if not os.path.isfile(pptx_path):
        return

    try:
        from pptx import Presentation
        prs = Presentation(pptx_path)
    except Exception as e:
        check("PPTX file readable", False, str(e))
        return

    slide_count = len(prs.slides)
    check("PPTX has >= 6 slides", slide_count >= 6, f"Got {slide_count} slides")

    # Build per-slide text and overall text
    slide_texts = []
    for slide in prs.slides:
        parts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                parts.append(shape.text_frame.text.lower())
        slide_texts.append(" ".join(parts))
    all_text = " ".join(slide_texts)

    # All 5 symbols (or full company names)
    sym_aliases = {
        "AMZN": ["amzn", "amazon"],
        "GOOGL": ["googl", "alphabet", "google"],
        "JNJ": ["jnj", "johnson"],
        "JPM": ["jpm", "morgan", "jp morgan"],
        "XOM": ["xom", "exxon"],
    }
    for sym, aliases in sym_aliases.items():
        ok = any(a in all_text for a in aliases)
        check(f"PPTX mentions {sym} (or alias)", ok,
              f"None of {aliases} found")

    # Each topic must appear on at least one DEDICATED slide (not aggregated across slides);
    # we require a single slide to contain at least one alt phrase, enforcing per-slide
    # structure so a one-slide-with-everything deck cannot satisfy multiple topics.
    topic_phrases = {
        "market overview": ["market overview", "companies and sectors"],
        "company profiles": ["company profile", "profiles", "key metrics"],
        "financial comparison": ["financial comparison", "annual financials", "revenue", "net income"],
        "scoring": ["scoring", "weighted score", "peer scoring"],
        "recommendations": ["recommendation", "recommended actions", "buy", "hold", "sell"],
    }
    used_slides = set()
    for topic, alts in topic_phrases.items():
        # Find first slide that matches and has not been claimed by another topic.
        matched_idx = None
        for idx, st in enumerate(slide_texts):
            if idx in used_slides:
                continue
            if any(a in st for a in alts):
                matched_idx = idx
                break
        ok = matched_idx is not None
        if ok:
            used_slides.add(matched_idx)
        check(
            f"PPTX has dedicated slide for '{topic}'",
            ok,
            f"alts={alts}; used_slides={sorted(used_slides)}",
        )


def check_emails():
    print("\n=== Checking Emails ===")
    try:
        conn = psycopg2.connect(**DB_CONFIG)
        cur = conn.cursor()
    except Exception as e:
        check("Email DB connection", False, str(e), db=True)
        return

    def _check_email(addr, expected_subject, body_phrases, label):
        cur.execute("""
            SELECT subject, to_addr, body_text FROM email.messages
            WHERE to_addr::text ILIKE %s
            ORDER BY id DESC LIMIT 5
        """, (f"%{addr}%",))
        rows = cur.fetchall()
        check(f"Email to {addr}", len(rows) > 0, "No email found", db=True)
        if not rows:
            check(f"{label} subject is exactly '{expected_subject}'", False,
                  "no email", db=True)
            check(f"{label} body has expected content", False,
                  "no email", db=True)
            return
        # Find the one with matching exact subject if possible
        chosen = None
        for r in rows:
            if r[0] and r[0].strip().lower() == expected_subject.lower():
                chosen = r
                break
        if chosen is None:
            chosen = rows[0]
        subj_ok = (chosen[0] or "").strip().lower() == expected_subject.lower()
        check(f"{label} subject is exactly '{expected_subject}'", subj_ok,
              f"Got: {chosen[0]!r}", db=True)
        body = (chosen[2] or "").lower()
        for phrase in body_phrases:
            ok = phrase.lower() in body
            check(f"{label} body mentions '{phrase}'", ok,
                  f"body excerpt: {body[:200]}", db=True)

    _check_email(
        "portfolio_managers@firm.com",
        "Peer Comparison Summary",
        body_phrases=["buy", "hold"],
        label="Portfolio Managers"
    )
    _check_email(
        "research_team@firm.com",
        "Peer Comparison Detailed Findings",
        body_phrases=["score", "revenue"],
        label="Research Team"
    )
    # Compliance/risk email: must mention beta + risk language and at least one
    # high-beta ticker that's specifically flagged in GT (AMZN beta=1.42 > 1).
    _check_email(
        "compliance@firm.com",
        "Peer Comparison Risk Review",
        body_phrases=["beta", "risk", "amzn"],
        label="Compliance"
    )

    cur.close()
    conn.close()


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--agent_workspace", required=False, default=".")
    parser.add_argument("--groundtruth_workspace", required=False, default=".")
    parser.add_argument("--launch_time", required=False)
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()

    gt_ws = args.groundtruth_workspace or os.path.join(
        os.path.dirname(__file__), "..", "groundtruth_workspace")

    check_excel(args.agent_workspace, gt_ws)
    check_pptx(args.agent_workspace)
    check_emails()

    total_pass = FILE_PASS + DB_PASS
    total_fail = FILE_FAIL + DB_FAIL
    overall_ok = FILE_FAIL == 0 and DB_FAIL == 0

    print(f"\n=== SUMMARY ===")
    print(f"  File checks - Passed: {FILE_PASS}, Failed: {FILE_FAIL}")
    print(f"  DB checks   - Passed: {DB_PASS}, Failed: {DB_FAIL}")
    print(f"  Overall: {'PASS' if overall_ok else 'FAIL'}")

    if args.res_log_file:
        result = {"passed": total_pass, "failed": total_fail, "success": overall_ok}
        with open(args.res_log_file, "w") as f:
            json.dump(result, f, indent=2)

    sys.exit(0 if overall_ok else 1)


if __name__ == "__main__":
    main()
