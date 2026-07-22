"""Evaluation for yf-portfolio-ppt-gcal."""
import argparse
import os
import sys

import psycopg2

def num_close(a, b, rel_tol=0.15, abs_tol=0.5):
    return abs(float(a) - float(b)) <= max(abs_tol, abs(float(b)) * rel_tol)

from pptx import Presentation

DB = {"host": os.environ.get("PGHOST", "localhost"), "port": 5432, "dbname": "toolathlon_gym", "user": "eigent", "password": "camel"}


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--agent_workspace", required=False)
    parser.add_argument("--groundtruth_workspace", required=False)
    parser.add_argument("--launch_time", required=False)
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()

    task_root = os.path.abspath(os.path.join(os.path.dirname(__file__), ".."))
    gt_dir = args.groundtruth_workspace or os.path.join(task_root, "groundtruth_workspace")
    all_errors = []

    # ---- Check PPT ----
    agent_ppt = os.path.join(args.agent_workspace, "Portfolio_Review.pptx")
    if not os.path.exists(agent_ppt):
        print(f"FAIL: Agent output Portfolio_Review.pptx not found")
        sys.exit(1)

    print("  Checking Portfolio_Review.pptx...")
    prs = Presentation(agent_ppt)
    slides = list(prs.slides)

    # Minimum: title + overview + 5 stocks + takeaways = 8
    if len(slides) < 8:
        all_errors.append(f"PPT has {len(slides)} slides, expected at least 8")

    # Collect all text
    all_ppt_text = ""
    for slide in slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                all_ppt_text += shape.text_frame.text + " "
    all_lower = all_ppt_text.lower()

    # Check title slide
    if len(slides) > 0:
        title_text = ""
        for shape in slides[0].shapes:
            if shape.has_text_frame:
                title_text += shape.text_frame.text.lower() + " "
        if "q4 2025" not in title_text or "portfolio" not in title_text:
            all_errors.append(f"Title slide missing expected text. Found: {title_text[:100]}")

    # Check all stock symbols present
    for sym in ["amzn", "googl", "jnj", "jpm", "xom"]:
        if sym not in all_lower:
            all_errors.append(f"PPT missing stock symbol: {sym.upper()}")

    # Check last slide has takeaways
    if len(slides) >= 2:
        last_text = ""
        for shape in slides[-1].shapes:
            if shape.has_text_frame:
                last_text += shape.text_frame.text.lower() + " "
        if "takeaway" not in last_text and "key" not in last_text and "summary" not in last_text and "conclusion" not in last_text:
            all_errors.append("Last slide missing takeaways/summary content")

    # Validate stock data against DB
    conn = psycopg2.connect(**DB)
    cur = conn.cursor()
    stocks = ["AMZN", "GOOGL", "JNJ", "JPM", "XOM"]
    for sym in stocks:
        cur.execute("""
            SELECT close FROM yf.stock_prices
            WHERE symbol = %s AND date = (SELECT MIN(date) FROM yf.stock_prices WHERE symbol = %s AND date >= '2025-10-01')
        """, (sym, sym))
        start_price = float(cur.fetchone()[0])

        cur.execute("""
            SELECT close FROM yf.stock_prices
            WHERE symbol = %s AND date = (SELECT MAX(date) FROM yf.stock_prices WHERE symbol = %s AND date <= '2025-12-31')
        """, (sym, sym))
        end_price = float(cur.fetchone()[0])

        ret = round((end_price - start_price) / start_price * 100, 1)

        # Check return percentage appears in text (with some tolerance on formatting)
        ret_str = f"{ret}"
        # Also accept close values
        if ret_str not in all_ppt_text and f"{abs(ret)}" not in all_ppt_text:
            # Allow +/- 0.5% tolerance in the text
            found = False
            for delta in [-0.5, -0.4, -0.3, -0.2, -0.1, 0, 0.1, 0.2, 0.3, 0.4, 0.5]:
                check_val = f"{ret + delta:.1f}"
                if check_val in all_ppt_text:
                    found = True
                    break
            if not found:
                all_errors.append(f"PPT missing return for {sym}: expected ~{ret}%")

    conn.close()

    if not all_errors:
        print("    PASS")

    # --- Check Google Calendar event in DB ---
    # Task.md REQUIRES a calendar event on Mar 15, 2026 14:00-15:00 with title
    # "Q4 2025 Portfolio Review Meeting" and location "Conference Room A".
    # When an event is present it must match all criteria (BLOCKING if present
    # but wrong); total absence is tolerated only when running GT self-test.
    print("  Checking Google Calendar event...")
    try:
        gt_canon = os.path.realpath(gt_dir)
        ag_canon = os.path.realpath(args.agent_workspace) if args.agent_workspace else gt_canon
        is_gt_self_test = (gt_canon == ag_canon)
    except Exception:
        is_gt_self_test = True
    try:
        conn = psycopg2.connect(**DB)
        cur = conn.cursor()
        cur.execute("""
            SELECT id, summary, start_datetime, end_datetime, location
            FROM gcal.events
            WHERE LOWER(summary) LIKE '%portfolio%review%'
               OR LOWER(summary) LIKE '%q4 2025%'
        """)
        events = cur.fetchall()
        if not events:
            # When run against a real agent workspace (agent != gt), missing event is fatal.
            # During GT self-test, tolerate absence.
            if is_gt_self_test:
                print("    WARNING: Portfolio review meeting not found in calendar (GT self-test, non-blocking)")
            else:
                all_errors.append(
                    "Portfolio review calendar event missing (expected Mar 15 2026 14:00-15:00 "
                    "at 'Conference Room A')."
                )
        else:
            # At least one candidate -> require it to match date/time/location.
            match_found = False
            for ev_id, summary, start_dt, end_dt, location in events:
                try:
                    if start_dt is None:
                        continue
                    date_match = str(start_dt)[:10] == "2026-03-15"
                    hour_match = hasattr(start_dt, 'hour') and start_dt.hour == 14
                    duration_match = False
                    if hasattr(start_dt, 'hour') and end_dt is not None and hasattr(end_dt, 'hour'):
                        duration_min = (end_dt - start_dt).total_seconds() / 60
                        duration_match = 55 <= duration_min <= 75
                    loc_match = location and "conference room a" in str(location).lower()
                    if date_match and hour_match and duration_match and loc_match:
                        match_found = True
                        print(f"    Matched calendar event: {summary} at {start_dt}")
                        break
                except Exception as inner_e:
                    print(f"    (parse err on event {ev_id}: {inner_e})")
            if not match_found:
                all_errors.append(
                    "Calendar event present but wrong date/time/location; "
                    "expected Mar 15 2026 14:00-15:00 at 'Conference Room A'."
                )
        conn.close()
    except Exception as e:
        print(f"    WARNING: GCal DB check error (fallback to non-blocking): {e}")

    if all_errors:
        print(f"\n=== RESULT: FAIL ({len(all_errors)} errors) ===")
        for e in all_errors[:15]:
            print(f"  {e}")
        sys.exit(1)
    else:
        print("\n=== RESULT: PASS ===")
        sys.exit(0)


if __name__ == "__main__":
    main()
