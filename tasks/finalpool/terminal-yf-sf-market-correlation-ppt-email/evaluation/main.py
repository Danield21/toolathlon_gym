"""Evaluation for terminal-yf-sf-market-correlation-ppt-email.
Checks:
1. Market_Correlation_Report.pptx with 6+ slides, all required tickers, segments, key revenue numbers
2. Email sent to cfo@company.com with exact subject + key findings in body
3. correlation_analysis.py script exists
"""
import argparse
import json
import os
import re
import sys
import psycopg2

DB_CONFIG = {
    "host": os.environ.get("PGHOST", "localhost"), "port": int(os.environ.get("PGPORT", "5432")),
    "dbname": os.environ.get("PGDATABASE", "toolathlon_gym"),
    "user": "eigent", "password": "camel",
}

PASS_COUNT = 0
FAIL_COUNT = 0

REQUIRED_TICKERS = ["AMZN", "GOOGL", "JNJ", "JPM", "XOM", "GC=F", "^DJI"]
REQUIRED_SEGMENTS = ["Consumer", "Enterprise", "Government", "SMB"]
EXPECTED_TOP_SEGMENT = "Consumer"
EXPECTED_SUBJECT = "Market and Sales Correlation Analysis Report"


def check(name, condition, detail=""):
    global PASS_COUNT, FAIL_COUNT
    if condition:
        PASS_COUNT += 1
        print(f"  [PASS] {name}")
    else:
        FAIL_COUNT += 1
        print(f"  [FAIL] {name}: {str(detail)[:200]}")


def _slide_text(slide):
    text = []
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text:
            text.append(shape.text)
        if getattr(shape, "has_table", False):
            try:
                for row in shape.table.rows:
                    for cell in row.cells:
                        text.append(cell.text)
            except Exception:
                pass
    return "\n".join(text)


def _slide_title(slide):
    """Extract title shape text only (not body), for stricter title validation."""
    try:
        if slide.shapes.title is not None and slide.shapes.title.text:
            return slide.shapes.title.text
    except Exception:
        pass
    return ""


def check_pptx(workspace):
    print("\n=== Check 1: Market_Correlation_Report.pptx ===")
    path = os.path.join(workspace, "Market_Correlation_Report.pptx")
    if not os.path.exists(path):
        check("PPTX file exists", False, f"Not found at {path}")
        # add subordinate failures so missing file isn't masked
        for label in ["Has at least 6 slides", "Slide 1 has correct title",
                      "All 7 required tickers appear", "All 4 customer segments mentioned",
                      "Has recommendations content", "Mentions revenue"]:
            check(label, False, "PPTX missing")
        return
    check("PPTX file exists", True)

    try:
        from pptx import Presentation
        prs = Presentation(path)
    except Exception as e:
        check("PPTX readable", False, str(e))
        return
    check("PPTX readable", True)

    slides = list(prs.slides)
    slide_count = len(slides)
    check("Has at least 6 slides", slide_count >= 6, f"Found {slide_count} slides")

    # Slide 1 title — restrict to the title shape (with fallback to slide text for robustness)
    slide1_title = _slide_title(slides[0]).lower() if slide_count >= 1 else ""
    slide1_text = _slide_text(slides[0]).lower() if slide_count >= 1 else ""
    title_match = (
        "market" in slide1_title and "sales" in slide1_title and "correlation" in slide1_title
    ) or (
        # Fallback: title shape might not be set; accept body if all keywords on slide 1
        slide1_title == "" and "market" in slide1_text and "sales" in slide1_text and "correlation" in slide1_text
    )
    check("Slide 1 has correct title (title shape preferred)",
          title_match,
          f"Slide 1 title: {slide1_title[:120]}; slide text: {slide1_text[:200]}")

    # All tickers must appear somewhere in pptx (case-insensitive)
    full_text = "\n".join(_slide_text(s) for s in slides)
    full_lower = full_text.lower()

    missing_tickers = []
    for t in REQUIRED_TICKERS:
        # special chars in tickers like ^DJI, GC=F: do simple substring match (case-insensitive)
        if t.lower() not in full_lower:
            missing_tickers.append(t)
    check(f"All {len(REQUIRED_TICKERS)} required tickers appear",
          len(missing_tickers) == 0,
          f"Missing: {missing_tickers}")

    # All 4 customer segments mentioned
    missing_segments = [s for s in REQUIRED_SEGMENTS if s.lower() not in full_lower]
    check("All 4 customer segments mentioned",
          len(missing_segments) == 0,
          f"Missing: {missing_segments}")

    # Has recommendations content — require the actual word 'recommendation' (or 'recommendations')
    check("Has 'recommendation' content (not just 'recommend')",
          ("recommendation" in full_lower) or ("recommendations" in full_lower),
          "")

    # Per-slide content validation: at least one slide mentions all 4 segments and at least
    # one slide mentions multiple tickers (as proxy for proper distribution of content)
    seg_slide_match = any(
        all(s.lower() in _slide_text(slide).lower() for s in REQUIRED_SEGMENTS)
        for slide in slides
    )
    check("At least one slide contains all 4 segments together",
          seg_slide_match,
          "Segments should be grouped on a 'Sales by Segment' slide")
    multi_ticker_slide = any(
        sum(1 for t in REQUIRED_TICKERS if t.lower() in _slide_text(slide).lower()) >= 3
        for slide in slides
    )
    check("At least one slide groups 3+ tickers",
          multi_ticker_slide,
          "Tickers should be grouped on a 'Market Indicators' slide")

    # Revenue keyword present
    check("Mentions revenue", "revenue" in full_lower, "")


def check_email():
    print("\n=== Check 2: Email to cfo@company.com ===")
    try:
        conn = psycopg2.connect(**DB_CONFIG)
        cur = conn.cursor()
        cur.execute("SELECT subject, from_addr, to_addr, body_text FROM email.messages")
        messages = cur.fetchall()
        cur.close()
        conn.close()
    except Exception as e:
        check("Email DB query OK", False, str(e))
        return

    matching = None
    for subject, from_addr, to_addr, body_text in messages:
        to_str = ""
        if isinstance(to_addr, list):
            to_str = " ".join(str(r).lower() for r in to_addr)
        elif isinstance(to_addr, str):
            try:
                parsed = json.loads(to_addr)
                to_str = " ".join(str(r).lower() for r in parsed) if isinstance(parsed, list) else str(to_addr).lower()
            except Exception:
                to_str = str(to_addr).lower()
        if "cfo@company.com" in to_str:
            matching = (subject, from_addr, to_addr, body_text)
            break

    check("Email sent to cfo@company.com", matching is not None,
          f"Total messages: {len(messages)}")

    if not matching:
        # subordinate failures
        for label in ["Email subject matches", "Email body mentions revenue",
                      "Email body mentions top segment", "Email body mentions presentation"]:
            check(label, False, "no email")
        return

    subject, _, _, body_text = matching
    subj_norm = re.sub(r'\s+', ' ', (subject or "").strip().lower())
    expected_norm = EXPECTED_SUBJECT.lower()
    check("Email subject matches",
          subj_norm == expected_norm,
          f"Subject: {subject!r}, expected: {EXPECTED_SUBJECT!r}")

    body_lower = (body_text or "").lower()
    check("Email body mentions revenue",
          "revenue" in body_lower,
          f"Body: {body_lower[:200]}")

    check("Email body mentions top segment",
          EXPECTED_TOP_SEGMENT.lower() in body_lower,
          f"Expected segment '{EXPECTED_TOP_SEGMENT}' in body")

    check("Email body mentions presentation",
          ("presentation" in body_lower or "powerpoint" in body_lower or
           ".pptx" in body_lower or "report" in body_lower),
          f"Body: {body_lower[:200]}")


def check_script(workspace):
    print("\n=== Check 3: correlation_analysis.py ===")
    path = os.path.join(workspace, "correlation_analysis.py")
    check("correlation_analysis.py exists", os.path.exists(path))


def check_reverse_validation(workspace):
    """Verify things that should NOT exist in the output."""
    print("\n=== Reverse Validation ===")
    try:
        conn = psycopg2.connect(**DB_CONFIG)
        cur = conn.cursor()
        cur.execute("""
            SELECT COUNT(*) FROM email.messages
            WHERE (lower(subject) LIKE '%%correlation%%' OR lower(subject) LIKE '%%market and sales%%')
              AND to_addr::text NOT ILIKE '%%cfo%%'
        """)
        wrong_emails = cur.fetchone()[0]
        check("No correlation emails to wrong recipients", wrong_emails == 0,
              f"Found {wrong_emails} misrouted emails")
        cur.close()
        conn.close()
    except Exception as e:
        check("Reverse-validation DB OK", False, str(e))

    path = os.path.join(workspace, "Market_Correlation_Report.pptx")
    if os.path.exists(path):
        try:
            from pptx import Presentation
            prs = Presentation(path)
            empty_slides = 0
            for slide in prs.slides:
                slide_text = _slide_text(slide)
                if not slide_text.strip():
                    empty_slides += 1
            check("No empty slides in PPTX", empty_slides == 0,
                  f"Found {empty_slides} empty slides")
        except Exception:
            pass


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--agent_workspace", required=False, default=".")
    parser.add_argument("--groundtruth_workspace", required=False, default=".")
    parser.add_argument("--launch_time", required=False)
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()

    check_pptx(args.agent_workspace)
    check_email()
    check_script(args.agent_workspace)
    check_reverse_validation(args.agent_workspace)

    total = PASS_COUNT + FAIL_COUNT
    print(f"\nOverall: {PASS_COUNT}/{total} checks passed")

    result = {"total_passed": PASS_COUNT, "total_checks": total, "failed": FAIL_COUNT}
    if args.res_log_file:
        with open(args.res_log_file, "w") as f:
            json.dump(result, f, indent=2)

    if FAIL_COUNT == 0:
        print("PASS")
        sys.exit(0)
    else:
        print("FAIL")
        sys.exit(1)


if __name__ == "__main__":
    main()
