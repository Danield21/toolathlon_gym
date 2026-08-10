import argparse
import json
import os
import re
import sys
import openpyxl
import psycopg2
from pptx import Presentation

DB_CONFIG = {
    "host": os.environ.get("PGHOST", "localhost"),
    "port": int(os.environ.get("PGPORT", "5432")),
    "dbname": os.environ.get("PGDATABASE", "toolathlon_gym"),
    "user": os.environ.get("PGUSER", "eigent"),
    "password": os.environ.get("PGPASSWORD", "camel"),
}

PASS_COUNT = 0
FAIL_COUNT = 0


def check(name, condition, detail=""):
    global PASS_COUNT, FAIL_COUNT
    if condition:
        PASS_COUNT += 1
        print(f"  [PASS] {name}")
    else:
        FAIL_COUNT += 1
        print(f"  [FAIL] {name}: {str(detail)[:200]}")


def _to_float(v):
    """Parse a value into a float, tolerating currency symbols, thousands
    separators, percent signs and surrounding whitespace.

    Returns None when the value cannot be parsed as a number (e.g. a
    formula string, plain text, or None)."""
    if v is None:
        return None
    if isinstance(v, bool):
        return None
    if isinstance(v, (int, float)):
        return float(v)
    s = str(v).strip()
    if s == "":
        return None
    if s.startswith("="):
        # Excel formula with no usable cached value here.
        return None
    # Normalize Unicode minus/dash characters (LLM output sometimes uses U+2212).
    s = s.replace("−", "-").replace("–", "-").replace("—", "-")
    s = s.replace(",", "").replace("$", "").replace("¥", "").replace("€", "")
    s = s.replace("%", "").replace(" ", "").replace("+", "")
    try:
        return float(s)
    except ValueError:
        return None


def _is_formula(v):
    return isinstance(v, str) and v.strip().startswith("=")


def num_close(a, b, tol=2.0):
    fa = _to_float(a)
    fb = _to_float(b)
    if fa is not None and fb is not None:
        return abs(fa - fb) <= tol
    if fa is None and fb is None:
        if a is None or b is None:
            return a is None and b is None
        return str(a).strip().lower() == str(b).strip().lower()
    return False


def _metric_key(s):
    """Normalize a Correlation_Analysis metric label for fuzzy matching:
    lowercase and drop every non-alphanumeric character (so 'P-Value',
    'P_Value', 'P Value' and 'pvalue' all map to 'pvalue')."""
    return re.sub(r"[^a-z0-9]", "", str(s).lower())


# Exact-label aliases for common statistical spellings that do not contain the
# canonical key as a substring (e.g. 'r' for the Pearson coefficient, 'P' for
# the p-value). Keys are canonical names, values are normalized alias strings.
_METRIC_ALIASES = {
    "Pearson_Correlation": {"pearson", "coefficient", "correlation", "r"},
    "P_Value": {"p", "significance", "prob"},
}


def _find_metric(metric_rows, canonical):
    """Search metric_rows (list of (raw_label, value)) for a row whose label
    matches the canonical metric name. Accepts common statistical aliases such
    as 'Pearson Correlation Coefficient' (canonical is a substring), 'P-Value'
    (normalizes to the canonical form), or short labels like 'Correlation',
    'r' and 'P'."""
    canon = _metric_key(canonical)
    if canon == "":
        return None
    alias = _METRIC_ALIASES.get(canonical, set())
    for k, v in metric_rows:
        kn = _metric_key(k)
        if kn == "":
            continue
        # canonical is a substring of the label (e.g. 'Pearson Correlation
        # Coefficient'), exact normalized equality ('P-Value'), a substantial
        # substring of the canonical (e.g. 'Correlation', 'Number'), or an
        # exact alias ('r', 'P').
        if canon in kn or kn == canon or (len(kn) >= 6 and kn in canon) or kn in alias:
            return v
    return None


def _gt_variants(gt):
    """Return the primary GT plus any accepted alternative interpretations.
    Alternatives let a model that follows a defensible but different reading of
    the instructions (e.g. counting only completed/processing orders) pass
    without being false-FAILed."""
    return [gt] + gt.get("_alts", [])


def load_groundtruth():
    """Load expected values - prefer static groundtruth_data.json, fallback to DB."""
    gt_path = os.path.join(
        os.path.dirname(os.path.dirname(os.path.abspath(__file__))),
        "groundtruth_workspace",
        "groundtruth_data.json",
    )
    gt = None
    if os.path.exists(gt_path):
        try:
            with open(gt_path) as f:
                gt = json.load(f)
        except Exception:
            pass
    if gt is None:
        gt = _compute_groundtruth_from_db()

    # Compute alternative GT interpretations so a model that follows a
    # defensible WooCommerce order-status reading (e.g. counting only
    # completed/processing orders as real purchases) is not false-FAILed.
    # The DB is guaranteed available during evaluation (email checks need it);
    # if it is not reachable we simply degrade to the primary GT.
    alts = []
    try:
        cp = _compute_groundtruth_from_db(
            status_filter=["completed", "processing"]
        )
        if cp.get("num_matched", 0) != gt.get("num_matched", 0):
            alts.append(cp)
    except Exception:
        alts = []
    gt["_alts"] = alts
    return gt


def _compute_groundtruth_from_db(status_filter=None):
    """Query canvas and wc schemas to compute expected values dynamically.

    When ``status_filter`` is a list of WooCommerce order statuses, only
    orders whose status is in that list contribute to the purchase totals
    (used to build an accepted-alternative GT interpretation)."""
    import csv as _csv

    conn = psycopg2.connect(**DB_CONFIG)
    cur = conn.cursor()
    try:
        # Read student_bookstore_registry.csv to get mappings
        registry_path = os.path.join(
            os.path.dirname(os.path.dirname(os.path.abspath(__file__))),
            "initial_workspace",
            "student_bookstore_registry.csv",
        )
        mappings = []
        with open(registry_path) as f:
            reader = _csv.DictReader(f)
            for row in reader:
                mappings.append({
                    "student_id": int(row["student_id"]),
                    "bookstore_customer_id": int(row["bookstore_customer_id"]),
                    "department": row.get("department", ""),
                })

        student_ids = [m["student_id"] for m in mappings]
        sid_to_bcid = {m["student_id"]: m["bookstore_customer_id"] for m in mappings}

        # Get enrollments for courses 1 and 2
        cur.execute("""
            SELECT DISTINCT user_id FROM canvas.enrollments
            WHERE course_id IN (1, 2)
        """)
        enrolled_ids = {r[0] for r in cur.fetchall()}

        # Filter to enrolled students in registry
        matched_sids = [sid for sid in student_ids if sid in enrolled_ids]

        # Get student names
        if matched_sids:
            cur.execute(
                "SELECT id, name FROM canvas.users WHERE id = ANY(%s)",
                (matched_sids,),
            )
            sid_to_name = dict(cur.fetchall())
        else:
            sid_to_name = {}

        # Get average scores per student for courses 1,2
        cur.execute("""
            SELECT s.user_id, AVG(s.score)
            FROM canvas.submissions s
            JOIN canvas.assignments a ON s.assignment_id = a.id
            WHERE a.course_id IN (1, 2)
              AND s.score IS NOT NULL
              AND s.user_id = ANY(%s)
            GROUP BY s.user_id
        """, (matched_sids,))
        sid_to_avg_score = {r[0]: round(float(r[1]), 2) for r in cur.fetchall()}

        # Get submission counts
        cur.execute("""
            SELECT s.user_id, COUNT(*)
            FROM canvas.submissions s
            JOIN canvas.assignments a ON s.assignment_id = a.id
            WHERE a.course_id IN (1, 2)
              AND s.score IS NOT NULL
              AND s.user_id = ANY(%s)
            GROUP BY s.user_id
        """, (matched_sids,))
        sid_to_sub_count = {r[0]: int(r[1]) for r in cur.fetchall()}

        # Get electronics spending from WC
        # Find Electronics/Cameras category IDs
        cur.execute("""
            SELECT id, name FROM wc.product_categories
            WHERE lower(name) IN ('electronics', 'cameras')
        """)
        elec_cats = {r[0]: r[1] for r in cur.fetchall()}
        elec_cat_names = [name.lower() for name in elec_cats.values()]

        # Get orders for matched bookstore customer IDs and compute electronics spend
        bcids = [sid_to_bcid[sid] for sid in matched_sids if sid in sid_to_bcid]
        bcid_spend = {}
        bcid_order_count = {}
        bcid_item_count = {}

        if bcids:
            if status_filter:
                cur.execute("""
                    SELECT id, customer_id, line_items FROM wc.orders
                    WHERE customer_id = ANY(%s)
                      AND status = ANY(%s)
                """, (bcids, list(status_filter)))
            else:
                cur.execute("""
                    SELECT id, customer_id, line_items FROM wc.orders
                    WHERE customer_id = ANY(%s)
                """, (bcids,))
            for order_id, cust_id, line_items in cur.fetchall():
                if not line_items:
                    continue
                items = line_items if isinstance(line_items, list) else json.loads(line_items)
                has_electronics = False
                for item in items:
                    # Check if product is in electronics/cameras category
                    product_id = item.get("product_id")
                    # Check categories in line item
                    cats = []
                    if "categories" in item:
                        cats = [c.get("name", "").lower() if isinstance(c, dict) else str(c).lower()
                                for c in item["categories"]]
                    # Also check via product_categories join
                    is_elec = any(c in elec_cat_names for c in cats)
                    if not is_elec and product_id:
                        # Look up product categories
                        cur.execute("""
                            SELECT categories FROM wc.products WHERE id = %s
                        """, (product_id,))
                        prod_row = cur.fetchone()
                        if prod_row and prod_row[0]:
                            prod_cats = prod_row[0] if isinstance(prod_row[0], list) else json.loads(str(prod_row[0]))
                            for pc in prod_cats:
                                if isinstance(pc, dict) and pc.get("name", "").lower() in elec_cat_names:
                                    is_elec = True
                                    break
                    if is_elec:
                        has_electronics = True
                        total = float(item.get("total", 0))
                        qty = int(item.get("quantity", 0))
                        bcid_spend[cust_id] = bcid_spend.get(cust_id, 0) + total
                        bcid_item_count[cust_id] = bcid_item_count.get(cust_id, 0) + qty
                if has_electronics:
                    bcid_order_count[cust_id] = bcid_order_count.get(cust_id, 0) + 1

        # Build matched_data
        matched_data = []
        for sid in matched_sids:
            if sid not in sid_to_avg_score:
                continue
            bcid = sid_to_bcid.get(sid)
            if bcid is None or bcid not in bcid_spend:
                continue
            matched_data.append({
                "student_name": sid_to_name.get(sid, f"Student_{sid}"),
                "student_id": sid,
                "avg_score": sid_to_avg_score[sid],
                "submission_count": sid_to_sub_count.get(sid, 0),
                "bookstore_customer_id": bcid,
                "total_electronics_spend": round(bcid_spend.get(bcid, 0), 2),
                "electronics_order_count": bcid_order_count.get(bcid, 0),
                "electronics_item_count": bcid_item_count.get(bcid, 0),
            })

        num_matched = len(matched_data)

        # Compute correlation
        if num_matched >= 2:
            scores = [m["avg_score"] for m in matched_data]
            spends = [m["total_electronics_spend"] for m in matched_data]
            mean_s = sum(scores) / len(scores)
            mean_sp = sum(spends) / len(spends)
            cov = sum((s - mean_s) * (sp - mean_sp) for s, sp in zip(scores, spends))
            std_s = (sum((s - mean_s) ** 2 for s in scores)) ** 0.5
            std_sp = (sum((sp - mean_sp) ** 2 for sp in spends)) ** 0.5
            correlation = round(cov / (std_s * std_sp), 4) if std_s > 0 and std_sp > 0 else 0
        else:
            scores = [m["avg_score"] for m in matched_data]
            spends = [m["total_electronics_spend"] for m in matched_data]
            mean_s = sum(scores) / len(scores) if scores else 0
            mean_sp = sum(spends) / len(spends) if spends else 0
            correlation = 0

        return {
            "matched_data": matched_data,
            "num_matched": num_matched,
            "correlation": correlation,
            "mean_score": round(mean_s, 2),
            "mean_spend": round(mean_sp, 2),
        }
    finally:
        cur.close()
        conn.close()


def _find_header_row(ws, required_col):
    """Find the 1-based row index that contains the given column name.

    Falls back to 1 when no header row is detected (e.g. a title row sits
    above the real header)."""
    target = required_col.strip().lower()
    for r in range(1, min(ws.max_row, 30) + 1):
        for c in range(1, ws.max_column + 1):
            cell = ws.cell(r, c).value
            if cell is not None and str(cell).strip().lower() == target:
                return r
    return 1


def _sheet_headers(ws, header_row):
    return [str(ws.cell(header_row, c).value or "").strip()
            for c in range(1, ws.max_column + 1)]


def _count_data_rows(ws, header_row, name_col):
    """Count non-empty rows below the header (extra blank/trailing rows ignored)."""
    count = 0
    for r in range(header_row + 1, ws.max_row + 1):
        cell = ws.cell(r, name_col).value
        if cell is not None and str(cell).strip() != "":
            count += 1
    return count


def check_excel(workspace, gt):
    print("\n=== Excel Checks ===")
    xlsx_path = os.path.join(workspace, "Student_Purchase_Analysis.xlsx")
    check("Excel file exists", os.path.exists(xlsx_path), xlsx_path)
    if not os.path.exists(xlsx_path):
        return

    wb = openpyxl.load_workbook(xlsx_path)
    sheets = wb.sheetnames

    # Sheet existence
    check("Sheet Student_Performance exists", "Student_Performance" in sheets, sheets)
    check("Sheet Purchase_Summary exists", "Purchase_Summary" in sheets, sheets)
    check("Sheet Correlation_Analysis exists", "Correlation_Analysis" in sheets, sheets)
    check("Sheet Recommendations exists", "Recommendations" in sheets, sheets)

    # Student_Performance content
    if "Student_Performance" in sheets:
        ws = wb["Student_Performance"]
        hr = _find_header_row(ws, "Student_Name")
        headers = _sheet_headers(ws, hr)
        check(
            "Student_Performance has Student_Name column",
            "Student_Name" in headers,
            headers,
        )
        check(
            "Student_Performance has Average_Score column",
            "Average_Score" in headers,
            headers,
        )
        if "Student_Name" in headers:
            name_col = headers.index("Student_Name") + 1
            data_rows = _count_data_rows(ws, hr, name_col)
            expected = [v["num_matched"] for v in _gt_variants(gt)]
            check(
                "Student_Performance has correct row count",
                any(data_rows >= v["num_matched"] for v in _gt_variants(gt)),
                f"got {data_rows} non-empty rows, expected at least one of {expected}",
            )

        # Spot check a few scores
        if "Average_Score" in headers and "Student_Name" in headers:
            score_col = headers.index("Average_Score") + 1
            name_col = headers.index("Student_Name") + 1
            gt_by_name = {m["student_name"]: m["avg_score"] for m in gt["matched_data"]}
            checked = 0
            for row in range(hr + 1, ws.max_row + 1):
                name = ws.cell(row, name_col).value
                score = ws.cell(row, score_col).value
                if name is not None and str(name).strip() in gt_by_name:
                    if _is_formula(score):
                        check(f"Score for {str(name).strip()} (formula cell)", True, score)
                    else:
                        check(
                            f"Score for {str(name).strip()}",
                            num_close(score, gt_by_name[str(name).strip()], 1.0),
                            f"got {score}, expected {gt_by_name[str(name).strip()]}",
                        )
                    checked += 1
                    if checked >= 3:
                        break

    # Purchase_Summary content
    if "Purchase_Summary" in sheets:
        ws = wb["Purchase_Summary"]
        hr = _find_header_row(ws, "Student_Name")
        headers = _sheet_headers(ws, hr)
        check(
            "Purchase_Summary has Total_Electronics_Spend",
            "Total_Electronics_Spend" in headers,
            headers,
        )
        check(
            "Purchase_Summary has Electronics_Order_Count",
            "Electronics_Order_Count" in headers,
            headers,
        )

        # Spot check spending
        if "Total_Electronics_Spend" in headers and "Student_Name" in headers:
            spend_col = headers.index("Total_Electronics_Spend") + 1
            name_col = headers.index("Student_Name") + 1
            # Accept the spend reported under any accepted GT interpretation
            # (full status set vs. completed/processing only).
            gt_spend_by_name = [
                {
                    m["student_name"]: m["total_electronics_spend"]
                    for m in v["matched_data"]
                }
                for v in _gt_variants(gt)
            ]
            checked = 0
            for row in range(hr + 1, ws.max_row + 1):
                name = ws.cell(row, name_col).value
                spend = ws.cell(row, spend_col).value
                if name is None:
                    continue
                nm = str(name).strip()
                if not any(nm in m for m in gt_spend_by_name):
                    continue
                if _is_formula(spend):
                    check(f"Spend for {nm} (formula cell)", True, spend)
                else:
                    matched = any(
                        nm in m and num_close(spend, m[nm], 5.0)
                        for m in gt_spend_by_name
                    )
                    check(
                        f"Spend for {nm}",
                        matched,
                        f"got {spend}, expected one of {[m.get(nm) for m in gt_spend_by_name]}",
                    )
                checked += 1
                if checked >= 3:
                    break

    # Correlation_Analysis content
    if "Correlation_Analysis" in sheets:
        ws = wb["Correlation_Analysis"]
        metric_rows = []
        for row in range(1, ws.max_row + 1):
            key = ws.cell(row, 1).value
            val = ws.cell(row, 2).value
            if key is None:
                continue
            k = str(key).strip()
            if k.lower() in ("metric", "value"):
                continue
            if k == "":
                continue
            metric_rows.append((k, val))
        metric_labels = [k for k, _ in metric_rows]

        # Fuzzy label matching: accept common statistical aliases such as
        # 'Pearson Correlation Coefficient' or 'P-Value'.
        pc = _find_metric(metric_rows, "Pearson_Correlation")
        pv = _find_metric(metric_rows, "P_Value")
        mas = _find_metric(metric_rows, "Mean_Academic_Score")
        mes = _find_metric(metric_rows, "Mean_Electronics_Spend")
        nom = _find_metric(metric_rows, "Number_of_Matched_Students")
        rec = _find_metric(metric_rows, "Recommendation")

        check(
            "Correlation has Pearson_Correlation",
            pc is not None,
            metric_labels,
        )
        if pc is not None:
            if _is_formula(pc):
                check("Pearson_Correlation value (formula cell)", True, pc)
            else:
                matched = any(
                    num_close(pc, v["correlation"], 0.1)
                    for v in _gt_variants(gt)
                )
                check(
                    "Pearson_Correlation value",
                    matched,
                    f"got {pc}, expected one of {[v['correlation'] for v in _gt_variants(gt)]}",
                )
        check(
            "Correlation has P_Value",
            pv is not None,
            metric_labels,
        )
        check(
            "Correlation has Mean_Academic_Score",
            mas is not None,
            metric_labels,
        )
        if mas is not None:
            if _is_formula(mas):
                check("Mean_Academic_Score value (formula cell)", True, mas)
            else:
                matched = any(
                    num_close(mas, v["mean_score"], 2.0)
                    for v in _gt_variants(gt)
                )
                check(
                    "Mean_Academic_Score value",
                    matched,
                    f"got {mas}, expected one of {[v['mean_score'] for v in _gt_variants(gt)]}",
                )
        check(
            "Correlation has Mean_Electronics_Spend",
            mes is not None,
            metric_labels,
        )
        if mes is not None:
            if _is_formula(mes):
                check("Mean_Electronics_Spend value (formula cell)", True, mes)
            else:
                # Tolerance is generous (a few excluded order statuses shift the
                # mean by roughly $25-35 while keeping every student matched);
                # the correlation check is the tighter discriminator.
                matched = any(
                    num_close(mes, v["mean_spend"], 40.0)
                    for v in _gt_variants(gt)
                )
                check(
                    "Mean_Electronics_Spend value",
                    matched,
                    f"got {mes}, expected one of {[v['mean_spend'] for v in _gt_variants(gt)]}",
                )
        check(
            "Correlation has Number_of_Matched_Students",
            nom is not None,
            metric_labels,
        )
        check(
            "Correlation has Recommendation",
            rec is not None,
            metric_labels,
        )
        if rec is not None:
            rec_str = str(rec).lower()
            # When the correlation is very close to zero, its sign (and hence the
            # natural recommendation direction) is sensitive to small rounding
            # differences in the model's own computation, so accept either.
            gt_corr = _to_float(gt["correlation"])
            near_zero = gt_corr is None or abs(gt_corr) < 0.15
            if near_zero:
                expected_keywords = ("discount", "expand")
            else:
                expected_keywords = ("discount",) if gt_corr < 0 else ("expand",)
            matched = any(
                re.search(rf"\b{re.escape(kw)}\b", rec_str) is not None
                for kw in expected_keywords
            )
            check(
                f"Recommendation matches correlation direction (keyword: {' or '.join(expected_keywords)})",
                matched,
                f"got '{rec_str[:80]}', expected one of {expected_keywords}",
            )

    # Recommendations sheet
    if "Recommendations" in sheets:
        ws = wb["Recommendations"]
        check(
            "Recommendations has at least 3 data rows",
            ws.max_row >= 4,
            f"rows={ws.max_row}",
        )


def check_pptx(workspace, gt):
    print("\n=== PowerPoint Checks ===")
    pptx_path = os.path.join(workspace, "Purchase_Behavior_Presentation.pptx")
    check("PPTX file exists", os.path.exists(pptx_path), pptx_path)
    if not os.path.exists(pptx_path):
        return

    prs = Presentation(pptx_path)
    slides = prs.slides
    check("PPTX has 5 slides", len(slides) == 5, f"got {len(slides)}")

    if len(slides) >= 1:
        title = slides[0].shapes.title
        if title:
            check(
                "Slide 1 title contains 'Purchase Behavior' or 'Electronics'",
                "purchase" in title.text.lower() or "electronics" in title.text.lower(),
                title.text,
            )

    expected_titles = [
        "Study Overview",
        "Student Performance Summary",
        "Correlation Findings",
        "Recommendations",
    ]
    for i, exp in enumerate(expected_titles):
        if i + 1 < len(slides):
            slide = slides[i + 1]
            t = slide.shapes.title
            if t:
                check(
                    f"Slide {i+2} title contains '{exp.split()[0]}'",
                    exp.split()[0].lower() in t.text.lower(),
                    t.text,
                )

    # Check slide content mentions key data
    all_text = ""
    for slide in slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                all_text += shape.text_frame.text + " "

    expected_counts = [str(v["num_matched"]) for v in _gt_variants(gt)]
    check(
        "PPTX mentions number of matched students",
        any(s in all_text for s in expected_counts),
        f"looking for one of {expected_counts}",
    )
    # Compare the correlation numerically rather than by exact string: models
    # commonly round the coefficient to 2-3 decimals in slides (e.g. -0.05, -0.051)
    # or, when it is essentially zero, present it as "0.00"/"-0.0". Tolerate both
    # by matching any integer or decimal token near any accepted GT correlation.
    # Unicode minus (U+2212) is normalized to ASCII '-'.
    scan_text = all_text.replace("−", "-").replace("–", "-").replace("—", "-")
    found_corr = False
    for tok in re.findall(r"[-+]?\d+(?:\.\d+)?", scan_text):
        f = _to_float(tok)
        if f is not None and any(
            num_close(f, v["correlation"], 0.06) for v in _gt_variants(gt)
        ):
            found_corr = True
            break
    check(
        "PPTX mentions correlation coefficient",
        found_corr,
        f"looking for a value within 0.06 of any of {[v['correlation'] for v in _gt_variants(gt)]}",
    )


def _recipient_matches(to_addr, needle):
    """Check whether a to_addr (jsonb array or string) contains the given address."""
    needle = needle.lower()
    if to_addr is None:
        return False
    if isinstance(to_addr, list):
        return any(needle in str(r).lower() for r in to_addr)
    if isinstance(to_addr, str):
        return needle in to_addr.lower()
    return False


def check_emails(gt):
    print("\n=== Email Checks ===")
    conn = psycopg2.connect(**DB_CONFIG)
    cur = conn.cursor()
    try:
        sent_folder = "(SELECT id FROM email.folders WHERE name = 'Sent' LIMIT 1)"

        # Check bookstore manager email. Use ILIKE (case-insensitive): email
        # addresses are case-insensitive, so a model may legitimately render the
        # address with capital letters.
        cur.execute(
            f"""SELECT subject, body_text, to_addr FROM email.messages
            WHERE folder_id = {sent_folder}
            AND to_addr::text ILIKE '%bookstore_manager@university.edu%'"""
        )
        mgr_emails = cur.fetchall()
        check(
            "Email sent to bookstore_manager@university.edu",
            len(mgr_emails) >= 1,
            f"found {len(mgr_emails)}",
        )
        if mgr_emails:
            subj_ok = any(
                any(kw in (e[0] or "").lower() for kw in ["purchase", "behavior", "report", "electronics"])
                for e in mgr_emails
            )
            body_ok = any(
                "correlation" in (e[1] or "").lower()
                or "matched" in (e[1] or "").lower()
                or str(gt["num_matched"]) in (e[1] or "")
                for e in mgr_emails
            )
            check(
                "Manager email subject mentions purchase/behavior/report",
                subj_ok,
                [e[0] for e in mgr_emails],
            )
            check(
                "Manager email body mentions correlation or matched students",
                body_ok,
                [ (e[1] or "")[:100] for e in mgr_emails ],
            )

        # Check academic affairs email
        cur.execute(
            f"""SELECT subject, body_text, to_addr FROM email.messages
            WHERE folder_id = {sent_folder}
            AND to_addr::text ILIKE '%academic_affairs@university.edu%'"""
        )
        acad_emails = cur.fetchall()
        check(
            "Email sent to academic_affairs@university.edu",
            len(acad_emails) >= 1,
            f"found {len(acad_emails)}",
        )
        if acad_emails:
            subj_ok = any(
                any(kw in (e[0] or "").lower() for kw in ["correlation", "academic", "performance"])
                for e in acad_emails
            )
            body_ok = any(
                "correlation" in (e[1] or "").lower()
                or "mean" in (e[1] or "").lower()
                or str(gt["mean_score"]) in (e[1] or "")
                for e in acad_emails
            )
            check(
                "Academic email subject mentions correlation or academic",
                subj_ok,
                [e[0] for e in acad_emails],
            )
            check(
                "Academic email body mentions mean score or correlation",
                body_ok,
                [ (e[1] or "")[:100] for e in acad_emails ],
            )
    finally:
        cur.close()
        conn.close()


def check_reverse_validation(workspace):
    print("\n=== Reverse Validation ===")
    conn = psycopg2.connect(**DB_CONFIG)
    cur = conn.cursor()
    try:
        # Check no emails sent to noise recipients (exact match, not substring)
        noise_recipients = [
            "newsletter@university.edu",
            "all-staff@university.edu",
            "alumni@university.edu",
            "admissions@university.edu",
            "events@university.edu",
            "marketing@university.edu",
            "inventory@university.edu",
        ]
        # Only check Sent folder
        cur.execute(
            "SELECT to_addr FROM email.messages WHERE folder_id = (SELECT id FROM email.folders WHERE name='Sent' LIMIT 1)"
        )
        sent_recips = set()
        for (to,) in cur.fetchall():
            if to is None:
                continue
            to_list = to if isinstance(to, list) else json.loads(to) if isinstance(to, str) else []
            for r in to_list:
                sent_recips.add(str(r).strip().lower())
        for addr in noise_recipients:
            check(f"No email sent to noise recipient {addr}",
                  addr not in sent_recips,
                  f"Sent recipients include {addr}")

        # Verify preprocess-injected noise emails are preserved (not deleted by agent)
        noise_subjects = [
            "Weekly Bookstore Inventory Update",
            "Campus Event: Tech Fair Next Month",
            "Student Discount Program Proposal",
        ]
        cur.execute("SELECT subject FROM email.messages WHERE folder_id IN (SELECT id FROM email.folders WHERE name IN ('INBOX','Sent'))")
        existing_subjects = {r[0] for r in cur.fetchall()}
        preserved = sum(1 for s in noise_subjects if s in existing_subjects)
        check("Reverse: noise emails preserved (not deleted)",
              preserved >= 2,
              f"Only {preserved}/3 noise emails remain")

        # Check Excel does not include non-Electronics categories in Purchase_Summary
        xlsx_path = os.path.join(workspace, "Student_Purchase_Analysis.xlsx")
        if os.path.exists(xlsx_path):
            wb = openpyxl.load_workbook(xlsx_path)
            if "Purchase_Summary" in wb.sheetnames:
                ws = wb["Purchase_Summary"]
                all_text = " ".join(
                    str(ws.cell(r, c).value) for r in range(2, ws.max_row + 1)
                    for c in range(1, ws.max_column + 1) if ws.cell(r, c).value
                ).lower()
                check("Purchase_Summary does not contain Clothing category data",
                      "clothing" not in all_text,
                      "Found 'clothing' in Purchase_Summary")
    except Exception as e:
        check("Reverse validation", False, str(e))
    finally:
        cur.close()
        conn.close()


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--agent_workspace", required=False, default=".")
    parser.add_argument("--groundtruth_workspace", required=False, default=".")
    parser.add_argument("--launch_time", required=False)
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()

    gt = load_groundtruth()

    check_excel(args.agent_workspace, gt)
    check_pptx(args.agent_workspace, gt)
    check_emails(gt)
    check_reverse_validation(args.agent_workspace)

    total = PASS_COUNT + FAIL_COUNT
    accuracy = PASS_COUNT / total * 100 if total > 0 else 0
    print(f"\nOverall: {PASS_COUNT}/{total} ({accuracy:.1f}%)")
    result = {
        "total_passed": PASS_COUNT,
        "total_checks": total,
        "accuracy": accuracy,
    }
    if args.res_log_file:
        with open(args.res_log_file, "w") as f:
            json.dump(result, f, indent=2)
    # Tightened: require all checks to pass (previously >=70%).
    sys.exit(0 if FAIL_COUNT == 0 else 1)


if __name__ == "__main__":
    main()
