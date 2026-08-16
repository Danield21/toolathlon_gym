"""
Evaluation for scholarly-arxiv-ppt-gsheet task.
Checks: PPT file, GSheet, Word document.
"""
import argparse
import os
import sys

import psycopg2
from docx import Document
from pptx import Presentation

DB_CONFIG = {
    "host": os.environ.get("PGHOST", "localhost"),
    "port": int(os.environ.get("PGPORT", "5432")),
    "dbname": "toolathlon_gym",
    "user": "eigent",
    "password": "camel",
}

PASS_COUNT = 0
FAIL_COUNT = 0

RELEVANT_KEYWORDS = ["chain-of-thought", "self-consistency", "tree of thoughts", "chain of thought", "reasoning"]
PAPER_TITLES_LOWER = [
    "chain-of-thought prompting elicits reasoning",
    "self-consistency improves chain of thought",
    "tree of thoughts: deliberate problem solving",
]


def record(name, passed, detail=""):
    global PASS_COUNT, FAIL_COUNT
    if passed:
        PASS_COUNT += 1
        print(f"  [PASS] {name}")
    else:
        FAIL_COUNT += 1
        msg = f": {detail[:300]}" if detail else ""
        print(f"  [FAIL] {name}{msg}")


def check_pptx(agent_workspace):
    print("\n=== Checking PowerPoint ===")
    ppt_path = os.path.join(agent_workspace, "LLM_Reasoning_Survey.pptx")
    if not os.path.isfile(ppt_path):
        record("PPT file LLM_Reasoning_Survey.pptx exists", False, f"Not found at: {ppt_path}")
        return
    record("PPT file LLM_Reasoning_Survey.pptx exists", True)

    try:
        prs = Presentation(ppt_path)
    except Exception as e:
        record("PPT file readable", False, str(e))
        return
    record("PPT file readable", True)

    num_slides = len(prs.slides)
    # Task implies 5 slides (title + 3 papers + summary). Allow 5-6 for tolerance, but bound noise.
    record("PPT has 5-6 slides", 5 <= num_slides <= 6, f"Found {num_slides} slides")

    # Per-slide text extraction
    slide_texts = []
    for slide in prs.slides:
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text.append(shape.text.lower())
        slide_texts.append(" ".join(slide_text))
    full_text = " ".join(slide_texts)

    # Check FIRST slide title contains 'survey' AND 'reasoning'
    if slide_texts:
        first = slide_texts[0]
        first_ok = ("survey" in first) and ("reasoning" in first) and ("llm" in first or "language" in first)
        record("PPT first slide titled 'Survey of LLM Reasoning Methods'", first_ok,
               f"First slide text: {first[:200]}")
    else:
        record("PPT first slide titled 'Survey of LLM Reasoning Methods'", False, "No slides")

    # Check paper content - must each appear on a dedicated slide
    paper_keywords = [
        ("chain-of-thought paper", ["chain-of-thought", "chain of thought"]),
        ("self-consistency paper", ["self-consistency", "self consistency"]),
        ("tree of thoughts paper", ["tree of thoughts", "tree of thought"]),
    ]
    for label, kws in paper_keywords:
        # at least one body slide (not first/last) contains this paper
        found_in_slide = False
        for st in slide_texts[1:-1] if len(slide_texts) >= 3 else slide_texts:
            if any(kw in st for kw in kws):
                found_in_slide = True
                break
        # fallback: also accept if appears in any non-first slide
        if not found_in_slide:
            for st in slide_texts[1:]:
                if any(kw in st for kw in kws):
                    found_in_slide = True
                    break
        record(f"PPT has dedicated slide for {label}", found_in_slide)

    # Check summary slide (last slide should mention summary/comparison/conclusion)
    if slide_texts:
        last = slide_texts[-1]
        has_summary_last = any(k in last for k in ["summary", "comparison", "conclusion", "compare"])
        record("PPT last slide is summary/comparison", has_summary_last,
               f"Last slide text: {last[:200]}")
    else:
        record("PPT last slide is summary/comparison", False, "No slides")


def check_gsheet():
    print("\n=== Checking Google Sheet ===")
    try:
        conn = psycopg2.connect(**DB_CONFIG)
        cur = conn.cursor()

        cur.execute("SELECT id, title FROM gsheet.spreadsheets")
        spreadsheets = cur.fetchall()

        target_ss = None
        for sid, title in spreadsheets:
            if title and ("llm" in title.lower() or "reasoning" in title.lower()) and "paper" in title.lower():
                target_ss = sid
                break

        record("GSheet 'LLM Reasoning Paper Tracker' exists",
               target_ss is not None,
               f"Found sheets: {[t for _, t in spreadsheets]}")

        if target_ss is None:
            conn.close()
            return

        # Find any sheet in this spreadsheet
        cur.execute("SELECT id, title FROM gsheet.sheets WHERE spreadsheet_id = %s", (target_ss,))
        sheets = cur.fetchall()
        record("GSheet has at least one sheet", len(sheets) > 0, f"Found: {sheets}")

        if not sheets:
            conn.close()
            return

        sheet_id = sheets[0][0]

        # Count data rows (exclude header). Task expects 3 papers; allow 3-4 only.
        cur.execute("""
            SELECT COUNT(DISTINCT row_index) FROM gsheet.cells
            WHERE spreadsheet_id = %s AND sheet_id = %s AND row_index > 0
        """, (target_ss, sheet_id))
        data_rows = cur.fetchone()[0]
        record("GSheet has 3-4 data rows (3 expected)", 3 <= data_rows <= 4, f"Found {data_rows} data rows")

        # Check column headers (row_index = 0)
        cur.execute("""
            SELECT LOWER(value) FROM gsheet.cells
            WHERE spreadsheet_id = %s AND sheet_id = %s AND row_index = 0
        """, (target_ss, sheet_id))
        header_values = [row[0] for row in cur.fetchall() if row[0]]
        header_text = " ".join(header_values)
        required_headers = ["title", "author", "arxiv", "publish", "key_method"]
        missing_headers = [h for h in required_headers if h not in header_text and h.replace("_", " ") not in header_text]
        record("GSheet has required column headers (Title/Authors/ArXiv_ID/Published_Date/Key_Method)",
               len(missing_headers) == 0,
               f"Missing: {missing_headers}; Found headers: {header_values}")

        # Check that paper titles appear in cells
        cur.execute("""
            SELECT LOWER(value) FROM gsheet.cells
            WHERE spreadsheet_id = %s AND sheet_id = %s
        """, (target_ss, sheet_id))
        cell_values = [row[0] for row in cur.fetchall() if row[0]]
        all_cells_text = " ".join(cell_values)

        has_chain = "chain-of-thought" in all_cells_text or "chain of thought" in all_cells_text
        has_self = "self-consistency" in all_cells_text or "self consistency" in all_cells_text
        has_tree = "tree of thoughts" in all_cells_text or "tree of thought" in all_cells_text
        record("GSheet contains chain-of-thought paper entry", has_chain)
        record("GSheet contains self-consistency paper entry", has_self)
        record("GSheet contains tree of thoughts paper entry", has_tree)

        # Verify GSheet contains correct ArXiv IDs (cross-reference with scholarly DB)
        # Only the 3 RELEVANT paper IDs should be present
        for arxiv_id, label in [("2301.00234", "chain-of-thought"),
                                 ("2302.11382", "self-consistency"),
                                 ("2305.10601", "tree of thoughts")]:
            record(f"GSheet contains arxiv id {arxiv_id} ({label})",
                   arxiv_id in all_cells_text)

        # Verify per-paper authors appear in GSheet
        for author, label in [("jason wei", "chain-of-thought author"),
                              ("xuezhi wang", "self-consistency author"),
                              ("shunyu yao", "tree of thoughts author")]:
            record(f"GSheet contains {label} '{author}'",
                   author in all_cells_text)

        # Verify NOISE arxiv IDs are NOT included (no irrelevant papers)
        noise_ids_present = [nid for nid in ("2301.11093", "2302.05543", "2304.09842")
                             if nid in all_cells_text]
        record("GSheet excludes noise paper arXiv IDs",
               len(noise_ids_present) == 0,
               f"unexpected noise: {noise_ids_present}")

        conn.close()
    except Exception as e:
        record("GSheet connection", False, str(e))


def check_word(agent_workspace):
    print("\n=== Checking Word Document ===")
    doc_path = os.path.join(agent_workspace, "LLM_Reasoning_Literature_Review.docx")
    if not os.path.isfile(doc_path):
        record("Word file LLM_Reasoning_Literature_Review.docx exists", False, f"Not found at: {doc_path}")
        return
    record("Word file LLM_Reasoning_Literature_Review.docx exists", True)

    try:
        doc = Document(doc_path)
    except Exception as e:
        record("Word file readable", False, str(e))
        return
    record("Word file readable", True)

    full_text = "\n".join(p.text for p in doc.paragraphs).lower()

    has_heading = "literature review" in full_text and ("llm" in full_text or "reasoning" in full_text)
    record("Word has 'Literature Review' heading with LLM/reasoning", has_heading)

    has_intro = len(full_text) > 300
    record("Word has substantial content (intro + sections)", has_intro, f"Text length: {len(full_text)}")

    has_chain = "chain-of-thought" in full_text or "chain of thought" in full_text
    has_self = "self-consistency" in full_text or "self consistency" in full_text
    has_tree = "tree of thoughts" in full_text or "tree of thought" in full_text
    record("Word mentions chain-of-thought method", has_chain)
    record("Word mentions self-consistency method", has_self)
    record("Word mentions tree of thoughts method", has_tree)


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--agent_workspace", required=True)
    parser.add_argument("--groundtruth_workspace", required=False)
    parser.add_argument("--res_log_file", required=False)
    parser.add_argument("--launch_time", required=False)
    args = parser.parse_args()

    check_pptx(args.agent_workspace)
    check_gsheet()
    check_word(args.agent_workspace)

    total = PASS_COUNT + FAIL_COUNT
    print(f"\n=== Results: {PASS_COUNT}/{total} passed ===")
    if FAIL_COUNT > 0:
        print(f"{FAIL_COUNT} checks failed")
        sys.exit(1)
    else:
        print("All checks passed!")
        sys.exit(0)


if __name__ == "__main__":
    main()
