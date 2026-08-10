"""
Evaluation for 12306-yf-investor-roadshow-excel-ppt-email task.

Checks:
1. Roadshow_Analysis.xlsx exists with Travel_Plan, Stock_Summary, Financial_Highlights sheets
2. Travel_Plan has Shanghai (G11 or G1) and Guangzhou (G105) rows, business class, prices, 2026-03-10
3. Stock_Summary has GOOGL and AMZN rows with tolerant numeric values
4. Financial_Highlights has GOOGL revenue/net income for the two most recent fiscal years
5. Investor_Roadshow.pptx exists with >= 5 slides and roadshow/travel/finance content
6. Emails sent to both investors@fundmanager.com and shanghai_partners@finance.com

Design notes (fix round):
- Tickers are GOOGL (primary) / AMZN (comparison) because the runtime yfinance mock
  DB (yf.* tables) only contains GOOGL, AMZN, JPM, JNJ, XOM, ^DJI, GC=F. AAPL/MSFT do
  not exist anywhere in the seed, so they were replaced per the R11 recipe.
- All DB connections read PGHOST/PGPORT/PGDATABASE/PGUSER/PGPASSWORD from env (R1).
- Workbooks are read with data_only=False so formula cells do not silently become None;
  cached values are used when present (R2).
- String comparisons normalize underscores/spaces/punctuation/case (fix for
  'Net_Income' vs 'Net Income', 'Apple Inc.' vs 'Apple Inc').
- Change_30D_Pct is runtime-only because the mock price history ends 2026-03-05 and
  the 'past 30 days' window depends on the runtime clock (underdetermined).
- Row comparison is content-driven (matched by ticker / metric+year) instead of
  positional, so extra/duplicate rows from the P=2 homogeneous swarm do not cause
  false FAILs (R10).

Design notes (fix round 2 - adversarial review):
- City matching is bilingual (_city_matches): the 12306 mock returns Chinese
  station/city names (上海/上海虹桥/广州/广州南), so the City column may be written
  in either language.
- Meeting date / fiscal-year parsing are format-tolerant (_has_mar10, to_year):
  2026/3/10, 2026年3月10日, FY2025, Dec 2025, 2025年, 2-digit years all accepted.
- Stock_Summary columns are matched by header name (_agent_col + HEADER_ALIASES),
  not position, so reordering value columns does not cause cross-column compares.
- Market_Cap_B tolerance widened to 12% so a price*floatShares computation
  (~10.5% from the GT marketCap) is accepted; the band still separates tickers.
- rec_keyword recognizes analyst rating phrasings (Outperform/Overweight/Neutral/
  Underperform/Reduce/...), not just buy/hold/sell.
- Formula cells without a cached value now FAIL numeric checks (they were previously
  skipped): formulas are banned by the task's formatting requirements, so a correct
  model never hits this, and empty/no-cache cells must not silently pass.
- to_float strips currency labels/symbols (元/￥/CNY/RMB/USD) and unit suffixes
  (B/bn/billion) so correctly-sourced values with units still parse.
"""
import datetime as _dt
import json
import os
import re
import sys
from argparse import ArgumentParser

import psycopg2
import openpyxl

DB_CONFIG = {
    "host": os.environ.get("PGHOST", "localhost"),
    "port": int(os.environ.get("PGPORT", "5432")),
    "dbname": os.environ.get("PGDATABASE", "toolathlon_gym"),
    "user": os.environ.get("PGUSER", "eigent"),
    "password": os.environ.get("PGPASSWORD", "camel"),
}

# Expected trains per route (Beijing -> city). Both G11 and G1 serve Beijing South ->
# Shanghai Hongqiao on 2026-03-10 (per train.trains in the seed); G105 is the only
# Beijing -> Guangzhou service that day. Business-seat prices come from train.train_seats.
SHANGHAI_TRAINS = {"g11", "g1"}
SHANGHAI_PRICE = 1748.5
GUANGZHOU_TRAIN = "g105"
GUANGZHOU_PRICE = 2631.5

# Tickers the task is built on (primary GOOGL, comparison AMZN).
PRIMARY_TICKER = "GOOGL"
COMPARISON_TICKER = "AMZN"

PASS_COUNT = 0
FAIL_COUNT = 0
BLOCKING_FAIL_COUNT = 0


def record(name, passed, detail="", runtime_only=False):
    global PASS_COUNT, FAIL_COUNT, BLOCKING_FAIL_COUNT
    if passed:
        PASS_COUNT += 1
        print(f"  [PASS] {name}")
    else:
        FAIL_COUNT += 1
        if not runtime_only:
            BLOCKING_FAIL_COUNT += 1
        msg = f": {detail[:300]}" if detail else ""
        suffix = " (runtime-only)" if runtime_only else ""
        print(f"  [FAIL] {name}{suffix}{msg}")


# ---------------------------------------------------------------- value helpers

def norm(s):
    """Normalize text for comparison: lower-case, strip, remove punctuation/space."""
    if s is None:
        return ""
    s = str(s).strip().lower()
    for ch in " _.,-()'\"/\\:*¥$€%":
        s = s.replace(ch, "")
    return s


def to_float(v):
    """Robustly parse a numeric value (str/int/float/None). Strips currency
    symbols/labels (incl. 元/￥/CNY/RMB/USD), thousands separators, '%',
    whitespace, and common unit suffixes ('B'/'bn'/'billion'). Formula cells
    (starting with '=') and unparseable strings return None."""
    if v is None:
        return None
    if isinstance(v, bool):
        return None
    if isinstance(v, (int, float)):
        return float(v)
    s = str(v).strip()
    if s.startswith("="):
        return None
    s = s.replace(",", "").replace("$", "").replace("¥", "").replace("€", "")
    s = s.replace("%", "").replace(" ", "").replace("'", "")
    s = s.replace("元", "").replace("￥", "")
    s = re.sub(r"(?i)^(usd|cny|rmb)", "", s)  # leading currency label, e.g. 'CNY 1748.5'
    for _ in range(3):
        s2 = re.sub(r"(?i)(usd|cny|rmb|bn|billion|b)$", "", s)  # trailing unit/currency
        if s2 == s:
            break
        s = s2
    try:
        return float(s)
    except ValueError:
        return None


def num_close(a, b, tol=1.0):
    """Compare two values numerically when both parse; otherwise fall back to a
    normalized string comparison (only when one side is not numeric)."""
    fa, fb = to_float(a), to_float(b)
    if fa is not None and fb is not None:
        return abs(fa - fb) <= tol
    if fa is None or fb is None:
        if a is None or b is None:
            return a is None and b is None
        return norm(a) == norm(b)
    return False


def rec_keyword(v):
    """Extract the consensus keyword (buy/hold/sell) from a recommendation string.

    Recognizes the plain words (buy/hold/sell) as well as common analyst rating
    phrasings: 'Outperform'/'Overweight'/'Accumulate' -> buy, 'Neutral'/
    'Equal-Weight'/'Market Perform'/'Sector Perform' -> hold, 'Underperform'/
    'Underweight'/'Reduce' -> sell.
    """
    n = norm(v)
    if not n:
        return None
    sell_terms = ("sell", "underperform", "underweight", "reduce")
    buy_terms = ("buy", "overweight", "outperform", "accumulate")
    hold_terms = ("hold", "neutral", "equalweight", "marketperform", "sectorperform", "inline", "in-line")
    if any(t in n for t in sell_terms):
        return "sell"
    if any(t in n for t in buy_terms):
        return "buy"
    if any(t in n for t in hold_terms):
        return "hold"
    return None


def _city_matches(v, aliases):
    """Whether a cell value refers to one of the given city aliases.

    Supports both English/pinyin ('Shanghai', 'Shanghai Hongqiao') and Chinese
    ('上海', '上海虹桥') spellings returned by the 12306 mock.
    """
    n = norm(v)
    if not n:
        return False
    return any(a in n for a in aliases)


def _train_ok(v, codes):
    """Whether a cell value denotes a train in the given code set.

    Accepts the display code ('G11', case/punctuation variants), plus the
    internal `train_no` form returned by the mock ('G11_260310_1' normalizes to
    'g112603101') via a leading-code + trailing-digit-run prefix match. The
    trailing digit-run must be >= 6 digits so 'g105' does not accidentally match
    'g1' (its leading 'g1' is followed by only '05').
    """
    n = norm(v)
    if not n:
        return False
    if n in codes:
        return True
    for c in sorted(codes, key=len, reverse=True):
        if n.startswith(c):
            rest = n[len(c):]
            if len(rest) >= 6 and rest.isdigit():
                return True
    return False


def _has_mar10(all_text, rows):
    """Whether the Travel_Plan sheet mentions the meeting date 2026-03-10 in any
    common representation (ISO date, slash date, Chinese date, 'March 10'...)."""
    for row in rows:
        for c in row:
            if isinstance(c, (_dt.datetime, _dt.date)):
                if c.year == 2026 and c.month == 3 and c.day == 10:
                    return True
    t = all_text.lower()
    # 2026-03-10 / 2026/3/10 / 2026-3-10 / 2026年3月10日 / 2026年03月10日
    if re.search(r"2026[-/年]0?3[-/月]10", t):
        return True
    # 3/10/2026 / 03/10/2026
    if re.search(r"0?3/10/2026", t):
        return True
    for kw in ("march 10", "mar 10", "10 march", "10 mar", "3月10", "march 10th", "mar 10th"):
        if kw in t:
            return True
    return False


# ---------------------------------------------------------------- workbook helpers

def _load_sheet_values(path):
    """Load a workbook and return {sheet_name: [ [row values...], ... ]}.

    Values are read with data_only=False so formula cells keep their source string;
    when a cell is a formula, the cached value (from a data_only=True load) is used if
    present, otherwise the raw formula string is returned and the numeric checks fail
    it (formulas are banned by the task's formatting requirements).
    """
    wb_raw = openpyxl.load_workbook(path, data_only=False)
    wb_cached = openpyxl.load_workbook(path, data_only=True)
    result = {}
    for ws_name in wb_raw.sheetnames:
        raw = wb_raw[ws_name]
        cached = wb_cached[ws_name] if ws_name in wb_cached.sheetnames else None
        rows = []
        raw_iter = raw.iter_rows()
        cached_iter = cached.iter_rows() if cached else None
        for i, raw_row in enumerate(raw_iter):
            cached_row = next(cached_iter) if cached_iter else None
            row = []
            for j, cell in enumerate(raw_row):
                v = cell.value
                if isinstance(v, str) and v.startswith("="):
                    if cached_row is not None and j < len(cached_row):
                        cv = cached_row[j].value
                        if cv is not None:
                            v = cv
                row.append(v)
            rows.append(row)
        result[ws_name] = rows
    wb_raw.close()
    wb_cached.close()
    return result


def _find_sheet(sheet_values, target):
    """Find a sheet by normalized name. Returns the sheet name or None."""
    t = norm(target)
    for name in sheet_values:
        if norm(name) == t:
            return name
    return None


def _data_rows(rows):
    """Rows after the header that contain at least one non-empty cell."""
    return [r for r in rows[1:] if any(c is not None and str(c).strip() != "" for c in r)]


def _header_of(rows):
    return [str(c).strip() if c is not None else "" for c in rows[0]]


# ---------------------------------------------------------------- excel checks

def check_excel(agent_workspace, groundtruth_workspace="."):
    print("\n=== Check 1: Excel Roadshow_Analysis.xlsx ===")

    xlsx_path = os.path.join(agent_workspace, "Roadshow_Analysis.xlsx")
    if not os.path.exists(xlsx_path):
        record("Roadshow_Analysis.xlsx exists", False, f"Not found at {xlsx_path}")
        return
    record("Roadshow_Analysis.xlsx exists", True)

    try:
        sheets = _load_sheet_values(xlsx_path)
    except Exception as e:
        record("Excel readable", False, str(e))
        return
    record("Excel readable", True)

    # --- Sheet presence (normalized name matching, so 'Travel Plan' == 'Travel_Plan') ---
    travel_sheet = _find_sheet(sheets, "Travel_Plan")
    stock_sheet = _find_sheet(sheets, "Stock_Summary")
    financial_sheet = _find_sheet(sheets, "Financial_Highlights")

    record("Excel has Travel_Plan sheet", travel_sheet is not None, f"Sheets: {list(sheets)}")
    record("Excel has Stock_Summary sheet", stock_sheet is not None, f"Sheets: {list(sheets)}")
    record("Excel has Financial_Highlights sheet", financial_sheet is not None, f"Sheets: {list(sheets)}")

    # --- Travel_Plan content ---
    if travel_sheet:
        rows = sheets[travel_sheet]
        data_rows = _data_rows(rows)
        record("Travel_Plan has >= 2 rows", len(data_rows) >= 2, f"Found {len(data_rows)} rows")

        all_text = " ".join(str(c) for row in rows for c in row if c is not None).lower()

        # City aliases: the 12306 mock returns Chinese station/city names
        # (上海 / 上海虹桥 / 广州 / 广州南); the agent may transcribe either the
        # English city or the Chinese name into the City column.
        SH_ALIASES = ("shanghai", "上海")
        GZ_ALIASES = ("guangzhou", "广州")

        # Shanghai leg: a data row whose city is Shanghai and train is G11 or G1;
        # fall back to "G11/G1 token present AND Shanghai mentioned" if the city column
        # is laid out differently.
        has_shanghai_train = any(
            len(r) > 1 and r[0] is not None and _city_matches(r[0], SH_ALIASES)
            and _train_ok(r[1], SHANGHAI_TRAINS)
            for r in data_rows
        ) or (any(_train_ok(c, SHANGHAI_TRAINS) for row in rows for c in row if c is not None)
              and any(r and r[0] and _city_matches(r[0], SH_ALIASES) for r in data_rows))
        record("Travel_Plan contains Shanghai train (G11 or G1)", has_shanghai_train,
               f"Trains found in sheet: {[r[1] for r in data_rows if len(r) > 1]}")

        has_guangzhou_train = any(
            len(r) > 1 and r[0] is not None and _city_matches(r[0], GZ_ALIASES)
            and _train_ok(r[1], {GUANGZHOU_TRAIN})
            for r in data_rows
        ) or (any(_train_ok(c, {GUANGZHOU_TRAIN}) for row in rows for c in row if c is not None)
              and any(r and r[0] and _city_matches(r[0], GZ_ALIASES) for r in data_rows))
        record(f"Travel_Plan contains Guangzhou train ({GUANGZHOU_TRAIN.upper()})",
               has_guangzhou_train, f"Trains found in sheet: {[r[1] for r in data_rows if len(r) > 1]}")

        # Ticket prices (from train.train_seats seed: G11/G1 business = 1748.5, G105 = 2631.5)
        numerics = [f for r in data_rows for c in r if (f := to_float(c)) is not None]
        has_sh_price = any(abs(v - SHANGHAI_PRICE) <= 5 for v in numerics)
        has_gz_price = any(abs(v - GUANGZHOU_PRICE) <= 5 for v in numerics)
        record(f"Travel_Plan has Shanghai ticket price ~{SHANGHAI_PRICE}", has_sh_price,
               f"Numerics: {[round(v, 1) for v in numerics]}")
        record(f"Travel_Plan has Guangzhou ticket price ~{GUANGZHOU_PRICE}", has_gz_price,
               f"Numerics: {[round(v, 1) for v in numerics]}")

        # Business class mentioned (Chinese '商务座' or English 'business')
        has_business = any(
            "商务" in str(c) or norm(c) == "business" or "business" in norm(c)
            for row in rows for c in row if c is not None
        )
        record("Travel_Plan mentions business class", has_business, f"Content sample: {all_text[:200]}")

        # Meeting date must mention 2026-03-10 (any common format)
        has_mar10 = _has_mar10(all_text, rows)
        record("Travel_Plan mentions 2026-03-10", has_mar10, f"Content sample: {all_text[:200]}")

    # --- Stock_Summary content ---
    if stock_sheet:
        rows = sheets[stock_sheet]
        data_rows = _data_rows(rows)
        record("Stock_Summary has >= 2 rows", len(data_rows) >= 2, f"Found {len(data_rows)} rows")

        all_text = " ".join(str(c) for row in rows for c in row if c is not None).upper()
        record(f"Stock_Summary contains {PRIMARY_TICKER}", PRIMARY_TICKER in all_text, f"Content: {all_text[:200]}")
        record(f"Stock_Summary contains {COMPARISON_TICKER}", COMPARISON_TICKER in all_text, f"Content: {all_text[:200]}")

    # --- Financial_Highlights content ---
    if financial_sheet:
        rows = sheets[financial_sheet]
        data_rows = _data_rows(rows)
        record("Financial_Highlights has >= 2 rows", len(data_rows) >= 2, f"Found {len(data_rows)} rows")

    # --- Groundtruth value comparison (content-driven, tolerant) ---
    gt_path = os.path.join(groundtruth_workspace, "Roadshow_Analysis.xlsx")
    if not os.path.isfile(gt_path):
        record("Groundtruth xlsx exists", False, gt_path)
        return
    try:
        gt_sheets = _load_sheet_values(gt_path)
    except Exception as e:
        record("Groundtruth xlsx readable", False, str(e))
        return
    record("Groundtruth xlsx readable", True)

    _compare_stock_summary(gt_sheets, sheets, stock_sheet)
    _compare_financial_highlights(gt_sheets, sheets, financial_sheet)


def _tol_for(header, gt_val):
    """Tolerance for a numeric GT value based on the column header."""
    h = norm(header)
    if h in ("currentprice",):
        return max(3.0, abs(gt_val) * 0.03)
    if h in ("marketcapb",):
        # 12% tolerance: a model that computes Market_Cap_B as price * floatShares
        # deviates ~10.5% (GOOGL) / ~9.1% (AMZN) from the stock_info marketCap that
        # the GT is built from. The tolerance band still separates the two tickers
        # and the other value columns, so the check stays discriminating.
        return max(40.0, abs(gt_val) * 0.12)
    if h in ("revenueb", "valueb", "value"):
        return max(10.0, abs(gt_val) * 0.05)
    if h in ("eps",):
        return max(0.5, abs(gt_val) * 0.1)
    if h in ("durationhours",):
        return max(0.5, abs(gt_val) * 0.1)
    if h in ("ticketpricecny", "ticketprice"):
        return 5.0
    if h in ("change30dpct",):
        return max(8.0, abs(gt_val) * 0.3)
    return max(1.0, abs(gt_val) * 0.1)


HEADER_ALIASES = {
    "ticker": ("ticker", "symbol", "stock", "tickersymbol"),
    "company": ("company", "companyname", "name", "compan"),
    "currentprice": ("currentprice", "currentpriceusd", "price", "lastprice", "priceusd", "last"),
    "change30dpct": ("change30dpct", "change30d", "chg30d", "change", "30dchange", "change30dpct"),
    "marketcapb": ("marketcapb", "marketcap", "marketcapbn", "marketcapbillion", "marketcapbillions", "marketcapusd"),
    "revenueb": ("revenueb", "revenue", "totalrevenue", "revenuebn", "revenuebillion", "revenuebillions", "annualrevenue"),
    "eps": ("eps", "earningspershare", "epsttm", "epsusd"),
    "recommendation": ("recommendation", "recommendationkey", "rating", "analystrating", "rec", "consensus"),
}


def _header_map(header):
    """Map normalized header text -> column index."""
    return {norm(h): i for i, h in enumerate(header) if str(h).strip()}


def _agent_col(gt_header_norm, agent_cols):
    """Resolve a GT column header to the agent's column index.

    Exact normalized match first; then, if the GT header is a known alias of a
    column base, look for the base or any of its aliases among the agent headers.
    This makes value-column comparison robust to column reordering and common
    header spellings, while still requiring the column to exist.
    """
    if gt_header_norm in agent_cols:
        return agent_cols[gt_header_norm]
    base = gt_header_norm
    for b, variants in HEADER_ALIASES.items():
        if gt_header_norm in variants:
            base = b
            break
    candidates = (base,) + tuple(HEADER_ALIASES.get(base, ()))
    for cand in candidates:
        if cand in agent_cols:
            return agent_cols[cand]
    return None


def _compare_stock_summary(gt_sheets, agent_sheets, agent_stock_sheet):
    print("\n  [GT] Stock_Summary values")
    gt_stock = _find_sheet(gt_sheets, "Stock_Summary")
    if gt_stock is None:
        record("GT Stock_Summary sheet present", False, "Missing in GT workbook")
        return
    gt_rows = _data_rows(gt_sheets[gt_stock])
    gt_header = _header_of(gt_sheets[gt_stock])
    if agent_stock_sheet is None:
        record("GT Stock_Summary exists in agent", False, "Agent has no Stock_Summary sheet")
        return
    agent_rows = _data_rows(agent_sheets[agent_stock_sheet])
    agent_header = _header_of(agent_sheets[agent_stock_sheet])

    # Column indexes resolved by header name (normalized), not by position, so
    # reordering value columns does not cause cross-column comparisons.
    gt_cols = _header_map(gt_header)
    ag_cols = _header_map(agent_header)
    ticker_col_gt = gt_cols.get("ticker", 1)
    ticker_col_ag = _agent_col("ticker", ag_cols) if ag_cols else 1
    if ticker_col_ag is None:
        ticker_col_ag = 1

    for gt_row in gt_rows:
        gt_ticker = norm(gt_row[ticker_col_gt]) if ticker_col_gt < len(gt_row) else ""
        # find an agent row with the same ticker
        agent_row = None
        for ar in agent_rows:
            if ticker_col_ag < len(ar) and norm(ar[ticker_col_ag]) == gt_ticker:
                agent_row = ar
                break
        if agent_row is None:
            record(f"GT Stock_Summary row for {gt_ticker.upper()} exists in agent",
                   False, f"No agent row with ticker {gt_ticker.upper()}")
            continue

        # compare each GT column, locating the matching agent column by header
        for header, col_idx in gt_cols.items():
            if header == "ticker":
                continue  # already anchored
            gt_val = gt_row[col_idx] if col_idx < len(gt_row) else None
            if gt_val is None:
                continue
            ag_col = _agent_col(header, ag_cols)
            if ag_col is None or ag_col >= len(agent_row):
                record(f"GT Stock_Summary {gt_ticker.upper()} {header} (column)",
                       False, f"Expected {gt_val}, agent sheet lacks column '{header}'")
                continue
            a_val = agent_row[ag_col]
            if header == "company":
                # Free text; the ticker already anchored this row, so this column
                # is informational.
                ok = True
            elif header == "recommendation":
                exp = rec_keyword(gt_val)
                got = rec_keyword(a_val)
                ok = exp is not None and exp == got
                record(f"GT Stock_Summary {gt_ticker.upper()} {header}",
                       ok, f"Expected '{gt_val}', got '{a_val}'")
                continue
            elif header == "change30dpct":
                # Underdetermined by the mock (depends on the 30-day window / clock).
                fa, fb = to_float(gt_val), to_float(a_val)
                ok = fa is not None and fb is not None and abs(fa - fb) <= _tol_for(header, fa)
                record(f"GT Stock_Summary {gt_ticker.upper()} {header}",
                       ok, f"Expected ~{gt_val}, got {a_val}", runtime_only=True)
                continue
            else:
                # Formula cells with no cached value fail here (they are banned by
                # the task's formatting requirements); a cached formula value was
                # already substituted by _load_sheet_values.
                ok = num_close(a_val, gt_val, _tol_for(header, gt_val))
            record(f"GT Stock_Summary {gt_ticker.upper()} {header}",
                   ok, f"Expected {gt_val}, got {a_val}")
    print("  [/GT] Stock_Summary")


def _compare_financial_highlights(gt_sheets, agent_sheets, agent_financial_sheet):
    print("\n  [GT] Financial_Highlights values")
    gt_fin = _find_sheet(gt_sheets, "Financial_Highlights")
    if gt_fin is None:
        record("GT Financial_Highlights sheet present", False, "Missing in GT workbook")
        return
    gt_rows = _data_rows(gt_sheets[gt_fin])
    if agent_financial_sheet is None:
        record("GT Financial_Highlights exists in agent", False, "Agent has no Financial_Highlights sheet")
        return
    agent_rows = _data_rows(agent_sheets[agent_financial_sheet])

    # Build agent lookup keyed by (norm_metric, int(year))
    def to_year(v):
        if v is None:
            return None
        if isinstance(v, (int, float)):
            return int(v)
        if isinstance(v, (_dt.datetime, _dt.date)):
            return v.year
        s = str(v).strip()
        # Full ISO dates: 2025-12-31 / 2025/12/31, and bare '2025'.
        for fmt in ("%Y-%m-%d", "%Y/%m/%d", "%Y"):
            try:
                return _dt.datetime.strptime(s[:10], fmt).year
            except (ValueError, IndexError):
                continue
        # Fiscal/marketing year spellings: 'FY2025', 'FY 2025', '2025年',
        # 'Dec 2025', 'Fiscal 2025'.
        m = re.search(r"(19|20)\d{2}", s)
        if m:
            return int(m.group(0))
        # 2-digit shorthand in a Year column ('25' -> 2025).
        m2 = re.fullmatch(r"\d{2}", s)
        if m2:
            return 2000 + int(m2.group(0))
        try:
            return int(float(s))
        except (ValueError, TypeError):
            return None

    agent_map = {}
    for r in agent_rows:
        if len(r) < 3:
            continue
        key = (norm(r[0]), to_year(r[1]))
        agent_map.setdefault(key, []).append(r)

    for gt_row in gt_rows:
        if len(gt_row) < 3:
            continue
        metric = norm(gt_row[0])
        year = to_year(gt_row[1])
        gt_val = gt_row[2]
        if year is None or gt_val is None:
            continue
        matches = agent_map.get((metric, year), [])
        # norm() already folds 'Net Income' / 'Net_Income' / 'NetIncome' to the same key,
        # so a direct map lookup covers metric-name variants. A formula cell with no
        # cached value fails here (formulas are banned by the task's formatting
        # requirements); a cached formula value was substituted by _load_sheet_values.
        ok = any(
            num_close(r[2], gt_val, _tol_for("value_b", gt_val))
            for r in matches
        )
        record(f"GT Financial_Highlights {gt_row[0]} {year}", ok,
               f"Expected {gt_val}, agent values: {[r[2] for k, rs in agent_map.items() for r in rs if k[1] == year]}")


# ---------------------------------------------------------------- pptx checks

def check_pptx(agent_workspace):
    print("\n=== Check 2: PPT Investor_Roadshow.pptx ===")

    pptx_path = os.path.join(agent_workspace, "Investor_Roadshow.pptx")
    if not os.path.exists(pptx_path):
        record("Investor_Roadshow.pptx exists", False, f"Not found at {pptx_path}")
        return
    record("Investor_Roadshow.pptx exists", True)

    try:
        from pptx import Presentation
        prs = Presentation(pptx_path)
        slide_count = len(prs.slides)
        record("PPT has >= 5 slides", slide_count >= 5, f"Found {slide_count} slides")

        all_text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    all_text += " " + shape.text
        all_text_lower = all_text.lower()

        has_roadshow = any(kw in all_text_lower for kw in
                           ["roadshow", "investor", "q1 2026", "financial", "路演", "投资"])
        record("PPT contains roadshow/investor/financial content", has_roadshow,
               f"Text sample: {all_text[:200]}")

        has_train = any(kw in all_text_lower for kw in
                        ["g11", "g105", "shanghai", "guangzhou", "train", "travel",
                         "上海", "广州", "旅行", "行程"])
        record("PPT mentions travel itinerary", has_train,
               f"Text sample: {all_text[:200]}")

        has_finance = any(kw in all_text_lower for kw in
                          ["googl", "amzn", "alphabet", "amazon", "revenue", "earnings", "stock", "price"])
        record("PPT mentions financial data", has_finance,
               f"Text sample: {all_text[:200]}")

    except ImportError:
        record("python-pptx available", False, "python-pptx not installed")
    except Exception as e:
        record("PPT readable", False, str(e))


# ---------------------------------------------------------------- email checks

def check_emails():
    print("\n=== Check 3: Emails sent ===")

    try:
        conn = psycopg2.connect(**DB_CONFIG)
    except Exception as e:
        record("Email database reachable", False, f"DB connect failed: {e}", runtime_only=True)
        print("  (db checks skipped - no database available on this host)")
        return

    try:
        cur = conn.cursor()
        cur.execute("SELECT subject, from_addr, to_addr, body_text FROM email.messages")
        messages = cur.fetchall()
        cur.close()
        conn.close()
    except Exception as e:
        record("Email table query", False, str(e), runtime_only=True)
        print("  (db checks skipped - query failed)")
        return

    def to_addresses(to_addr):
        if isinstance(to_addr, list):
            return " ".join(str(r).lower() for r in to_addr)
        elif to_addr:
            try:
                parsed = json.loads(str(to_addr))
                return " ".join(str(r).lower() for r in parsed) if isinstance(parsed, list) else str(to_addr).lower()
            except Exception:
                return str(to_addr).lower()
        return ""

    to_fundmanager = [m for m in messages if "investors@fundmanager.com" in to_addresses(m[2])]
    to_shanghai = [m for m in messages if "shanghai_partners@finance.com" in to_addresses(m[2])]

    record("Email sent to investors@fundmanager.com", len(to_fundmanager) >= 1,
           f"Total messages: {len(messages)}", runtime_only=True)
    record("Email sent to shanghai_partners@finance.com", len(to_shanghai) >= 1,
           f"Total messages: {len(messages)}", runtime_only=True)

    if to_fundmanager:
        subj, _, _, body = to_fundmanager[0]
        content = ((subj or "") + " " + (body or "")).lower()
        has_finance = any(kw in content for kw in ["roadshow", "financial", "presentation", "schedule", "路演"])
        record("Fundmanager email mentions roadshow/financial content", has_finance,
               f"Subject: {subj}", runtime_only=True)

    # Shanghai email body must include the roadshow arrival date (March 10 / 2026-03-10)
    if to_shanghai:
        subj_s, _, _, body_s = to_shanghai[0]
        content_s = ((subj_s or "") + " " + (body_s or "")).lower()
        has_march_10 = (
            "march 10" in content_s
            or "2026-03-10" in content_s
            or "mar 10" in content_s
            or "3/10" in content_s
            or "3月10" in content_s
        )
        record("Shanghai email mentions arrival date March 10", has_march_10,
               f"Subject: {subj_s}", runtime_only=True)


# ---------------------------------------------------------------- main

def main():
    parser = ArgumentParser()
    parser.add_argument("--agent_workspace", required=False, default=".")
    parser.add_argument("--groundtruth_workspace", required=False, default=".")
    parser.add_argument("--launch_time", required=False)
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()

    check_excel(args.agent_workspace, args.groundtruth_workspace)
    check_pptx(args.agent_workspace)
    check_emails()

    total = PASS_COUNT + FAIL_COUNT
    if total == 0:
        print("\nFAIL: No checks were performed.")
        sys.exit(1)

    accuracy = PASS_COUNT / total * 100
    print(f"\nOverall: {PASS_COUNT}/{total} checks passed ({accuracy:.1f}%) ; blocking_fail={BLOCKING_FAIL_COUNT}")

    result = {
        "total_passed": PASS_COUNT,
        "total_checks": total,
        "accuracy": accuracy,
        "blocking_fail_count": BLOCKING_FAIL_COUNT,
    }

    if args.res_log_file:
        with open(args.res_log_file, "w") as f:
            json.dump(result, f, indent=2)

    # Blocking-fail gate: any non-runtime (i.e. local file) failure is fatal.
    # Runtime-only (email / Change_30D_Pct) checks are still counted but excluded from the gate.
    if BLOCKING_FAIL_COUNT == 0:
        print("PASS")
        sys.exit(0)
    else:
        print("FAIL")
        sys.exit(1)


if __name__ == "__main__":
    main()
