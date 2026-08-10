"""Evaluation for sf-workforce-diversity-review.

Reads the agent's ``Diversity_Analysis.xlsx`` and ``DEI_Board_Report.pptx`` and
compares them against the groundtruth workspace.

Robustness design (a model that follows the task always passes):
  - Workbooks are read with ``data_only=False`` so literal values AND formulas
    are both visible. A numeric cell that holds a formula is evaluated
    best-effort (simple arithmetic + SUM/AVERAGE/ROUND/... including
    whole-column ranges); if a formula cannot be evaluated the numeric check
    FAILS rather than being silently skipped, so the arithmetic checks stay
    binding (the task explicitly asks for literal values, see docs/task.md).
  - Number parsing tolerates '%', thousand separators, currency symbols and
    whitespace, and a 0-1 vs 0-100 percent-scale difference.
  - Labels (departments, age groups, education levels, metric names) are matched
    through a normalizer so phrasing variants ("30 to 39" vs "30-39",
    "50 plus" vs "50+", "Bachelor's" vs "Bachelors") do not cause false
    negatives.
"""
import argparse
import os
import re
import sys

import openpyxl

# ---------------------------------------------------------------------------
# parsing helpers
# ---------------------------------------------------------------------------


def _to_float(v):
    """Robustly convert a cell value to float (or None).

    Handles int/float and numeric strings with thousand separators, currency
    symbols, percent signs and whitespace. Returns None for None, unparseable
    strings and formula strings (formulas are handled by the caller).
    """
    if v is None:
        return None
    if isinstance(v, bool):
        return float(v)
    if isinstance(v, (int, float)):
        return float(v)
    s = str(v).strip()
    if not s or s.startswith("="):
        return None
    s = (s.replace(",", "").replace("$", "").replace("¥", "")
         .replace("€", "").replace("%", "").replace(" ", "")
         .replace("USD", "").replace("usd", ""))
    # strip a trailing unit word ("34.7 years" -> "34.7", "8.3yrs" -> "8.3")
    s = re.sub(r"[a-zA-Z]+$", "", s)
    try:
        return float(s)
    except ValueError:
        return None


def str_match(a, b):
    if a is None or b is None:
        return a is None and b is None
    return str(a).strip().lower() == str(b).strip().lower()


def _norm_label(s):
    """Canonical key for labels/categories.

    Maps phrasing/punctuation variants ("30 to 39" -> "30-39", "50 plus" ->
    "50+", "Under 30" -> "under30", "Marital Status" -> "marital",
    "Total Employees" -> "totalemployees") so that equivalent labels match.
    It never merges semantically distinct categories.
    """
    s = str(s).strip().lower()
    # unicode dashes (en/em dash, minus sign) -> ASCII hyphen
    s = s.replace("–", "-").replace("—", "-").replace("−", "-")
    # department aliases so semantically identical names match ('R&D' vs
    # 'Research & Development', 'HR' vs 'Human Resources')
    s = s.replace("research & development", "rd")
    s = s.replace("research and development", "rd")
    s = s.replace("r & d", "rd").replace("r&d", "rd")
    s = s.replace("human resources", "hr")
    # age-group phrasing variants ('Over 50'/'Aged 30-39')
    s = s.replace("over 50", "50+").replace("above 50", "50+")
    s = s.replace("aged ", "")
    s = s.replace(" to ", "-")
    s = s.replace(" and over", "+").replace(" and above", "+")
    s = s.replace(" or over", "+").replace(" or older", "+").replace(" or more", "+")
    s = s.replace("50 plus", "50+").replace("50plus", "50+")
    s = s.replace("under 30", "under30")
    s = s.replace("< 30", "under30").replace("<30", "under30")
    s = s.replace("'", "").replace(".", "")
    s = s.replace("degrees", "").replace("degree", "")
    s = s.replace("education level", "education")
    s = s.replace("education distribution", "education")
    s = s.replace("age group", "age")
    s = s.replace("age distribution", "age")
    s = s.replace("marital status", "marital")
    s = s.replace("marital distribution", "marital")
    s = s.replace("_", " ").replace("-", " ")
    s = s.replace(" years", "").replace(" year", "").replace(" yrs", "")
    s = s.replace(" ", "")
    return s


def _norm_edu(s):
    """Canonical key for education-level labels (also strips trailing plural s)."""
    k = _norm_label(s)
    if len(k) > 3 and k.endswith("s") and not k.endswith("us"):
        k = k[:-1]
    return k


# ---------------------------------------------------------------------------
# lightweight Excel-formula evaluation (best-effort; failure yields None)
# ---------------------------------------------------------------------------

_TOK = re.compile(
    r"""\s*(?:(?P<num>\d+(?:\.\d+)?(?:[eE][+-]?\d+)?)
    |(?P<cell>\$?[A-Za-z]{1,3}\$?\d{1,7})
    |(?P<func>[A-Za-z_][A-Za-z0-9_.]*)
    |(?P<op>[+\-*/^(),:]))""",
    re.VERBOSE,
)


def _tokenize_formula(s):
    toks = []
    i = 0
    n = len(s)
    while i < n:
        m = _TOK.match(s, i)
        if not m:
            return None
        i = m.end()
        toks.append((m.lastgroup, m.group()))
    return toks


def _expand_column_ranges(s, max_row):
    """Expand whole-column references (``C:C``, ``$C:$C``) to ``C1:C<max_row>``.

    Excel allows ``SUM(C:C)``; our tokenizer only understands explicit ranges,
    so whole-column refs are rewritten before tokenizing. Cell ranges such as
    ``C2:C6`` are left untouched (they contain row digits).
    """
    def _rep(m):
        col = m.group(1).replace("$", "")
        return f"{col}1:{col}{max_row}"
    return re.sub(r"(\$?[A-Za-z]{1,3})\s*:\s*\1", _rep, s)


def _eval_formula(formula, get_cell, depth=0, max_row=10000):
    """Evaluate a simple Excel formula to a float, or None if unsupported.

    Supports numeric literals, cell references (relative/absolute), parentheses,
    + - * / ^, and the functions SUM/AVERAGE/ROUND/MAX/MIN/COUNT over values or
    ranges. ``get_cell(row, col)`` resolves a referenced cell to float/None.
    Anything outside that set (IF/VLOOKUP/string operands ...) yields None; the
    caller then treats the cell as unparseable (an error, not a silent pass).
    """
    if depth > 12:
        return None
    s = formula.strip()
    if s.startswith("="):
        s = s[1:]
    s = _expand_column_ranges(s, max_row)
    toks = _tokenize_formula(s)
    if not toks:
        return None
    pos = [0]
    val = _parse_addsub(toks, pos, get_cell, depth)
    if val is None or pos[0] != len(toks):
        return None
    return val


def _parse_addsub(toks, pos, get_cell, depth):
    v = _parse_muldiv(toks, pos, get_cell, depth)
    while v is not None and pos[0] < len(toks):
        kind, op = toks[pos[0]]
        if kind == "op" and op in ("+", "-"):
            pos[0] += 1
            r = _parse_muldiv(toks, pos, get_cell, depth)
            if r is None:
                return None
            v = v + r if op == "+" else v - r
        else:
            break
    return v


def _parse_muldiv(toks, pos, get_cell, depth):
    v = _parse_unary(toks, pos, get_cell, depth)
    while v is not None and pos[0] < len(toks):
        kind, op = toks[pos[0]]
        if kind == "op" and op in ("*", "/"):
            pos[0] += 1
            r = _parse_unary(toks, pos, get_cell, depth)
            if r is None:
                return None
            if op == "*":
                v = v * r
            else:
                if r == 0:
                    return None
                v = v / r
        else:
            break
    return v


def _parse_unary(toks, pos, get_cell, depth):
    if pos[0] < len(toks) and toks[pos[0]][0] == "op" and toks[pos[0]][1] == "-":
        pos[0] += 1
        v = _parse_unary(toks, pos, get_cell, depth)
        return None if v is None else -v
    return _parse_primary(toks, pos, get_cell, depth)


def _parse_primary(toks, pos, get_cell, depth):
    if pos[0] >= len(toks):
        return None
    kind, val = toks[pos[0]]
    if kind == "num":
        pos[0] += 1
        return float(val)
    if kind == "cell":
        pos[0] += 1
        return _resolve_ref(val, get_cell, depth)
    if kind == "op" and val == "(":
        pos[0] += 1
        inner = _parse_addsub(toks, pos, get_cell, depth)
        if inner is None or pos[0] >= len(toks) or toks[pos[0]] != ("op", ")"):
            return None
        pos[0] += 1
        return inner
    if kind == "func":
        pos[0] += 1
        if pos[0] >= len(toks) or toks[pos[0]] != ("op", "("):
            return None
        pos[0] += 1
        args = _parse_args(toks, pos, get_cell, depth)
        if args is None or pos[0] >= len(toks) or toks[pos[0]] != ("op", ")"):
            return None
        pos[0] += 1
        return _apply_func(val.upper(), args)
    return None


def _parse_args(toks, pos, get_cell, depth):
    args = []
    while True:
        if pos[0] >= len(toks):
            return None
        if toks[pos[0]] == ("op", ")"):
            break
        # top-level range like C2:C6
        if (pos[0] + 2 < len(toks) and toks[pos[0]][0] == "cell"
                and toks[pos[0] + 1] == ("op", ":") and toks[pos[0] + 2][0] == "cell"):
            rng = _resolve_range(toks[pos[0]][1], toks[pos[0] + 2][1], get_cell, depth)
            pos[0] += 3
            if rng is None:
                return None
            args.append(("range", rng))
        else:
            v = _parse_addsub(toks, pos, get_cell, depth)
            if v is None:
                return None
            args.append(("num", v))
        if pos[0] < len(toks) and toks[pos[0]] == ("op", ","):
            pos[0] += 1
            continue
        break
    return args


def _cell_key(ref):
    ref = ref.replace("$", "")
    m = re.match(r"^([A-Za-z]{1,3})(\d{1,7})$", ref)
    if not m:
        return None
    col = 0
    for ch in m.group(1):
        col = col * 26 + (ord(ch.upper()) - 64)
    return int(m.group(2)), col


def _resolve_ref(ref, get_cell, depth):
    k = _cell_key(ref)
    if k is None:
        return None
    row, col = k
    return get_cell(row, col, depth)


def _resolve_range(ref1, ref2, get_cell, depth):
    k1, k2 = _cell_key(ref1), _cell_key(ref2)
    if k1 is None or k2 is None:
        return None
    vals = []
    r1, c1 = k1
    r2, c2 = k2
    for r in range(min(r1, r2), max(r1, r2) + 1):
        for c in range(min(c1, c2), max(c1, c2) + 1):
            v = get_cell(r, c, depth)
            if v is not None:
                vals.append(v)
    return vals


def _apply_func(name, args):
    if name == "SUM":
        tot = 0.0
        for kind, v in args:
            if kind == "num":
                if v is None:
                    return None
                tot += v
            else:
                tot += sum(v)
        return tot
    if name in ("AVERAGE", "AVERAGEA"):
        vals = []
        for kind, v in args:
            if kind == "num":
                if v is None:
                    return None
                vals.append(v)
            else:
                vals.extend(v)
        return sum(vals) / len(vals) if vals else None
    if name == "ROUND":
        if len(args) != 2 or args[0][1] is None or args[1][1] is None:
            return None
        return round(args[0][1], int(args[1][1]))
    if name in ("MAX", "MIN", "COUNT", "COUNTA"):
        vals = []
        for kind, v in args:
            if kind == "num":
                if v is None:
                    return None
                vals.append(v)
            else:
                vals.extend(v)
        if name == "COUNT":
            return float(len(vals))
        return max(vals) if name == "MAX" else (min(vals) if vals else None)
    return None


# ---------------------------------------------------------------------------
# workbook access
# ---------------------------------------------------------------------------


class _Cell(object):
    __slots__ = ("raw", "cached", "num", "is_formula", "unparsed")

    def __init__(self, raw, cached, num, is_formula, unparsed):
        self.raw = raw          # value from data_only=False view
        self.cached = cached    # value from data_only=True view
        self.num = num          # resolved float or None
        self.is_formula = is_formula
        self.unparsed = unparsed  # formula that could not be resolved


class WorkbookPair(object):
    """Two views of a workbook: data_only=False (raw/formulas) and data_only=True (cached)."""

    def __init__(self, wb_vals, wb_cached):
        self.wb_vals = wb_vals
        self.wb_cached = wb_cached

    def sheet_name(self, wanted):
        for n in self.wb_vals.sheetnames:
            if n.strip().lower() == wanted.strip().lower():
                return n
        return None

    def _ref_num(self, sheet, row, col, depth):
        cell = self._cell(sheet, row, col, depth)
        return cell.num

    def _cell(self, sheet, row, col, depth=0):
        raw = self.wb_vals[sheet].cell(row=row, column=col).value
        cached = self.wb_cached[sheet].cell(row=row, column=col).value
        if isinstance(raw, str) and raw.strip().startswith("="):
            num = _to_float(cached)
            if num is None:
                def _g(rr, cc, _d=depth):
                    return self._ref_num(sheet, rr, cc, _d + 1)
                num = _eval_formula(raw, _g, depth,
                                    max_row=self.wb_vals[sheet].max_row)
            return _Cell(raw, cached, num, True, num is None)
        return _Cell(raw, cached, _to_float(raw), False, False)

    def rows(self, sheet):
        ws = self.wb_vals[sheet]
        out = []
        for r in ws.iter_rows():
            out.append([self._cell(sheet, c.row, c.column) for c in r])
        return out


def load_pair(path):
    wb_vals = openpyxl.load_workbook(path, data_only=False)
    wb_cached = openpyxl.load_workbook(path, data_only=True)
    return WorkbookPair(wb_vals, wb_cached)


def _find_file(ws_dir, name):
    """Find a file in ws_dir, case-insensitively tolerant."""
    p = os.path.join(ws_dir, name)
    if os.path.isfile(p):
        return p
    if os.path.isdir(ws_dir):
        for fn in os.listdir(ws_dir):
            if fn.lower() == name.lower() and os.path.isfile(os.path.join(ws_dir, fn)):
                return os.path.join(ws_dir, fn)
    return p


def _looks_header(row):
    if not row or row[0] is None or row[0].raw is None:
        return False
    k = _norm_label(row[0].raw)
    if k in ("department", "dept", "metric", "college", "position", "title"):
        return True
    if row[0].num is None and row[1] is not None and row[1].num is None and row[1].raw is not None:
        k2 = _norm_label(row[1].raw)
        if k2 in ("educationlevel", "agegroup", "category", "educationlevels"):
            return True
    return False


def _data_rows(rows):
    """Return the data rows, skipping a leading header/title row."""
    if len(rows) > 1 and _looks_header(rows[0]):
        rows = rows[1:]
    if rows and _looks_header(rows[0]):
        rows = rows[1:]
    return rows


def num_close_cell(ac, gc, tol):
    """Compare two _Cell objects. Returns True when the agent's numeric cell is
    within tolerance of the groundtruth. A formula the evaluator could not
    resolve leaves ``ac.num`` as None; that is treated as an error, never a
    silent pass, so the numeric checks stay binding (the task explicitly asks
    for literal values).
    """
    a, b = ac.num, gc.num
    if a is not None and b is not None:
        if abs(a - b) <= tol:
            return True
        # 0-1 fraction vs 0-100 percent scale tolerance
        if tol > 0 and a > 0 and b > 0:
            lo, hi = sorted((a, b))
            if lo <= 1.5 and abs(hi / lo - 100.0) < 1e-6:
                return True
        return False
    if a is None and b is None:
        return True
    if b is None:
        # groundtruth cell is not numeric; compare as strings
        return str_match(ac.raw, gc.raw)
    # agent cell is empty or holds an unresolvable formula while the groundtruth
    # is numeric -> error (never silently skip the arithmetic check)
    return False


_STATUS_BOUND = 3.0
_STATUS_EPS = 0.25


def status_ok(ac_int, ac_bench, ac_status, gc_status):
    """Validate an Overall Scorecard Status cell.

    The classification boundary is ±3.0 points. The groundtruth rounds the
    internal percentages to one decimal (Married 55.014 -> 55.0, gap 3.0 ->
    Aligned); a model that keeps full precision computes a gap of 3.014 and,
    following the task's literal rule, writes Above/Below. Both are correct, so
    besides an exact match with the groundtruth status we accept any status that
    is consistent with the agent's own Internal_Value / Benchmark_Value within a
    small boundary tolerance (the task does not specify how many decimals the
    gap is rounded to).
    """
    a_raw = ac_status.raw if ac_status is not None else None
    g_raw = gc_status.raw if gc_status is not None else None
    a = str(a_raw).strip().lower() if a_raw is not None else None
    g = str(g_raw).strip().lower() if g_raw is not None else None
    if a is None or g is None:
        return False
    if a == g:
        return True
    if (ac_int is not None and ac_int.num is not None
            and ac_bench is not None and ac_bench.num is not None):
        gap = ac_int.num - ac_bench.num
        if gap > _STATUS_BOUND + _STATUS_EPS:
            return a == "above"
        if gap < -_STATUS_BOUND - _STATUS_EPS:
            return a == "below"
        # Boundary band (|gap| within ~0.25 of 3.0): Aligned is always
        # defensible here. A directional status is also defensible when the
        # stored gap is at or just past the 3.0 boundary (the task does not say
        # how many decimals to keep: full precision gives Married 55.014 ->
        # gap 3.014 -> Above, while the groundtruth rounds to 55.0 -> gap 3.0 ->
        # Aligned; both are correct).
        if a == "aligned":
            return True
        if gap >= _STATUS_BOUND:
            return a == "above"
        if gap <= -_STATUS_BOUND:
            return a == "below"
    return False


# ---------------------------------------------------------------------------
# main
# ---------------------------------------------------------------------------


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--agent_workspace", required=False)
    parser.add_argument("--groundtruth_workspace", required=False)
    parser.add_argument("--launch_time", required=False)
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()

    task_root = os.path.abspath(os.path.join(os.path.dirname(__file__), ".."))
    gt_dir = args.groundtruth_workspace or os.path.join(task_root, "groundtruth_workspace")
    agent_ws = args.agent_workspace or gt_dir

    agent_file = _find_file(agent_ws, "Diversity_Analysis.xlsx")
    gt_file = os.path.join(gt_dir, "Diversity_Analysis.xlsx")

    if not os.path.exists(agent_file):
        print(f"FAIL: Agent output not found: {agent_file}")
        sys.exit(1)
    if not os.path.exists(gt_file):
        print(f"FAIL: Groundtruth not found: {gt_file}")
        sys.exit(1)

    all_errors = []
    warnings = []

    # ---------------------------------------------------------------
    # Education Distribution
    # ---------------------------------------------------------------
    print("  Checking Education Distribution...")
    a_pair = load_pair(agent_file)
    g_pair = load_pair(gt_file)

    a_name = a_pair.sheet_name("Education Distribution")
    g_name = g_pair.sheet_name("Education Distribution")
    if a_name is None:
        all_errors.append("Sheet 'Education Distribution' not found")
    elif g_name is None:
        all_errors.append("Sheet 'Education Distribution' not found in groundtruth")
    else:
        a_data = _data_rows(a_pair.rows(a_name))
        g_data = _data_rows(g_pair.rows(g_name))
        a_lookup = {}
        for r in a_data:
            if len(r) >= 2 and r[0].raw is not None and r[1].raw is not None:
                k = f"{_norm_label(r[0].raw)}|{_norm_edu(r[1].raw)}"
                a_lookup[k] = r
        errors = []
        for g_row in g_data:
            if not g_row or g_row[0].raw is None:
                continue
            key = f"{_norm_label(g_row[0].raw)}|{_norm_edu(g_row[1].raw)}"
            a_row = a_lookup.get(key)
            if a_row is None:
                errors.append(f"Missing: {g_row[0].raw}|{g_row[1].raw}")
                continue
            if len(a_row) > 2 and len(g_row) > 2:
                if not num_close_cell(a_row[2], g_row[2], 0):
                    errors.append(f"{key}.Count: {a_row[2].raw} vs {g_row[2].raw}")
            if len(a_row) > 3 and len(g_row) > 3:
                if not num_close_cell(a_row[3], g_row[3], 0.5):
                    errors.append(f"{key}.Internal_Pct: {a_row[3].raw} vs {g_row[3].raw}")
            if len(a_row) > 4 and len(g_row) > 4:
                if not num_close_cell(a_row[4], g_row[4], 0.5):
                    errors.append(f"{key}.Benchmark_Pct: {a_row[4].raw} vs {g_row[4].raw}")
            if len(a_row) > 5 and len(g_row) > 5:
                if not num_close_cell(a_row[5], g_row[5], 0.5):
                    errors.append(f"{key}.Gap: {a_row[5].raw} vs {g_row[5].raw}")
        if errors:
            all_errors.extend(errors)
            print(f"    ERRORS: {len(errors)}")
            for e in errors[:5]:
                print(f"      {e}")
        else:
            print("    PASS")

    # ---------------------------------------------------------------
    # Age Distribution
    # ---------------------------------------------------------------
    print("  Checking Age Distribution...")
    a_name = a_pair.sheet_name("Age Distribution")
    g_name = g_pair.sheet_name("Age Distribution")
    if a_name is None:
        all_errors.append("Sheet 'Age Distribution' not found")
    elif g_name is None:
        all_errors.append("Sheet 'Age Distribution' not found in groundtruth")
    else:
        a_data = _data_rows(a_pair.rows(a_name))
        g_data = _data_rows(g_pair.rows(g_name))
        a_lookup = {}
        for r in a_data:
            if len(r) >= 2 and r[0].raw is not None and r[1].raw is not None:
                k = f"{_norm_label(r[0].raw)}|{_norm_label(r[1].raw)}"
                a_lookup[k] = r
        errors = []
        for g_row in g_data:
            if not g_row or g_row[0].raw is None:
                continue
            key = f"{_norm_label(g_row[0].raw)}|{_norm_label(g_row[1].raw)}"
            a_row = a_lookup.get(key)
            if a_row is None:
                errors.append(f"Missing: {g_row[0].raw}|{g_row[1].raw}")
                continue
            if len(a_row) > 2 and len(g_row) > 2:
                if not num_close_cell(a_row[2], g_row[2], 0):
                    errors.append(f"{key}.Count: {a_row[2].raw} vs {g_row[2].raw}")
            if len(a_row) > 3 and len(g_row) > 3:
                if not num_close_cell(a_row[3], g_row[3], 0.5):
                    errors.append(f"{key}.Internal_Pct: {a_row[3].raw} vs {g_row[3].raw}")
            if len(a_row) > 4 and len(g_row) > 4:
                if not num_close_cell(a_row[4], g_row[4], 0.5):
                    errors.append(f"{key}.Benchmark_Pct: {a_row[4].raw} vs {g_row[4].raw}")
            if len(a_row) > 5 and len(g_row) > 5:
                if not num_close_cell(a_row[5], g_row[5], 0.5):
                    errors.append(f"{key}.Gap: {a_row[5].raw} vs {g_row[5].raw}")
        if errors:
            all_errors.extend(errors)
            print(f"    ERRORS: {len(errors)}")
            for e in errors[:5]:
                print(f"      {e}")
        else:
            print("    PASS")

    # ---------------------------------------------------------------
    # Overall Scorecard
    # ---------------------------------------------------------------
    print("  Checking Overall Scorecard...")
    a_name = a_pair.sheet_name("Overall Scorecard")
    g_name = g_pair.sheet_name("Overall Scorecard")
    if a_name is None:
        all_errors.append("Sheet 'Overall Scorecard' not found")
    elif g_name is None:
        all_errors.append("Sheet 'Overall Scorecard' not found in groundtruth")
    else:
        a_data = _data_rows(a_pair.rows(a_name))
        g_data = _data_rows(g_pair.rows(g_name))
        a_lookup = {}
        for r in a_data:
            if len(r) >= 2 and r[0].raw is not None and r[1].raw is not None:
                k = f"{_norm_label(r[0].raw)}|{_norm_edu(r[1].raw)}"
                a_lookup[k] = r
        errors = []
        for g_row in g_data:
            if not g_row or g_row[0].raw is None:
                continue
            key = f"{_norm_label(g_row[0].raw)}|{_norm_edu(g_row[1].raw)}"
            a_row = a_lookup.get(key)
            if a_row is None:
                errors.append(f"Missing: {g_row[0].raw}|{g_row[1].raw}")
                continue
            if len(a_row) > 2 and len(g_row) > 2:
                if not num_close_cell(a_row[2], g_row[2], 0.5):
                    errors.append(f"{key}.Internal_Value: {a_row[2].raw} vs {g_row[2].raw}")
            if len(a_row) > 3 and len(g_row) > 3:
                if not num_close_cell(a_row[3], g_row[3], 0.5):
                    errors.append(f"{key}.Benchmark_Value: {a_row[3].raw} vs {g_row[3].raw}")
            if len(a_row) > 4 and len(g_row) > 4:
                if not num_close_cell(a_row[4], g_row[4], 0.5):
                    errors.append(f"{key}.Gap: {a_row[4].raw} vs {g_row[4].raw}")
            if len(a_row) > 5 and len(g_row) > 5:
                if not status_ok(a_row[2], a_row[3], a_row[5], g_row[5]):
                    errors.append(f"{key}.Status: '{a_row[5].raw}' vs '{g_row[5].raw}'")
        if errors:
            all_errors.extend(errors)
            print(f"    ERRORS: {len(errors)}")
            for e in errors[:5]:
                print(f"      {e}")
        else:
            print("    PASS")

    # ---------------------------------------------------------------
    # Summary
    # ---------------------------------------------------------------
    print("  Checking Summary...")
    a_name = a_pair.sheet_name("Summary")
    g_name = g_pair.sheet_name("Summary")
    if a_name is None:
        all_errors.append("Sheet 'Summary' not found")
    elif g_name is None:
        all_errors.append("Sheet 'Summary' not found in groundtruth")
    else:
        a_data = _data_rows(a_pair.rows(a_name))
        g_data = _data_rows(g_pair.rows(g_name))
        a_lookup = {}
        for r in a_data:
            if r and r[0].raw is not None and str(r[0].raw).strip():
                a_lookup[_norm_label(r[0].raw)] = r
        errors = []
        for g_row in g_data:
            if not g_row or g_row[0].raw is None:
                continue
            key = _norm_label(g_row[0].raw)
            a_row = a_lookup.get(key)
            if a_row is None:
                errors.append(f"Missing: {g_row[0].raw}")
                continue
            if len(a_row) > 1 and len(g_row) > 1:
                if key in ("totalemployees", "totaldepartments", "educationlevels"):
                    if not num_close_cell(a_row[1], g_row[1], 0):
                        errors.append(f"{key}: {a_row[1].raw} vs {g_row[1].raw}")
                else:
                    if not num_close_cell(a_row[1], g_row[1], 0.2):
                        errors.append(f"{key}: {a_row[1].raw} vs {g_row[1].raw}")
        if errors:
            all_errors.extend(errors)
            print(f"    ERRORS: {len(errors)}")
            for e in errors[:5]:
                print(f"      {e}")
        else:
            print("    PASS")

    # ---------------------------------------------------------------
    # PowerPoint
    # ---------------------------------------------------------------
    print("  Checking DEI_Board_Report.pptx...")
    pptx_file = _find_file(agent_ws, "DEI_Board_Report.pptx")
    if not os.path.exists(pptx_file):
        all_errors.append("DEI_Board_Report.pptx not found")
        print("    FAIL: file not found")
    else:
        try:
            from pptx import Presentation
        except ImportError as e:
            warnings.append(f"PPTX check skipped (python-pptx unavailable): {e}")
            print("    SKIPPED: python-pptx not installed")
            pptx_file = None

        if pptx_file and os.path.exists(pptx_file):
            try:
                prs = Presentation(pptx_file)
                slide_count = len(prs.slides)
                if slide_count < 5:
                    all_errors.append(f"PowerPoint has only {slide_count} slides, expected 5+")
                    print(f"    FAIL: only {slide_count} slides")

                slide_texts = []
                slide_line_counts = []
                slide_titles = []
                for slide in prs.slides:
                    t = ""
                    nlines = 0
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            t += shape.text_frame.text + " "
                            nlines += len(shape.text_frame.paragraphs)
                    slide_texts.append(t.lower())
                    slide_line_counts.append(nlines)
                    # Slide title: use the title placeholder when present,
                    # otherwise fall back to the first text shape's first line.
                    title = ""
                    if slide.shapes.title is not None:
                        try:
                            title = slide.shapes.title.text_frame.text.strip()
                        except Exception:
                            title = ""
                    if not title:
                        for shape in slide.shapes:
                            if shape.has_text_frame and shape.text_frame.text.strip():
                                title = shape.text_frame.text.split("\n")[0].strip()
                                break
                    slide_titles.append(title.lower())
                full = " ".join(slide_texts)

                # Detect the recommendations/actions slide. A correct model may
                # title it "Recommendations", "Action Plan", "Next Steps",
                # "Improving Diversity", "Key Recommendations", etc., so we match
                # a broad set of recommendation vocabulary. The slide TITLE is the
                # primary signal; the body text is only a fallback so a slide
                # whose title is generic but whose content is clearly
                # recommendations is still detected. The candidate with the most
                # action items wins, so an earlier sparser slide cannot be
                # mis-identified as the recommendations slide.
                rec_keywords = (
                    "recommendation", "recommend", "action item", "action plan",
                    "action", "next step", "improve", "improving", "improvement",
                    "suggest", "suggestion", "initiative", "opportunit",
                    "priority", "way forward", "going forward",
                )
                rec_idx = None
                rec_bullets = -1
                candidates = list(zip(slide_titles, slide_texts, slide_line_counts))
                for i, (title, st, nlines) in enumerate(candidates):
                    if any(k in title for k in rec_keywords):
                        lines = [ln.strip() for ln in st.split("\n") if ln.strip()]
                        bullets = max(st.count("•"), st.count("-"), len(lines), nlines)
                        if bullets > rec_bullets:
                            rec_bullets = bullets
                            rec_idx = i
                if rec_idx is None:
                    for i, (title, st, nlines) in enumerate(candidates):
                        if any(k in st for k in rec_keywords):
                            lines = [ln.strip() for ln in st.split("\n") if ln.strip()]
                            bullets = max(st.count("•"), st.count("-"), len(lines), nlines)
                            if bullets > rec_bullets:
                                rec_bullets = bullets
                                rec_idx = i

                required_topics = {
                    "education": "education",
                    "age": "age",
                    # The recommendations/actions slide may legitimately be titled
                    # "Recommendations", "Action Plan", "Next Steps",
                    # "Improving Diversity", "Key Initiatives", etc., and its
                    # bullets are often verb phrases that contain none of those
                    # words, so the topic is detected from the slide title/text
                    # with a broad keyword set (not just 'recommendation').
                    "recommendations or actions": rec_idx is not None,
                }
                for topic, present in required_topics.items():
                    if isinstance(present, bool):
                        if not present:
                            all_errors.append(f"PPTX missing topic '{topic}'")
                    else:
                        if present not in full:
                            all_errors.append(f"PPTX missing topic '{topic}'")

                # Recommendations slide should carry at least 4 action items.
                # Count text-frame paragraphs (robust to bullets-as-shapes) and
                # fall back to newline/bullet-char counting. If several slides
                # match the vocabulary, use the one with the most action items so
                # an earlier, sparser slide cannot be mis-identified as the
                # recommendations slide.
                if rec_idx is not None and rec_bullets < 4:
                    all_errors.append(
                        f"Recommendations slide has only {rec_bullets} action items (need >=4)")

                if len(full.strip()) < 200:
                    all_errors.append(f"PPTX total content too sparse ({len(full)} chars)")

                if not any("DEI_Board_Report" in e or "PPTX" in e for e in all_errors):
                    print("    PASS")
                else:
                    print("    Some PPTX checks failed.")
            except Exception as e:
                all_errors.append(f"PPTX read error: {e}")
                print(f"    ERROR: {e}")

    if warnings:
        for w in warnings:
            print(f"  (note: {w})")

    if all_errors:
        print(f"\n=== RESULT: FAIL ({len(all_errors)} errors) ===")
        for e in all_errors[:10]:
            print(f"  {e}")
        sys.exit(1)
    else:
        print("\n=== RESULT: PASS ===")
        sys.exit(0)


if __name__ == "__main__":
    main()
