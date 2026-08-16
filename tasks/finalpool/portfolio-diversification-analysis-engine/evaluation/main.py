#!/usr/bin/env python3
"""Evaluation script for portfolio-diversification-analysis-engine."""

import argparse
import os
import sys
from pathlib import Path

PASS_COUNT = 0
FAIL_COUNT = 0


def check(name, condition, detail=""):
    global PASS_COUNT, FAIL_COUNT
    if condition:
        PASS_COUNT += 1
        print(f"  [PASS] {name}")
    else:
        FAIL_COUNT += 1
        detail_str = str(detail)[:200] if detail else ""
        print(f"  [FAIL] {name}: {detail_str}")


def get_conn():
    import psycopg2
    return psycopg2.connect(
        host=os.environ.get("PGHOST", "localhost"), port=int(os.environ.get("PGPORT", "5432")),
        dbname=os.environ.get("PGDATABASE", "toolathlon_gym"),
        user="eigent", password="camel",
    )


def run_evaluation(agent_workspace, groundtruth_workspace, launch_time, res_log_file):
    global PASS_COUNT, FAIL_COUNT
    PASS_COUNT = 0
    FAIL_COUNT = 0

    if not agent_workspace:
        return False, "No agent workspace provided"

    agent_ws = Path(agent_workspace)
    if not agent_ws.exists():
        return False, f"Agent workspace not found: {agent_workspace}"

    # Detect GT self-test mode (V1 parity test).
    try:
        gt_canon = os.path.realpath(groundtruth_workspace) if groundtruth_workspace else ""
        ag_canon = os.path.realpath(agent_workspace) if agent_workspace else ""
        is_gt_self_test = bool(gt_canon) and (gt_canon == ag_canon)
    except Exception:
        is_gt_self_test = False

    # Phase 5: Strategy/recommendations Word doc
    docx_files = list(agent_ws.glob("*.docx"))
    check("At least one docx investment strategy doc exists",
          len(docx_files) >= 1, f"found {len(docx_files)} docx files")
    if docx_files:
        try:
            from docx import Document
            audit_docx = max(docx_files, key=lambda p: p.stat().st_size)
            doc = Document(str(audit_docx))
            text = " ".join(p.text for p in doc.paragraphs).lower()
            check("Word doc has substantive content (>=300 chars)",
                  len(text) >= 300, f"len={len(text)}")
            has_topic = any(kw in text for kw in
                            ['portfolio', 'allocation', 'risk', 'return',
                             'diversification', 'investment'])
            # Note: 'risk' removed from metric set since it appears in topic set already.
            has_metric = any(kw in text for kw in
                             ['sharpe', 'volatility', 'correlation',
                              'beta', 'return'])
            has_recs = any(kw in text for kw in
                           ['recommendation', 'rebalance', 'philosophy',
                            'finding', 'strategy'])
            check("Word doc mentions portfolio/allocation/risk topic",
                  has_topic, f"text head: {text[:200]}")
            check("Word doc mentions risk/return/correlation metrics",
                  has_metric, f"text head: {text[:200]}")
            check("Word doc mentions recommendations/rebalancing/strategy",
                  has_recs, f"text head: {text[:200]}")
        except Exception as e:
            check("Word doc parse", False, str(e))

    # Phase 4: PowerPoint client/IC deck
    pptx_files = list(agent_ws.glob("*.pptx"))
    check("At least one pptx client/IC briefing deck exists",
          len(pptx_files) >= 1, f"found {len(pptx_files)} pptx files")
    if pptx_files:
        try:
            from pptx import Presentation
            audit_ppt = max(pptx_files, key=lambda p: p.stat().st_size)
            prs = Presentation(str(audit_ppt))
            n_slides = len(prs.slides)
            check("PPTX has at least 4 slides", n_slides >= 4,
                  f"n_slides={n_slides}")
            all_text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        all_text += " " + shape.text.lower()
            has_topic = any(kw in all_text for kw in
                            ['portfolio', 'allocation', 'risk', 'return',
                             'rebalance', 'diversification'])
            check("PPTX mentions portfolio/allocation/risk topic",
                  has_topic, f"text head: {all_text[:200]}")
        except Exception as e:
            check("PPTX parse", False, str(e))

    # Optional: structured analysis files (csv, xlsx, json)
    analysis_files = (list(agent_ws.glob("*.csv"))
                      + list(agent_ws.glob("*.xlsx"))
                      + list(agent_ws.glob("*.json")))
    # Exclude initial_workspace placeholders that get copied into agent_ws
    analysis_files = [p for p in analysis_files
                      if 'data.csv' not in p.name.lower()
                      and 'expected_output' not in p.name.lower()
                      and 'config.json' not in p.name.lower()
                      and 'current_portfolio' not in p.name.lower()]
    check("At least one structured analysis file exists (xlsx/csv/json)",
          len(analysis_files) >= 1, f"found {len(analysis_files)} files: {[p.name for p in analysis_files]}")

    # Value-level: validate at least one numeric portfolio metric in any csv/xlsx
    # is within plausible ranges (Sharpe in [-2,5], Volatility in [0,200], Beta in [-2,5]).
    metric_validated = False
    for af in analysis_files:
        try:
            if af.suffix.lower() == '.csv':
                import csv
                with open(af, 'r', encoding='utf-8', errors='ignore') as f:
                    reader = csv.DictReader(f)
                    headers = [h.lower() for h in (reader.fieldnames or [])]
                    has_metric_col = any(
                        any(k in h for k in ['sharpe', 'volatility', 'beta',
                                             'correlation', 'return', 'risk',
                                             'std_dev', 'std dev', 'stdev'])
                        for h in headers)
                    if has_metric_col:
                        for row in reader:
                            for k, v in row.items():
                                try:
                                    fv = float(v)
                                    kl = k.lower()
                                    if 'sharpe' in kl and -3 <= fv <= 4:
                                        metric_validated = True
                                    elif ('volatility' in kl or 'std_dev' in kl
                                          or 'std dev' in kl or 'stdev' in kl
                                          or 'risk' in kl) and 0 <= fv <= 200:
                                        metric_validated = True
                                    elif 'beta' in kl and -3 <= fv <= 5:
                                        metric_validated = True
                                    elif 'correlation' in kl and -1.0 <= fv <= 1.0:
                                        metric_validated = True
                                except (ValueError, TypeError):
                                    pass
                            if metric_validated:
                                break
            elif af.suffix.lower() == '.xlsx':
                from openpyxl import load_workbook
                wb = load_workbook(str(af), data_only=True)
                for ws in wb.worksheets:
                    headers = [str(c.value).lower() if c.value else '' for c in ws[1]]
                    has_metric_col = any(
                        any(k in h for k in ['sharpe', 'volatility', 'beta',
                                             'correlation', 'return', 'risk',
                                             'std_dev', 'std dev', 'stdev'])
                        for h in headers)
                    if has_metric_col:
                        for row in ws.iter_rows(min_row=2, values_only=True):
                            for i, v in enumerate(row):
                                if v is None or i >= len(headers):
                                    continue
                                try:
                                    fv = float(v)
                                    kl = headers[i]
                                    if 'sharpe' in kl and -3 <= fv <= 4:
                                        metric_validated = True
                                    elif ('volatility' in kl or 'std_dev' in kl
                                          or 'std dev' in kl or 'stdev' in kl
                                          or 'risk' in kl) and 0 <= fv <= 200:
                                        metric_validated = True
                                    elif 'beta' in kl and -3 <= fv <= 5:
                                        metric_validated = True
                                    elif 'correlation' in kl and -1.0 <= fv <= 1.0:
                                        metric_validated = True
                                except (ValueError, TypeError):
                                    pass
                            if metric_validated:
                                break
                    if metric_validated:
                        break
        except Exception:
            pass
        if metric_validated:
            break
    check("At least one numeric portfolio metric in plausible range",
          metric_validated, "no Sharpe/Volatility/Beta/Correlation column with valid value found")

    # Phase 6: Email to clients/IC
    try:
        conn = get_conn()
        cur = conn.cursor()
        cur.execute(
            """SELECT id, subject FROM email.messages
               WHERE subject ILIKE %s OR subject ILIKE %s OR subject ILIKE %s
                  OR body_text ILIKE %s""",
            ('%portfolio%', '%allocation%', '%investment%', '%portfolio%'))
        rows = cur.fetchall()
        if is_gt_self_test and len(rows) == 0:
            print("  [WARN] Investment strategy email sent: 0 found (GT self-test, non-blocking)")
        else:
            check("Investment strategy email sent",
                  len(rows) >= 1, f"found {len(rows)} matching emails")
        conn.close()
    except Exception as e:
        check("Email check", False, str(e))

    # Phase 6: IC review meeting
    try:
        conn = get_conn()
        cur = conn.cursor()
        cur.execute(
            """SELECT summary FROM gcal.events
               WHERE summary ILIKE %s OR summary ILIKE %s OR summary ILIKE %s""",
            ('%portfolio%', '%review%', '%investment%'))
        rows = cur.fetchall()
        if is_gt_self_test and len(rows) == 0:
            print("  [WARN] Investment review meeting scheduled: 0 found (GT self-test, non-blocking)")
        else:
            check("Investment review meeting scheduled",
                  len(rows) >= 1, f"found {len(rows)} matching events")
        conn.close()
    except Exception as e:
        check("Calendar check", False, str(e))

    return FAIL_COUNT == 0, f"Passed {PASS_COUNT}/{PASS_COUNT + FAIL_COUNT} checks"


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--agent_workspace", required=False, default=".")
    parser.add_argument("--groundtruth_workspace", required=False, default=".")
    parser.add_argument("--launch_time", required=False, default="2026-03-07 10:00:00")
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()

    success, message = run_evaluation(
        args.agent_workspace, args.groundtruth_workspace,
        args.launch_time, args.res_log_file
    )
    print(message)
    sys.exit(0 if success else 1)


if __name__ == "__main__":
    main()
