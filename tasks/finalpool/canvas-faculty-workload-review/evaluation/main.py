"""Evaluation for canvas-faculty-workload-review."""
import argparse
import os
import sys
from collections import defaultdict

import openpyxl
import psycopg2

DB_CONFIG = {
    "host": os.environ.get("PGHOST", "localhost"),
    "port": 5432,
    "dbname": os.environ.get("PGDATABASE", "toolathlon_gym"),
    "user": "eigent",
    "password": "camel",
}

PASS_COUNT = 0
FAIL_COUNT = 0


def check(name, condition, detail=""):
    global PASS_COUNT, FAIL_COUNT
    if condition:
        PASS_COUNT += 1
        print(f"  [PASS] {name}")
    else:
        FAIL_COUNT += 1
        msg = f": {str(detail)[:300]}" if detail else ""
        print(f"  [FAIL] {name}{msg}")


def num_close(a, b, tol=2.0):
    if a is None or b is None:
        return False
    try:
        return abs(float(a) - float(b)) <= tol
    except (TypeError, ValueError):
        return False


def load_sheet_rows(wb, sheet_name):
    for name in wb.sheetnames:
        if name.strip().lower() == sheet_name.strip().lower():
            return [[cell.value for cell in row] for row in wb[name].iter_rows()]
    return None


def get_expected_instructors():
    """Get expected instructor data from DB."""
    conn = psycopg2.connect(**DB_CONFIG)
    cur = conn.cursor()
    cur.execute("""
        SELECT u.name as instructor,
               COUNT(DISTINCT c.id) as courses_count,
               SUM(sc.cnt) as total_students,
               SUM(ac.cnt) as total_assignments,
               ROUND(SUM(sc.cnt) * 0.5, 1) as est_grading_hours
        FROM canvas.enrollments e
        JOIN canvas.users u ON u.id = e.user_id
        JOIN canvas.courses c ON c.id = e.course_id
        LEFT JOIN (SELECT course_id, COUNT(*) as cnt FROM canvas.enrollments WHERE type='StudentEnrollment' GROUP BY course_id) sc ON sc.course_id = c.id
        LEFT JOIN (SELECT course_id, COUNT(*) as cnt FROM canvas.assignments GROUP BY course_id) ac ON ac.course_id = c.id
        WHERE e.type = 'TeacherEnrollment'
        GROUP BY u.name
        ORDER BY u.name
    """)
    rows = cur.fetchall()
    cur.close()
    conn.close()
    return rows


def get_expected_departments():
    """Expected Department Summary from DB, matching the aggregation defined in
    task.md (using initial_workspace/department_mapping.csv):
      - department = subject area of each course's base name
      - Instructor_Count = distinct instructors teaching >=1 course in the dept
      - Course_Count = number of course offerings in the dept (all offerings,
        regardless of whether the course has an instructor assigned)
      - Total_Students = sum of StudentEnrollment counts across the dept's offerings
    """
    DEPT_MAP = {
        "applied analytics & algorithms": "Analytics",
        "biochemistry & bioinformatics": "Biochemistry",
        "creative computing & culture": "Computing",
        "data-driven design": "Design",
        "environmental economics & ethics": "Economics",
        "foundations of finance": "Finance",
        "global governance & geopolitics": "Governance",
    }

    def dept_of(name):
        base = str(name).split("(")[0].strip().lower()
        return DEPT_MAP.get(base)

    conn = psycopg2.connect(**DB_CONFIG)
    cur = conn.cursor()
    cur.execute("SELECT id, name FROM canvas.courses")
    courses = cur.fetchall()
    cur.execute("""SELECT course_id, COUNT(*) FROM canvas.enrollments
        WHERE type='StudentEnrollment' GROUP BY course_id""")
    stu_counts = dict(cur.fetchall())
    cur.execute("""SELECT course_id, user_id FROM canvas.enrollments
        WHERE type='TeacherEnrollment'""")
    dept_instr = defaultdict(set)
    dept_courses = defaultdict(list)  # dept -> [student count per offering]
    for cid, uid in cur.fetchall():
        cname = next((n for i, n in courses if i == cid), None)
        d = dept_of(cname) if cname is not None else None
        if d is not None:
            dept_instr[d].add(uid)
    cur.close()
    conn.close()

    for cid, name in courses:
        d = dept_of(name)
        if d is not None:
            dept_courses[d].append(stu_counts.get(cid, 0))

    rows = []
    for d in sorted(dept_courses, key=str.lower):
        rows.append((d, len(dept_instr[d]), len(dept_courses[d]), sum(dept_courses[d])))
    return rows


def check_excel(agent_workspace):
    print("\n=== Checking Excel ===")
    xlsx_path = os.path.join(agent_workspace, "Faculty_Workload.xlsx")
    if not os.path.isfile(xlsx_path):
        check("Faculty_Workload.xlsx exists", False, f"Not found: {xlsx_path}")
        return
    check("Faculty_Workload.xlsx exists", True)

    try:
        wb = openpyxl.load_workbook(xlsx_path, data_only=True)
    except Exception as e:
        check("Excel readable", False, str(e))
        return
    check("Excel readable", True)

    expected = get_expected_instructors()

    # Instructor Load sheet
    il_rows = load_sheet_rows(wb, "Instructor Load")
    if il_rows is None:
        check("Sheet 'Instructor Load' exists", False, f"Available: {wb.sheetnames}")
    else:
        check("Sheet 'Instructor Load' exists", True)
        data_rows = il_rows[1:] if len(il_rows) > 1 else []
        check(f"Instructor Load has {len(expected)} rows",
              abs(len(data_rows) - len(expected)) <= 1,
              f"Found {len(data_rows)}")

        header = il_rows[0] if il_rows else []
        header_lower = [str(h).lower().replace(" ", "_") if h else "" for h in header]
        for col in ["instructor", "courses_count", "total_students",
                    "total_assignments", "est_grading_hours",
                    "weekly_grading_hours", "overloaded_yn"]:
            check(f"Column '{col}' present", any(col in h for h in header_lower),
                  f"Header: {header}")

        # Build column-index map
        def col_idx(target):
            for i, h in enumerate(header_lower):
                if target in h:
                    return i
            return None

        i_instr = col_idx("instructor")
        i_courses = col_idx("courses_count") if col_idx("courses_count") is not None else col_idx("course")
        i_students = col_idx("total_students") if col_idx("total_students") is not None else col_idx("student")
        i_assignments = col_idx("total_assignments") if col_idx("total_assignments") is not None else col_idx("assignment")
        i_est = col_idx("est_grading_hours") if col_idx("est_grading_hours") is not None else col_idx("est_")
        i_weekly = col_idx("weekly_grading_hours") if col_idx("weekly_grading_hours") is not None else col_idx("weekly")
        i_overload = col_idx("overloaded") if col_idx("overloaded") is not None else (len(header_lower) - 1)

        # Build agent map: instructor_name -> row data
        agent_map = {}
        if i_instr is not None:
            for row in data_rows:
                if row and i_instr < len(row) and row[i_instr]:
                    agent_map[str(row[i_instr]).strip().lower()] = row

        # Validate EVERY expected instructor (not just 2 spot checks)
        for inst_row in expected:
            instr_name, courses_count, total_students, total_assignments, est_hrs = inst_row
            key = str(instr_name).strip().lower()
            actual_row = agent_map.get(key)
            if actual_row is None:
                check(f"Instructor '{instr_name}' present", False,
                      f"agent rows: {list(agent_map.keys())[:10]}")
                continue
            check(f"Instructor '{instr_name}' present", True)
            # courses_count
            if i_courses is not None and i_courses < len(actual_row):
                check(f"'{instr_name}' courses_count == {courses_count}",
                      num_close(actual_row[i_courses], courses_count, 0),
                      f"Got {actual_row[i_courses]}")
            # total_students
            if i_students is not None and i_students < len(actual_row):
                check(f"'{instr_name}' total_students == {total_students}",
                      num_close(actual_row[i_students], total_students, 5),
                      f"Got {actual_row[i_students]}")
            # total_assignments
            if i_assignments is not None and i_assignments < len(actual_row):
                check(f"'{instr_name}' total_assignments == {total_assignments}",
                      num_close(actual_row[i_assignments], total_assignments, 0),
                      f"Got {actual_row[i_assignments]}")
            # est_grading_hours (total semester) = total_students * grading_hours_per_student (0.5)
            if i_est is not None and i_est < len(actual_row):
                check(f"'{instr_name}' est_grading_hours == {est_hrs}",
                      num_close(actual_row[i_est], est_hrs, 0.5),
                      f"Got {actual_row[i_est]}")
            # weekly_grading_hours = est_grading_hours / semester_weeks (16)
            weekly_hours = float(total_students) * 0.5 / 16
            if i_weekly is not None and i_weekly < len(actual_row):
                check(f"'{instr_name}' weekly_grading_hours == {round(weekly_hours, 1)}",
                      num_close(actual_row[i_weekly], round(weekly_hours, 1), 0.5),
                      f"Got {actual_row[i_weekly]}")
            # overloaded computed
            expected_overloaded = (courses_count > 4) or (weekly_hours > 40)
            if i_overload is not None and i_overload < len(actual_row):
                actual_y = str(actual_row[i_overload]).strip().lower() in ("yes", "y", "true", "1")
                check(f"'{instr_name}' overloaded_yn == {'Yes' if expected_overloaded else 'No'}",
                      actual_y == expected_overloaded,
                      f"Got '{actual_row[i_overload]}'; weekly_hours={weekly_hours:.1f}")

    # Department Summary sheet
    ds_rows = load_sheet_rows(wb, "Department Summary")
    if ds_rows is None:
        check("Sheet 'Department Summary' exists", False, f"Available: {wb.sheetnames}")
    else:
        check("Sheet 'Department Summary' exists", True)
        data_rows = ds_rows[1:] if len(ds_rows) > 1 else []
        expected_deps = get_expected_departments()
        check(f"Department Summary has {len(expected_deps)} departments",
              len(data_rows) == len(expected_deps),
              f"Found {len(data_rows)}, expected {len(expected_deps)}")
        # Verify required columns
        if ds_rows:
            ds_header = [str(h).lower().replace(" ", "_") if h else "" for h in ds_rows[0]]
            for col in ["department", "instructor_count", "course_count", "total_students"]:
                check(f"Department Summary has '{col}'",
                      any(col in h for h in ds_header), f"Got: {ds_rows[0]}")

        # Validate every department's values (aggregation per department_mapping.csv)
        if ds_rows:
            ds_header = [str(h).lower().replace(" ", "_") if h else "" for h in ds_rows[0]]
            def col_idx(target):
                for i, h in enumerate(ds_header):
                    if target in h:
                        return i
                return None
            i_dep = col_idx("department")
            i_ic = col_idx("instructor_count")
            i_cc = col_idx("course_count")
            i_ts = col_idx("total_students")

            agent_deps = {}
            for row in data_rows:
                if row and i_dep is not None and i_dep < len(row) and row[i_dep]:
                    agent_deps[str(row[i_dep]).strip().lower()] = row

            for dname, ic, cc, ts in expected_deps:
                row = agent_deps.get(str(dname).strip().lower())
                if row is None:
                    check(f"Department '{dname}' present", False,
                          f"Agent departments: {sorted(agent_deps.keys())}")
                    continue
                check(f"Department '{dname}' present", True)
                if i_ic is not None and i_ic < len(row):
                    check(f"'{dname}' instructor_count == {ic}",
                          num_close(row[i_ic], ic, 0), f"Got {row[i_ic]}")
                if i_cc is not None and i_cc < len(row):
                    check(f"'{dname}' course_count == {cc}",
                          num_close(row[i_cc], cc, 0), f"Got {row[i_cc]}")
                if i_ts is not None and i_ts < len(row):
                    check(f"'{dname}' total_students == {ts}",
                          num_close(row[i_ts], ts, 5), f"Got {row[i_ts]}")


def check_pptx(agent_workspace):
    print("\n=== Checking PowerPoint ===")
    pptx_path = os.path.join(agent_workspace, "Workload_Review.pptx")
    if not os.path.isfile(pptx_path):
        check("Workload_Review.pptx exists", False, f"Not found: {pptx_path}")
        return
    check("Workload_Review.pptx exists", True)

    try:
        from pptx import Presentation
        prs = Presentation(pptx_path)
        check("PPTX has at least 3 slides", len(prs.slides) >= 3,
              f"Found {len(prs.slides)} slides")
        all_text = " ".join(
            shape.text for slide in prs.slides for shape in slide.shapes if shape.has_text_frame
        ).lower()
        check("PPTX mentions workload or overloaded",
              "workload" in all_text or "overload" in all_text,
              f"Sample: {all_text[:200]}")
    except ImportError:
        check("python-pptx available", False)


def check_gsheet():
    print("\n=== Checking Google Sheet ===")
    try:
        conn = psycopg2.connect(**DB_CONFIG)
        cur = conn.cursor()
        cur.execute("""
            SELECT id, title FROM gsheet.spreadsheets
            WHERE title ILIKE '%%workload%%' OR title ILIKE '%%faculty%%'
        """)
        sheets = cur.fetchall()
        check("Google Sheet created for workload data", len(sheets) >= 1,
              "No matching spreadsheet found")
        if sheets:
            sid = sheets[0][0]
            cur.execute("SELECT COUNT(*) FROM gsheet.cells WHERE spreadsheet_id = %s", (sid,))
            cell_count = cur.fetchone()[0]
            check("Google Sheet has data (cells)", cell_count > 10,
                  f"Found {cell_count} cells")
        cur.close()
        conn.close()
    except Exception as e:
        check("Google Sheet check", False, str(e))


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--agent_workspace", required=True)
    parser.add_argument("--groundtruth_workspace", required=False)
    parser.add_argument("--launch_time", required=False)
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()

    check_excel(args.agent_workspace)
    check_pptx(args.agent_workspace)
    check_gsheet()

    total = PASS_COUNT + FAIL_COUNT
    print(f"\n=== Results: {PASS_COUNT}/{total} passed ===")
    if FAIL_COUNT > 0:
        print(f"{FAIL_COUNT} checks failed")
        sys.exit(1)
    else:
        print("All checks passed!")
        sys.exit(0)


if __name__ == "__main__":
    main()
