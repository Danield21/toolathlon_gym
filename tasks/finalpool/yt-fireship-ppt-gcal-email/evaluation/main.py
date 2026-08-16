"""
Evaluation for yt-fireship-ppt-gcal-email task.

Checks:
1. Fireship_2025_Digest.pptx exists and has 7 slides (title + 6 monthly)
2. Title slide contains 'Fireship 2025 Monthly Tech Digest'
3. Monthly slides contain month names and video data
4. GCal has 6 monthly review events in Jan-Jun 2025
5. GCal event titles match pattern 'Monthly Tech Digest Review - [Month Year]'
6. Email sent to team@company.com with subject 'Fireship 2025 Monthly Digest Ready'
7. Email body lists top videos from multiple months
"""
import json
import os
import sys
from argparse import ArgumentParser

import psycopg2
from zoneinfo import ZoneInfo

# ──────────────────────────────────────────────────────────────────────────
# EVALUATION GROUND TRUTH SPEC (gcal tz root-fix v3, case-study 2026-08-13)
# Source of truth: docs/task.md. The Calendar MCP writes events as UTC
# instants into the DB (`TIMESTAMPTZ`), so all comparisons anchor to the
# timezone that task.md declares for the event's wall-clock semantics.
# Never use bare `sd.hour` / `sd.date()` on rows read from `gcal.events`:
# psycopg2 returns them in the PG session `TimeZone` (case-study: compute
# node default was Asia/Shanghai, masking 9 AM ET as 21:00+08). Use
# `utils.evaluation.gcal_helpers` instead (session-tz-independent).
# ──────────────────────────────────────────────────────────────────────────
# task.md line 7: "2025-01-07 from 09:00 to 09:30" ... "2025-02-04 from
# 09:00 to 09:30" (six monthly review events, 09:00-09:30).
EXPECTED_TIMEZONE = ZoneInfo("America/New_York")
EXPECTED_REVIEW_HOUR = 9
EXPECTED_REVIEW_MINUTE = 0
EXPECTED_REVIEW_DURATION_MIN = 30

from utils.evaluation.gcal_helpers import get_zone_components  # noqa: E402

DB_CONFIG = {
    "host": os.environ.get("PGHOST", "localhost"),
    "port": int(os.environ.get("PGPORT", "5432")),
    "dbname": "toolathlon_gym",
    "user": "eigent",
    "password": "camel",
}

PASS_COUNT = 0
FAIL_COUNT = 0


def record(name, passed, detail=""):
    global PASS_COUNT, FAIL_COUNT
    if passed:
        PASS_COUNT += 1
        print(f"  [PASS] {name}")
    else:
        FAIL_COUNT += 1
        msg = f": {detail[:300]}" if detail else ""
        print(f"  [FAIL] {name}{msg}")


def _fetch_top_videos():
    """Query DB for top video per month + total monthly count."""
    try:
        conn = psycopg2.connect(**DB_CONFIG)
        cur = conn.cursor()
        cur.execute("""
            WITH ranked AS (
                SELECT TO_CHAR(published_at, 'YYYY-MM') as ym, title, view_count, duration,
                       ROW_NUMBER() OVER (PARTITION BY TO_CHAR(published_at, 'YYYY-MM')
                                          ORDER BY view_count DESC) as rn
                FROM youtube.videos
                WHERE channel_title ILIKE '%fireship%'
                  AND published_at >= '2025-01-01' AND published_at < '2025-07-01'
            )
            SELECT ym, title, view_count, duration FROM ranked WHERE rn = 1 ORDER BY ym
        """)
        top_per_month = cur.fetchall()
        cur.execute("""
            SELECT TO_CHAR(published_at, 'YYYY-MM'), COUNT(*)
            FROM youtube.videos
            WHERE channel_title ILIKE '%fireship%'
              AND published_at >= '2025-01-01' AND published_at < '2025-07-01'
            GROUP BY 1 ORDER BY 1
        """)
        monthly_total = {r[0]: r[1] for r in cur.fetchall()}
        cur.close()
        conn.close()
        # Build dict by ym
        result = {}
        for ym, title, vc, dur in top_per_month:
            try:
                dur_min = round(float(dur) / 60.0, 1)
            except (TypeError, ValueError):
                dur_min = None
            result[ym] = {"title": title, "view_count": int(vc), "duration_min": dur_min}
        return result, monthly_total
    except Exception as e:
        print(f"[WARN] _fetch_top_videos: {e}")
        return None, None


def check_pptx(agent_workspace, top_per_month, monthly_total):
    print("\n=== Check 1: Fireship_2025_Digest.pptx ===")
    pptx_path = os.path.join(agent_workspace, "Fireship_2025_Digest.pptx")
    if not os.path.exists(pptx_path):
        record("Fireship_2025_Digest.pptx exists", False, f"Not found at {pptx_path}")
        return
    record("Fireship_2025_Digest.pptx exists", True)

    try:
        from pptx import Presentation
        prs = Presentation(pptx_path)
    except Exception as e:
        record("PPTX file readable", False, str(e))
        return
    record("PPTX file readable", True)

    slide_count = len(prs.slides)
    record("PPTX has at least 7 slides (title + 6 months)", slide_count >= 7,
           f"Found {slide_count} slides")

    # Per-slide text (use both .text attr and text_frame.text where available)
    slide_texts = []
    for slide in prs.slides:
        stext = ""
        for shape in slide.shapes:
            t = getattr(shape, "text", "") or ""
            if not t and getattr(shape, "has_text_frame", False):
                try:
                    t = shape.text_frame.text or ""
                except Exception:
                    t = ""
            if t:
                stext += " " + t
        slide_texts.append(stext)
    full_text = " ".join(slide_texts).lower()

    record("Title slide contains 'Fireship 2025 Monthly Tech Digest'",
           "fireship 2025 monthly tech digest" in slide_texts[0].lower() if slide_texts else False,
           "Title text not found")

    # Check ALL 6 months present
    expected_months = [("january", "01"), ("february", "02"), ("march", "03"),
                       ("april", "04"), ("may", "05"), ("june", "06")]
    months_found = 0
    for mname, _ in expected_months:
        if f"{mname} 2025" in full_text:
            months_found += 1
    record("All 6 month names appear in slides", months_found == 6,
           f"Found {months_found}/6")

    # If we have DB data, validate per-month content
    if top_per_month and monthly_total:
        for mname, mnum in expected_months:
            ym = f"2025-{mnum}"
            tv = top_per_month.get(ym)
            mt = monthly_total.get(ym)
            if not tv:
                continue
            # Find slide containing this month name (case-insensitive)
            slide_idx = None
            for i, st in enumerate(slide_texts):
                if f"{mname} 2025".lower() in st.lower() and i > 0:
                    slide_idx = i
                    break
            if slide_idx is None:
                record(f"Slide for {mname.title()} 2025", False, "no slide with month title")
                continue
            stxt = slide_texts[slide_idx]
            stxt_lower = stxt.lower()
            # Top video title: check for distinctive keyword (>=5 alphabetic chars)
            import re as _re
            title_lower = (tv["title"] or "").lower()
            # Extract alphabetic words (>=5 chars), strip punctuation
            words = [w for w in _re.findall(r"[a-z]+", title_lower) if len(w) >= 5]
            keyword_found = False
            for w in words[:5]:
                if w in stxt_lower:
                    keyword_found = True
                    break
            # Fallback: check first 15 chars of original
            if not keyword_found and len(tv["title"]) >= 15:
                if tv["title"][:15].lower() in stxt_lower:
                    keyword_found = True
            record(f"Slide for {mname.title()} 2025 has top video title",
                   keyword_found, f"Looking for keyword from '{(tv['title'] or '')[:30]}' in slide")
            # View count (formatted with commas or plain)
            vc_str = str(tv["view_count"])
            vc_comma = f"{tv['view_count']:,}"
            record(f"Slide for {mname.title()} 2025 has view count {tv['view_count']}",
                   vc_str in stxt or vc_comma in stxt,
                   f"Looking for {vc_str} or {vc_comma}")
            # Duration in minutes (1 decimal)
            if tv["duration_min"] is not None:
                dur_str = f"{tv['duration_min']:.1f}"
                record(f"Slide for {mname.title()} 2025 has duration {dur_str} min",
                       dur_str in stxt or str(int(tv["duration_min"])) in stxt,
                       f"Looking for {dur_str}")
            # Total monthly count
            if mt is not None:
                record(f"Slide for {mname.title()} 2025 has monthly count {mt}",
                       str(mt) in stxt,
                       f"Looking for {mt}")


def check_gcal(top_per_month):
    print("\n=== Check 2: Google Calendar events ===")
    conn = psycopg2.connect(**DB_CONFIG)
    cur = conn.cursor()
    cur.execute("""
        SELECT summary, description, start_datetime, end_datetime
        FROM gcal.events
        WHERE start_datetime >= '2025-01-01' AND start_datetime < '2025-07-01'
        ORDER BY start_datetime
    """)
    events = cur.fetchall()
    cur.close()
    conn.close()

    digest_events = [
        e for e in events
        if "monthly tech digest review" in (e[0] or "").lower()
    ]

    # Strict count == 6
    record("Exactly 6 Monthly Tech Digest Review events",
           len(digest_events) == 6,
           f"Found {len(digest_events)} matching events out of {len(events)} total")

    # Assert DB returned 6 months of top videos so per-month checks can't silently skip.
    if top_per_month is not None:
        record(
            "DB has top videos for all 6 months (2025-01..2025-06)",
            len(top_per_month) == 6,
            f"got months: {sorted(top_per_month.keys()) if top_per_month else []}",
        )

    # Expected exact dates and event titles
    expected_events = [
        ("2025-01-07", "January 2025"),
        ("2025-02-04", "February 2025"),
        ("2025-03-04", "March 2025"),
        ("2025-04-01", "April 2025"),
        ("2025-05-06", "May 2025"),
        ("2025-06-03", "June 2025"),
    ]

    # Build event lookup by date
    by_date = {}
    for ev in digest_events:
        if ev[2]:
            # gcal.events.start_datetime is TIMESTAMPTZ (UTC instant). The
            # event date must be extracted in EXPECTED_TIMEZONE
            # (America/New_York per task.md), NOT via bare sd.date() which
            # silently follows the PG session tz (case-study 2026-08-13).
            ev_date, _, _ = get_zone_components(ev[2], EXPECTED_TIMEZONE)
            if ev_date is not None:
                by_date[ev_date.isoformat()] = ev

    for exp_date, exp_month_year in expected_events:
        ev = by_date.get(exp_date)
        if ev is None:
            record(f"GCal event on {exp_date}", False,
                   f"available dates: {list(by_date.keys())}")
            continue
        # Title pattern - exact equality (case-insensitive, trimmed) to prevent
        # extra prefixes/suffixes from sneaking through.
        summary = (ev[0] or "")
        expected_title = f"Monthly Tech Digest Review - {exp_month_year}"
        record(f"Event on {exp_date} title is exactly '{expected_title}'",
               summary.strip().lower() == expected_title.strip().lower(),
               f"got '{summary}'")
        # Time 09:00 start (in EXPECTED_TIMEZONE per task.md).
        ev_date, ev_hour, ev_minute = get_zone_components(ev[2], EXPECTED_TIMEZONE)
        if ev_hour is None:
            record(f"Event on {exp_date} starts at {EXPECTED_REVIEW_HOUR:02d}:00", False,
                   "start_datetime missing/unparseable")
        else:
            record(f"Event on {exp_date} starts at {EXPECTED_REVIEW_HOUR:02d}:00",
                   ev_hour == EXPECTED_REVIEW_HOUR and ev_minute == EXPECTED_REVIEW_MINUTE,
                   f"got {ev_hour:02d}:{ev_minute:02d}")
        # Duration 30 min (timedelta arithmetic on UTC instants is tz-safe).
        if ev[2] and ev[3]:
            duration_min = (ev[3] - ev[2]).total_seconds() / 60
            record(f"Event on {exp_date} is {EXPECTED_REVIEW_DURATION_MIN} minutes",
                   EXPECTED_REVIEW_DURATION_MIN - 5 <= duration_min <= EXPECTED_REVIEW_DURATION_MIN + 5,
                   f"duration {duration_min:.0f} min")
        # Description mentions that month's top video title
        if top_per_month:
            import re as _re2
            ym = exp_date[:7]
            tv = top_per_month.get(ym)
            if tv:
                desc = (ev[1] or "").lower()
                title_words = [w for w in _re2.findall(r"[a-z]+", (tv["title"] or "").lower())
                               if len(w) >= 5]
                found = any(w in desc for w in title_words[:5]) or (tv["title"] or "")[:20].lower() in desc
                record(f"Event on {exp_date} description mentions top video",
                       found, f"description: {desc[:120]}")


def check_email(top_per_month):
    print("\n=== Check 3: Email to team@company.com ===")
    conn = psycopg2.connect(**DB_CONFIG)
    cur = conn.cursor()
    cur.execute("""
        SELECT to_addr, subject, body_text FROM email.messages
        WHERE to_addr::text ILIKE '%team@company.com%'
        ORDER BY id DESC LIMIT 10
    """)
    emails = cur.fetchall()
    cur.close()
    conn.close()

    # Find email with exact subject
    target = None
    for to, subj, body in emails:
        if subj and "fireship 2025 monthly digest ready" in subj.lower():
            target = (to, subj, body)
            break

    record("Email with subject 'Fireship 2025 Monthly Digest Ready' to team@company.com exists",
           target is not None,
           f"subjects: {[e[1] for e in emails]}")

    if target:
        to_addr, subject, body = target
        body_lower = (body or "").lower()
        # All 6 months mentioned
        for m in ["january", "february", "march", "april", "may", "june"]:
            record(f"Email body mentions {m.title()}",
                   m in body_lower, "missing month name")
        # All 6 top videos mentioned (by keyword)
        if top_per_month:
            import re as _re3
            for ym, tv in top_per_month.items():
                title_words = [w for w in _re3.findall(r"[a-z]+", (tv["title"] or "").lower())
                               if len(w) >= 5]
                found = any(w in body_lower for w in title_words[:5]) or (tv["title"] or "")[:20].lower() in body_lower
                record(f"Email body mentions top video for {ym}",
                       found, f"title keyword from '{(tv['title'] or '')[:30]}'")
                # View count mention
                vc = tv["view_count"]
                vc_str = str(vc)
                vc_comma = f"{vc:,}"
                record(f"Email body mentions view count for {ym}",
                       vc_str in (body or "") or vc_comma in (body or ""),
                       f"view count {vc}")


def main():
    parser = ArgumentParser()
    parser.add_argument("--agent_workspace", required=False, default=".")
    parser.add_argument("--groundtruth_workspace", required=False, default=".")
    parser.add_argument("--launch_time", required=False)
    parser.add_argument("--res_log_file", required=False)
    args = parser.parse_args()

    top_per_month, monthly_total = _fetch_top_videos()
    check_pptx(args.agent_workspace, top_per_month, monthly_total)
    check_gcal(top_per_month)
    check_email(top_per_month)

    total = PASS_COUNT + FAIL_COUNT
    if total == 0:
        print("\nFAIL: No checks were performed.")
        sys.exit(1)

    accuracy = PASS_COUNT / total * 100
    print(f"\nOverall: {PASS_COUNT}/{total} checks passed ({accuracy:.1f}%)")

    result = {
        "total_passed": PASS_COUNT,
        "total_checks": total,
        "accuracy": accuracy,
    }

    if args.res_log_file:
        with open(args.res_log_file, "w") as f:
            json.dump(result, f, indent=2)

    # Require all checks pass
    if FAIL_COUNT == 0:
        print("PASS")
        sys.exit(0)
    else:
        print("FAIL")
        sys.exit(1)


if __name__ == "__main__":
    main()
